"""
dual_run.py — Run both v1 and v2 transformers on the same PPTX and compare.

Orchestrates:
  1. Copy input PPTX twice → v1 working copy, v2 working copy
  2. Run Phase 2 (masters/layouts) on both
  3. Run Phase 3 v1 (SlideContentTransformer) on copy 1
  4. Run Phase 3 v2 (SlideContentTransformerV2) on copy 2
  5. Diff the two outputs shape-by-shape
  6. Classify diffs
  7. Generate markdown report
"""

from __future__ import annotations

import json
import logging
import os
import shutil
import tempfile
import time
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Callable, Dict, List, Optional

from pptx import Presentation

from .structural_differ import (
    EMU_PER_INCH,
    ShapeDiff,
    SlideDiff,
    StructuralDiffer,
)
from .diff_classifier import (
    ClassifiedDiff,
    ClassifiedSlide,
    DiffClassifier,
    DiffLabel,
    GoldenShape,
    GoldenSlide,
    classify_all,
)

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# Config & result types
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class DualRunConfig:
    """Configuration for a dual-run comparison."""
    input_pptx: str
    output_dir: str
    translate_fn: Optional[Callable] = None
    skip_translation: bool = True
    golden_slides: Optional[Dict[int, GoldenSlide]] = None
    render_images: bool = False
    render_dpi: int = 150


@dataclass
class DualRunResult:
    """Full result of a dual run."""
    input_path: str
    v1_output_path: str
    v2_output_path: str
    slide_diffs: List[SlideDiff] = field(default_factory=list)
    classified: List[ClassifiedSlide] = field(default_factory=list)
    duration_ms: int = 0

    @property
    def total_regressions(self) -> int:
        return sum(cs.regressions for cs in self.classified)

    @property
    def total_improvements(self) -> int:
        return sum(cs.improvements for cs in self.classified)

    @property
    def total_neutrals(self) -> int:
        return sum(cs.neutrals for cs in self.classified)

    @property
    def total_unknowns(self) -> int:
        return sum(
            sum(1 for d in cs.diffs if d.label == DiffLabel.UNKNOWN)
            for cs in self.classified
        )


@dataclass
class DualRunSummary:
    """Lightweight summary for serialization."""
    input_file: str
    slide_count: int
    total_shapes: int
    regressions: int
    improvements: int
    neutrals: int
    unknowns: int
    per_slide: List[dict] = field(default_factory=list)

    def to_dict(self) -> dict:
        return {
            'input_file': self.input_file,
            'slide_count': self.slide_count,
            'total_shapes': self.total_shapes,
            'regressions': self.regressions,
            'improvements': self.improvements,
            'neutrals': self.neutrals,
            'unknowns': self.unknowns,
            'per_slide': self.per_slide,
        }


# ─────────────────────────────────────────────────────────────────────────────
# Dummy translation helper
# ─────────────────────────────────────────────────────────────────────────────

def _dummy_translate(text: str) -> str:
    """Prefix text with Arabic marker for testing without API keys."""
    if not text or not text.strip():
        return text
    return f'\u0639\u0631\u0628\u064a: {text}'


# ─────────────────────────────────────────────────────────────────────────────
# DualRunHarness
# ─────────────────────────────────────────────────────────────────────────────

class DualRunHarness:
    """Run v1 and v2 on the same input and compare."""

    def __init__(self, config: DualRunConfig):
        self.config = config
        self._translate_fn = config.translate_fn or _dummy_translate

    def run(self) -> DualRunResult:
        """Execute the dual run."""
        t0 = time.time()
        cfg = self.config
        out = Path(cfg.output_dir)
        out.mkdir(parents=True, exist_ok=True)

        # Copy input for v1 and v2
        v1_path = str(out / 'v1_output.pptx')
        v2_path = str(out / 'v2_output.pptx')
        shutil.copy2(cfg.input_pptx, v1_path)
        shutil.copy2(cfg.input_pptx, v2_path)

        # --- Phase 2: Masters & Layouts (same for both) ---
        v1_prs = Presentation(v1_path)
        v2_prs = Presentation(v2_path)

        try:
            from slidearabi.rtl_transforms import MasterLayoutTransformer
            mlt_v1 = MasterLayoutTransformer(v1_prs)
            mlt_v1.transform_all_masters()
            mlt_v1.transform_all_layouts()

            mlt_v2 = MasterLayoutTransformer(v2_prs)
            mlt_v2.transform_all_masters()
            mlt_v2.transform_all_layouts()
            logger.info('Phase 2 complete for both copies')
        except ImportError:
            logger.warning('MasterLayoutTransformer not available; skipping Phase 2')

        # --- Build translations ---
        translations = {}
        if cfg.skip_translation:
            for slide in v1_prs.slides:
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                            txt = shape.text_frame.text
                            if txt and txt.strip():
                                translations[txt] = self._translate_fn(txt)
                    except Exception:
                        pass

        # --- Phase 3 v1: SlideContentTransformer ---
        v1_report = None
        try:
            from slidearabi.rtl_transforms import SlideContentTransformer
            sct_v1 = SlideContentTransformer(
                v1_prs, translations=translations)
            v1_report = sct_v1.transform_all_slides()
            logger.info('Phase 3 v1: %d changes', v1_report.total_changes)
        except ImportError:
            logger.warning('SlideContentTransformer not available')
        except Exception as e:
            logger.error('Phase 3 v1 failed: %s', e)

        # --- Phase 3 v2: SlideContentTransformerV2 ---
        v2_report = None
        try:
            from slidearabi_v2.rtl_transforms_v2 import SlideContentTransformerV2
            sct_v2 = SlideContentTransformerV2(
                v2_prs, translations=translations)
            v2_report = sct_v2.transform_all_slides()
            v2_total = sum(v2_report.counts.values()) if hasattr(v2_report, 'counts') else 0
            logger.info('Phase 3 v2: %d changes', v2_total)
        except ImportError:
            logger.warning('SlideContentTransformerV2 not available; v2 output = Phase 2 only')
        except Exception as e:
            logger.error('Phase 3 v2 failed: %s', e)

        # Save both outputs
        v1_prs.save(v1_path)
        v2_prs.save(v2_path)

        # --- Diff ---
        differ = StructuralDiffer(v1_path, v2_path)
        slide_diffs = differ.diff_all_slides()

        # --- Classify ---
        classified = classify_all(slide_diffs, cfg.golden_slides)

        duration_ms = int((time.time() - t0) * 1000)

        result = DualRunResult(
            input_path=cfg.input_pptx,
            v1_output_path=v1_path,
            v2_output_path=v2_path,
            slide_diffs=slide_diffs,
            classified=classified,
            duration_ms=duration_ms,
        )

        # --- Generate outputs ---
        self._write_report(result, out)
        self._write_json(result, out)

        if cfg.render_images:
            self._render_images(v1_path, v2_path, out, cfg.render_dpi)

        return result

    # ─── Report generation ────────────────────────────────────────────────

    def _write_report(self, result: DualRunResult, out: Path) -> None:
        """Write markdown comparison report."""
        lines = [
            '# Dual-Run Comparison Report',
            '',
            f'**Input**: `{os.path.basename(result.input_path)}`',
            f'**Duration**: {result.duration_ms}ms',
            f'**Slides**: {len(result.classified)}',
            '',
            '## Summary',
            '',
            f'| Label | Count |',
            f'|-------|-------|',
            f'| REGRESSION | {result.total_regressions} |',
            f'| IMPROVEMENT | {result.total_improvements} |',
            f'| NEUTRAL | {result.total_neutrals} |',
            f'| UNKNOWN | {result.total_unknowns} |',
            '',
        ]

        for cs in result.classified:
            lines.append(f'## Slide {cs.slide_number} ({cs.layout_type})')
            lines.append('')
            lines.append(f'R={cs.regressions} I={cs.improvements} '
                         f'N={cs.neutrals}')
            lines.append('')

            if not cs.diffs:
                lines.append('No shapes on this slide.')
                lines.append('')
                continue

            lines.append('| ID | Shape | Label | Severity | '
                         'v1 Left | v2 Left | ΔLeft | Reason |')
            lines.append('|----|-------|-------|----------|'
                         '---------|---------|-------|--------|')

            for cd in cs.diffs:
                sd = cd.shape_diff
                v1l = f'{sd.v1_left / EMU_PER_INCH:.2f}"'
                v2l = f'{sd.v2_left / EMU_PER_INCH:.2f}"'
                dl = f'{sd.left_delta}'
                if sd.v1_only:
                    v2l = '—'
                    dl = 'V1 ONLY'
                elif sd.v2_only:
                    v1l = '—'
                    dl = 'V2 ONLY'
                lines.append(
                    f'| {sd.shape_id} | {sd.shape_name} | '
                    f'{cd.label.value} | {cd.severity} | '
                    f'{v1l} | {v2l} | {dl} | {cd.reason} |'
                )
            lines.append('')

        report_path = out / 'dual_run_report.md'
        report_path.write_text('\n'.join(lines), encoding='utf-8')
        logger.info('Report written to %s', report_path)

    def _write_json(self, result: DualRunResult, out: Path) -> None:
        """Write structured JSON summary."""
        summary = DualRunSummary(
            input_file=os.path.basename(result.input_path),
            slide_count=len(result.classified),
            total_shapes=sum(len(cs.diffs) for cs in result.classified),
            regressions=result.total_regressions,
            improvements=result.total_improvements,
            neutrals=result.total_neutrals,
            unknowns=result.total_unknowns,
            per_slide=[
                {
                    'slide': cs.slide_number,
                    'layout': cs.layout_type,
                    'regressions': cs.regressions,
                    'improvements': cs.improvements,
                    'neutrals': cs.neutrals,
                    'diffs': [
                        {
                            'shape_id': cd.shape_diff.shape_id,
                            'shape_name': cd.shape_diff.shape_name,
                            'label': cd.label.value,
                            'severity': cd.severity,
                            'reason': cd.reason,
                            'left_delta': cd.shape_diff.left_delta,
                            'v1_left': cd.shape_diff.v1_left,
                            'v2_left': cd.shape_diff.v2_left,
                        }
                        for cd in cs.diffs
                    ],
                }
                for cs in result.classified
            ],
        )
        json_path = out / 'dual_run_summary.json'
        json_path.write_text(
            json.dumps(summary.to_dict(), indent=2, ensure_ascii=False),
            encoding='utf-8',
        )
        logger.info('JSON written to %s', json_path)

    def _render_images(
        self, v1_path: str, v2_path: str, out: Path, dpi: int
    ) -> None:
        """Render both PPTX files to images for visual comparison."""
        import subprocess

        for label, pptx_path in [('v1', v1_path), ('v2', v2_path)]:
            img_dir = out / f'{label}_images'
            img_dir.mkdir(exist_ok=True)
            try:
                # Convert PPTX → PDF
                subprocess.run(
                    ['soffice', '--headless', '--convert-to', 'pdf',
                     '--outdir', str(img_dir), pptx_path],
                    capture_output=True, timeout=120,
                )
                pdf_name = Path(pptx_path).stem + '.pdf'
                pdf_path = img_dir / pdf_name
                if pdf_path.exists():
                    subprocess.run(
                        ['pdftoppm', '-png', '-r', str(dpi),
                         str(pdf_path), str(img_dir / 'slide')],
                        capture_output=True, timeout=120,
                    )
                    logger.info('Rendered %s to %s', label, img_dir)
            except Exception as e:
                logger.warning('Rendering %s failed: %s', label, e)


# ─────────────────────────────────────────────────────────────────────────────
# Convenience
# ─────────────────────────────────────────────────────────────────────────────

def dual_run(
    input_pptx: str,
    output_dir: str,
    golden_slides: Optional[Dict[int, GoldenSlide]] = None,
    render: bool = False,
) -> DualRunResult:
    """One-liner: dual-run comparison."""
    cfg = DualRunConfig(
        input_pptx=input_pptx,
        output_dir=output_dir,
        golden_slides=golden_slides,
        render_images=render,
    )
    return DualRunHarness(cfg).run()


# ─────────────────────────────────────────────────────────────────────────────
# CLI
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == '__main__':
    import sys
    logging.basicConfig(level=logging.INFO, format='%(levelname)s: %(message)s')

    if len(sys.argv) < 3:
        print(f'Usage: {sys.argv[0]} <input.pptx> <output_dir> [--render]')
        sys.exit(1)

    render = '--render' in sys.argv
    result = dual_run(sys.argv[1], sys.argv[2], render=render)

    print(f'\nDual run complete in {result.duration_ms}ms')
    print(f'  Slides:       {len(result.classified)}')
    print(f'  REGRESSION:   {result.total_regressions}')
    print(f'  IMPROVEMENT:  {result.total_improvements}')
    print(f'  NEUTRAL:      {result.total_neutrals}')
    print(f'  UNKNOWN:      {result.total_unknowns}')
    print(f'\nv1 output: {result.v1_output_path}')
    print(f'v2 output: {result.v2_output_path}')

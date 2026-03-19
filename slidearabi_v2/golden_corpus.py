"""
golden_corpus.py — Freeze known-good shape positions as golden baselines.

Loads v1 patched output and records expected shape positions so the
DiffClassifier can compare v2 output against a known-good reference.
"""

from __future__ import annotations

import json
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation

from .diff_classifier import GoldenShape, GoldenSlide, BOUNDS_MARGIN_EMU
from .structural_differ import EMU_PER_INCH, POSITION_TOLERANCE_EMU

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# GoldenCorpus
# ─────────────────────────────────────────────────────────────────────────────

class GoldenCorpus:
    """Collection of golden slides from one or more presentations."""

    def __init__(self):
        self.slides: Dict[int, GoldenSlide] = {}

    def add_slide(
        self,
        slide_number: int,
        shapes: Dict[int, GoldenShape],
    ) -> None:
        """Add a golden slide with expected shape positions."""
        self.slides[slide_number] = GoldenSlide(
            slide_number=slide_number,
            shapes=shapes,
        )

    def add_from_pptx(
        self,
        pptx_path: str,
        slide_numbers: Optional[List[int]] = None,
        tolerance: int = POSITION_TOLERANCE_EMU,
    ) -> None:
        """Extract golden positions from an already-transformed PPTX.

        Args:
            pptx_path: Path to a transformed (v1 patched) PPTX file.
            slide_numbers: Which slides to freeze (1-indexed). None = all.
            tolerance: EMU tolerance for position matching.
        """
        prs = Presentation(pptx_path)
        total = len(prs.slides)

        if slide_numbers is None:
            slide_numbers = list(range(1, total + 1))

        for sn in slide_numbers:
            if sn < 1 or sn > total:
                logger.warning('Slide %d out of range (1-%d), skipping', sn, total)
                continue

            slide = prs.slides[sn - 1]
            shapes: Dict[int, GoldenShape] = {}

            for shape in slide.shapes:
                sid = shape.shape_id
                left = int(shape.left) if shape.left is not None else 0
                width = int(shape.width) if shape.width is not None else 0
                shapes[sid] = GoldenShape(
                    shape_id=sid,
                    expected_left=left,
                    expected_width=width,
                    tolerance=tolerance,
                )

            self.add_slide(sn, shapes)
            logger.debug('Frozen slide %d: %d shapes', sn, len(shapes))

        logger.info('Golden corpus: %d slides from %s', len(self.slides), pptx_path)

    def save(self, path: str) -> None:
        """Save golden corpus to JSON."""
        data = {}
        for sn, gs in sorted(self.slides.items()):
            data[str(sn)] = {
                str(sid): {
                    'shape_id': g.shape_id,
                    'expected_left': g.expected_left,
                    'expected_width': g.expected_width,
                    'tolerance': g.tolerance,
                }
                for sid, g in gs.shapes.items()
            }
        Path(path).write_text(
            json.dumps(data, indent=2), encoding='utf-8')
        logger.info('Saved golden corpus (%d slides) to %s',
                     len(self.slides), path)

    @classmethod
    def load(cls, path: str) -> 'GoldenCorpus':
        """Load golden corpus from JSON."""
        corpus = cls()
        data = json.loads(Path(path).read_text(encoding='utf-8'))
        for sn_str, shapes_data in data.items():
            sn = int(sn_str)
            shapes = {}
            for sid_str, g in shapes_data.items():
                sid = int(sid_str)
                shapes[sid] = GoldenShape(
                    shape_id=g['shape_id'],
                    expected_left=g['expected_left'],
                    expected_width=g['expected_width'],
                    tolerance=g.get('tolerance', POSITION_TOLERANCE_EMU),
                )
            corpus.add_slide(sn, shapes)
        logger.info('Loaded golden corpus: %d slides from %s',
                     len(corpus.slides), path)
        return corpus


# ─────────────────────────────────────────────────────────────────────────────
# Lukas corpus
# ─────────────────────────────────────────────────────────────────────────────

def freeze_lukas_corpus(
    v1_patched_path: str = '/home/user/workspace/lukas_patched_output.pptx',
    output_path: Optional[str] = None,
) -> GoldenCorpus:
    """Freeze the Lukas school presentation golden slides.

    Uses the v1 patched output (which passed all 13 verification checks)
    as the golden baseline.  Focuses on the 5 critical slides but includes
    all 13 for completeness.

    Args:
        v1_patched_path: Path to the v1 patched PPTX.
        output_path: Optional path to save the JSON corpus.

    Returns:
        GoldenCorpus with all 13 slides frozen.
    """
    corpus = GoldenCorpus()

    # The critical slides are 1, 7, 8, 12, 13 — but freeze all for completeness
    corpus.add_from_pptx(v1_patched_path)

    if output_path:
        corpus.save(output_path)

    return corpus


# ─────────────────────────────────────────────────────────────────────────────
# CLI
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == '__main__':
    import sys
    logging.basicConfig(level=logging.INFO, format='%(levelname)s: %(message)s')

    if len(sys.argv) < 2:
        print(f'Usage: {sys.argv[0]} <patched.pptx> [output.json]')
        print(f'       {sys.argv[0]} --lukas [output.json]')
        sys.exit(1)

    if sys.argv[1] == '--lukas':
        out = sys.argv[2] if len(sys.argv) > 2 else 'lukas_golden.json'
        corpus = freeze_lukas_corpus(output_path=out)
    else:
        out = sys.argv[2] if len(sys.argv) > 2 else 'golden_corpus.json'
        corpus = GoldenCorpus()
        corpus.add_from_pptx(sys.argv[1])
        corpus.save(out)

    print(f'Golden corpus: {len(corpus.slides)} slides')
    for sn, gs in sorted(corpus.slides.items()):
        print(f'  Slide {sn}: {len(gs.shapes)} shapes')

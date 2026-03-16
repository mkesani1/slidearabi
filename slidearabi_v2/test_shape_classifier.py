"""
Tests for SlideArabi v2 Shape Classification Layer.

Uses the LUKAS-SCHOOL-PRESENTATION.pptx to verify classification of
5 broken slides (1, 7, 8, 12, 13) — the slides that exhibited RTL
mirroring bugs in v1 due to misclassification of shape intent.

Run with:
    pytest slidearabi_v2/test_shape_classifier.py -v
"""

from __future__ import annotations

import os
from pathlib import Path

import pytest
from pptx import Presentation

from slidearabi_v2.shape_classifier import (
    ShapeClassifier,
    ShapeRole,
    ShapeClassification,
    SlideClassificationResult,
    SplitPanelInfo,
    SlideContext,
    classify_slide,
    _ROLE_ACTIONS,
)


# ═══════════════════════════════════════════════════════════════════════════════
# Fixtures
# ═══════════════════════════════════════════════════════════════════════════════

PPTX_PATH = Path(__file__).resolve().parent.parent / 'LUKAS-SCHOOL-PRESENTATION.pptx'


@pytest.fixture(scope='module')
def prs():
    """Load the test presentation once per module."""
    if not PPTX_PATH.exists():
        pytest.skip(f'Test PPTX not found: {PPTX_PATH}')
    return Presentation(str(PPTX_PATH))


@pytest.fixture(scope='module')
def slide_dims(prs):
    """Return (slide_width, slide_height) in EMU."""
    return int(prs.slide_width), int(prs.slide_height)


@pytest.fixture(scope='module')
def classifier(slide_dims):
    """Create a ShapeClassifier for the test presentation."""
    sw, sh = slide_dims
    return ShapeClassifier(sw, sh)


def _classify(prs, slide_dims, slide_number: int) -> SlideClassificationResult:
    """Classify a slide using the module-level convenience function."""
    sw, sh = slide_dims
    slide = prs.slides[slide_number - 1]
    layout_type = slide.slide_layout._element.get('type', 'cust')
    return classify_slide(slide, slide_number, sw, sh, layout_type)


def _get_shape_by_name(slide, name: str):
    """Find a shape on a slide by its name."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise ValueError(f'Shape {name!r} not found on slide')


def _role_for(result: SlideClassificationResult, slide, name: str) -> ShapeRole:
    """Get the ShapeRole for a named shape."""
    shape = _get_shape_by_name(slide, name)
    return result.get(shape).role


# ═══════════════════════════════════════════════════════════════════════════════
# Unit tests — ShapeRole and _ROLE_ACTIONS
# ═══════════════════════════════════════════════════════════════════════════════

class TestShapeRoleEnum:
    """Verify the ShapeRole enum is complete and well-formed."""

    def test_all_roles_present(self):
        expected = {
            'PLACEHOLDER', 'TABLE', 'CHART', 'CONNECTOR', 'DIRECTIONAL',
            'BACKGROUND', 'BLEED', 'FOOTER', 'BADGE', 'LOGO', 'OVERLAY',
            'PANEL_LEFT', 'PANEL_RIGHT', 'DECORATIVE', 'GROUP',
            'CONTENT_IMAGE', 'CONTENT_TEXT', 'UNKNOWN',
        }
        actual = {r.name for r in ShapeRole}
        assert actual == expected

    def test_role_actions_cover_all_roles(self):
        """Every ShapeRole must have an entry in _ROLE_ACTIONS."""
        for role in ShapeRole:
            assert role in _ROLE_ACTIONS, f'{role.name} missing from _ROLE_ACTIONS'

    def test_role_actions_have_required_keys(self):
        for role, actions in _ROLE_ACTIONS.items():
            assert 'position' in actions, f'{role.name} missing position action'
            assert 'text' in actions, f'{role.name} missing text action'
            assert 'direction' in actions, f'{role.name} missing direction action'

    def test_position_actions_are_valid(self):
        valid = {'mirror', 'keep', 'swap', 'inherit', 'reposition'}
        for role, actions in _ROLE_ACTIONS.items():
            assert actions['position'] in valid, (
                f'{role.name} has invalid position action: {actions["position"]}'
            )

    def test_text_actions_are_valid(self):
        valid = {'translate_rtl', 'rtl_only', 'none'}
        for role, actions in _ROLE_ACTIONS.items():
            assert actions['text'] in valid, (
                f'{role.name} has invalid text action: {actions["text"]}'
            )

    def test_direction_actions_are_valid(self):
        valid = {'remove_flip', 'toggle_flipH', 'swap_preset', 'none'}
        for role, actions in _ROLE_ACTIONS.items():
            assert actions['direction'] in valid, (
                f'{role.name} has invalid direction action: {actions["direction"]}'
            )


# ═══════════════════════════════════════════════════════════════════════════════
# Unit tests — ShapeClassification properties
# ═══════════════════════════════════════════════════════════════════════════════

class TestShapeClassificationProperties:

    def test_should_mirror(self):
        cls = ShapeClassification(
            role=ShapeRole.CONTENT_IMAGE,
            position_action='mirror',
            text_action='none',
            direction_action='remove_flip',
        )
        assert cls.should_mirror is True
        assert cls.should_swap is False
        assert cls.should_keep is False

    def test_should_swap(self):
        cls = ShapeClassification(
            role=ShapeRole.PANEL_LEFT,
            position_action='swap',
            text_action='translate_rtl',
            direction_action='remove_flip',
        )
        assert cls.should_swap is True
        assert cls.should_mirror is False

    def test_should_keep(self):
        cls = ShapeClassification(
            role=ShapeRole.BACKGROUND,
            position_action='keep',
            text_action='rtl_only',
            direction_action='remove_flip',
        )
        assert cls.should_keep is True
        assert cls.should_translate is True  # 'rtl_only' counts as translate

    def test_should_translate_covers_both(self):
        for action in ('translate_rtl', 'rtl_only'):
            cls = ShapeClassification(
                role=ShapeRole.UNKNOWN,
                position_action='mirror',
                text_action=action,
                direction_action='remove_flip',
            )
            assert cls.should_translate is True

    def test_should_not_translate_none(self):
        cls = ShapeClassification(
            role=ShapeRole.CONTENT_IMAGE,
            position_action='mirror',
            text_action='none',
            direction_action='remove_flip',
        )
        assert cls.should_translate is False


# ═══════════════════════════════════════════════════════════════════════════════
# Integration tests — Slide 1 (layout=picTx — "Picture with Caption")
#
# Shapes:
#   Picture 33  (100%w, 100%h)  → BACKGROUND (full-slide photo)
#   Title 3     (placeholder)    → PLACEHOLDER (title ph)
#   Picture 8   (50%w, 100%h)   → CONTENT_IMAGE (half-slide photo)
#
# v1 bug: Picture 33 was mirrored (nudged off-slide) and Picture 8's
# large half-slide image was treated as a panel anchor, causing the
# title to shift incorrectly.
# ═══════════════════════════════════════════════════════════════════════════════

class TestSlide1:
    """Slide 1 — picTx layout with full-background + half-page image."""

    @pytest.fixture(autouse=True)
    def setup(self, prs, slide_dims):
        self.slide = prs.slides[0]
        self.result = _classify(prs, slide_dims, 1)

    def test_layout_type(self):
        assert self.result.layout_type == 'picTx'

    def test_no_split_panel(self):
        assert self.result.has_split_panel is False

    def test_no_map_overlay(self):
        assert self.result.has_map_overlay is False

    def test_shape_count(self):
        assert len(self.result.classifications) == 3

    def test_picture_33_is_background(self):
        """Full-slide image (100%w × 100%h) → BACKGROUND."""
        role = _role_for(self.result, self.slide, 'Picture 33')
        assert role == ShapeRole.BACKGROUND

    def test_picture_33_action_is_keep(self):
        shape = _get_shape_by_name(self.slide, 'Picture 33')
        cls = self.result.get(shape)
        assert cls.position_action == 'keep'

    def test_title_3_is_placeholder(self):
        """Title placeholder → PLACEHOLDER."""
        role = _role_for(self.result, self.slide, 'Title 3')
        assert role == ShapeRole.PLACEHOLDER

    def test_title_3_action_is_inherit(self):
        shape = _get_shape_by_name(self.slide, 'Title 3')
        cls = self.result.get(shape)
        assert cls.position_action == 'inherit'
        assert cls.placeholder_type == 'title'

    def test_picture_8_is_content_image(self):
        """Half-page photo (50%w × 100%h) → CONTENT_IMAGE."""
        role = _role_for(self.result, self.slide, 'Picture 8')
        assert role == ShapeRole.CONTENT_IMAGE

    def test_picture_8_action_is_mirror(self):
        shape = _get_shape_by_name(self.slide, 'Picture 8')
        cls = self.result.get(shape)
        assert cls.position_action == 'mirror'


# ═══════════════════════════════════════════════════════════════════════════════
# Integration tests — Slide 7 (layout=title — split panel)
#
# Shapes:
#   Title 7     (39%w, 55%h, left)   → PANEL_LEFT (text title on left)
#   Graphic 11  (45%w, 80%h, right)  → PANEL_RIGHT (large image on right)
#
# v1 bug: Panel anchor detected but shapes on the "wrong" side were
# mirrored instead of swapped, breaking the spatial relationship.
# ═══════════════════════════════════════════════════════════════════════════════

class TestSlide7:
    """Slide 7 — title layout with detected split panel."""

    @pytest.fixture(autouse=True)
    def setup(self, prs, slide_dims):
        self.slide = prs.slides[6]
        self.result = _classify(prs, slide_dims, 7)

    def test_layout_type(self):
        assert self.result.layout_type == 'title'

    def test_has_split_panel(self):
        assert self.result.has_split_panel is True

    def test_shape_count(self):
        assert len(self.result.classifications) == 2

    def test_title_7_is_panel_left(self):
        """Title placeholder on left side of split → PANEL_LEFT."""
        role = _role_for(self.result, self.slide, 'Title 7')
        assert role == ShapeRole.PANEL_LEFT

    def test_title_7_action_is_swap(self):
        shape = _get_shape_by_name(self.slide, 'Title 7')
        cls = self.result.get(shape)
        assert cls.position_action == 'swap'

    def test_graphic_11_is_panel_right(self):
        """Large image on right side of split → PANEL_RIGHT."""
        role = _role_for(self.result, self.slide, 'Graphic 11')
        assert role == ShapeRole.PANEL_RIGHT

    def test_graphic_11_action_is_swap(self):
        shape = _get_shape_by_name(self.slide, 'Graphic 11')
        cls = self.result.get(shape)
        assert cls.position_action == 'swap'

    def test_panel_shift_delta_positive(self):
        """Panel shift delta must be positive (right panel is to the right)."""
        panel = self.result.context.split_panel
        assert panel is not None
        assert panel.shift_delta > 0

    def test_panel_groups_are_disjoint(self):
        """Left and right panel shape sets must not overlap."""
        panel = self.result.context.split_panel
        assert panel.left_shape_ids & panel.right_shape_ids == frozenset()


# ═══════════════════════════════════════════════════════════════════════════════
# Integration tests — Slide 8 (layout=title — complex split panel with groups)
#
# Shapes:
#   Title 7       (47%w, right)     → PANEL_RIGHT
#   Group 219     (40%w, right, bleeds top) → PANEL_RIGHT
#   Picture 6     (39%w, right)     → PANEL_RIGHT
#   Rectangle 301 (6%w, left)       → PANEL_LEFT
#   Group 303     (54%w, left, bleeds left) → PANEL_LEFT
#   Picture 2     (45%w, left)      → PANEL_LEFT
#
# v1 bug: Groups with negative positions confused the mirror logic.
# The bleed on Group 219 (top) and Group 303 (left) caused double-
# mirroring in v1.
# ═══════════════════════════════════════════════════════════════════════════════

class TestSlide8:
    """Slide 8 — title layout with groups and bleeds in split panel."""

    @pytest.fixture(autouse=True)
    def setup(self, prs, slide_dims):
        self.slide = prs.slides[7]
        self.result = _classify(prs, slide_dims, 8)

    def test_has_split_panel(self):
        assert self.result.has_split_panel is True

    def test_shape_count(self):
        assert len(self.result.classifications) == 6

    def test_title_7_is_panel_right(self):
        role = _role_for(self.result, self.slide, 'Title 7')
        assert role == ShapeRole.PANEL_RIGHT

    def test_group_219_is_panel_right(self):
        """Group with top bleed on right side → PANEL_RIGHT (not BLEED)."""
        role = _role_for(self.result, self.slide, 'Group 219')
        assert role == ShapeRole.PANEL_RIGHT

    def test_picture_6_is_panel_right(self):
        role = _role_for(self.result, self.slide, 'Picture 6')
        assert role == ShapeRole.PANEL_RIGHT

    def test_rectangle_301_is_panel_left(self):
        role = _role_for(self.result, self.slide, 'Rectangle 301')
        assert role == ShapeRole.PANEL_LEFT

    def test_group_303_is_panel_left(self):
        """Group with left bleed on left side → PANEL_LEFT (not BLEED)."""
        role = _role_for(self.result, self.slide, 'Group 303')
        assert role == ShapeRole.PANEL_LEFT

    def test_picture_2_is_panel_left(self):
        role = _role_for(self.result, self.slide, 'Picture 2')
        assert role == ShapeRole.PANEL_LEFT

    def test_all_panel_actions_are_swap(self):
        """Every shape on a split panel slide gets position_action='swap'."""
        for shape in self.slide.shapes:
            cls = self.result.get(shape)
            if cls.role in (ShapeRole.PANEL_LEFT, ShapeRole.PANEL_RIGHT):
                assert cls.position_action == 'swap', (
                    f'{shape.name}: expected swap, got {cls.position_action}'
                )

    def test_left_panel_shapes(self):
        """Left panel must contain Rectangle 301, Group 303, Picture 2."""
        panel = self.result.context.split_panel
        left_names = set()
        for shape in self.slide.shapes:
            if shape.shape_id in panel.left_shape_ids:
                left_names.add(shape.name)
        assert 'Rectangle 301' in left_names
        assert 'Group 303' in left_names
        assert 'Picture 2' in left_names

    def test_right_panel_shapes(self):
        """Right panel must contain Title 7, Group 219, Picture 6."""
        panel = self.result.context.split_panel
        right_names = set()
        for shape in self.slide.shapes:
            if shape.shape_id in panel.right_shape_ids:
                right_names.add(shape.name)
        assert 'Title 7' in right_names
        assert 'Group 219' in right_names
        assert 'Picture 6' in right_names


# ═══════════════════════════════════════════════════════════════════════════════
# Integration tests — Slide 12 (layout=title — single placeholder)
#
# Shapes:
#   Title 7  (59%w, 54%h) → PLACEHOLDER
#
# v1 bug: Single-shape title slides had their placeholder mirrored
# instead of inheriting from the layout, causing text misalignment.
# ═══════════════════════════════════════════════════════════════════════════════

class TestSlide12:
    """Slide 12 — title layout with a single text placeholder."""

    @pytest.fixture(autouse=True)
    def setup(self, prs, slide_dims):
        self.slide = prs.slides[11]
        self.result = _classify(prs, slide_dims, 12)

    def test_layout_type(self):
        assert self.result.layout_type == 'title'

    def test_no_split_panel(self):
        assert self.result.has_split_panel is False

    def test_shape_count(self):
        assert len(self.result.classifications) == 1

    def test_title_7_is_placeholder(self):
        """Sole title placeholder → PLACEHOLDER (not CONTENT_TEXT or UNKNOWN)."""
        role = _role_for(self.result, self.slide, 'Title 7')
        assert role == ShapeRole.PLACEHOLDER

    def test_title_7_action_is_inherit(self):
        """Placeholder should inherit position from layout master."""
        shape = _get_shape_by_name(self.slide, 'Title 7')
        cls = self.result.get(shape)
        assert cls.position_action == 'inherit'

    def test_title_7_placeholder_type(self):
        shape = _get_shape_by_name(self.slide, 'Title 7')
        cls = self.result.get(shape)
        assert cls.placeholder_type == 'ctrTitle'

    def test_title_7_text_action(self):
        shape = _get_shape_by_name(self.slide, 'Title 7')
        cls = self.result.get(shape)
        assert cls.text_action == 'translate_rtl'


# ═══════════════════════════════════════════════════════════════════════════════
# Integration tests — Slide 13 (layout=title — background + bleed image)
#
# Shapes:
#   Picture 17 (100%w, 100%h)           → BACKGROUND
#   Title 7    (placeholder)             → PLACEHOLDER
#   Picture 3  (56%w, left_bleed=1.1M)   → BLEED
#
# v1 bug: Picture 3 bleeds ~1.2" off the left edge. v1 mirrored it
# using the standard formula but didn't preserve the bleed offset,
# causing it to be clamped to x=0 on the right side instead of
# extending proportionally past the right edge.
# ═══════════════════════════════════════════════════════════════════════════════

class TestSlide13:
    """Slide 13 — title layout with background + bleeding image."""

    @pytest.fixture(autouse=True)
    def setup(self, prs, slide_dims):
        self.slide = prs.slides[12]
        self.result = _classify(prs, slide_dims, 13)

    def test_layout_type(self):
        assert self.result.layout_type == 'title'

    def test_no_split_panel(self):
        assert self.result.has_split_panel is False

    def test_shape_count(self):
        assert len(self.result.classifications) == 3

    def test_picture_17_is_background(self):
        """Full-slide background image (100%w × 100%h) → BACKGROUND."""
        role = _role_for(self.result, self.slide, 'Picture 17')
        assert role == ShapeRole.BACKGROUND

    def test_picture_17_action_is_keep(self):
        shape = _get_shape_by_name(self.slide, 'Picture 17')
        cls = self.result.get(shape)
        assert cls.position_action == 'keep'

    def test_title_7_is_placeholder(self):
        role = _role_for(self.result, self.slide, 'Title 7')
        assert role == ShapeRole.PLACEHOLDER

    def test_picture_3_is_bleed(self):
        """Image extending 1.1M EMU past left edge → BLEED."""
        role = _role_for(self.result, self.slide, 'Picture 3')
        assert role == ShapeRole.BLEED

    def test_picture_3_action_is_mirror(self):
        shape = _get_shape_by_name(self.slide, 'Picture 3')
        cls = self.result.get(shape)
        assert cls.position_action == 'mirror'

    def test_picture_3_bleed_metadata(self):
        """BLEED classification should record left bleed amount."""
        shape = _get_shape_by_name(self.slide, 'Picture 3')
        cls = self.result.get(shape)
        # Picture 3 has left=-1110996, so bleed_left should be ~1110996
        assert cls.bleed_left > 1_000_000, (
            f'Expected significant left bleed, got {cls.bleed_left}'
        )


# ═══════════════════════════════════════════════════════════════════════════════
# Cross-slide tests — consistency and invariants
# ═══════════════════════════════════════════════════════════════════════════════

class TestClassifierInvariants:
    """Verify classification invariants across all 5 test slides."""

    SLIDE_NUMBERS = [1, 7, 8, 12, 13]

    @pytest.fixture(autouse=True)
    def setup(self, prs, slide_dims):
        self.prs = prs
        self.slide_dims = slide_dims
        self.results = {
            n: _classify(prs, slide_dims, n) for n in self.SLIDE_NUMBERS
        }

    def test_every_shape_classified(self):
        """Every shape on every slide must have a classification."""
        for n, result in self.results.items():
            slide = self.prs.slides[n - 1]
            for shape in slide.shapes:
                cls = result.get(shape)
                assert cls.role != ShapeRole.UNKNOWN or cls.rule_name == 'fallback', (
                    f'Slide {n}, {shape.name}: unexpected UNKNOWN without fallback'
                )

    def test_no_duplicate_shape_ids(self):
        """No two shapes on a slide should share the same classification key."""
        for n, result in self.results.items():
            keys = list(result.classifications.keys())
            assert len(keys) == len(set(keys)), (
                f'Slide {n}: duplicate keys in classifications'
            )

    def test_classification_count_matches_shape_count(self):
        """Number of classifications must match number of shapes."""
        for n, result in self.results.items():
            slide = self.prs.slides[n - 1]
            shape_count = len(list(slide.shapes))
            assert len(result.classifications) == shape_count, (
                f'Slide {n}: {len(result.classifications)} classifications '
                f'for {shape_count} shapes'
            )

    def test_panel_shapes_have_swap_action(self):
        """All PANEL_LEFT/PANEL_RIGHT shapes must have position_action='swap'."""
        for n, result in self.results.items():
            for sid, cls in result.classifications.items():
                if cls.role in (ShapeRole.PANEL_LEFT, ShapeRole.PANEL_RIGHT):
                    assert cls.position_action == 'swap', (
                        f'Slide {n}: panel shape {sid} has '
                        f'position_action={cls.position_action!r}'
                    )

    def test_background_shapes_have_keep_action(self):
        """All BACKGROUND shapes must have position_action='keep'."""
        for n, result in self.results.items():
            for sid, cls in result.classifications.items():
                if cls.role == ShapeRole.BACKGROUND:
                    assert cls.position_action == 'keep', (
                        f'Slide {n}: bg shape {sid} has '
                        f'position_action={cls.position_action!r}'
                    )

    def test_placeholder_shapes_have_inherit_action(self):
        """All PLACEHOLDER shapes must have position_action='inherit'."""
        for n, result in self.results.items():
            for sid, cls in result.classifications.items():
                if cls.role == ShapeRole.PLACEHOLDER:
                    assert cls.position_action == 'inherit', (
                        f'Slide {n}: placeholder shape {sid} has '
                        f'position_action={cls.position_action!r}'
                    )


# ═══════════════════════════════════════════════════════════════════════════════
# API tests — classify_slide convenience function
# ═══════════════════════════════════════════════════════════════════════════════

class TestClassifySlideConvenience:
    """Verify the module-level classify_slide function."""

    def test_returns_result_type(self, prs, slide_dims):
        result = _classify(prs, slide_dims, 1)
        assert isinstance(result, SlideClassificationResult)

    def test_result_has_context(self, prs, slide_dims):
        result = _classify(prs, slide_dims, 1)
        assert isinstance(result.context, SlideContext)

    def test_slide_number_recorded(self, prs, slide_dims):
        result = _classify(prs, slide_dims, 7)
        assert result.slide_number == 7

    def test_get_unknown_shape_returns_default(self, prs, slide_dims):
        """Calling get() with a shape not on the classified slide → UNKNOWN."""
        result = _classify(prs, slide_dims, 1)
        # Use a shape from slide 12 — not in slide 1's classification
        other_shape = list(prs.slides[11].shapes)[0]
        cls = result.get(other_shape)
        assert cls.role == ShapeRole.UNKNOWN

    def test_shapes_with_role(self, prs, slide_dims):
        result = _classify(prs, slide_dims, 7)
        left_ids = result.shapes_with_role(ShapeRole.PANEL_LEFT)
        right_ids = result.shapes_with_role(ShapeRole.PANEL_RIGHT)
        assert len(left_ids) >= 1
        assert len(right_ids) >= 1

    def test_get_by_id(self, prs, slide_dims):
        """get_by_id should work with XML shape_id values."""
        result = _classify(prs, slide_dims, 12)
        slide = prs.slides[11]
        shape = list(slide.shapes)[0]  # Title 7
        cls_by_get = result.get(shape)
        cls_by_id = result.get_by_id(shape.shape_id)
        assert cls_by_get.role == cls_by_id.role

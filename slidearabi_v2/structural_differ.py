"""
structural_differ.py — Shape-by-shape comparison between two PPTX files.

Compares the output of v1 and v2 transformers by matching shapes on each slide
via their shape_id (cNvPr id attribute) and computing positional, textual, and
directional deltas.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Tuple

from lxml import etree
from pptx import Presentation

logger = logging.getLogger(__name__)

EMU_PER_INCH = 914_400
POSITION_TOLERANCE_EMU = 50_000  # ~0.055" — sub-pixel, treated as NEUTRAL


# ─────────────────────────────────────────────────────────────────────────────
# Data classes
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class ShapeDiff:
    """Difference between v1 and v2 output for a single shape."""
    shape_name: str
    shape_id: int
    slide_number: int

    # Position
    v1_left: int = 0
    v2_left: int = 0
    left_delta: int = 0  # v2 - v1
    v1_width: int = 0
    v2_width: int = 0
    width_delta: int = 0
    v1_top: int = 0
    v2_top: int = 0
    top_delta: int = 0
    v1_height: int = 0
    v2_height: int = 0
    height_delta: int = 0

    # Text
    v1_text: str = ''
    v2_text: str = ''
    text_changed: bool = False

    # Direction
    v1_flipH: bool = False
    v2_flipH: bool = False
    flipH_changed: bool = False

    # Metadata
    v2_role: str = ''  # ShapeRole name from v2 classifier (set externally)
    is_placeholder: bool = False
    placeholder_type: str = ''

    # Existence flags — True means shape only exists in one version
    v1_only: bool = False
    v2_only: bool = False

    @property
    def position_changed(self) -> bool:
        return abs(self.left_delta) > POSITION_TOLERANCE_EMU

    @property
    def v1_right(self) -> int:
        return self.v1_left + self.v1_width

    @property
    def v2_right(self) -> int:
        return self.v2_left + self.v2_width


@dataclass
class SlideDiff:
    """All shape diffs for one slide."""
    slide_number: int
    layout_type: str
    shape_diffs: List[ShapeDiff] = field(default_factory=list)
    slide_width: int = 0
    slide_height: int = 0

    @property
    def shapes_moved(self) -> int:
        return sum(1 for d in self.shape_diffs if d.position_changed)

    @property
    def shapes_text_changed(self) -> int:
        return sum(1 for d in self.shape_diffs if d.text_changed)

    @property
    def total_position_delta(self) -> int:
        return sum(abs(d.left_delta) for d in self.shape_diffs)

    @property
    def shapes_v1_only(self) -> int:
        return sum(1 for d in self.shape_diffs if d.v1_only)

    @property
    def shapes_v2_only(self) -> int:
        return sum(1 for d in self.shape_diffs if d.v2_only)


# ─────────────────────────────────────────────────────────────────────────────
# Shape snapshot
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class _ShapeSnapshot:
    """Immutable snapshot of a shape's state."""
    shape_id: int
    name: str
    left: int
    top: int
    width: int
    height: int
    text: str
    flipH: bool
    is_placeholder: bool
    placeholder_type: str


def _snapshot_shape(shape) -> _ShapeSnapshot:
    """Take a snapshot of shape geometry, text, and flipH."""
    x = int(shape.left) if shape.left is not None else 0
    y = int(shape.top) if shape.top is not None else 0
    w = int(shape.width) if shape.width is not None else 0
    h = int(shape.height) if shape.height is not None else 0

    # Text
    text = ''
    try:
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            text = shape.text_frame.text or ''
    except Exception:
        pass

    # flipH — check the XML directly
    flipH = False
    try:
        sp_el = shape._element
        xfrm = sp_el.find('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
        if xfrm is not None:
            flipH = xfrm.get('flipH') == '1'
    except Exception:
        pass

    # Placeholder
    is_ph = False
    ph_type = ''
    try:
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            is_ph = True
            ph_type = str(shape.placeholder_format.type).split('.')[-1].lower()
    except Exception:
        pass

    return _ShapeSnapshot(
        shape_id=shape.shape_id,
        name=shape.name,
        left=x, top=y, width=w, height=h,
        text=text,
        flipH=flipH,
        is_placeholder=is_ph,
        placeholder_type=ph_type,
    )


def _get_layout_type(slide) -> str:
    """Extract layout type from slide layout name."""
    try:
        name = slide.slide_layout.name.lower()
        if 'title' in name and 'content' not in name:
            return 'title'
        if 'section' in name or 'sechead' in name:
            return 'secHead'
        if 'blank' in name:
            return 'blank'
        if 'pic' in name:
            return 'picTx'
        if 'two' in name:
            return 'twoCol'
        return name
    except Exception:
        return 'cust'


# ─────────────────────────────────────────────────────────────────────────────
# StructuralDiffer
# ─────────────────────────────────────────────────────────────────────────────

class StructuralDiffer:
    """Compares two PPTX files shape-by-shape."""

    def __init__(self, v1_pptx_path: str, v2_pptx_path: str):
        self.v1_prs = Presentation(v1_pptx_path)
        self.v2_prs = Presentation(v2_pptx_path)
        self.slide_width = int(self.v1_prs.slide_width)
        self.slide_height = int(self.v1_prs.slide_height)

        v1_count = len(self.v1_prs.slides)
        v2_count = len(self.v2_prs.slides)
        if v1_count != v2_count:
            logger.warning('Slide count mismatch: v1=%d, v2=%d', v1_count, v2_count)

    def diff_all_slides(self) -> List[SlideDiff]:
        """Compare every slide, matching shapes by shape_id."""
        results = []
        slide_count = min(len(self.v1_prs.slides), len(self.v2_prs.slides))
        for i in range(slide_count):
            results.append(self.diff_slide(i + 1))
        return results

    def diff_slide(self, slide_number: int) -> SlideDiff:
        """Compare a specific slide (1-indexed)."""
        v1_slide = self.v1_prs.slides[slide_number - 1]
        v2_slide = self.v2_prs.slides[slide_number - 1]
        layout_type = _get_layout_type(v1_slide)

        # Snapshot all shapes
        v1_map: Dict[int, _ShapeSnapshot] = {}
        for s in v1_slide.shapes:
            snap = _snapshot_shape(s)
            v1_map[snap.shape_id] = snap

        v2_map: Dict[int, _ShapeSnapshot] = {}
        for s in v2_slide.shapes:
            snap = _snapshot_shape(s)
            v2_map[snap.shape_id] = snap

        all_ids = sorted(set(v1_map.keys()) | set(v2_map.keys()))
        diffs = []

        for sid in all_ids:
            v1 = v1_map.get(sid)
            v2 = v2_map.get(sid)

            if v1 and not v2:
                # Shape only in v1
                diffs.append(ShapeDiff(
                    shape_name=v1.name, shape_id=sid, slide_number=slide_number,
                    v1_left=v1.left, v1_width=v1.width, v1_top=v1.top,
                    v1_height=v1.height, v1_text=v1.text, v1_flipH=v1.flipH,
                    is_placeholder=v1.is_placeholder,
                    placeholder_type=v1.placeholder_type,
                    v1_only=True,
                ))
                continue

            if v2 and not v1:
                # Shape only in v2
                diffs.append(ShapeDiff(
                    shape_name=v2.name, shape_id=sid, slide_number=slide_number,
                    v2_left=v2.left, v2_width=v2.width, v2_top=v2.top,
                    v2_height=v2.height, v2_text=v2.text, v2_flipH=v2.flipH,
                    is_placeholder=v2.is_placeholder,
                    placeholder_type=v2.placeholder_type,
                    v2_only=True,
                ))
                continue

            # Both present — compute deltas
            diffs.append(ShapeDiff(
                shape_name=v1.name,
                shape_id=sid,
                slide_number=slide_number,
                v1_left=v1.left, v2_left=v2.left,
                left_delta=v2.left - v1.left,
                v1_width=v1.width, v2_width=v2.width,
                width_delta=v2.width - v1.width,
                v1_top=v1.top, v2_top=v2.top,
                top_delta=v2.top - v1.top,
                v1_height=v1.height, v2_height=v2.height,
                height_delta=v2.height - v1.height,
                v1_text=v1.text, v2_text=v2.text,
                text_changed=v1.text != v2.text,
                v1_flipH=v1.flipH, v2_flipH=v2.flipH,
                flipH_changed=v1.flipH != v2.flipH,
                is_placeholder=v1.is_placeholder,
                placeholder_type=v1.placeholder_type,
            ))

        return SlideDiff(
            slide_number=slide_number,
            layout_type=layout_type,
            shape_diffs=diffs,
            slide_width=self.slide_width,
            slide_height=self.slide_height,
        )


# ─────────────────────────────────────────────────────────────────────────────
# Convenience
# ─────────────────────────────────────────────────────────────────────────────

def diff_pptx(v1_path: str, v2_path: str) -> List[SlideDiff]:
    """One-liner: diff two PPTX files, return slide diffs."""
    return StructuralDiffer(v1_path, v2_path).diff_all_slides()


if __name__ == '__main__':
    import sys
    if len(sys.argv) != 3:
        print(f'Usage: {sys.argv[0]} <v1.pptx> <v2.pptx>')
        sys.exit(1)
    diffs = diff_pptx(sys.argv[1], sys.argv[2])
    for sd in diffs:
        moved = sd.shapes_moved
        if moved == 0:
            continue
        print(f'Slide {sd.slide_number} ({sd.layout_type}): {moved} shapes moved, '
              f'total delta={sd.total_position_delta} EMU')
        for d in sd.shape_diffs:
            if d.position_changed or d.v1_only or d.v2_only:
                flag = ''
                if d.v1_only:
                    flag = ' [V1 ONLY]'
                elif d.v2_only:
                    flag = ' [V2 ONLY]'
                print(f'  [{d.shape_id:3d}] {d.shape_name:20s} '
                      f'left: {d.v1_left} → {d.v2_left} (Δ={d.left_delta}){flag}')

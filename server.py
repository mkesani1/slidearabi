# === STARTUP DIAGNOSTIC — remove after Railway health check is resolved ===
import sys as _sys, os as _os
print(f"[BOOT] Python {_sys.version}, PID={_os.getpid()}, PORT={_os.environ.get('PORT', 'NOT SET')}", flush=True)
# === END DIAGNOSTIC ===

import base64
import hashlib
import json
import logging
import math
import os
import re
import secrets
import shutil
import subprocess
import tempfile
import threading
import time
import uuid
from dataclasses import dataclass, field
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional

import stripe
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from PIL import Image, ImageDraw, ImageFont
from pydantic import BaseModel
from pptx import Presentation as PptxPresentation

from slidearabi.llm_translator import DualLLMTranslator, TranslatorConfig
from slidearabi.pipeline import PipelineConfig, SlideArabiPipeline
import slidearabi.pipeline as pipeline_module


logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s [%(name)s] %(message)s",
)
logger = logging.getLogger("slidearabi.server")


MAX_FILE_SIZE = 150 * 1024 * 1024  # 150 MB
BASE_DIR = Path("/tmp/slideshift_jobs")
BASE_DIR.mkdir(parents=True, exist_ok=True)
JOB_TTL_HOURS = 24

# Validate critical environment variables at startup
for _env_var in ["OPENAI_API_KEY", "ANTHROPIC_API_KEY"]:
    _val = os.getenv(_env_var, "")
    if not _val:
        logger.critical("MISSING: %s environment variable is not set!", _env_var)
    elif len(_val) < 10:
        logger.critical("INVALID: %s appears too short (%d chars)", _env_var, len(_val))
del _env_var, _val

PIPELINE_SEMAPHORE = threading.Semaphore(1)
JOBS_LOCK = threading.Lock()

# ── Pending job queue: jobs wait here when the semaphore is busy ─────────────
import queue as _queue_mod
PENDING_QUEUE: _queue_mod.Queue = _queue_mod.Queue(maxsize=10)

def _queue_dispatcher() -> None:
    """Background thread that drains the pending queue sequentially.
    
    Uses a BLOCKING semaphore acquire so queued jobs wait their turn
    instead of failing with 'Pipeline busy'.
    """
    while True:
        job_id = PENDING_QUEUE.get()  # blocks until a job appears
        try:
            # Block until the pipeline is free — this is the whole point
            # of the queue. Jobs wait here, not fail.
            PIPELINE_SEMAPHORE.acquire(blocking=True)
            try:
                _run_pipeline_worker_inner(job_id)
            finally:
                PIPELINE_SEMAPHORE.release()
        except Exception as exc:
            logger.exception("Queue dispatcher failed for %s: %s", job_id, exc)
        finally:
            PENDING_QUEUE.task_done()

_dispatcher_thread = threading.Thread(target=_queue_dispatcher, daemon=True)
_dispatcher_thread.start()

PHASE_EXTRACTING = "extracting"
PHASE_TRANSLATING = "translating"
PHASE_RTL = "rtl_transforms"
PHASE_QC = "quality_check"
PHASE_PREVIEW = "generating_preview"
PHASE_DONE = "done"


@dataclass
class JobState:
    job_id: str
    status: str = "queued"
    progress_pct: int = 0
    current_phase: str = PHASE_EXTRACTING
    input_path: str = ""
    output_path: str = ""
    total_slides: int = 0
    paid: bool = False
    created_at: datetime = field(default_factory=lambda: datetime.now(timezone.utc))
    error: Optional[str] = None
    preview_slides: List[dict] = field(default_factory=list)
    preview_origin: str = ""  # "output" or "fallback_input" — tracks which file the preview was rendered from


JOBS: Dict[str, JobState] = {}


class CheckoutRequest(BaseModel):
    job_id: str
    slide_count: int


class VerifyPaymentRequest(BaseModel):
    session_id: str
    job_id: str


class PromoCodeRequest(BaseModel):
    job_id: str
    code: str


# ── Valid promo codes that bypass Stripe payment ──
VALID_PROMO_CODES: Dict[str, str] = {
    "SLIDETEST2026": "Internal testing",
    "FOUNDER": "Founder access",
    "DEMO": "Demo / sales",
}


app = FastAPI(title="SlideArabi API", version="1.1.4")


# === STARTUP DIAGNOSTIC — remove after Railway deploy is stable ===
@app.on_event("startup")
async def _startup_diagnostic():
    try:
        import resource
        peak_mb = resource.getrusage(resource.RUSAGE_SELF).ru_maxrss / 1024
        logger.info("[BOOT] FastAPI ready — peak RSS: %.0f MB, PORT=%s", peak_mb, os.environ.get("PORT", "NOT SET"))
    except Exception:
        logger.info("[BOOT] FastAPI ready")
# === END DIAGNOSTIC ===


app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "https://slidearabi.com",
        "https://www.slidearabi.com",
        "http://localhost:3000",
    ],
    allow_credentials=False,
    allow_methods=["GET", "POST", "OPTIONS", "DELETE", "PUT"],
    allow_headers=["*"],
    expose_headers=["Content-Length", "Content-Type"],
    max_age=3600,
)

# ─── Agent Stack: REST API v1 + MCP + Stripe Credits ────────────────────────
try:
    from slidearabi.api_gateway import api_router
    app.include_router(api_router, prefix="/v1")
    logger.info("[AGENT] REST API gateway mounted at /v1")
except Exception as exc:
    logger.warning("[AGENT] REST API gateway not mounted: %s", exc)

try:
    from slidearabi.stripe_credits import stripe_router
    app.include_router(stripe_router)
    logger.info("[AGENT] Stripe credits webhook mounted")
except Exception as exc:
    logger.warning("[AGENT] Stripe credits not mounted: %s", exc)

_mcp_error = None
try:
    from slidearabi.mcp_server import create_mcp_app
    app.mount("/mcp", create_mcp_app())
    logger.info("[AGENT] MCP server mounted at /mcp")
except Exception as exc:
    _mcp_error = str(exc)
    logger.warning("[AGENT] MCP server not mounted: %s", exc)
# ─── End Agent Stack ────────────────────────────────────────────────────────


def _cleanup_expired_jobs() -> None:
    cutoff = datetime.now(timezone.utc) - timedelta(hours=JOB_TTL_HOURS)
    remove_ids = []
    with JOBS_LOCK:
        for job_id, job in JOBS.items():
            if job.created_at < cutoff:
                remove_ids.append(job_id)

        for job_id in remove_ids:
            job = JOBS.pop(job_id, None)
            if job:
                job_dir = BASE_DIR / job_id
                try:
                    if job_dir.exists():
                        shutil.rmtree(job_dir, ignore_errors=True)
                except Exception as exc:
                    logger.warning("Failed to remove expired job dir for %s: %s", job_id, exc)

    if remove_ids:
        logger.info("Cleaned up %d expired jobs", len(remove_ids))


def _job_cleanup_loop() -> None:
    while True:
        try:
            _cleanup_expired_jobs()
        except Exception as exc:
            logger.exception("Cleanup loop error: %s", exc)
        time.sleep(3600)


threading.Thread(target=_job_cleanup_loop, daemon=True).start()


def _count_slides(pptx_path: Path) -> int:
    try:
        import zipfile
        import xml.etree.ElementTree as ET

        with zipfile.ZipFile(pptx_path, "r") as zf:
            with zf.open("ppt/presentation.xml") as f:
                tree = ET.parse(f)
                root = tree.getroot()
                ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
                sld_id_lst = root.find("p:sldIdLst", ns)
                if sld_id_lst is None:
                    return 0
                return len(sld_id_lst.findall("p:sldId", ns))
    except Exception as exc:
        logger.exception("Failed counting slides for %s: %s", pptx_path, exc)
        raise HTTPException(status_code=400, detail="Unable to read PPTX slide structure")


# ---------------------------------------------------------------------------
# Preview rendering constants
# ---------------------------------------------------------------------------
THUMB_WIDTH_PX = 800
JPEG_QUALITY = 85
LIBREOFFICE_TIMEOUT = 300
PDFTOPPM_TIMEOUT = 120

ARABIC_RE = re.compile(r"[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF\uFB50-\uFDFF\uFE70-\uFEFF]")


def _apply_watermark(
    image_path: Path,
    watermark_text: str = "SlideArabi Preview",
) -> None:
    """Apply a semi-transparent diagonal watermark to a JPEG preview thumbnail.

    Uses Pillow to draw repeated diagonal text so the watermark is visible
    regardless of slide content. Overwrites the original file in-place.

    Watermark density reduced in v1.1.3 — sparse enough for quality evaluation,
    visible enough for IP protection.
    """
    watermark_mode = os.getenv("WATERMARK_MODE", "light")

    with Image.open(image_path) as base:
        base = base.convert("RGBA")
        width, height = base.size

        txt_layer = Image.new("RGBA", base.size, (255, 255, 255, 0))
        draw = ImageDraw.Draw(txt_layer)

        # Dynamic font size: ~5% of image width (was 8%)
        font_size = max(20, int(width * 0.05))

        # Try several font paths (Docker / macOS / Linux fallbacks)
        font = None
        for font_path in [
            "/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
        ]:
            try:
                font = ImageFont.truetype(font_path, font_size)
                break
            except (IOError, OSError):
                continue
        if font is None:
            font = ImageFont.load_default()

        # Measure text
        bbox = draw.textbbox((0, 0), watermark_text, font=font)
        text_w = bbox[2] - bbox[0]
        text_h = bbox[3] - bbox[1]

        shadow_color = (0, 0, 0, 28)    # Was 60
        text_color = (255, 255, 255, 55) # Was 100

        if watermark_mode == "light":
            # Light mode: 3 fixed diagonal positions instead of dense tiling
            positions = [
                (int(width * 0.10), int(height * 0.20)),
                (int(width * 0.38), int(height * 0.48)),
                (int(width * 0.66), int(height * 0.76)),
            ]
            for x, y in positions:
                draw.text((x + 1, y + 1), watermark_text, font=font, fill=shadow_color)
                draw.text((x, y), watermark_text, font=font, fill=text_color)
        else:
            # Heavy mode — tiling with improved spacing
            spacing_x = int(text_w * 1.8)
            spacing_y = int(height * 0.45)

            for y in range(0, height * 2, spacing_y):
                for x in range(0, width * 2, spacing_x):
                    draw.text((x + 1, y + 1), watermark_text, font=font, fill=shadow_color)
                    draw.text((x, y), watermark_text, font=font, fill=text_color)

        # Rotate the text layer for diagonal effect
        watermark_layer = txt_layer.rotate(
            28, resample=Image.BICUBIC, expand=False,  # Was 30
            center=(width // 2, height // 2),
        )

        watermarked = Image.alpha_composite(base, watermark_layer)
        final = watermarked.convert("RGB")
        final.save(str(image_path), "JPEG", quality=JPEG_QUALITY)


def _natural_sort_key(path: Path) -> int:
    """Extract trailing integer from filename for natural sort order."""
    match = re.search(r"(\d+)\.\w+$", path.name)
    return int(match.group(1)) if match else 0


def _render_preview_slides(
    input_path: Path,
    preview_dir: Path,
    max_slides: int = 100,
    thumb_width: int = THUMB_WIDTH_PX,
) -> List[dict]:
    """Render individual slide thumbnails from a PPTX file.

    Pipeline:
        1. LibreOffice converts PPTX → PDF (one page per slide)
        2. pdftoppm converts each PDF page → JPEG thumbnail

    Returns a list of dicts matching the frontend PreviewSlide interface:
        { index: int, image_url: str, title: str }
    """
    preview_dir.mkdir(parents=True, exist_ok=True)

    # ── Step 1: PPTX → PDF via LibreOffice ────────────────────────────
    with tempfile.TemporaryDirectory(prefix="pptx_pdf_") as tmp_dir:
        pdf_out_dir = Path(tmp_dir)

        cmd_pdf = [
            "libreoffice",
            "--headless",
            "--norestore",
            "--convert-to", "pdf",
            "--outdir", str(pdf_out_dir),
            str(input_path),
        ]

        logger.info("Converting PPTX to PDF: %s", " ".join(cmd_pdf))
        try:
            subprocess.run(
                cmd_pdf,
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=LIBREOFFICE_TIMEOUT,
            )
        except subprocess.TimeoutExpired:
            logger.error("LibreOffice timed out after %ds", LIBREOFFICE_TIMEOUT)
            raise RuntimeError(
                f"LibreOffice PDF conversion timed out after {LIBREOFFICE_TIMEOUT}s"
            )
        except subprocess.CalledProcessError as exc:
            logger.error(
                "LibreOffice failed (rc=%d): %s",
                exc.returncode,
                exc.stderr.decode(errors="replace"),
            )
            raise RuntimeError(
                f"LibreOffice PDF conversion failed: {exc.stderr.decode(errors='replace')}"
            )

        # LibreOffice names the output <stem>.pdf in --outdir
        pdf_path = pdf_out_dir / f"{input_path.stem}.pdf"
        if not pdf_path.exists():
            pdfs = list(pdf_out_dir.glob("*.pdf"))
            if not pdfs:
                raise FileNotFoundError(
                    f"LibreOffice did not produce a PDF in {pdf_out_dir}"
                )
            pdf_path = pdfs[0]
            logger.warning("Expected PDF name mismatch; using %s", pdf_path.name)

        logger.info("PDF created: %s (%d bytes)", pdf_path, pdf_path.stat().st_size)

        # ── Step 2: PDF → per-page JPEGs via pdftoppm ────────────────
        output_prefix = str(preview_dir / "slide")

        cmd_jpg = [
            "pdftoppm",
            "-jpeg",
            "-jpegopt", f"quality={JPEG_QUALITY}",
            "-scale-to-x", str(thumb_width),
            "-scale-to-y", "-1",
        ]

        if max_slides and max_slides > 0:
            cmd_jpg += ["-l", str(max_slides)]

        cmd_jpg += [str(pdf_path), output_prefix]

        logger.info("Rendering slide thumbnails: %s", " ".join(cmd_jpg))
        try:
            subprocess.run(
                cmd_jpg,
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=PDFTOPPM_TIMEOUT,
            )
        except subprocess.TimeoutExpired:
            logger.error("pdftoppm timed out after %ds", PDFTOPPM_TIMEOUT)
            raise RuntimeError(
                f"pdftoppm timed out after {PDFTOPPM_TIMEOUT}s"
            )
        except subprocess.CalledProcessError as exc:
            logger.error(
                "pdftoppm failed (rc=%d): %s",
                exc.returncode,
                exc.stderr.decode(errors="replace"),
            )
            raise RuntimeError(
                f"pdftoppm failed: {exc.stderr.decode(errors='replace')}"
            )

    # ── Step 3: Apply watermark + collect as base64 data URIs ─────────
    generated = sorted(preview_dir.glob("slide-*.jpg"), key=_natural_sort_key)

    if not generated:
        logger.warning("No slide thumbnails were generated in %s", preview_dir)
        return []

    logger.info("Generated %d slide thumbnails — applying watermarks", len(generated))

    # Apply watermark to every preview thumbnail
    for img_path in generated:
        try:
            _apply_watermark(img_path)
        except Exception as wm_exc:
            logger.warning("Watermark failed for %s: %s", img_path, wm_exc)

    previews: List[dict] = []
    for idx, img_path in enumerate(generated[:max_slides]):
        with open(img_path, "rb") as f:
            encoded = base64.b64encode(f.read()).decode("ascii")
        previews.append(
            {
                "index": idx,
                "image_url": f"data:image/jpeg;base64,{encoded}",
                "title": f"Slide {idx + 1}",
            }
        )

    return previews


def _build_translate_fn():
    """Build a translate_fn matching the pipeline's signature: List[str] → Dict[str, str]."""
    openai_key = os.getenv("OPENAI_API_KEY", "")
    anthropic_key = os.getenv("ANTHROPIC_API_KEY", "")

    if not openai_key:
        raise RuntimeError(
            "OPENAI_API_KEY is not configured. "
            "Set it in Railway environment variables."
        )
    if not anthropic_key:
        logger.warning("ANTHROPIC_API_KEY not set — Claude QA pass will be skipped")

    config = TranslatorConfig(
        openai_api_key=openai_key,
        anthropic_api_key=anthropic_key,
    )
    translator = DualLLMTranslator(config)

    def translate_fn(texts: list) -> dict:
        return translator.translate(texts)

    return translate_fn


def _set_job_state(job_id: str, **kwargs) -> None:
    with JOBS_LOCK:
        job = JOBS.get(job_id)
        if not job:
            return
        for k, v in kwargs.items():
            setattr(job, k, v)


def _write_preview_debug(
    job_dir: Path, job_id: str, source_pptx: str, origin: str,
) -> None:
    """Write a tiny sidecar JSON so /debug/{job_id} can prove which file the preview used."""
    try:
        data = {
            "job_id": job_id,
            "source_pptx": source_pptx,
            "origin": origin,
            "timestamp_utc": datetime.now(timezone.utc).isoformat(),
        }
        (job_dir / "preview_debug.json").write_text(
            json.dumps(data, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
    except Exception as exc:
        logger.warning("Failed to write preview_debug.json: %s", exc)


def _run_pipeline_worker(job_id: str) -> None:
    """Direct-start path: acquire semaphore non-blocking, run if free.
    If busy, the caller should enqueue instead — never fail with 'busy'.
    """
    acquired = PIPELINE_SEMAPHORE.acquire(blocking=False)
    if not acquired:
        # Don't fail — enqueue for the dispatcher to process
        try:
            PENDING_QUEUE.put_nowait(job_id)
            logger.info("Job %s re-queued from worker (semaphore busy, queue depth: %d)",
                       job_id, PENDING_QUEUE.qsize())
        except _queue_mod.Full:
            _set_job_state(
                job_id,
                status="failed",
                error="Server overloaded. Please retry in a minute.",
                progress_pct=0,
            )
            logger.error("Queue full, job %s rejected", job_id)
        return

    try:
        _run_pipeline_worker_inner(job_id)
    finally:
        PIPELINE_SEMAPHORE.release()
        # After releasing, kick the next queued job if any
        if not PENDING_QUEUE.empty():
            logger.info("Semaphore released, queue has %d pending", PENDING_QUEUE.qsize())


def _run_pipeline_worker_inner(job_id: str) -> None:
    """Core pipeline execution — assumes semaphore is already held."""
    try:
        with JOBS_LOCK:
            job = JOBS.get(job_id)
            if not job:
                return

        _set_job_state(job_id, status="processing", current_phase=PHASE_EXTRACTING, progress_pct=5)

        input_path = Path(job.input_path)
        output_path = Path(job.output_path)

        _set_job_state(job_id, progress_pct=10)

        # ── Translation pipeline ──────────────────────────────────────
        translate_fn = _build_translate_fn()
        cfg = PipelineConfig(
            input_path=str(input_path),
            output_path=str(output_path),
            translate_fn=translate_fn,
        )

        pipeline = SlideArabiPipeline(config=cfg)

        _set_job_state(job_id, current_phase=PHASE_TRANSLATING, progress_pct=25)
        result = pipeline.run()

        if not result.success:
            error_msg = result.error or "Pipeline failed (unknown internal error)"
            logger.error(
                "Pipeline returned failure for job %s: %s | phase_reports=%s",
                job_id,
                error_msg,
                list(result.phase_reports.keys()) if result.phase_reports else [],
            )
            raise RuntimeError(error_msg)

        # Persist phase reports for diagnostic endpoint
        try:
            phase_report_path = output_path.parent / "pipeline_phase_reports.json"
            phase_report_path.write_text(
                json.dumps(result.phase_reports or {}, ensure_ascii=False, indent=2, default=str),
                encoding="utf-8",
            )
        except Exception as pr_exc:
            logger.warning("Failed to persist phase reports: %s", pr_exc)

        _set_job_state(job_id, current_phase=PHASE_RTL, progress_pct=70)
        _set_job_state(job_id, current_phase=PHASE_QC, progress_pct=90)

        if not output_path.exists():
            raise RuntimeError("Pipeline reported success but output file missing")

        # ── Preview from TRANSLATED Arabic output ─────────────────────
        _set_job_state(job_id, current_phase=PHASE_PREVIEW, progress_pct=92)
        preview_dir = output_path.parent / "preview"
        # Clear any stale preview files from previous runs
        if preview_dir.exists():
            for old_file in preview_dir.glob("slide-*.jpg"):
                old_file.unlink(missing_ok=True)

        try:
            logger.info(
                "Generating Arabic preview for job %s from %s (%d bytes)",
                job_id,
                output_path,
                output_path.stat().st_size,
            )
            preview_slides = _render_preview_slides(
                output_path, preview_dir, max_slides=50
            )
            if not preview_slides:
                raise RuntimeError("_render_preview_slides returned empty list")
            _set_job_state(job_id, preview_slides=preview_slides, preview_origin="output")
            # Write preview sidecar metadata for diagnostics
            _write_preview_debug(
                output_path.parent, job_id,
                source_pptx=str(output_path), origin="output",
            )
            _set_job_state(job_id, progress_pct=96)  # Signal: preview done, finalizing
            logger.info(
                "Arabic preview generated for job %s (%d slides)",
                job_id,
                len(preview_slides),
            )
        except Exception as exc:
            logger.error(
                "ARABIC PREVIEW FAILED for job %s: %s — falling back to English input",
                job_id,
                exc,
                exc_info=True,
            )
            # Fallback: render from original English so user sees *something*
            try:
                fallback_dir = input_path.parent / "preview_fallback"
                preview_slides = _render_preview_slides(
                    input_path, fallback_dir, max_slides=50
                )
                _set_job_state(job_id, preview_slides=preview_slides, preview_origin="fallback_input")
                _write_preview_debug(
                    input_path.parent, job_id,
                    source_pptx=str(input_path), origin="fallback_input",
                )
                logger.warning(
                    "ENGLISH FALLBACK preview generated for job %s (%d slides)",
                    job_id,
                    len(preview_slides),
                )
            except Exception as fallback_exc:
                logger.error(
                    "Fallback preview also failed for %s: %s",
                    job_id,
                    fallback_exc,
                    exc_info=True,
                )

        _set_job_state(
            job_id,
            status="completed",
            progress_pct=100,
            current_phase=PHASE_DONE,
            error=None,
        )
        logger.info("Job %s completed", job_id)

    except Exception as exc:
        logger.exception("Pipeline failed for job %s: %s", job_id, exc)
        _set_job_state(
            job_id,
            status="failed",
            error=str(exc),
        )
    # Note: semaphore release is handled by the CALLER
    # (_run_pipeline_worker or _queue_dispatcher), not here.


def _start_pipeline_if_allowed(job_id: str) -> None:
    with JOBS_LOCK:
        job = JOBS.get(job_id)
        if not job:
            raise HTTPException(status_code=404, detail="Job not found")
        if job.status in {"processing", "completed"}:
            return

    t = threading.Thread(target=_run_pipeline_worker, args=(job_id,), daemon=True)
    t.start()


def _is_pptx(upload: UploadFile) -> bool:
    filename = (upload.filename or "").lower()
    return filename.endswith(".pptx")


@app.post("/convert")
async def convert(file: UploadFile = File(...)):
    try:
        if not _is_pptx(file):
            raise HTTPException(status_code=400, detail="Only .pptx files are allowed")

        content = await file.read()
        if len(content) > MAX_FILE_SIZE:
            raise HTTPException(status_code=413, detail="File too large (max 150MB)")

        # Reject only if the queue is completely full (10 pending + 1 active)
        if PENDING_QUEUE.full() and PIPELINE_SEMAPHORE._value <= 0:
            raise HTTPException(status_code=503, detail="Server busy — too many queued jobs. Please retry in a minute.", headers={"Retry-After": "60"})

        job_id = str(uuid.uuid4())
        job_dir = BASE_DIR / job_id
        job_dir.mkdir(parents=True, exist_ok=True)

        input_path = job_dir / "input.pptx"
        output_path = job_dir / "output_ar.pptx"
        with open(input_path, "wb") as f:
            f.write(content)

        slide_count = _count_slides(input_path)

        # Do NOT block on LibreOffice preview here — respond immediately.
        # Preview generation happens in the background pipeline worker.
        job = JobState(
            job_id=job_id,
            status="queued",
            progress_pct=0,
            current_phase=PHASE_EXTRACTING,
            input_path=str(input_path),
            output_path=str(output_path),
            total_slides=slide_count,
            paid=False,
            preview_slides=[],
        )

        with JOBS_LOCK:
            JOBS[job_id] = job

        # Try direct start; if semaphore is busy, enqueue for the dispatcher
        if PIPELINE_SEMAPHORE._value > 0:
            _start_pipeline_if_allowed(job_id)
        else:
            try:
                PENDING_QUEUE.put_nowait(job_id)
                logger.info("Job %s queued (pipeline busy, queue depth: %d)", job_id, PENDING_QUEUE.qsize())
            except _queue_mod.Full:
                _set_job_state(job_id, status="failed", error="Server overloaded. Please retry in a minute.")

        return {"job_id": job_id, "status": "queued", "slide_count": slide_count, "total_slides": slide_count}
    except HTTPException:
        raise
    except Exception as exc:
        logger.exception("/convert failed: %s", exc)
        raise HTTPException(status_code=500, detail="Internal server error")


@app.get("/status/{job_id}")
def status(job_id: str):
    try:
        with JOBS_LOCK:
            job = JOBS.get(job_id)
            if not job:
                raise HTTPException(status_code=404, detail="Job not found")

            return {
                "status": job.status,
                "progress_pct": max(0, min(100, job.progress_pct)),
                "current_phase": job.current_phase,
                "total_slides": job.total_slides,
                "error": job.error if job.status == "failed" else None,
            }
    except HTTPException:
        raise
    except Exception as exc:
        logger.exception("/status failed for %s: %s", job_id, exc)
        raise HTTPException(status_code=500, detail="Internal server error")


@app.get("/preview/{job_id}")
def preview(job_id: str):
    try:
        with JOBS_LOCK:
            job = JOBS.get(job_id)
            if not job:
                raise HTTPException(status_code=404, detail="Job not found")
            if not job.preview_slides:
                raise HTTPException(status_code=404, detail="Preview not ready")
            return {
                "preview_slides": job.preview_slides,
                "total_slides": job.total_slides,
            }
    except HTTPException:
        raise
    except Exception as exc:
        logger.exception("/preview failed for %s: %s", job_id, exc)
        raise HTTPException(status_code=500, detail="Internal server error")


@app.get("/download/{job_id}")
def download(job_id: str):
    try:
        with JOBS_LOCK:
            job = JOBS.get(job_id)
            if not job:
                raise HTTPException(status_code=404, detail="Job not found")
            if job.status != "completed":
                raise HTTPException(status_code=400, detail="Job not completed")
            output_path = Path(job.output_path)

        if not output_path.exists():
            raise HTTPException(status_code=404, detail="Output not found")

        return FileResponse(
            path=str(output_path),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"slideshift_{job_id}.pptx",
        )
    except HTTPException:
        raise
    except Exception as exc:
        logger.exception("/download failed for %s: %s", job_id, exc)
        raise HTTPException(status_code=500, detail="Internal server error")


def _country_from_ip() -> str:
    try:
        cmd = [
            "curl",
            "-s",
            "https://ipapi.co/json/",
        ]
        out = subprocess.check_output(cmd, text=True, timeout=10)
        data = json.loads(out)
        return str(data.get("country_code", "")).upper()
    except Exception as exc:
        logger.warning("Geo lookup failed, defaulting to USD: %s", exc)
        return ""


# Geo-pricing — per-slide in smallest currency unit (cents/halalas/fils)
GEO_PRICING = {
    "SA": ("sar", 500),     # SAR 5/slide
    "AE": ("aed", 500),     # AED 5/slide
    "EG": ("egp", 5000),    # EGP 50/slide
    # GCC countries default to AED
    "BH": ("aed", 500),
    "KW": ("aed", 500),
    "OM": ("aed", 500),
    "QA": ("aed", 500),
}


def _price_for_country(country_code: str, slide_count: int):
    cc = (country_code or "").upper()
    if cc in GEO_PRICING:
        return GEO_PRICING[cc]
    return "usd", 100       # Default: $1/slide


@app.post("/create-checkout-session")
def create_checkout_session(payload: CheckoutRequest):
    try:
        with JOBS_LOCK:
            job = JOBS.get(payload.job_id)
            if not job:
                raise HTTPException(status_code=404, detail="Job not found")

        if payload.slide_count <= 0:
            raise HTTPException(status_code=400, detail="slide_count must be positive")

        stripe_key = os.getenv("STRIPE_SECRET_KEY", "")
        if not stripe_key:
            raise HTTPException(status_code=500, detail="Stripe key not configured")

        stripe.api_key = stripe_key

        country = _country_from_ip()
        currency, unit_amount = _price_for_country(country, payload.slide_count)

        session = stripe.checkout.Session.create(
            mode="payment",
            payment_method_types=["card"],
            line_items=[
                {
                    "quantity": payload.slide_count,
                    "price_data": {
                        "currency": currency,
                        "unit_amount": unit_amount,
                        "product_data": {
                            "name": "SlideArabi PPT Translation",
                            "description": f"{payload.slide_count} slides",
                        },
                    },
                }
            ],
            success_url="https://slidearabi.com/success?session_id={CHECKOUT_SESSION_ID}",
            cancel_url="https://slidearabi.com/cancel",
            metadata={"job_id": payload.job_id},
        )

        return {"checkout_url": session.url, "session_id": session.id}
    except HTTPException:
        raise
    except Exception as exc:
        logger.exception("/create-checkout-session failed: %s", exc)
        raise HTTPException(status_code=500, detail="Internal server error")


@app.post("/verify-payment")
def verify_payment(payload: VerifyPaymentRequest):
    try:
        stripe_key = os.getenv("STRIPE_SECRET_KEY", "")
        if not stripe_key:
            raise HTTPException(status_code=500, detail="Stripe key not configured")

        stripe.api_key = stripe_key

        with JOBS_LOCK:
            job = JOBS.get(payload.job_id)
            if not job:
                raise HTTPException(status_code=404, detail="Job not found")

        session = stripe.checkout.Session.retrieve(payload.session_id)
        paid = session.payment_status == "paid"

        if paid:
            _set_job_state(payload.job_id, paid=True)
            _start_pipeline_if_allowed(payload.job_id)
            return {"verified": True, "status": "processing"}

        return {"verified": False, "status": "unpaid"}
    except HTTPException:
        raise
    except Exception as exc:
        logger.exception("/verify-payment failed: %s", exc)
        raise HTTPException(status_code=500, detail="Internal server error")


@app.post("/apply-promo")
def apply_promo(payload: PromoCodeRequest):
    """Validate a promo code and mark the job as paid (bypasses Stripe)."""
    try:
        code = payload.code.strip().upper()
        if code not in VALID_PROMO_CODES:
            raise HTTPException(status_code=400, detail="Invalid promo code")

        with JOBS_LOCK:
            job = JOBS.get(payload.job_id)
            if not job:
                raise HTTPException(status_code=404, detail="Job not found")
            if job.status != "completed":
                raise HTTPException(status_code=400, detail="Job is not ready for download")

        _set_job_state(payload.job_id, paid=True)
        logger.info("Promo code %s applied to job %s (%s)", code, payload.job_id, VALID_PROMO_CODES[code])

        return {"success": True, "message": f"Promo code applied. You can now download your file."}
    except HTTPException:
        raise
    except Exception as exc:
        logger.exception("/apply-promo failed: %s", exc)
        raise HTTPException(status_code=500, detail="Internal server error")


@app.get("/health")
def health():
    try:
        # Detect mounted agent stack components
        route_paths = [getattr(r, 'path', '') for r in app.routes]
        mount_paths = [getattr(r, 'path', '') for r in app.routes if hasattr(r, 'app')]
        agent_stack = {
            "rest_api": any(p.startswith('/v1') for p in route_paths),
            "mcp": '/mcp' in mount_paths,
            "stripe_webhooks": any('stripe' in p.lower() or 'webhook' in p.lower() for p in route_paths),
        }
        if _mcp_error:
            agent_stack["mcp_error"] = _mcp_error
        return {"status": "ok", "version": "1.1.4", "agent_stack": agent_stack}
    except Exception as exc:
        logger.exception("/health failed: %s", exc)
        raise HTTPException(status_code=500, detail="Internal server error")


# ─────────────────────────────────────────────────────────────────────────────
# TEMPORARY DEBUG ENDPOINTS  — remove after root cause is confirmed
# Gated by SLIDEARABI_DEBUG_TOKEN env var (fail-closed: 403 if unset)
# ─────────────────────────────────────────────────────────────────────────────

DEBUG_TOKEN = os.getenv("SLIDEARABI_DEBUG_TOKEN", "")


def _require_debug_auth(token: str) -> None:
    """Validate debug token. 403 if unset or mismatched."""
    if not DEBUG_TOKEN or not secrets.compare_digest(token, DEBUG_TOKEN):
        raise HTTPException(status_code=403, detail="Forbidden")

def _sha256_file(path: Path) -> Optional[str]:
    """SHA-256 hash of a file, or None if it doesn't exist."""
    if not path.exists():
        return None
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def _count_arabic_chars_in_pptx(pptx_path: Path) -> Dict[str, Any]:
    """Scan a PPTX for Arabic script evidence. Returns a diagnostic dict."""
    stats: Dict[str, Any] = {
        "slides": 0,
        "text_shapes": 0,
        "nonempty_paragraphs": 0,
        "paragraphs_with_arabic": 0,
        "arabic_char_count": 0,
        "sample_arabic": [],
        "sample_english": [],
    }
    try:
        prs = PptxPresentation(str(pptx_path))
        stats["slides"] = len(prs.slides)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not (getattr(shape, "has_text_frame", False) and shape.has_text_frame):
                    continue
                stats["text_shapes"] += 1
                for para in shape.text_frame.paragraphs:
                    txt = (para.text or "").strip()
                    if not txt:
                        continue
                    stats["nonempty_paragraphs"] += 1
                    matches = ARABIC_RE.findall(txt)
                    if matches:
                        stats["paragraphs_with_arabic"] += 1
                        stats["arabic_char_count"] += len(matches)
                        if len(stats["sample_arabic"]) < 5:
                            stats["sample_arabic"].append(txt[:180])
                    else:
                        if len(stats["sample_english"]) < 5:
                            stats["sample_english"].append(txt[:180])
    except Exception as exc:
        stats["error"] = str(exc)
    return stats


@app.get("/debug/imports")
def debug_imports(token: str = ""):
    """Check which pipeline modules loaded and what system tools are available."""
    _require_debug_auth(token)
    try:
        flags = {
            "HAS_PROPERTY_RESOLVER": pipeline_module.HAS_PROPERTY_RESOLVER,
            "HAS_LAYOUT_ANALYZER": pipeline_module.HAS_LAYOUT_ANALYZER,
            "HAS_TEMPLATE_REGISTRY": pipeline_module.HAS_TEMPLATE_REGISTRY,
            "HAS_RTL_TRANSFORMS": pipeline_module.HAS_RTL_TRANSFORMS,
            "HAS_RTL_TRANSFORMS_V2": pipeline_module.HAS_RTL_TRANSFORMS_V2,
            "ENGINE_VERSION": pipeline_module.ENGINE_VERSION,
            "HAS_TYPOGRAPHY": pipeline_module.HAS_TYPOGRAPHY,
            "HAS_STRUCTURAL_VALIDATOR": pipeline_module.HAS_STRUCTURAL_VALIDATOR,
            "HAS_VQA": getattr(pipeline_module, "HAS_VQA", False),
            "HAS_LLM_TRANSLATOR": getattr(pipeline_module, "HAS_LLM_TRANSLATOR", False),
        }

        def _check_cmd(cmd: List[str]) -> Dict[str, Any]:
            exe = shutil.which(cmd[0])
            if not exe:
                return {"installed": False, "path": None}
            try:
                proc = subprocess.run(cmd, capture_output=True, text=True, timeout=10)
                return {"installed": True, "path": exe, "stdout": proc.stdout.strip()[:200]}
            except Exception as e:
                return {"installed": True, "path": exe, "error": str(e)}

        # Arabic fonts check
        arabic_fonts_raw = _check_cmd(["fc-list", ":lang=ar", "family"])
        arabic_font_lines = [
            line.strip()
            for line in arabic_fonts_raw.get("stdout", "").splitlines()
            if line.strip()
        ]

        return {
            "debug": True,
            "timestamp_utc": datetime.now(timezone.utc).isoformat(),
            "import_flags": flags,
            "system_tools": {
                "libreoffice": _check_cmd(["libreoffice", "--version"]),
                "pdftoppm": _check_cmd(["pdftoppm", "-v"]),
            },
            "arabic_fonts": {
                "count": len(arabic_font_lines),
                "families": arabic_font_lines[:50],
            },
            "env": {
                "OPENAI_API_KEY_set": bool(os.getenv("OPENAI_API_KEY")),
                "ANTHROPIC_API_KEY_set": bool(os.getenv("ANTHROPIC_API_KEY")),
            },
            "interpretation": {
                "likely_phase0_failure": not pipeline_module.HAS_PROPERTY_RESOLVER,
                "likely_translation_disabled": not getattr(pipeline_module, "HAS_LLM_TRANSLATOR", False),
            },
        }
    except Exception as exc:
        logger.exception("/debug/imports failed: %s", exc)
        raise HTTPException(status_code=500, detail=f"debug failed: {exc}")


@app.get("/debug/{job_id}")
def debug_job(job_id: str, token: str = ""):
    """Comprehensive diagnostic for a specific job.

    Reveals: module flags, file hashes, Arabic text scan, translation probe,
    preview source metadata, and phase reports.  Read-only — safe to hit
    repeatedly without side effects.
    """
    _require_debug_auth(token)
    try:
        with JOBS_LOCK:
            job = JOBS.get(job_id)
            if not job:
                raise HTTPException(status_code=404, detail="Job not found")

        input_path = Path(job.input_path)
        output_path = Path(job.output_path)

        # File metadata + hashing
        in_hash = _sha256_file(input_path)
        out_hash = _sha256_file(output_path)

        files = {
            "input_exists": input_path.exists(),
            "output_exists": output_path.exists(),
            "input_size": input_path.stat().st_size if input_path.exists() else 0,
            "output_size": output_path.stat().st_size if output_path.exists() else 0,
            "input_sha256": in_hash,
            "output_sha256": out_hash,
            "files_identical": (in_hash == out_hash) if (in_hash and out_hash) else None,
        }

        # Scan both files for Arabic text
        text_scan = {
            "input": _count_arabic_chars_in_pptx(input_path) if input_path.exists() else None,
            "output": _count_arabic_chars_in_pptx(output_path) if output_path.exists() else None,
        }

        # Module flags
        module_flags = {
            "HAS_PROPERTY_RESOLVER": pipeline_module.HAS_PROPERTY_RESOLVER,
            "HAS_LAYOUT_ANALYZER": pipeline_module.HAS_LAYOUT_ANALYZER,
            "HAS_TEMPLATE_REGISTRY": pipeline_module.HAS_TEMPLATE_REGISTRY,
            "HAS_RTL_TRANSFORMS": pipeline_module.HAS_RTL_TRANSFORMS,
            "HAS_RTL_TRANSFORMS_V2": pipeline_module.HAS_RTL_TRANSFORMS_V2,
            "ENGINE_VERSION": pipeline_module.ENGINE_VERSION,
            "HAS_TYPOGRAPHY": pipeline_module.HAS_TYPOGRAPHY,
            "HAS_STRUCTURAL_VALIDATOR": pipeline_module.HAS_STRUCTURAL_VALIDATOR,
            "HAS_VQA": getattr(pipeline_module, "HAS_VQA", False),
            "HAS_LLM_TRANSLATOR": getattr(pipeline_module, "HAS_LLM_TRANSLATOR", False),
        }

        # Phase reports (persisted by worker)
        phase_reports = {}
        phase_report_path = output_path.parent / "pipeline_phase_reports.json"
        if phase_report_path.exists():
            try:
                phase_reports = json.loads(phase_report_path.read_text(encoding="utf-8"))
            except Exception:
                phase_reports = {"error": "Failed reading pipeline_phase_reports.json"}

        # Preview sidecar
        preview_debug = {}
        preview_debug_path = output_path.parent / "preview_debug.json"
        if not preview_debug_path.exists():
            preview_debug_path = input_path.parent / "preview_debug.json"
        if preview_debug_path.exists():
            try:
                preview_debug = json.loads(preview_debug_path.read_text(encoding="utf-8"))
            except Exception:
                preview_debug = {"error": "Failed reading preview_debug.json"}

        preview_info = {
            "preview_slides_count": len(job.preview_slides or []),
            "preview_origin": job.preview_origin or "unknown",
            "sidecar_metadata": preview_debug,
        }

        # Diagnosis summary
        diagnosis: List[str] = []
        if files["files_identical"]:
            diagnosis.append("OUTPUT_IDENTICAL_TO_INPUT — pipeline was a no-op")
        if text_scan.get("output") and text_scan["output"].get("arabic_char_count", 0) == 0:
            diagnosis.append("NO_ARABIC_IN_OUTPUT — translation/apply failed")
        if text_scan.get("output") and text_scan["output"].get("arabic_char_count", 0) > 0:
            diagnosis.append("ARABIC_PRESENT_IN_OUTPUT — translation succeeded")
        if not module_flags["HAS_PROPERTY_RESOLVER"]:
            diagnosis.append("PROPERTY_RESOLVER_DISABLED — Phase 0 would skip")
        if not module_flags["HAS_RTL_TRANSFORMS"]:
            diagnosis.append("RTL_TRANSFORMS_DISABLED — Phase 3 would skip")
        if job.preview_origin == "fallback_input":
            diagnosis.append("PREVIEW_FROM_ENGLISH_FALLBACK")

        return {
            "debug": True,
            "timestamp_utc": datetime.now(timezone.utc).isoformat(),
            "job": {
                "job_id": job.job_id,
                "status": job.status,
                "progress_pct": job.progress_pct,
                "current_phase": job.current_phase,
                "error": job.error,
                "total_slides": job.total_slides,
            },
            "module_flags": module_flags,
            "files": files,
            "text_scan": text_scan,
            "preview": preview_info,
            "phase_reports": phase_reports,
            "env": {
                "OPENAI_API_KEY_set": bool(os.getenv("OPENAI_API_KEY")),
                "ANTHROPIC_API_KEY_set": bool(os.getenv("ANTHROPIC_API_KEY")),
            },
            "diagnosis": diagnosis,
        }
    except HTTPException:
        raise
    except Exception as exc:
        logger.exception("/debug/job failed for %s: %s", job_id, exc)
        raise HTTPException(status_code=500, detail=f"debug failed: {exc}")

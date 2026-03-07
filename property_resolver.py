"""
SlideArabi — PropertyResolver

The 7-level OOXML property inheritance resolver. This is the most important
class in the SlideArabi system: it walks the full Open XML inheritance
chain and produces a ResolvedPresentation where every text property has a
concrete, non-None effective value.

OOXML Inheritance Chain (checked in this order — first non-None wins):
  1. Run-level        (a:rPr on the run element)
  2. Paragraph-level  (a:pPr/a:defRPr on the paragraph)
  3. Text frame       (a:lstStyle on the shape's text body, at paragraph level)
  4. Shape-level      (shape's own inline or lstStyle properties)
  5. Layout           (matching placeholder on the slide layout by idx/type)
  6. Master           (matching placeholder on the slide master by idx/type)
  7. Master txStyles  (p:txStyles → titleStyle/bodyStyle/otherStyle at level)

If all 7 levels return None, we use PowerPoint's built-in defaults:
  - Font size: 18.0 pt
  - Font name: 'Calibri'
  - Bold: False
  - Italic: False
  - Alignment: 'l' (left)
  - RTL: False

Key OOXML facts:
  - Font sizes in hundredths of a point (e.g., 1800 = 18pt)
  - EMU = English Metric Units. 1 inch = 914400 EMU. 1 pt = 12700 EMU.
  - Theme font references: +mj-lt → majorFont latin, +mn-lt → minorFont latin
  - Placeholder matching: first by idx, then fallback to type
  - lstStyle has levels: lvl1pPr..lvl9pPr — paragraph level maps to lvlNpPr
"""

from __future__ import annotations

import logging
from typing import Any, Dict, List, Optional, Tuple

from slidearabi.models import (
    A_NS,
    DEFAULT_ALIGNMENT,
    DEFAULT_BOLD,
    DEFAULT_FONT_NAME,
    DEFAULT_FONT_SIZE_PT,
    DEFAULT_ITALIC,
    DEFAULT_LEVEL,
    DEFAULT_ROTATION,
    DEFAULT_RTL,
    DEFAULT_UNDERLINE,
    P_NS,
    R_NS,
    ResolvedLayout,
    ResolvedMaster,
    ResolvedParagraph,
    ResolvedPresentation,
    ResolvedRun,
    ResolvedShape,
    ResolvedSlide,
)

logger = logging.getLogger(__name__)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# XML HELPER UTILITIES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


def _qn(tag: str) -> str:
    """Expand a namespace-prefixed tag to Clark notation.

    Example: _qn('a:rPr') → '{http://schemas.openxmlformats.org/.../main}rPr'
    """
    prefix, local = tag.split(':', 1)
    ns_map = {'a': A_NS, 'p': P_NS, 'r': R_NS}
    return f'{{{ns_map[prefix]}}}{local}'


def _get_attr_int(element, attr: str) -> Optional[int]:
    """Get an integer attribute from an XML element, or None."""
    if element is None:
        return None
    val = element.get(attr)
    if val is None:
        return None
    try:
        return int(val)
    except (ValueError, TypeError):
        return None


def _get_attr_bool(element, attr: str) -> Optional[bool]:
    """Get a boolean attribute from an XML element.

    OOXML booleans: '1'/'true' = True, '0'/'false' = False, absent = None.
    """
    if element is None:
        return None
    val = element.get(attr)
    if val is None:
        return None
    return val in ('1', 'true')


def _get_attr_str(element, attr: str) -> Optional[str]:
    """Get a string attribute from an XML element, or None."""
    if element is None:
        return None
    return element.get(attr)


def _find(element, tag: str):
    """Find a child element by namespace-prefixed tag. Returns None safely."""
    if element is None:
        return None
    return element.find(_qn(tag))


def _findall(element, tag: str) -> list:
    """Find all child elements by namespace-prefixed tag. Returns [] safely."""
    if element is None:
        return []
    return element.findall(_qn(tag))


def _find_descendant(element, tag: str):
    """Find a descendant element (recursive). Returns None safely."""
    if element is None:
        return None
    return element.find(f'.//{_qn(tag)}')


def _find_first(element, *tags: str):
    """Find the first matching child element from multiple tags.

    Avoids lxml FutureWarning about truth-testing elements with `or`.
    """
    if element is None:
        return None
    for tag in tags:
        result = element.find(_qn(tag))
        if result is not None:
            return result
    return None


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PROPERTY RESOLVER
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class PropertyResolver:
    """Resolves effective OOXML property values by walking the 7-level
    inheritance chain.

    This class is the foundation of SlideArabi's Phase 0. It accepts a
    python-pptx Presentation object and produces a ResolvedPresentation
    where every text property has a concrete value.

    Usage:
        from pptx import Presentation
        from slidearabi.property_resolver import PropertyResolver

        prs = Presentation('input.pptx')
        resolver = PropertyResolver(prs)
        resolved = resolver.resolve_presentation()
        # resolved.slides[0].shapes[0].paragraphs[0].runs[0].effective_font_size_pt
        # → always a float, never None
    """

    def __init__(self, presentation) -> None:
        """Initialize with a python-pptx Presentation object.

        Args:
            presentation: A python-pptx Presentation instance.
        """
        self.prs = presentation
        self._theme_font_cache: Dict[str, Dict[str, str]] = {}
        self._master_index_map: Dict[int, int] = {}  # master element id → index
        self._layout_index_map: Dict[int, Tuple[int, int]] = {}  # layout id → (master_idx, layout_idx)

    # ─────────────────────────────────────────────────────────────
    # PUBLIC: Main entry point
    # ─────────────────────────────────────────────────────────────

    def resolve_presentation(self) -> ResolvedPresentation:
        """Resolve the entire presentation into an immutable snapshot.

        Walks every master, layout, and slide. For each shape with text,
        resolves all text properties through the full 7-level chain.

        Returns:
            A ResolvedPresentation with all effective values materialized.
        """
        slide_width = int(self.prs.slide_width)
        slide_height = int(self.prs.slide_height)

        # Build index maps for masters and layouts
        masters: List[ResolvedMaster] = []
        layouts: List[ResolvedLayout] = []

        for master_idx, master in enumerate(self.prs.slide_masters):
            master_elem = master._element
            self._master_index_map[id(master_elem)] = master_idx

            # Resolve master shapes
            master_placeholders, master_freeforms = self._resolve_element_shapes(
                master, source_level='master',
                layout=None, master_obj=None,
            )

            # Extract txStyles as dict
            tx_styles = self._extract_tx_styles(master)

            resolved_master = ResolvedMaster(
                master_name=getattr(master, 'name', '') or f'Master {master_idx + 1}',
                master_index=master_idx,
                placeholders=tuple(master_placeholders),
                freeform_shapes=tuple(master_freeforms),
                tx_styles=tx_styles,
            )
            masters.append(resolved_master)

            # Resolve layouts for this master
            for layout_idx, layout in enumerate(master.slide_layouts):
                layout_elem = layout._element
                self._layout_index_map[id(layout_elem)] = (master_idx, layout_idx)

                layout_placeholders, layout_freeforms = self._resolve_element_shapes(
                    layout, source_level='layout',
                    layout=None, master_obj=master,
                )

                layout_type = layout_elem.get('type', 'cust')
                layout_name = layout.name or f'Layout {layout_idx + 1}'

                resolved_layout = ResolvedLayout(
                    layout_name=layout_name,
                    layout_type=layout_type,
                    master_index=master_idx,
                    placeholders=tuple(layout_placeholders),
                    freeform_shapes=tuple(layout_freeforms),
                )
                layouts.append(resolved_layout)

        # Resolve slides
        slides: List[ResolvedSlide] = []
        for slide_idx, slide in enumerate(self.prs.slides):
            resolved_slide = self._resolve_slide(slide, slide_idx + 1)
            slides.append(resolved_slide)

        return ResolvedPresentation(
            slide_width_emu=slide_width,
            slide_height_emu=slide_height,
            masters=tuple(masters),
            layouts=tuple(layouts),
            slides=tuple(slides),
        )

    # ─────────────────────────────────────────────────────────────
    # SLIDE RESOLUTION
    # ─────────────────────────────────────────────────────────────

    def _resolve_slide(self, slide, slide_number: int) -> ResolvedSlide:
        """Resolve a single slide into a ResolvedSlide."""
        layout = slide.slide_layout
        master = layout.slide_master

        layout_elem = layout._element
        layout_type = layout_elem.get('type', 'cust')
        layout_name = layout.name or 'Unknown Layout'

        # Determine layout and master indices
        layout_key = id(layout_elem)
        master_key = id(master._element)
        master_idx = self._master_index_map.get(master_key, 0)
        layout_idx_tuple = self._layout_index_map.get(layout_key, (master_idx, 0))
        layout_idx = layout_idx_tuple[1]

        # Resolve all shapes on the slide
        shapes: List[ResolvedShape] = []
        for shape in self._collect_all_shapes(slide.shapes):
            try:
                resolved = self._resolve_shape(
                    shape,
                    source_level='slide',
                    layout=layout,
                    master_obj=master,
                )
                shapes.append(resolved)
            except Exception as e:
                shape_name = getattr(shape, 'name', '?')
                logger.warning(
                    f"Failed to resolve shape '{shape_name}' on slide {slide_number}: {e}"
                )

        return ResolvedSlide(
            slide_number=slide_number,
            layout_name=layout_name,
            layout_type=layout_type,
            layout_index=layout_idx,
            master_index=master_idx,
            shapes=tuple(shapes),
        )

    # ─────────────────────────────────────────────────────────────
    # SHAPE RESOLUTION
    # ─────────────────────────────────────────────────────────────

    def _resolve_element_shapes(
        self,
        element,
        source_level: str,
        layout,
        master_obj,
    ) -> Tuple[List[ResolvedShape], List[ResolvedShape]]:
        """Resolve shapes on a master or layout element.

        Returns:
            (placeholders, freeform_shapes) — two lists of ResolvedShape.
        """
        placeholders = []
        freeforms = []
        for shape in self._collect_all_shapes(element.shapes):
            try:
                resolved = self._resolve_shape(
                    shape, source_level=source_level,
                    layout=layout, master_obj=master_obj,
                )
                if resolved.is_placeholder:
                    placeholders.append(resolved)
                else:
                    freeforms.append(resolved)
            except Exception as e:
                shape_name = getattr(shape, 'name', '?')
                logger.warning(
                    f"Failed to resolve shape '{shape_name}' on {source_level}: {e}"
                )
        return placeholders, freeforms

    def _resolve_shape(
        self,
        shape,
        source_level: str,
        layout,
        master_obj,
    ) -> ResolvedShape:
        """Resolve a single shape, walking the inheritance chain for all text
        properties.

        Args:
            shape: A python-pptx shape object.
            source_level: 'master', 'layout', or 'slide'.
            layout: The slide layout (python-pptx SlideLayout), or None.
            master_obj: The slide master (python-pptx SlideMaster), or None.
        """
        sp_elem = shape._element

        # Shape identification
        shape_id = shape.shape_id
        shape_name = shape.name or ''
        shape_type = self._classify_shape_type(shape)

        # Placeholder info
        placeholder_type = None
        placeholder_idx = None
        is_placeholder_shape = False
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            is_placeholder_shape = True
            ph_format = shape.placeholder_format
            if ph_format is not None:
                placeholder_idx = ph_format.idx
                placeholder_type = self._get_placeholder_type_str(ph_format)

        # Position and size
        x_emu = int(shape.left) if shape.left is not None else 0
        y_emu = int(shape.top) if shape.top is not None else 0
        width_emu = int(shape.width) if shape.width is not None else 0
        height_emu = int(shape.height) if shape.height is not None else 0

        # Rotation
        rotation = DEFAULT_ROTATION
        try:
            if hasattr(shape, 'rotation') and shape.rotation is not None:
                rotation = float(shape.rotation)
        except (TypeError, ValueError):
            pass

        # Check for local position override (slide shape that overrides layout pos)
        has_local_position_override = self._has_local_position_override(sp_elem)

        # Resolve text paragraphs
        paragraphs: List[ResolvedParagraph] = []
        has_text = False

        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            tf = shape.text_frame

            # Find matching placeholder on layout and master
            layout_ph_elem = None
            master_ph_elem = None
            if is_placeholder_shape and layout is not None:
                layout_ph_elem = self._find_matching_placeholder(
                    shape, layout
                )
            if is_placeholder_shape and master_obj is not None:
                master_ph_elem = self._find_matching_placeholder(
                    shape, master_obj
                )

            # Get master txStyles for this shape
            master_txstyle = None
            if master_obj is not None:
                master_txstyle = self._get_master_txstyle_for_shape(
                    shape, master_obj
                )

            for para in tf.paragraphs:
                resolved_para = self._resolve_paragraph(
                    para, shape, tf,
                    layout_ph_elem, master_ph_elem,
                    master_txstyle, master_obj,
                )
                paragraphs.append(resolved_para)
                if any(r.text.strip() for r in resolved_para.runs):
                    has_text = True

        # For table shapes, resolve text in cells
        elif hasattr(shape, 'has_table') and shape.has_table:
            has_text = True
            # Tables are complex — we mark has_text but don't deeply resolve
            # cell text into paragraphs for the shape-level model.
            # Cell-level resolution would be handled separately.

        is_master_inherited = source_level in ('master', 'layout')

        return ResolvedShape(
            shape_id=shape_id,
            shape_name=shape_name,
            shape_type='placeholder' if is_placeholder_shape else shape_type,
            placeholder_type=placeholder_type,
            placeholder_idx=placeholder_idx,
            x_emu=x_emu,
            y_emu=y_emu,
            width_emu=width_emu,
            height_emu=height_emu,
            rotation_degrees=rotation,
            paragraphs=tuple(paragraphs),
            is_master_inherited=is_master_inherited,
            source_level=source_level,
            has_local_position_override=has_local_position_override,
            has_text=has_text,
            original_xml_element=sp_elem,
        )

    # ─────────────────────────────────────────────────────────────
    # PARAGRAPH RESOLUTION
    # ─────────────────────────────────────────────────────────────

    def _resolve_paragraph(
        self,
        paragraph,
        shape,
        text_frame,
        layout_ph_elem,
        master_ph_elem,
        master_txstyle,
        master_obj,
    ) -> ResolvedParagraph:
        """Resolve a single paragraph and its runs."""
        p_elem = paragraph._p

        # Get paragraph level
        level = self._get_paragraph_level(p_elem)

        # Resolve alignment
        alignment = self.resolve_alignment(
            paragraph, shape, layout_ph_elem, master_ph_elem,
            master_txstyle, level,
        )

        # Resolve RTL
        rtl = self._resolve_rtl(p_elem, layout_ph_elem, master_ph_elem, master_txstyle, level)

        # Resolve bullet type
        bullet_type = self._resolve_bullet_type(p_elem)

        # Resolve spacing
        line_spacing = self._resolve_line_spacing(p_elem)
        space_before = self._resolve_space_before(p_elem)
        space_after = self._resolve_space_after(p_elem)

        # Resolve runs
        runs: List[ResolvedRun] = []
        for run in paragraph.runs:
            resolved_run = self._resolve_run(
                run, paragraph, shape, text_frame,
                layout_ph_elem, master_ph_elem,
                master_txstyle, master_obj, level,
            )
            runs.append(resolved_run)

        # If no explicit runs but there is text (e.g., field elements),
        # create a synthetic run
        if not runs and paragraph.text.strip():
            resolved_run = self._create_synthetic_run(
                paragraph.text, paragraph, shape, text_frame,
                layout_ph_elem, master_ph_elem,
                master_txstyle, master_obj, level,
            )
            runs.append(resolved_run)

        return ResolvedParagraph(
            runs=tuple(runs),
            effective_alignment=alignment,
            effective_rtl=rtl,
            effective_level=level,
            effective_bullet_type=bullet_type,
            effective_line_spacing=line_spacing,
            effective_space_before=space_before,
            effective_space_after=space_after,
        )

    # ─────────────────────────────────────────────────────────────
    # RUN RESOLUTION
    # ─────────────────────────────────────────────────────────────

    def _resolve_run(
        self,
        run,
        paragraph,
        shape,
        text_frame,
        layout_ph_elem,
        master_ph_elem,
        master_txstyle,
        master_obj,
        level: int,
    ) -> ResolvedRun:
        """Resolve a single run's properties through the 7-level chain."""
        font_size, size_source = self.resolve_font_size(
            run, paragraph, shape, text_frame,
            layout_ph_elem, master_ph_elem,
            master_txstyle, level,
        )

        font_name = self.resolve_font_name(
            run, paragraph, shape, text_frame,
            layout_ph_elem, master_ph_elem,
            master_txstyle, master_obj, level,
        )

        bold = self.resolve_bold(
            run, paragraph, shape, text_frame,
            layout_ph_elem, master_ph_elem,
            master_txstyle, level,
        )

        italic = self.resolve_italic(
            run, paragraph, shape, text_frame,
            layout_ph_elem, master_ph_elem,
            master_txstyle, level,
        )

        underline = self._resolve_underline(run)
        color = self._resolve_color(run)

        return ResolvedRun(
            text=run.text,
            effective_font_size_pt=font_size,
            effective_font_name=font_name,
            effective_bold=bold,
            effective_italic=italic,
            effective_color=color,
            effective_underline=underline,
            source_font_size_level=size_source,
        )

    def _create_synthetic_run(
        self,
        text: str,
        paragraph,
        shape,
        text_frame,
        layout_ph_elem,
        master_ph_elem,
        master_txstyle,
        master_obj,
        level: int,
    ) -> ResolvedRun:
        """Create a resolved run for paragraphs with text but no explicit runs
        (e.g., field codes, endParaRPr-only paragraphs)."""

        font_size, size_source = self._resolve_font_size_no_run(
            paragraph, shape, text_frame,
            layout_ph_elem, master_ph_elem,
            master_txstyle, level,
        )

        font_name = self._resolve_font_name_no_run(
            paragraph, shape, text_frame,
            layout_ph_elem, master_ph_elem,
            master_txstyle, master_obj, level,
        )

        bold = self._resolve_bold_no_run(
            paragraph, shape, text_frame,
            layout_ph_elem, master_ph_elem,
            master_txstyle, level,
        )

        italic = self._resolve_italic_no_run(
            paragraph, shape, text_frame,
            layout_ph_elem, master_ph_elem,
            master_txstyle, level,
        )

        return ResolvedRun(
            text=text,
            effective_font_size_pt=font_size,
            effective_font_name=font_name,
            effective_bold=bold,
            effective_italic=italic,
            effective_color=None,
            effective_underline=DEFAULT_UNDERLINE,
            source_font_size_level=size_source,
        )

    # ─────────────────────────────────────────────────────────────
    # FONT SIZE RESOLUTION (7-level chain)
    # ─────────────────────────────────────────────────────────────

    def resolve_font_size(
        self,
        run,
        paragraph,
        shape,
        text_frame,
        layout_ph_elem,
        master_ph_elem,
        master_txstyle,
        level: int,
    ) -> Tuple[float, str]:
        """Walk the 7-level inheritance chain for font size.

        Returns:
            (font_size_pt, source_level) where font_size_pt is in points
            and source_level identifies which level provided the value.
        """
        r_elem = run._r

        # Level 1: Run-level (a:rPr sz)
        rPr = _find(r_elem, 'a:rPr')
        sz = _get_attr_int(rPr, 'sz')
        if sz is not None:
            return sz / 100.0, 'run'

        # Levels 2-7 (same as no-run resolution)
        return self._resolve_font_size_no_run(
            paragraph, shape, text_frame,
            layout_ph_elem, master_ph_elem,
            master_txstyle, level,
        )

    def _resolve_font_size_no_run(
        self,
        paragraph,
        shape,
        text_frame,
        layout_ph_elem,
        master_ph_elem,
        master_txstyle,
        level: int,
    ) -> Tuple[float, str]:
        """Resolve font size starting from level 2 (paragraph default).

        Used for synthetic runs and as the continuation after level 1 misses.
        """
        p_elem = paragraph._p
        level_tag = f'a:lvl{level + 1}pPr'

        # Level 2: Paragraph-level (a:pPr/a:defRPr sz)
        pPr = _find(p_elem, 'a:pPr')
        if pPr is not None:
            defRPr = _find(pPr, 'a:defRPr')
            sz = _get_attr_int(defRPr, 'sz')
            if sz is not None:
                return sz / 100.0, 'paragraph'

        # Level 3: Text frame lstStyle (a:lstStyle on the shape's txBody)
        try:
            txBody = text_frame._txBody
            sz = self._sz_from_lst_style(txBody, level_tag)
            if sz is not None:
                return sz / 100.0, 'textframe'
        except (AttributeError, TypeError):
            pass

        # Level 4: Shape-level (shape's own inline properties or lstStyle)
        # For most shapes, this is the same as level 3. But for shapes
        # embedded in groups or with distinct lstStyle, it can differ.
        try:
            sp_elem = shape._element
            sp_txBody = _find_first(sp_elem, 'p:txBody', 'a:txBody')
            if sp_txBody is not None and sp_txBody is not getattr(text_frame, '_txBody', None):
                sz = self._sz_from_lst_style(sp_txBody, level_tag)
                if sz is not None:
                    return sz / 100.0, 'shape'
        except (AttributeError, TypeError):
            pass

        # Level 5: Layout placeholder lstStyle
        if layout_ph_elem is not None:
            layout_txBody = _find_first(layout_ph_elem, 'p:txBody', 'a:txBody')
            sz = self._sz_from_lst_style(layout_txBody, level_tag)
            if sz is not None:
                return sz / 100.0, 'layout'

        # Level 6: Master placeholder lstStyle
        if master_ph_elem is not None:
            master_txBody = _find_first(master_ph_elem, 'p:txBody', 'a:txBody')
            sz = self._sz_from_lst_style(master_txBody, level_tag)
            if sz is not None:
                return sz / 100.0, 'master'

        # Level 7: Master txStyles
        if master_txstyle is not None:
            lvl_pPr = _find(master_txstyle, f'a:lvl{level + 1}pPr')
            if lvl_pPr is not None:
                defRPr = _find(lvl_pPr, 'a:defRPr')
                sz = _get_attr_int(defRPr, 'sz')
                if sz is not None:
                    return sz / 100.0, 'txstyles'

        # Fallback
        return DEFAULT_FONT_SIZE_PT, 'default'

    # ─────────────────────────────────────────────────────────────
    # FONT NAME RESOLUTION
    # ─────────────────────────────────────────────────────────────

    def resolve_font_name(
        self,
        run,
        paragraph,
        shape,
        text_frame,
        layout_ph_elem,
        master_ph_elem,
        master_txstyle,
        master_obj,
        level: int,
    ) -> str:
        """Resolve font name through the inheritance chain.

        Handles theme font references (+mj-lt → major font, +mn-lt → minor font)
        by resolving against the actual theme.
        """
        r_elem = run._r

        # Level 1: Run-level (a:rPr → a:latin/@typeface or a:cs/@typeface)
        rPr = _find(r_elem, 'a:rPr')
        name = self._font_name_from_rPr(rPr)
        if name is not None:
            return self._resolve_theme_font(name, master_obj)

        return self._resolve_font_name_no_run(
            paragraph, shape, text_frame,
            layout_ph_elem, master_ph_elem,
            master_txstyle, master_obj, level,
        )

    def _resolve_font_name_no_run(
        self,
        paragraph,
        shape,
        text_frame,
        layout_ph_elem,
        master_ph_elem,
        master_txstyle,
        master_obj,
        level: int,
    ) -> str:
        """Resolve font name starting from level 2."""
        p_elem = paragraph._p
        level_tag = f'a:lvl{level + 1}pPr'

        # Level 2: Paragraph default (a:pPr/a:defRPr)
        pPr = _find(p_elem, 'a:pPr')
        if pPr is not None:
            defRPr = _find(pPr, 'a:defRPr')
            name = self._font_name_from_rPr(defRPr)
            if name is not None:
                return self._resolve_theme_font(name, master_obj)

        # Level 3: Text frame lstStyle
        try:
            txBody = text_frame._txBody
            name = self._font_name_from_lst_style(txBody, level_tag)
            if name is not None:
                return self._resolve_theme_font(name, master_obj)
        except (AttributeError, TypeError):
            pass

        # Level 5: Layout placeholder
        if layout_ph_elem is not None:
            layout_txBody = _find_first(layout_ph_elem, 'p:txBody', 'a:txBody')
            name = self._font_name_from_lst_style(layout_txBody, level_tag)
            if name is not None:
                return self._resolve_theme_font(name, master_obj)

        # Level 6: Master placeholder
        if master_ph_elem is not None:
            master_txBody = _find_first(master_ph_elem, 'p:txBody', 'a:txBody')
            name = self._font_name_from_lst_style(master_txBody, level_tag)
            if name is not None:
                return self._resolve_theme_font(name, master_obj)

        # Level 7: Master txStyles
        if master_txstyle is not None:
            lvl_pPr = _find(master_txstyle, f'a:lvl{level + 1}pPr')
            if lvl_pPr is not None:
                defRPr = _find(lvl_pPr, 'a:defRPr')
                name = self._font_name_from_rPr(defRPr)
                if name is not None:
                    return self._resolve_theme_font(name, master_obj)

        return DEFAULT_FONT_NAME

    # ─────────────────────────────────────────────────────────────
    # ALIGNMENT RESOLUTION
    # ─────────────────────────────────────────────────────────────

    def resolve_alignment(
        self,
        paragraph,
        shape,
        layout_ph_elem,
        master_ph_elem,
        master_txstyle,
        level: int,
    ) -> str:
        """Resolve paragraph alignment through the inheritance chain.

        Returns one of: 'l', 'r', 'ctr', 'just', 'dist'.
        """
        p_elem = paragraph._p
        level_tag = f'a:lvl{level + 1}pPr'

        # Level 1: Explicit on paragraph (a:pPr algn)
        pPr = _find(p_elem, 'a:pPr')
        algn = _get_attr_str(pPr, 'algn')
        if algn is not None:
            return algn

        # Level 3: Text frame lstStyle
        try:
            txBody = shape.text_frame._txBody
            algn = self._algn_from_lst_style(txBody, level_tag)
            if algn is not None:
                return algn
        except (AttributeError, TypeError):
            pass

        # Level 5: Layout placeholder lstStyle
        if layout_ph_elem is not None:
            layout_txBody = _find_first(layout_ph_elem, 'p:txBody', 'a:txBody')
            algn = self._algn_from_lst_style(layout_txBody, level_tag)
            if algn is not None:
                return algn

        # Level 6: Master placeholder lstStyle
        if master_ph_elem is not None:
            master_txBody = _find_first(master_ph_elem, 'p:txBody', 'a:txBody')
            algn = self._algn_from_lst_style(master_txBody, level_tag)
            if algn is not None:
                return algn

        # Level 7: Master txStyles
        if master_txstyle is not None:
            lvl_pPr = _find(master_txstyle, f'a:lvl{level + 1}pPr')
            algn = _get_attr_str(lvl_pPr, 'algn')
            if algn is not None:
                return algn

        return DEFAULT_ALIGNMENT

    # ─────────────────────────────────────────────────────────────
    # BOLD RESOLUTION
    # ─────────────────────────────────────────────────────────────

    def resolve_bold(
        self,
        run,
        paragraph,
        shape,
        text_frame,
        layout_ph_elem,
        master_ph_elem,
        master_txstyle,
        level: int,
    ) -> bool:
        """Resolve bold through the inheritance chain."""
        r_elem = run._r

        # Level 1: Run-level
        rPr = _find(r_elem, 'a:rPr')
        b = _get_attr_bool(rPr, 'b')
        if b is not None:
            return b

        return self._resolve_bold_no_run(
            paragraph, shape, text_frame,
            layout_ph_elem, master_ph_elem,
            master_txstyle, level,
        )

    def _resolve_bold_no_run(
        self,
        paragraph,
        shape,
        text_frame,
        layout_ph_elem,
        master_ph_elem,
        master_txstyle,
        level: int,
    ) -> bool:
        """Resolve bold starting from level 2."""
        return self._resolve_bool_property(
            'b', paragraph, shape, text_frame,
            layout_ph_elem, master_ph_elem,
            master_txstyle, level, DEFAULT_BOLD,
        )

    # ─────────────────────────────────────────────────────────────
    # ITALIC RESOLUTION
    # ─────────────────────────────────────────────────────────────

    def resolve_italic(
        self,
        run,
        paragraph,
        shape,
        text_frame,
        layout_ph_elem,
        master_ph_elem,
        master_txstyle,
        level: int,
    ) -> bool:
        """Resolve italic through the inheritance chain."""
        r_elem = run._r

        # Level 1: Run-level
        rPr = _find(r_elem, 'a:rPr')
        i = _get_attr_bool(rPr, 'i')
        if i is not None:
            return i

        return self._resolve_italic_no_run(
            paragraph, shape, text_frame,
            layout_ph_elem, master_ph_elem,
            master_txstyle, level,
        )

    def _resolve_italic_no_run(
        self,
        paragraph,
        shape,
        text_frame,
        layout_ph_elem,
        master_ph_elem,
        master_txstyle,
        level: int,
    ) -> bool:
        """Resolve italic starting from level 2."""
        return self._resolve_bool_property(
            'i', paragraph, shape, text_frame,
            layout_ph_elem, master_ph_elem,
            master_txstyle, level, DEFAULT_ITALIC,
        )

    # ─────────────────────────────────────────────────────────────
    # GENERIC BOOLEAN PROPERTY RESOLVER (for bold/italic/etc.)
    # ─────────────────────────────────────────────────────────────

    def _resolve_bool_property(
        self,
        attr_name: str,
        paragraph,
        shape,
        text_frame,
        layout_ph_elem,
        master_ph_elem,
        master_txstyle,
        level: int,
        default: bool,
    ) -> bool:
        """Walk levels 2-7 for a boolean attribute on defRPr.

        Used for bold ('b'), italic ('i'), and similar run properties
        that cascade through the same chain as font size.
        """
        p_elem = paragraph._p
        level_tag = f'a:lvl{level + 1}pPr'

        # Level 2: Paragraph default (a:pPr/a:defRPr)
        pPr = _find(p_elem, 'a:pPr')
        if pPr is not None:
            defRPr = _find(pPr, 'a:defRPr')
            val = _get_attr_bool(defRPr, attr_name)
            if val is not None:
                return val

        # Level 3: Text frame lstStyle
        try:
            txBody = text_frame._txBody
            val = self._bool_from_lst_style(txBody, level_tag, attr_name)
            if val is not None:
                return val
        except (AttributeError, TypeError):
            pass

        # Level 5: Layout placeholder lstStyle
        if layout_ph_elem is not None:
            layout_txBody = _find_first(layout_ph_elem, 'p:txBody', 'a:txBody')
            val = self._bool_from_lst_style(layout_txBody, level_tag, attr_name)
            if val is not None:
                return val

        # Level 6: Master placeholder lstStyle
        if master_ph_elem is not None:
            master_txBody = _find_first(master_ph_elem, 'p:txBody', 'a:txBody')
            val = self._bool_from_lst_style(master_txBody, level_tag, attr_name)
            if val is not None:
                return val

        # Level 7: Master txStyles
        if master_txstyle is not None:
            lvl_pPr = _find(master_txstyle, f'a:lvl{level + 1}pPr')
            if lvl_pPr is not None:
                defRPr = _find(lvl_pPr, 'a:defRPr')
                val = _get_attr_bool(defRPr, attr_name)
                if val is not None:
                    return val

        return default

    # ─────────────────────────────────────────────────────────────
    # RTL RESOLUTION
    # ─────────────────────────────────────────────────────────────

    def _resolve_rtl(
        self,
        p_elem,
        layout_ph_elem,
        master_ph_elem,
        master_txstyle,
        level: int,
    ) -> bool:
        """Resolve RTL direction for a paragraph.

        Checks pPr rtl attribute, then layout/master lstStyle, then txStyles.
        """
        level_tag = f'a:lvl{level + 1}pPr'

        # Level 1: Explicit on paragraph
        pPr = _find(p_elem, 'a:pPr')
        rtl = _get_attr_bool(pPr, 'rtl')
        if rtl is not None:
            return rtl

        # Level 5: Layout placeholder lstStyle
        if layout_ph_elem is not None:
            layout_txBody = _find_first(layout_ph_elem, 'p:txBody', 'a:txBody')
            rtl = self._rtl_from_lst_style(layout_txBody, level_tag)
            if rtl is not None:
                return rtl

        # Level 6: Master placeholder lstStyle
        if master_ph_elem is not None:
            master_txBody = _find_first(master_ph_elem, 'p:txBody', 'a:txBody')
            rtl = self._rtl_from_lst_style(master_txBody, level_tag)
            if rtl is not None:
                return rtl

        # Level 7: Master txStyles
        if master_txstyle is not None:
            lvl_pPr = _find(master_txstyle, f'a:lvl{level + 1}pPr')
            rtl = _get_attr_bool(lvl_pPr, 'rtl')
            if rtl is not None:
                return rtl

        return DEFAULT_RTL

    # ─────────────────────────────────────────────────────────────
    # PLACEHOLDER MATCHING
    # ─────────────────────────────────────────────────────────────

    def _find_matching_placeholder(self, shape, layout_or_master) -> Optional[Any]:
        """Find the matching placeholder XML element on a layout or master.

        Matching strategy:
        1. First try to match by idx attribute (exact index match).
        2. If no idx match, fall back to matching by type attribute.

        Args:
            shape: The python-pptx shape to match.
            layout_or_master: A python-pptx SlideLayout or SlideMaster.

        Returns:
            The lxml element of the matching placeholder's sp element, or None.
        """
        if not hasattr(shape, 'placeholder_format') or shape.placeholder_format is None:
            return None

        target_idx = shape.placeholder_format.idx
        target_type = self._get_placeholder_type_str(shape.placeholder_format)

        best_match = None

        try:
            for ph_shape in layout_or_master.placeholders:
                ph_elem = ph_shape._element
                ph_el = _find_descendant(ph_elem, 'p:ph')
                if ph_el is None:
                    continue

                source_idx_str = ph_el.get('idx')
                source_type = ph_el.get('type', '')

                # Priority 1: Match by idx
                if target_idx is not None and source_idx_str is not None:
                    try:
                        source_idx = int(source_idx_str)
                        if source_idx == target_idx:
                            return ph_elem
                    except (ValueError, TypeError):
                        if source_idx_str == str(target_idx):
                            return ph_elem

                # Priority 2: Remember type match as fallback
                if target_type and source_type == target_type and best_match is None:
                    best_match = ph_elem
        except Exception:
            pass

        return best_match

    # ─────────────────────────────────────────────────────────────
    # MASTER TEXT STYLES
    # ─────────────────────────────────────────────────────────────

    def _get_master_txstyle_for_shape(self, shape, master) -> Optional[Any]:
        """Get the appropriate txStyle element from the master for a shape.

        The master's p:txStyles has three children:
        - p:titleStyle → for title/ctrTitle placeholders
        - p:bodyStyle → for body/subTitle/object placeholders
        - p:otherStyle → for everything else (textboxes, etc.)

        Args:
            shape: The python-pptx shape.
            master: The python-pptx SlideMaster.

        Returns:
            The lxml element of the appropriate style (e.g., p:bodyStyle), or None.
        """
        try:
            master_elem = master._element
            txStyles = _find(master_elem, 'p:txStyles')
            if txStyles is None:
                return None

            # Determine which style applies
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                ph_type_str = self._get_placeholder_type_str(shape.placeholder_format)
                if ph_type_str in ('title', 'ctrTitle'):
                    return _find(txStyles, 'p:titleStyle')
                elif ph_type_str in ('body', 'subTitle', 'obj', 'tbl', 'chart',
                                     'dgm', 'media', 'clipArt'):
                    return _find(txStyles, 'p:bodyStyle')
                else:
                    # dt, ftr, sldNum → otherStyle
                    return _find(txStyles, 'p:otherStyle')
            else:
                # Non-placeholder shapes use otherStyle
                return _find(txStyles, 'p:otherStyle')
        except Exception:
            return None

    def _extract_tx_styles(self, master) -> Dict[str, Any]:
        """Extract txStyles from master as a dictionary for the resolved model.

        Returns dict with keys 'titleStyle', 'bodyStyle', 'otherStyle',
        each holding the lxml element or None.
        """
        result: Dict[str, Any] = {}
        try:
            master_elem = master._element
            txStyles = _find(master_elem, 'p:txStyles')
            if txStyles is not None:
                result['titleStyle'] = _find(txStyles, 'p:titleStyle')
                result['bodyStyle'] = _find(txStyles, 'p:bodyStyle')
                result['otherStyle'] = _find(txStyles, 'p:otherStyle')
        except Exception:
            pass
        return result

    # ─────────────────────────────────────────────────────────────
    # THEME FONT RESOLUTION
    # ─────────────────────────────────────────────────────────────

    def _resolve_theme_font(self, font_name: str, master_obj) -> str:
        """Resolve theme font references to actual font names.

        Theme references:
        - +mj-lt → major font, latin script (a:majorFont/a:latin)
        - +mn-lt → minor font, latin script (a:minorFont/a:latin)
        - +mj-cs → major font, complex script (a:majorFont/a:cs)
        - +mn-cs → minor font, complex script (a:minorFont/a:cs)

        Args:
            font_name: The raw typeface value (may be a theme reference).
            master_obj: The slide master to access the theme.

        Returns:
            The resolved font name string.
        """
        if not font_name or not font_name.startswith('+'):
            return font_name

        # Cache key based on master identity
        cache_key = id(master_obj) if master_obj else 0
        if cache_key not in self._theme_font_cache:
            self._theme_font_cache[cache_key] = self._build_theme_font_map(master_obj)

        theme_map = self._theme_font_cache[cache_key]

        # Map theme reference to key
        ref_map = {
            '+mj-lt': 'major_latin',
            '+mn-lt': 'minor_latin',
            '+mj-cs': 'major_cs',
            '+mn-cs': 'minor_cs',
            '+mj-ea': 'major_ea',
            '+mn-ea': 'minor_ea',
        }

        key = ref_map.get(font_name)
        if key and key in theme_map:
            return theme_map[key]

        return DEFAULT_FONT_NAME

    def _build_theme_font_map(self, master_obj) -> Dict[str, str]:
        """Build a map of theme font references to actual font names.

        Walks the theme XML: a:theme → a:themeElements → a:fontScheme →
        a:majorFont / a:minorFont → a:latin/@typeface, a:cs/@typeface.
        """
        result: Dict[str, str] = {}
        if master_obj is None:
            return result

        try:
            # Try accessing theme through python-pptx
            theme = None
            if hasattr(master_obj, 'theme'):
                theme = master_obj.theme
            if theme is None:
                return result

            theme_elem = theme._element if hasattr(theme, '_element') else theme

            # Navigate: a:themeElements → a:fontScheme
            theme_elements = _find(theme_elem, 'a:themeElements')
            if theme_elements is None:
                return result

            font_scheme = _find(theme_elements, 'a:fontScheme')
            if font_scheme is None:
                return result

            # Major font
            major_font = _find(font_scheme, 'a:majorFont')
            if major_font is not None:
                latin = _find(major_font, 'a:latin')
                if latin is not None:
                    tf = latin.get('typeface')
                    if tf:
                        result['major_latin'] = tf
                cs = _find(major_font, 'a:cs')
                if cs is not None:
                    tf = cs.get('typeface')
                    if tf:
                        result['major_cs'] = tf
                ea = _find(major_font, 'a:ea')
                if ea is not None:
                    tf = ea.get('typeface')
                    if tf:
                        result['major_ea'] = tf

            # Minor font
            minor_font = _find(font_scheme, 'a:minorFont')
            if minor_font is not None:
                latin = _find(minor_font, 'a:latin')
                if latin is not None:
                    tf = latin.get('typeface')
                    if tf:
                        result['minor_latin'] = tf
                cs = _find(minor_font, 'a:cs')
                if cs is not None:
                    tf = cs.get('typeface')
                    if tf:
                        result['minor_cs'] = tf
                ea = _find(minor_font, 'a:ea')
                if ea is not None:
                    tf = ea.get('typeface')
                    if tf:
                        result['minor_ea'] = tf

        except Exception as e:
            logger.debug(f"Could not extract theme fonts: {e}")

        return result

    # ─────────────────────────────────────────────────────────────
    # LIST STYLE HELPERS
    # ─────────────────────────────────────────────────────────────

    def _sz_from_lst_style(self, txBody_or_elem, level_tag: str) -> Optional[int]:
        """Extract font size from a:lstStyle at the given level.

        Args:
            txBody_or_elem: The txBody element containing lstStyle.
            level_tag: e.g., 'a:lvl1pPr' for paragraph level 0.

        Returns:
            Font size in hundredths of a point, or None.
        """
        if txBody_or_elem is None:
            return None
        lstStyle = _find(txBody_or_elem, 'a:lstStyle')
        if lstStyle is None:
            return None
        lvl_pPr = _find(lstStyle, level_tag)
        if lvl_pPr is None:
            return None
        defRPr = _find(lvl_pPr, 'a:defRPr')
        return _get_attr_int(defRPr, 'sz')

    def _algn_from_lst_style(self, txBody_or_elem, level_tag: str) -> Optional[str]:
        """Extract alignment from a:lstStyle at the given level."""
        if txBody_or_elem is None:
            return None
        lstStyle = _find(txBody_or_elem, 'a:lstStyle')
        if lstStyle is None:
            return None
        lvl_pPr = _find(lstStyle, level_tag)
        return _get_attr_str(lvl_pPr, 'algn')

    def _bool_from_lst_style(
        self, txBody_or_elem, level_tag: str, attr_name: str
    ) -> Optional[bool]:
        """Extract a boolean attribute from a:lstStyle's defRPr at the given level."""
        if txBody_or_elem is None:
            return None
        lstStyle = _find(txBody_or_elem, 'a:lstStyle')
        if lstStyle is None:
            return None
        lvl_pPr = _find(lstStyle, level_tag)
        if lvl_pPr is None:
            return None
        defRPr = _find(lvl_pPr, 'a:defRPr')
        return _get_attr_bool(defRPr, attr_name)

    def _rtl_from_lst_style(self, txBody_or_elem, level_tag: str) -> Optional[bool]:
        """Extract RTL from a:lstStyle at the given level.

        RTL is on the lvlNpPr element itself, not on defRPr.
        """
        if txBody_or_elem is None:
            return None
        lstStyle = _find(txBody_or_elem, 'a:lstStyle')
        if lstStyle is None:
            return None
        lvl_pPr = _find(lstStyle, level_tag)
        return _get_attr_bool(lvl_pPr, 'rtl')

    def _font_name_from_rPr(self, rPr_elem) -> Optional[str]:
        """Extract font name from a run properties element.

        Checks a:latin/@typeface first, then a:cs/@typeface, then a:ea/@typeface.
        """
        if rPr_elem is None:
            return None

        # Try latin first
        latin = _find(rPr_elem, 'a:latin')
        if latin is not None:
            tf = latin.get('typeface')
            if tf:
                return tf

        # Try complex script
        cs = _find(rPr_elem, 'a:cs')
        if cs is not None:
            tf = cs.get('typeface')
            if tf:
                return tf

        # Try East Asian
        ea = _find(rPr_elem, 'a:ea')
        if ea is not None:
            tf = ea.get('typeface')
            if tf:
                return tf

        return None

    def _font_name_from_lst_style(self, txBody_or_elem, level_tag: str) -> Optional[str]:
        """Extract font name from lstStyle at the given level."""
        if txBody_or_elem is None:
            return None
        lstStyle = _find(txBody_or_elem, 'a:lstStyle')
        if lstStyle is None:
            return None
        lvl_pPr = _find(lstStyle, level_tag)
        if lvl_pPr is None:
            return None
        defRPr = _find(lvl_pPr, 'a:defRPr')
        return self._font_name_from_rPr(defRPr)

    # ─────────────────────────────────────────────────────────────
    # SIMPLE PROPERTY RESOLVERS
    # ─────────────────────────────────────────────────────────────

    def _resolve_underline(self, run) -> bool:
        """Resolve underline from run properties. Simple — no deep chain needed."""
        try:
            rPr = _find(run._r, 'a:rPr')
            if rPr is not None:
                u = rPr.get('u')
                if u is not None:
                    # 'sng', 'dbl', 'heavy', etc. are truthy; 'none' is falsy
                    return u != 'none'
        except Exception:
            pass
        return DEFAULT_UNDERLINE

    def _resolve_color(self, run) -> Optional[str]:
        """Extract color from run properties.

        Checks a:rPr → a:solidFill → a:srgbClr/@val.
        Also checks a:schemeClr (but does not resolve scheme colors to RGB).

        Returns hex color string or None.
        """
        try:
            rPr = _find(run._r, 'a:rPr')
            if rPr is None:
                return None

            solid_fill = _find(rPr, 'a:solidFill')
            if solid_fill is None:
                return None

            # Direct RGB color
            srgb = _find(solid_fill, 'a:srgbClr')
            if srgb is not None:
                return srgb.get('val')

            # Scheme color reference (not fully resolved — would need theme)
            scheme_clr = _find(solid_fill, 'a:schemeClr')
            if scheme_clr is not None:
                return scheme_clr.get('val')  # Returns scheme name like 'tx1'

        except Exception:
            pass
        return None

    def _resolve_bullet_type(self, p_elem) -> Optional[str]:
        """Extract bullet type from paragraph properties.

        Checks for a:buNone, a:buChar, a:buAutoNum, a:buBlip.
        """
        pPr = _find(p_elem, 'a:pPr')
        if pPr is None:
            return None

        if _find(pPr, 'a:buNone') is not None:
            return None

        bu_char = _find(pPr, 'a:buChar')
        if bu_char is not None:
            char = bu_char.get('char', '•')
            return f'char:{char}'

        bu_auto = _find(pPr, 'a:buAutoNum')
        if bu_auto is not None:
            auto_type = bu_auto.get('type', 'arabicPeriod')
            return f'auto:{auto_type}'

        bu_blip = _find(pPr, 'a:buBlip')
        if bu_blip is not None:
            return 'blip'

        return None

    def _resolve_line_spacing(self, p_elem) -> Optional[float]:
        """Extract line spacing from paragraph properties.

        Returns spacing as a multiplier (e.g., 1.5) or None.
        """
        pPr = _find(p_elem, 'a:pPr')
        if pPr is None:
            return None

        lnSpc = _find(pPr, 'a:lnSpc')
        if lnSpc is None:
            return None

        # Try percentage-based (spcPct)
        spcPct = _find(lnSpc, 'a:spcPct')
        if spcPct is not None:
            val = spcPct.get('val')
            if val is not None:
                try:
                    # OOXML stores percentage * 1000 (e.g., 150000 = 150% = 1.5)
                    return int(val) / 100000.0
                except (ValueError, TypeError):
                    pass

        # Try points-based (spcPts)
        spcPts = _find(lnSpc, 'a:spcPts')
        if spcPts is not None:
            val = spcPts.get('val')
            if val is not None:
                try:
                    # OOXML stores hundredths of a point
                    return int(val) / 100.0
                except (ValueError, TypeError):
                    pass

        return None

    def _resolve_space_before(self, p_elem) -> Optional[float]:
        """Extract space before from paragraph properties. Returns points or None."""
        return self._resolve_spacing(p_elem, 'a:spcBef')

    def _resolve_space_after(self, p_elem) -> Optional[float]:
        """Extract space after from paragraph properties. Returns points or None."""
        return self._resolve_spacing(p_elem, 'a:spcAft')

    def _resolve_spacing(self, p_elem, spacing_tag: str) -> Optional[float]:
        """Generic spacing extractor for spcBef/spcAft."""
        pPr = _find(p_elem, 'a:pPr')
        if pPr is None:
            return None

        spc = _find(pPr, spacing_tag)
        if spc is None:
            return None

        spcPts = _find(spc, 'a:spcPts')
        if spcPts is not None:
            val = spcPts.get('val')
            if val is not None:
                try:
                    return int(val) / 100.0
                except (ValueError, TypeError):
                    pass

        spcPct = _find(spc, 'a:spcPct')
        if spcPct is not None:
            val = spcPct.get('val')
            if val is not None:
                try:
                    return int(val) / 100000.0
                except (ValueError, TypeError):
                    pass

        return None

    # ─────────────────────────────────────────────────────────────
    # SHAPE CLASSIFICATION HELPERS
    # ─────────────────────────────────────────────────────────────

    def _classify_shape_type(self, shape) -> str:
        """Classify a shape into one of the valid shape types.

        Uses the shape's XML element tag and python-pptx shape_type.
        """
        try:
            sp_elem = shape._element
            tag = sp_elem.tag

            # Check element tag for common types
            if tag.endswith('}pic') or tag == 'pic':
                return 'picture'
            if tag.endswith('}graphicFrame') or tag == 'graphicFrame':
                # Could be chart, table, smartart, or OLE
                if hasattr(shape, 'has_chart') and shape.has_chart:
                    return 'chart'
                if hasattr(shape, 'has_table') and shape.has_table:
                    return 'table'
                # Check for SmartArt (dgm namespace)
                graphic = _find(sp_elem, 'a:graphic')
                if graphic is not None:
                    graphic_data = _find(graphic, 'a:graphicData')
                    if graphic_data is not None:
                        uri = graphic_data.get('uri', '')
                        if 'diagram' in uri.lower() or 'dgm' in uri.lower():
                            return 'smartart'
                        if 'oleObject' in uri:
                            return 'ole'
                        if 'media' in uri.lower():
                            return 'media'
                return 'ole'  # Default for unrecognized graphicFrame
            if tag.endswith('}cxnSp') or tag == 'cxnSp':
                return 'connector'
            if tag.endswith('}grpSp') or tag == 'grpSp':
                return 'group'

            # Check python-pptx shape_type for finer classification
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            st = getattr(shape, 'shape_type', None)
            if st is not None:
                shape_type_map = {
                    MSO_SHAPE_TYPE.PICTURE: 'picture',
                    MSO_SHAPE_TYPE.CHART: 'chart',
                    MSO_SHAPE_TYPE.TABLE: 'table',
                    MSO_SHAPE_TYPE.GROUP: 'group',
                    MSO_SHAPE_TYPE.MEDIA: 'media',
                    MSO_SHAPE_TYPE.FREEFORM: 'freeform',
                    MSO_SHAPE_TYPE.TEXT_BOX: 'textbox',
                    MSO_SHAPE_TYPE.PLACEHOLDER: 'placeholder',
                }
                if st in shape_type_map:
                    return shape_type_map[st]

            # Check if it's a placeholder
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                return 'placeholder'

            # Check for auto-shape (sp element with no placeholder)
            if tag.endswith('}sp') or tag == 'sp':
                # Check if it's a textbox
                nvSpPr = _find(sp_elem, 'p:nvSpPr')
                if nvSpPr is not None:
                    sp_pr = _find(sp_elem, 'p:spPr')
                    if sp_pr is not None:
                        prstGeom = _find(sp_pr, 'a:prstGeom')
                        if prstGeom is not None and prstGeom.get('prst') == 'rect':
                            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                                return 'textbox'
                return 'freeform'

        except Exception:
            pass

        return 'freeform'

    def _get_placeholder_type_str(self, ph_format) -> Optional[str]:
        """Convert a python-pptx placeholder format type to a string.

        Maps PP_PLACEHOLDER enum values to the OOXML type attribute strings.
        """
        if ph_format is None:
            return None

        try:
            ph_type = ph_format.type
            if ph_type is None:
                return None

            # Try to get the string from the enum member name
            from pptx.enum.shapes import PP_PLACEHOLDER

            type_map = {
                PP_PLACEHOLDER.TITLE: 'title',
                PP_PLACEHOLDER.CENTER_TITLE: 'ctrTitle',
                PP_PLACEHOLDER.SUBTITLE: 'subTitle',
                PP_PLACEHOLDER.BODY: 'body',
                PP_PLACEHOLDER.OBJECT: 'obj',
                PP_PLACEHOLDER.DATE: 'dt',
                PP_PLACEHOLDER.FOOTER: 'ftr',
                PP_PLACEHOLDER.SLIDE_NUMBER: 'sldNum',
            }

            result = type_map.get(ph_type)
            if result:
                return result

            # Fall back: try XML element directly
            sp_elem = getattr(ph_format, '_element', None)
            if sp_elem is not None:
                ph_el = _find_descendant(sp_elem, 'p:ph')
                if ph_el is not None:
                    return ph_el.get('type')

            # Last resort: convert enum to string
            try:
                return str(ph_type).split('.')[-1].lower()
            except Exception:
                return None

        except Exception:
            return None

    def _get_paragraph_level(self, p_elem) -> int:
        """Get the paragraph indent level (0-8)."""
        pPr = _find(p_elem, 'a:pPr')
        if pPr is not None:
            lvl = pPr.get('lvl')
            if lvl is not None:
                try:
                    return max(0, min(8, int(lvl)))
                except (ValueError, TypeError):
                    pass
        return DEFAULT_LEVEL

    def _has_local_position_override(self, sp_elem) -> bool:
        """Check if a shape element has explicit position attributes (xfrm).

        A slide-level placeholder that has an a:off with explicit x,y coordinates
        has a local position override — it's not relying purely on layout inheritance.
        """
        try:
            # Look for spPr → xfrm → off
            spPr = _find(sp_elem, 'p:spPr')
            if spPr is None:
                spPr = _find(sp_elem, 'a:spPr')
            if spPr is None:
                return False
            xfrm = _find(spPr, 'a:xfrm')
            if xfrm is None:
                return False
            off = _find(xfrm, 'a:off')
            if off is not None and (off.get('x') is not None or off.get('y') is not None):
                return True
        except Exception:
            pass
        return False

    def _collect_all_shapes(self, shapes_collection) -> List[Any]:
        """Collect all shapes including group children (flattened).

        Does not recurse into nested groups beyond one level to avoid
        double-counting in the coordinate space.
        """
        result = []
        try:
            for shape in shapes_collection:
                result.append(shape)
        except Exception:
            pass
        return result

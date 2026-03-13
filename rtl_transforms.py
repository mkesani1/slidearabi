"""
rtl_transforms.py — Deterministic RTL transformation functions.

SlideArabi: Template-First Deterministic RTL Transformation Engine.

Phase 2: MasterLayoutTransformer  — transforms slide masters and layouts FIRST
Phase 3: SlideContentTransformer  — transforms individual content slides

Design principles (from 4-model architectural consensus):
1. Masters are style sheets — set direction/language defaults, mirror logos only.
   NEVER set algn at master/layout level (context-sensitive, slide-level concern).
   NEVER apply flipH (corrupts logos and gradients).
2. Layouts own placeholder geometry — mirror and swap placeholder positions.
3. Slides own content — placeholders inherit from transformed layout by deleting
   local position overrides; freeform shapes are mirrored explicitly.
4. Alignment (algn) is ALWAYS written explicitly at the paragraph level on slides.
"""

from __future__ import annotations

import logging
from copy import deepcopy
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Tuple

from lxml import etree

from .utils import (
    A_NS, P_NS, R_NS,
    mirror_x,
    swap_positions,
    has_arabic,
    compute_script_ratio,
    qn,
    ensure_pPr,
    set_rtl_on_paragraph,
    set_alignment_on_paragraph,
    get_placeholder_info,
    get_placeholder_info_from_xml,
    set_body_pr_rtl_col,
    set_defRPr_lang,
    iter_paragraphs,
    iter_runs,
    get_run_text,
    bounds_check_emu,
    clamp_emu,
)

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# Data classes
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class TransformReport:
    """Summary report produced by each transformation phase."""
    phase: str  # 'master', 'layout', 'slide', 'typography'
    total_changes: int = 0
    changes_by_type: Dict[str, int] = field(default_factory=dict)
    warnings: List[str] = field(default_factory=list)
    errors: List[str] = field(default_factory=list)

    def add(self, change_type: str, count: int = 1) -> None:
        """Increment a change-type counter and the total."""
        self.changes_by_type[change_type] = (
            self.changes_by_type.get(change_type, 0) + count
        )
        self.total_changes += count

    def warn(self, msg: str) -> None:
        self.warnings.append(msg)
        logger.warning('[%s] %s', self.phase, msg)

    def error(self, msg: str) -> None:
        self.errors.append(msg)
        logger.error('[%s] %s', self.phase, msg)

    def merge(self, other: 'TransformReport') -> None:
        """Merge another report's counters into this one."""
        self.total_changes += other.total_changes
        for k, v in other.changes_by_type.items():
            self.changes_by_type[k] = self.changes_by_type.get(k, 0) + v
        self.warnings.extend(other.warnings)
        self.errors.extend(other.errors)


# ─────────────────────────────────────────────────────────────────────────────
# Internal constants
# ─────────────────────────────────────────────────────────────────────────────

# Logo detection: image width must be less than this fraction of slide width
_LOGO_MAX_WIDTH_FRACTION = 0.20

# Position tolerance: skip mirror if the change would be smaller than this (EMU)
_POSITION_TOLERANCE_EMU = 50_000  # ≈ 0.055 inches

# Placeholder type strings that are "title-like" (keep centered or right-align)
_TITLE_PH_TYPES = frozenset({'title', 'ctrTitle', 'center_title'})

# Placeholder type strings that should always be left-aligned (footers, dates)
_FOOTER_PH_TYPES = frozenset({'ftr', 'sldNum', 'dt', 'footer', 'slideNumber', 'date_time'})


# ─────────────────────────────────────────────────────────────────────────────
# MasterLayoutTransformer — Phase 2
# ─────────────────────────────────────────────────────────────────────────────

class MasterLayoutTransformer:
    """
    Transforms slide masters and layouts for RTL (Phase 2).

    Operates BEFORE any content slides are touched so that placeholders
    on slides can inherit the correct RTL positions from their layouts.

    What this does:
    - Masters: set RTL text-direction defaults (rtl, rtlCol, lang); mirror logos.
    - Layouts: mirror/swap placeholder X positions; set RTL defaults.

    What this does NOT do (by design):
    - Set algn on masters/layouts (context-sensitive; set at slide paragraph level).
    - Apply flipH to any shape (corrupts logos, inverts brand gradients).
    - Modify shape heights or fonts (typography phase).
    """

    def __init__(self, presentation, template_registry=None):
        """
        Args:
            presentation: python-pptx Presentation object.
            template_registry: Optional TemplateRegistry instance with layout rules.
                When None, a built-in default ruleset is used.
        """
        self.prs = presentation
        self.template_registry = template_registry
        self._slide_width = int(presentation.slide_width)
        self._slide_height = int(presentation.slide_height)

    # ─────────────────────────────────────────────────────────────────────
    # Public entry points
    # ─────────────────────────────────────────────────────────────────────

    def transform_all_masters(self) -> TransformReport:
        """
        Transform all slide masters for RTL.

        Returns:
            TransformReport summarising all changes made.
        """
        report = TransformReport(phase='master')
        for idx, master in enumerate(self.prs.slide_masters):
            try:
                count = self._transform_master(master)
                report.add('master_transformed', count)
            except Exception as exc:
                report.error(f'master[{idx}]: {exc}')
        return report

    def transform_all_layouts(self) -> TransformReport:
        """
        Transform all slide layouts for RTL.

        Returns:
            TransformReport summarising all changes made.
        """
        report = TransformReport(phase='layout')
        for m_idx, master in enumerate(self.prs.slide_masters):
            for l_idx, layout in enumerate(master.slide_layouts):
                try:
                    count = self._transform_layout(layout)
                    report.add('layout_transformed', count)
                except Exception as exc:
                    report.error(f'master[{m_idx}].layout[{l_idx}]: {exc}')
        return report

    # ─────────────────────────────────────────────────────────────────────
    # Master transformation
    # ─────────────────────────────────────────────────────────────────────

    def _transform_master(self, master) -> int:
        """
        Transform a single slide master:

        1. Set RTL text-direction defaults in txStyles (titleStyle, bodyStyle,
           otherStyle): rtl='1', rtlCol='1', lang='ar-SA'.
        2. Mirror small logo/brand image positions (Step E whitelist).
        3. Mirror small brand text elements (company names, lettermarks).
        4. Set RTL defaults on all bodyPr, lstStyle, defPPr, defRPr elements.

        Does NOT set algn at master level (violates design principle).
        Does NOT apply flipH to any shape.

        Returns:
            Count of XML attribute writes performed.
        """
        changes = 0
        xml_el = master._element

        # 1. RTL text-direction defaults on all text body elements
        changes += self._apply_rtl_direction_defaults(xml_el)

        # 2. Arabic language defaults on all defRPr elements
        changes += self._apply_arabic_lang_defaults(xml_el)

        # 3. Set RTL in master txStyles (paragraph level defaults per style)
        changes += self._set_master_text_styles_rtl(master)

        # 4. Mirror logo images (position only — NO flipH)
        changes += self._mirror_logo_images(master)

        # 5. Mirror small brand text elements (company names, lettermarks)
        changes += self._mirror_brand_elements(master)

        return changes

    def _apply_rtl_direction_defaults(self, xml_element) -> int:
        """
        Set default RTL direction on master/layout elements.

        - bodyPr → rtlCol="1" (column order)
        - defPPr → rtl="1" (default paragraph direction)
        - lstStyle/lvlNpPr → rtl="1" (list level paragraph direction)

        CRITICAL: Do NOT set algn here — alignment is context-sensitive
        and must be set per-paragraph at slide level.
        """
        changes = 0

        # bodyPr → rtlCol="1" (safe — controls column order, not text direction)
        for body_pr in xml_element.iter(f'{{{A_NS}}}bodyPr'):
            body_pr.set('rtlCol', '1')
            changes += 1

        # defPPr → rtl="1" (default paragraph direction, but NOT algn)
        for def_ppr in xml_element.iter(f'{{{A_NS}}}defPPr'):
            def_ppr.set('rtl', '1')
            changes += 1

        # lstStyle/lvlNpPr → rtl="1" (list level direction, but NOT algn)
        for lst_style in xml_element.iter(f'{{{A_NS}}}lstStyle'):
            for child in lst_style:
                if child.tag.startswith(f'{{{A_NS}}}lvl') and child.tag.endswith('pPr'):
                    child.set('rtl', '1')
                    changes += 1

        return changes

    def _apply_arabic_lang_defaults(self, xml_element) -> int:
        """Set lang='ar-SA' on all <a:defRPr> elements."""
        changes = 0
        for def_rPr in xml_element.iter(f'{{{A_NS}}}defRPr'):
            def_rPr.set('lang', 'ar-SA')
            changes += 1
        return changes

    def _set_master_text_styles_rtl(self, master) -> int:
        """
        Set RTL-related defaults in the master's txStyles element.

        Sets rtl='1' on txStyles lvlNpPr for default RTL direction,
        and lang='ar-SA' on defRPr for Arabic font selection.
        Does NOT set algn — alignment is per-paragraph at slide level.
        """
        changes = 0
        try:
            xml_el = master._element
            tx_styles = xml_el.find(f'{{{P_NS}}}txStyles')
            if tx_styles is None:
                return 0

            for style_name in ('titleStyle', 'bodyStyle', 'otherStyle'):
                style_elem = tx_styles.find(f'{{{A_NS}}}{style_name}')
                if style_elem is None:
                    continue
                for level in range(1, 10):
                    for lvl_pPr in style_elem.findall(f'{{{A_NS}}}lvl{level}pPr'):
                        # Set rtl='1' for default paragraph direction (not algn)
                        lvl_pPr.set('rtl', '1')
                        changes += 1
                        # Set language for Arabic font selection
                        defRPr = lvl_pPr.find(f'{{{A_NS}}}defRPr')
                        if defRPr is not None:
                            defRPr.set('lang', 'ar-SA')
                            changes += 1
        except Exception as exc:
            logger.warning('_set_master_text_styles_rtl: %s', exc)
        return changes

    def _mirror_brand_elements(self, element) -> int:
        """
        Mirror small non-placeholder text shapes on masters/layouts.

        These are brand elements like company names ("MACHANI ROBOTICS"),
        lettermarks ("M"), and similar decorative text that should mirror
        position for RTL. Handles rotated shapes and negative positions.

        Whitelist criteria:
        - Non-placeholder shape
        - Has a text frame with non-empty text
        - Effective width (accounting for rotation) < 30% of slide width
        - Is not a background/full-width shape

        For rotated shapes with negative positions (e.g., vertical sidebar
        text extending off the left edge), the mirror accounts for the
        visual bounding box width rather than the shape's logical width.

        Returns:
            Count of shapes mirrored.
        """
        mirrored = 0
        try:
            for shape in element.shapes:
                # Skip placeholders
                sp_el = shape._element
                ph = sp_el.find(f'.//{{{P_NS}}}ph')
                if ph is not None:
                    continue

                # Must have text content
                if not (getattr(shape, 'has_text_frame', False) and shape.has_text_frame):
                    continue
                text = shape.text_frame.text or ''
                if not text.strip():
                    continue

                left = shape.left
                width = shape.width
                height = shape.height
                if left is None or width is None:
                    continue

                # Check rotation to determine effective visual width
                rotation_deg = 0
                for xfrm in sp_el.iter(f'{{{A_NS}}}xfrm'):
                    rot_val = xfrm.get('rot', '0')
                    try:
                        rotation_deg = int(rot_val) / 60000
                    except (ValueError, TypeError):
                        pass
                    break

                # For 90°/270° rotated shapes, visual width = logical height
                is_rotated_90 = abs(rotation_deg) in (90, 270)
                if is_rotated_90:
                    # Rotated shape — its visual footprint on the X axis
                    # is the height, not the width. The 'left' and 'width'
                    # in the XML refer to the pre-rotation bounding box.
                    # The shape's visual right edge = left + height (rotated).
                    visual_width = height if height else width
                else:
                    visual_width = width

                # Skip full-width or large shapes (backgrounds)
                if visual_width > self._slide_width * 0.30:
                    continue

                # Compute mirrored position
                new_left = mirror_x(left, width, self._slide_width)

                # For shapes with negative original positions (intentional
                # partial off-screen placement), preserve the same bleed
                # amount on the opposite side
                if left < 0:
                    # Original bleed: how far off-screen on the left
                    bleed = abs(left)
                    # Mirror the visible portion: right_edge = left + width
                    visible_right_edge = left + width
                    if visible_right_edge <= 0:
                        continue  # Entirely off-screen — skip
                    # New position: same bleed off the right edge
                    new_left = self._slide_width - width + bleed

                # Skip if change is negligible
                if abs(new_left - left) < _POSITION_TOLERANCE_EMU:
                    continue

                shape.left = new_left
                mirrored += 1
                logger.debug('Brand element mirrored: %s (%r) %d → %d',
                             shape.name, text[:30], left, new_left)

        except Exception as exc:
            logger.warning('_mirror_brand_elements: %s', exc)

        return mirrored

    def _mirror_logo_images(self, element) -> int:
        """
        Mirror the horizontal position of small logo images on a master or layout.

        Strict whitelist (4-model consensus — see master_slide_agent.py R35):
        - Shape is a picture element (<p:pic>)
        - Has no text frame
        - Width < 20% of slide width (small — not a background)
        - Is not a placeholder
        - Has an actual image relationship (blipFill with rEmbed)

        Applies position translation ONLY — no flipH.

        Returns:
            Count of shapes mirrored.
        """
        mirrored = 0
        try:
            for shape in element.shapes:
                if not self._is_logo_shape(shape, self._slide_width):
                    continue

                original_left = shape.left
                original_width = shape.width
                if original_left is None or original_width is None:
                    continue

                new_left = mirror_x(original_left, original_width, self._slide_width)

                # Bounds check
                if not bounds_check_emu(new_left, self._slide_width):
                    logger.debug('Logo mirror out of bounds: %s new_left=%d', shape.name, new_left)
                    continue

                # Skip if change is negligible (centred/symmetric shape)
                if abs(new_left - original_left) < _POSITION_TOLERANCE_EMU:
                    continue

                shape.left = new_left
                mirrored += 1
                logger.debug('Logo mirrored: %s  %d → %d', shape.name, original_left, new_left)

        except Exception as exc:
            logger.warning('_mirror_logo_images: %s', exc)

        return mirrored

    def _is_logo_shape(self, shape, slide_width: int) -> bool:
        """
        Detect if a shape is likely a logo (small image in a corner of the slide).

        Returns True only if ALL conditions are met:
        1. Is a picture element (<p:pic>) — not a text shape, group, chart, etc.
        2. Has no text frame (pure image).
        3. Width < 20% of slide width.
        4. Is not a placeholder.
        5. Has an actual image relationship (blipFill with rId embed).
        """
        try:
            sp_el = shape._element
            tag = sp_el.tag

            # Condition 1: Must be a picture element
            if not (tag.endswith('}pic') or tag == 'pic'):
                try:
                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                    if getattr(shape, 'shape_type', None) != MSO_SHAPE_TYPE.PICTURE:
                        return False
                except ImportError:
                    return False

            # Condition 2: Must not have a text frame
            if getattr(shape, 'has_text_frame', False) and shape.has_text_frame:
                return False

            # Condition 3: Width < 20% of slide width
            shape_width = getattr(shape, 'width', None)
            if shape_width is None or shape_width >= slide_width * _LOGO_MAX_WIDTH_FRACTION:
                return False

            # Condition 4: Must not be a placeholder
            nv_pic_pr = sp_el.find(f'{{{P_NS}}}nvPicPr')
            if nv_pic_pr is not None:
                ph = nv_pic_pr.find(f'.//{{{P_NS}}}ph')
                if ph is not None:
                    return False

            # Condition 5: Must have actual image (blipFill with rId)
            blip_fill = sp_el.find(f'{{{P_NS}}}blipFill')
            if blip_fill is None:
                blip_fill = sp_el.find(f'{{{A_NS}}}blipFill')
            if blip_fill is not None:
                blip = blip_fill.find(f'{{{A_NS}}}blip')
                if blip is not None and blip.get(f'{{{R_NS}}}embed'):
                    return True

            return False
        except Exception:
            return False

    # ─────────────────────────────────────────────────────────────────────
    # Layout transformation
    # ─────────────────────────────────────────────────────────────────────

    def _transform_layout(self, layout) -> int:
        """
        Transform a single slide layout for RTL.

        1. Look up the layout type (ST_SlideLayoutType) from the OOXML attribute.
        2. Mirror all placeholder X positions (new_x = slide_width - x - w).
        3. Swap X positions of column pairs for two-column layouts (twoColTx,
           txAndChart, chartAndTx, twoObj, twoTxTwoObj, etc.).
        4. Set RTL direction defaults on all text properties (same rules as master,
           but at layout scope).

        Returns:
            Count of XML attribute writes performed.
        """
        changes = 0
        xml_el = layout._element
        layout_type = xml_el.get('type', 'cust')

        # Apply direction defaults first
        changes += self._apply_rtl_direction_defaults(xml_el)
        changes += self._apply_arabic_lang_defaults(xml_el)

        # Mirror placeholder positions
        changes += self._mirror_layout_placeholders(layout, layout_type)

        return changes

    def _mirror_layout_placeholders(self, layout, layout_type: str) -> int:
        """
        Mirror placeholder positions in a layout according to RTL rules.

        For standard two-column layout types, the two content placeholders
        are *swapped* (left column moves to right side, right column to left side)
        so the reading order is preserved in RTL.

        For all other placeholders (title, body, subtitle, etc.), positions
        are *mirrored* using the standard formula: new_x = slide_width - x - w.

        Column-swap layout types:
            twoColTx, twoObj, twoTxTwoObj, txAndChart, chartAndTx, picTx

        Args:
            layout: python-pptx SlideLayout object.
            layout_type: ST_SlideLayoutType string from OOXML.

        Returns:
            Count of position changes applied.
        """
        changes = 0
        slide_width = self._slide_width

        # Identify two-column layout types that need column swapping
        two_column_types = frozenset({
            'twoColTx', 'twoObj', 'twoTxTwoObj',
            'txAndChart', 'chartAndTx', 'picTx',
        })

        if layout_type in two_column_types:
            changes += self._swap_two_column_placeholders(layout, slide_width)
        else:
            # Mirror all placeholders individually
            for shape in layout.placeholders:
                try:
                    left = shape.left
                    width = shape.width
                    if left is None or width is None:
                        continue
                    new_left = mirror_x(left, width, slide_width)
                    if not bounds_check_emu(new_left, slide_width):
                        logger.debug('Layout placeholder mirror OOB: %s', shape.name)
                        continue
                    if abs(new_left - left) < _POSITION_TOLERANCE_EMU:
                        continue
                    # Fix 21: Ensure full xfrm geometry is written when mirroring
                    # layout placeholders. python-pptx's shape.left setter creates
                    # a partial xfrm (off with only x) when the layout had no local
                    # xfrm. This causes y/cy to be lost, collapsing all placeholders
                    # to y=0. We write the complete xfrm explicitly.
                    self._set_layout_placeholder_position(
                        shape, new_left, shape.top, shape.width, shape.height
                    )
                    changes += 1
                except Exception as exc:
                    logger.debug('_mirror_layout_placeholders: %s', exc)

        return changes

    def _set_layout_placeholder_position(
        self, shape, new_left, top, width, height
    ) -> None:
        """
        Fix 21: Write complete xfrm geometry on a layout placeholder.

        When a layout placeholder has no local xfrm (inherits all geometry
        from the slide master), python-pptx's `shape.left = value` setter
        creates a minimal xfrm with only `<a:off x="..."/>`, omitting y, cx,
        and cy. This causes PowerPoint/LibreOffice to position ALL placeholders
        at y=0 with zero height — producing title/body overlap.

        This method reads the inherited geometry, writes a complete xfrm with
        the mirrored x and all original y/cx/cy values preserved.
        """
        sp_el = shape._element
        sp_pr = sp_el.find(f'{{{P_NS}}}spPr')
        if sp_pr is None:
            sp_pr = etree.SubElement(sp_el, f'{{{P_NS}}}spPr')

        # Get or create xfrm
        xfrm = sp_pr.find(f'{{{A_NS}}}xfrm')
        if xfrm is None:
            xfrm = etree.SubElement(sp_pr, f'{{{A_NS}}}xfrm')
            # Insert xfrm as the first child of spPr
            sp_pr.insert(0, xfrm)

        # Get or create off element
        off = xfrm.find(f'{{{A_NS}}}off')
        if off is None:
            off = etree.SubElement(xfrm, f'{{{A_NS}}}off')

        # Get or create ext element
        ext = xfrm.find(f'{{{A_NS}}}ext')
        if ext is None:
            ext = etree.SubElement(xfrm, f'{{{A_NS}}}ext')

        # Write all four position values
        off.set('x', str(int(new_left)))
        off.set('y', str(int(top) if top is not None else 0))
        ext.set('cx', str(int(width) if width is not None else 0))
        ext.set('cy', str(int(height) if height is not None else 0))

        logger.debug(
            'Fix 21: wrote full xfrm on layout PH "%s": x=%s y=%s cx=%s cy=%s',
            getattr(shape, 'name', '?'),
            off.get('x'), off.get('y'), ext.get('cx'), ext.get('cy'),
        )

    def _swap_two_column_placeholders(self, layout, slide_width: int) -> int:
        """
        Swap the horizontal positions of the two content-area placeholders
        in a two-column layout.

        Identifies the left and right column placeholders by comparing their
        X positions, then swaps using the mirror formula so each ends up on
        the opposite side at the correct RTL position.

        The title placeholder (idx=0) is mirrored separately, not swapped.

        Returns:
            Count of position changes applied.
        """
        changes = 0
        content_placeholders = []
        title_placeholders = []

        for shape in layout.placeholders:
            ph_info = get_placeholder_info(shape)
            if ph_info is None:
                continue
            ph_type, ph_idx = ph_info
            if ph_type in _TITLE_PH_TYPES or ph_idx == 0:
                title_placeholders.append(shape)
            else:
                content_placeholders.append(shape)

        # Mirror title placeholders individually
        for shape in title_placeholders:
            try:
                left, width = shape.left, shape.width
                if left is None or width is None:
                    continue
                new_left = mirror_x(left, width, slide_width)
                if bounds_check_emu(new_left, slide_width) and abs(new_left - left) >= _POSITION_TOLERANCE_EMU:
                    shape.left = new_left
                    changes += 1
            except Exception as exc:
                logger.debug('title mirror: %s', exc)

        # Sort content placeholders left-to-right
        content_placeholders.sort(key=lambda s: getattr(s, 'left', 0) or 0)

        if len(content_placeholders) >= 2:
            # Swap the leftmost and rightmost
            left_ph = content_placeholders[0]
            right_ph = content_placeholders[-1]
            try:
                new_x_left, new_x_right = swap_positions(
                    left_ph.left, left_ph.width,
                    right_ph.left, right_ph.width,
                    slide_width,
                )
                left_ph.left = clamp_emu(new_x_left, slide_width)
                right_ph.left = clamp_emu(new_x_right, slide_width)
                changes += 2
            except Exception as exc:
                logger.warning('_swap_two_column_placeholders: %s', exc)

            # Mirror any remaining content placeholders individually
            for shape in content_placeholders[1:-1]:
                try:
                    left, width = shape.left, shape.width
                    if left is None or width is None:
                        continue
                    new_left = mirror_x(left, width, slide_width)
                    if bounds_check_emu(new_left, slide_width):
                        shape.left = new_left
                        changes += 1
                except Exception as exc:
                    logger.debug('extra content mirror: %s', exc)

        elif len(content_placeholders) == 1:
            # Single content placeholder — just mirror it
            shape = content_placeholders[0]
            try:
                left, width = shape.left, shape.width
                if left is not None and width is not None:
                    new_left = mirror_x(left, width, slide_width)
                    if bounds_check_emu(new_left, slide_width):
                        shape.left = new_left
                        changes += 1
            except Exception as exc:
                logger.debug('single content mirror: %s', exc)

        return changes

    def _mirror_shape_position(self, shape, slide_width_emu: int) -> bool:
        """
        Mirror a shape's X position using new_x = slide_width - (left + width).

        Args:
            shape: python-pptx Shape.
            slide_width_emu: Slide width in EMU.

        Returns:
            True if the position was changed, False otherwise.
        """
        try:
            left = shape.left
            width = shape.width
            if left is None or width is None:
                return False
            new_left = mirror_x(left, width, slide_width_emu)
            if not bounds_check_emu(new_left, slide_width_emu):
                return False
            if abs(new_left - left) < _POSITION_TOLERANCE_EMU:
                return False
            shape.left = new_left
            return True
        except Exception as exc:
            logger.debug('_mirror_shape_position: %s', exc)
            return False

    def _swap_column_placeholders(
        self, layout, left_idx: int, right_idx: int, slide_width_emu: int
    ) -> bool:
        """
        Swap two placeholders' X positions given their placeholder indices.

        Args:
            layout: python-pptx SlideLayout.
            left_idx: Placeholder idx of the left column.
            right_idx: Placeholder idx of the right column.
            slide_width_emu: Slide width in EMU.

        Returns:
            True if the swap was applied, False otherwise.
        """
        left_ph = None
        right_ph = None
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == left_idx:
                left_ph = ph
            elif ph.placeholder_format.idx == right_idx:
                right_ph = ph

        if left_ph is None or right_ph is None:
            return False

        try:
            new_x_left, new_x_right = swap_positions(
                left_ph.left, left_ph.width,
                right_ph.left, right_ph.width,
                slide_width_emu,
            )
            left_ph.left = clamp_emu(new_x_left, slide_width_emu)
            right_ph.left = clamp_emu(new_x_right, slide_width_emu)
            return True
        except Exception as exc:
            logger.warning('_swap_column_placeholders(%d, %d): %s', left_idx, right_idx, exc)
            return False


# ─────────────────────────────────────────────────────────────────────────────
# SlideContentTransformer — Phase 3
# ─────────────────────────────────────────────────────────────────────────────

class SlideContentTransformer:
    """
    Transforms individual content slides for RTL (Phase 3).

    This phase runs AFTER MasterLayoutTransformer has finished so that layout
    placeholder positions are already in their final RTL coordinates.

    For each slide:
    - PLACEHOLDER shapes: remove local position overrides (spPr xfrm) so
      they re-inherit the updated positions from the transformed layout.
    - NON-PLACEHOLDER (freeform) shapes: mirror X position explicitly.
    - All text: apply translated Arabic text with RTL paragraph properties.
    - Tables: reverse column order.
    - Charts: reverse category axis and mirror legend.
    """

    def __init__(
        self,
        presentation,
        template_registry=None,
        layout_classifications: Optional[Dict] = None,
        translations: Optional[Dict[str, str]] = None,
    ):
        """
        Args:
            presentation: python-pptx Presentation object.
            template_registry: Optional TemplateRegistry with layout rules.
            layout_classifications: Optional dict {slide_number: layout_type_str}
                from LayoutAnalyzer. If absent, layout type is read directly from
                the slide layout XML.
            translations: Optional dict {english_text: arabic_text} translation map.
                If absent, no text replacement is performed.
        """
        self.prs = presentation
        self.template_registry = template_registry
        self.layout_classifications = layout_classifications or {}
        self.translations = translations or {}
        self._slide_width = int(presentation.slide_width)
        self._slide_height = int(presentation.slide_height)

        # Pre-build lowercase translation index for O(1) case-insensitive lookups
        self._translations_lower: Dict[str, str] = {}
        for key, val in self.translations.items():
            lower_key = key.strip().lower()
            if lower_key not in self._translations_lower:
                self._translations_lower[lower_key] = val

    # ─────────────────────────────────────────────────────────────────────
    # Public entry point
    # ─────────────────────────────────────────────────────────────────────

    def transform_all_slides(self) -> TransformReport:
        """
        Transform all content slides.

        Returns:
            TransformReport summarising all changes made.
        """
        report = TransformReport(phase='slide')
        for slide_idx, slide in enumerate(self.prs.slides):
            slide_number = slide_idx + 1
            try:
                count = self._transform_slide(slide, slide_number)
                report.add('slide_transformed', count)
            except Exception as exc:
                report.error(f'slide[{slide_number}]: {exc}')
        return report

    # ─────────────────────────────────────────────────────────────────────
    # Slide-level transformation
    # ─────────────────────────────────────────────────────────────────────

    def _transform_slide(self, slide, slide_number: int) -> int:
        """
        Transform a single slide:

        1. For PLACEHOLDER shapes: remove local position overrides so the
           shape inherits from the (now RTL-transformed) layout.
        2. For NON-PLACEHOLDER shapes: mirror X position.
        3. Insert translated Arabic text and set RTL paragraph properties.
        4. Transform tables (reverse column order).
        5. Transform charts (reverse category axis, mirror legend).

        Returns:
            Total count of changes made.
        """
        changes = 0
        layout = slide.slide_layout
        layout_type = layout._element.get('type', 'cust')

        # Override with classifier result if available
        if slide_number in self.layout_classifications:
            layout_type = self.layout_classifications[slide_number]

        all_shapes = self._collect_all_shapes(slide.shapes)

        # ── Pre-mirroring: detect split-panel layout BEFORE individual mirroring ──
        # If a split-panel is detected (image on one side, content on other),
        # swap panels as units and track which shapes were handled.
        panel_handled_shapes = set()  # shape ids that were repositioned by panel swap
        panel_swap_delta = self._pre_mirror_split_panel_swap(
            all_shapes, slide_number, panel_handled_shapes
        )
        changes += panel_swap_delta

        # ── Pre-mirroring: detect geographic/map overlay slides ───────────────
        # Shapes overlaid on a large map image should NOT be position-mirrored
        # (they are geographically anchored). Only their text gets RTL treatment.
        changes += self._exempt_map_overlay_shapes(
            all_shapes, slide_number, panel_handled_shapes
        )

        for shape in all_shapes:
            try:
                # Fix 5: Handle group shapes as a single unit
                if hasattr(shape, 'shapes') and not getattr(shape, 'is_placeholder', False):
                    # Mirror the group's position as a unit (unless already handled by panel swap)
                    if id(shape) not in panel_handled_shapes:
                        if self._should_mirror_shape(shape, layout_type):
                            if self._mirror_freeform_shape(shape, self._slide_width):
                                changes += 1
                    # Process text in group children (no position mirroring)
                    for child in self._collect_text_shapes_from_group(shape):
                        changes += self._apply_translation(child, self.translations)
                        changes += self._set_rtl_alignment_unconditional(child)
                    # Check for tables/charts inside group
                    for child in self._collect_text_shapes_from_group(shape):
                        if getattr(child, 'has_table', False) and child.has_table:
                            changes += self._transform_table_rtl(child)
                        if getattr(child, 'has_chart', False) and child.has_chart:
                            changes += self._transform_chart_rtl(child)
                    # Fix 4B: Directional reversal for group children (arrows, connectors)
                    if hasattr(shape, 'shapes'):
                        for child in shape.shapes:
                            changes += self._reverse_directional_shape(child)
                            changes += self._reverse_connector_direction(child)
                            changes += self._reverse_line_arrowheads(child)
                    continue

                # Fix 4A: Handle connector shapes (cxnSp)
                sp_el = shape._element
                if sp_el.tag.endswith('}cxnSp') or sp_el.tag == 'cxnSp':
                    changes += self._reverse_connector_direction(shape)
                    # Also mirror position
                    if self._should_mirror_shape(shape, layout_type):
                        if self._mirror_freeform_shape(shape, self._slide_width):
                            changes += 1
                    continue  # Connectors don't need text treatment

                if getattr(shape, 'is_placeholder', False):
                    # Placeholder: remove local position override to inherit
                    # from the transformed layout
                    removed = self._remove_local_position_override(shape, layout)
                    if removed:
                        changes += 1
                        logger.debug('Slide %d: removed position override on %s',
                                     slide_number, shape.name)
                elif id(shape) in panel_handled_shapes:
                    # Shape was already repositioned by pre-mirror panel swap
                    pass
                else:
                    # Freeform shape: determine whether to mirror
                    if self._should_mirror_shape(shape, layout_type):
                        if self._mirror_freeform_shape(shape, self._slide_width):
                            changes += 1
                    # Fix 10: Reverse directional shapes (chevrons, arrows)
                    changes += self._reverse_directional_shape(shape)

                # Apply text translation and RTL properties
                if getattr(shape, 'has_text_frame', False) and shape.has_text_frame:
                    changes += self._apply_translation(shape, self.translations)
                    # Fix 1: Unconditional RTL alignment on ALL paragraphs
                    changes += self._set_rtl_alignment_unconditional(shape)

                # Fix 8: Validate text box width for short text (page numbers)
                self._validate_textbox_width(shape)

                # Transform tables
                if getattr(shape, 'has_table', False) and shape.has_table:
                    changes += self._transform_table_rtl(shape)

                # Transform charts
                if getattr(shape, 'has_chart', False) and shape.has_chart:
                    changes += self._transform_chart_rtl(shape)

            except Exception as exc:
                logger.warning('Slide %d shape "%s": %s', slide_number,
                               getattr(shape, 'name', '?'), exc)

        # Fix 11: Fix wrap="none" text boxes for Arabic text
        for shape in all_shapes:
            try:
                if getattr(shape, 'has_text_frame', False):
                    changes += self._fix_wrap_none_for_arabic(shape)
            except Exception as exc:
                logger.debug('Fix 11 on slide %d: %s', slide_number, exc)

        # Fix 12: Resolve title-body vertical overlap for RTL
        changes += self._fix_title_body_overlap(all_shapes, slide_number)

        # Fix 22: Apply normAutofit to all Arabic-containing placeholder text frames.
        # Arabic text is 20-40% taller/wider than English at the same font size.
        # Without autofit, text overflows the placeholder box and overlaps adjacent shapes.
        changes += self._apply_arabic_autofit(all_shapes, slide_number)

        # Fix 14: Right-anchor cover title text for photo-background covers
        changes += self._fix_cover_title_anchor(all_shapes, slide_number)

        # Fix 15: Mirror split-panel layouts (image-left/text-right → image-right/text-left)
        changes += self._mirror_split_panel_layout(all_shapes, slide_number)

        # Fix 16: Reverse timeline alternation pattern for RTL
        changes += self._reverse_timeline_alternation(all_shapes, slide_number)

        # Fix 17: Reverse logo row ordering for RTL
        changes += self._reverse_logo_row_order(all_shapes, slide_number)

        # Fix 18: Center text in circular/bounded container shapes
        changes += self._center_text_in_container_shapes(all_shapes, slide_number)

        # Fix 19: Ensure bidi base direction on mixed Arabic/English titles
        changes += self._fix_bidi_base_direction(all_shapes, slide_number)

        # Fix 20: Move slide-number badge to top-left in RTL output
        changes += self._reposition_slide_number_badge(all_shapes, slide_number)

        # Fix 9: Collision detection (log warnings for overlapping shapes)
        self._detect_collisions(all_shapes, slide_number)

        # v1.1.3: Per-slide translation coverage telemetry
        translated_shapes = 0
        total_text_shapes = 0
        for s in all_shapes:
            if getattr(s, 'has_text_frame', False) and any(
                p.text.strip() for p in s.text_frame.paragraphs if p.text
            ):
                total_text_shapes += 1
                if any(has_arabic(p.text) for p in s.text_frame.paragraphs if p.text):
                    translated_shapes += 1
        if total_text_shapes > 0 and translated_shapes == 0:
            logger.error(
                'Slide %d: 0/%d text shapes translated — possible extraction miss',
                slide_number, total_text_shapes,
            )

        return changes

    def _collect_all_shapes(self, shapes) -> List:
        """
        Collect top-level shapes from the slide. Group shapes are included
        as single units — their children are NOT flattened here.

        Fix 5: Groups must be mirrored as a single unit (position only).
        Child shapes within groups get text treatment but NOT position mirroring,
        because child coordinates are relative to the group's local space.

        Args:
            shapes: python-pptx ShapeCollection (slide.shapes).

        Returns:
            List of top-level shapes (groups included as single items).
        """
        result = []
        for shape in shapes:
            result.append(shape)
        return result

    def _collect_text_shapes_from_group(self, group_shape) -> List:
        """
        Fix 5: Recursively collect all text-bearing shapes inside a group
        for text translation and RTL alignment (but NOT position mirroring).
        """
        result = []
        try:
            for child in group_shape.shapes:
                if getattr(child, 'has_text_frame', False) and child.has_text_frame:
                    result.append(child)
                if hasattr(child, 'shapes'):
                    result.extend(self._collect_text_shapes_from_group(child))
        except Exception:
            pass
        return result

    # ─────────────────────────────────────────────────────────────────────
    # Placeholder position inheritance
    # ─────────────────────────────────────────────────────────────────────

    def _remove_local_position_override(self, shape, layout) -> bool:
        """
        Remove explicit x/y/width/height from a placeholder shape's spPr/xfrm
        element so it re-inherits its position from the (already RTL-transformed)
        layout placeholder.

        Round 4 refinement: Before removing, check if the inherited position
        would overlap with large non-placeholder shapes on the slide. If so,
        mirror the placeholder's position explicitly instead of removing the
        xfrm (to avoid title-picture overlaps like R6_17 slide 1).

        Only removes the xfrm element if:
        1. The shape is a placeholder (has ph element).
        2. There is a matching placeholder in the layout (same idx).
        3. The shape has a local spPr/xfrm element (meaning it overrides position).
        4. The inherited position would NOT overlap with large freeform shapes.

        Args:
            shape: python-pptx Shape on a slide.
            layout: python-pptx SlideLayout for this slide.

        Returns:
            True if the position was modified (either removed or mirrored).
        """
        try:
            sp_el = shape._element

            # Verify shape is a placeholder
            ph_info = get_placeholder_info_from_xml(sp_el)
            if ph_info is None:
                return False
            _, ph_idx = ph_info

            # Verify layout has a matching placeholder
            layout_ph = None
            for lph in layout.placeholders:
                if lph.placeholder_format.idx == ph_idx:
                    layout_ph = lph
                    break

            if layout_ph is None:
                return False

            # Check if local xfrm exists
            sp_pr = sp_el.find(f'{{{P_NS}}}spPr')
            if sp_pr is None:
                sp_pr = sp_el.find(f'{{{A_NS}}}spPr')
            if sp_pr is None:
                return False

            xfrm = sp_pr.find(f'{{{A_NS}}}xfrm')
            if xfrm is None:
                return False  # No local position override to remove

            # Round 4: Check if inheriting layout position would cause overlap
            # with non-placeholder shapes on this slide.
            #
            # Two-tier check:
            # Tier 1: Large shapes (>30% slide width) — any overlap >20% of PH area
            # Tier 2: Image shapes (pics) at top of slide — check if logo/image
            #         would collide with the inherited title position after both
            #         are mirrored. This catches logo-title collisions (R6_17 s8).
            layout_left = layout_ph.left
            layout_width = layout_ph.width
            layout_top = layout_ph.top
            layout_height = layout_ph.height

            if layout_left is not None and layout_width is not None:
                # Check overlap with freeform shapes
                has_dangerous_overlap = False
                try:
                    sp_tree = sp_el.getparent()
                    if sp_tree is not None:
                        lx, ly = int(layout_left), int(layout_top or 0)
                        lw, lh = int(layout_width), int(layout_height or 0)

                        for sibling in sp_tree:
                            if sibling is sp_el:
                                continue
                            tag = sibling.tag
                            if not (tag.endswith('}sp') or tag.endswith('}pic')):
                                continue
                            sib_ph = sibling.find(f'.//{{{P_NS}}}ph')
                            if sib_ph is not None:
                                continue
                            sib_xfrm = sibling.find(f'.//{{{A_NS}}}xfrm')
                            if sib_xfrm is None:
                                continue
                            sib_off = sib_xfrm.find(f'{{{A_NS}}}off')
                            sib_ext = sib_xfrm.find(f'{{{A_NS}}}ext')
                            if sib_off is None or sib_ext is None:
                                continue
                            try:
                                sx = int(sib_off.get('x', 0))
                                sy = int(sib_off.get('y', 0))
                                sw = int(sib_ext.get('cx', 0))
                                sh = int(sib_ext.get('cy', 0))
                            except (ValueError, TypeError):
                                continue

                            # --- Tier 1: Large shapes (>30% width) ---
                            if sw >= self._slide_width * 0.30:
                                if (lx < sx + sw and lx + lw > sx and
                                    ly < sy + sh and ly + lh > sy):
                                    overlap_w = min(lx + lw, sx + sw) - max(lx, sx)
                                    overlap_h = min(ly + lh, sy + sh) - max(ly, sy)
                                    ph_area = max(lw * lh, 1)
                                    overlap_area = overlap_w * overlap_h
                                    if overlap_area > ph_area * 0.20:
                                        has_dangerous_overlap = True
                                        break

                            # --- Tier 2: Image shapes (logos) that will be mirrored ---
                            # Only for pic elements (images) in the title zone (top 25%)
                            is_pic = tag.endswith('}pic')
                            in_title_zone = (sy < self._slide_height * 0.25)
                            if is_pic and in_title_zone and sw < self._slide_width * 0.30:
                                # Predict where this image will end up after freeform mirror
                                mirrored_sx = self._slide_width - sx - sw
                                # Check if mirrored image position overlaps with
                                # inherited layout position
                                if (lx < mirrored_sx + sw and lx + lw > mirrored_sx and
                                    ly < sy + sh and ly + lh > sy):
                                    overlap_w = min(lx + lw, mirrored_sx + sw) - max(lx, mirrored_sx)
                                    overlap_h = min(ly + lh, sy + sh) - max(ly, sy)
                                    if overlap_w > 0 and overlap_h > 0:
                                        # For logos, any overlap with the title is dangerous
                                        has_dangerous_overlap = True
                                        logger.debug('Logo-title overlap predicted: image at %d->%d, title at %d',
                                                     sx, mirrored_sx, lx)
                                        break

                except Exception as exc:
                    logger.debug('overlap check: %s', exc)

                if has_dangerous_overlap:
                    # Don't remove xfrm — adjust placeholder position/size to avoid overlap.
                    # Strategy depends on whether mirroring alone resolves the collision:
                    current_left = shape.left
                    current_width = shape.width
                    if current_left is not None and current_width is not None:
                        new_left = mirror_x(current_left, current_width, self._slide_width)
                        new_left = clamp_emu(new_left, self._slide_width)

                        # Collect ALL overlapping logos (there may be one on each side)
                        # and compute a safe title zone between them.
                        overlapping_logos = []  # list of (mirrored_x, width) tuples
                        margin = 250000  # 0.27in — generous to prevent visual clipping
                        try:
                            for sibling in sp_tree:
                                if sibling is sp_el:
                                    continue
                                tag = sibling.tag
                                if not tag.endswith('}pic'):
                                    continue
                                sib_ph2 = sibling.find(f'.//{{{P_NS}}}ph')
                                if sib_ph2 is not None:
                                    continue
                                sib_xfrm2 = sibling.find(f'.//{{{A_NS}}}xfrm')
                                if sib_xfrm2 is None:
                                    continue
                                sib_off2 = sib_xfrm2.find(f'{{{A_NS}}}off')
                                sib_ext2 = sib_xfrm2.find(f'{{{A_NS}}}ext')
                                if sib_off2 is None or sib_ext2 is None:
                                    continue
                                try:
                                    s2x = int(sib_off2.get('x', 0))
                                    s2y = int(sib_off2.get('y', 0))
                                    s2w = int(sib_ext2.get('cx', 0))
                                    s2h = int(sib_ext2.get('cy', 0))
                                except (ValueError, TypeError):
                                    continue
                                if s2y > self._slide_height * 0.25:
                                    continue
                                if s2w >= self._slide_width * 0.30:
                                    continue  # Not a logo — too large
                                # Predict mirrored logo position
                                m_s2x = self._slide_width - s2x - s2w
                                # Check vertical overlap with title zone
                                if (int(layout_top or 0) < s2y + s2h and
                                    int(layout_top or 0) + int(layout_height or 0) > s2y):
                                    overlapping_logos.append((m_s2x, s2w))
                        except Exception:
                            pass

                        if overlapping_logos:
                            # Compute safe zone: the title must fit between all logos
                            # Find the rightmost left-side logo and leftmost right-side logo
                            title_center = new_left + current_width // 2
                            safe_left = new_left
                            safe_right = new_left + current_width

                            for (logo_x, logo_w) in overlapping_logos:
                                logo_right = logo_x + logo_w
                                logo_center = logo_x + logo_w // 2
                                # Check if logo actually overlaps horizontally with title
                                if not (new_left < logo_x + logo_w and new_left + current_width > logo_x):
                                    continue
                                if logo_center < title_center:
                                    # Logo on LEFT — push title left edge right
                                    safe_left = max(safe_left, logo_right + margin)
                                else:
                                    # Logo on RIGHT — pull title right edge left
                                    safe_right = min(safe_right, logo_x - margin)

                            safe_width = safe_right - safe_left
                            if safe_width >= current_width * 0.50 and safe_width > 0:
                                shape.left = safe_left
                                shape.width = safe_width
                                # Enable auto-fit (font scaling) when width is significantly reduced.
                                # Arabic text is typically wider than English for equivalent content,
                                # so title boxes that were wide enough for English may overflow in Arabic.
                                shrink_ratio = safe_width / current_width
                                if shrink_ratio < 0.85:
                                    self._enable_autofit(shape, shrink_ratio)
                                logger.debug('Placeholder %d: shrunk to avoid %d logos %d,%d -> %d,%d',
                                             ph_idx, len(overlapping_logos),
                                             current_left, current_width,
                                             safe_left, safe_width)
                            else:
                                # Shrinkage too aggressive — just mirror position
                                shape.left = new_left
                                logger.debug('Placeholder %d: explicit mirror (shrink too aggressive) %d -> %d',
                                             ph_idx, current_left, new_left)
                        else:
                            shape.left = new_left
                            logger.debug('Placeholder %d: explicit mirror (overlap risk) %d -> %d',
                                         ph_idx, current_left, new_left)
                    return True

            # ── Size-divergence guard ──────────────────────────────────
            # Before removing xfrm (which makes the shape inherit the layout's
            # geometry), check whether the slide-level shape's width/height
            # differ significantly from the layout placeholder's dimensions.
            # If they do, removing xfrm would collapse (or inflate) the shape
            # to the layout's incompatible box.  Instead, mirror the shape's
            # position explicitly while PRESERVING its original width/height.
            #
            # Symmetric threshold: if either dimension diverges by >30% in
            # EITHER direction (larger or smaller), keep explicit geometry.
            # ─────────────────────────────────────────────────────────────
            _SIZE_DIVERGENCE_THRESHOLD = 0.30  # 30% in either direction

            current_left = shape.left
            current_width = shape.width
            current_height = shape.height

            # ── Fallback: layout has zero/None dimensions ────────────────
            # If we can't read the layout's geometry, it's unsafe to inherit
            # from it.  Mirror the slide shape's position explicitly.
            if (layout_left is None or layout_width is None
                    or layout_height is None
                    or int(layout_width or 0) <= 0
                    or int(layout_height or 0) <= 0):
                if current_left is not None and current_width is not None:
                    slide_w = int(current_width)
                    new_left = self._slide_width - int(current_left) - slide_w
                    new_left = max(0, min(new_left, self._slide_width - slide_w))
                    shape.left = new_left
                    logger.debug(
                        'Placeholder %d: layout dims missing/zero — '
                        'explicit mirror %d→%d, keeping size.',
                        ph_idx, int(current_left), new_left,
                    )
                return True

            # ── Normal path: compare slide vs layout dimensions ──────────
            if (current_left is not None and current_width is not None
                    and current_height is not None):

                slide_w = int(current_width)
                slide_h = int(current_height)
                lay_w = int(layout_width)
                lay_h = int(layout_height)

                width_ratio = slide_w / lay_w
                height_ratio = slide_h / lay_h

                width_diverges = abs(width_ratio - 1.0) > _SIZE_DIVERGENCE_THRESHOLD
                height_diverges = abs(height_ratio - 1.0) > _SIZE_DIVERGENCE_THRESHOLD

                if width_diverges or height_diverges:
                    # Slide shape dimensions differ significantly from layout
                    # — do NOT remove xfrm.  Mirror position explicitly,
                    # keep original width/height.
                    new_left = self._slide_width - int(current_left) - slide_w
                    new_left = max(0, min(new_left, self._slide_width - slide_w))
                    shape.left = new_left
                    # Preserve top, width, height — do NOT touch them
                    logger.debug(
                        'Placeholder %d: size-divergence guard — '
                        'slide %dx%d vs layout %dx%d (ratio %.2fx%.2f). '
                        'Mirroring position %d→%d, keeping size.',
                        ph_idx, slide_w, slide_h, lay_w, lay_h,
                        width_ratio, height_ratio,
                        int(current_left), new_left,
                    )
                    return True

            # Dimensions are compatible — safe to remove xfrm and inherit
            sp_pr.remove(xfrm)
            return True

        except Exception as exc:
            logger.debug('_remove_local_position_override: %s', exc)
            return False

    # ─────────────────────────────────────────────────────────────────────
    # Freeform shape mirroring
    # ─────────────────────────────────────────────────────────────────────

    def _mirror_freeform_shape(self, shape, slide_width_emu: int) -> bool:
        """
        Mirror a non-placeholder shape's X position.

        Applies the standard RTL mirror formula: new_x = slide_width - (x + w).
        Skips shapes with negligible position change to avoid touching
        centred shapes.

        Round 2 fix: Clamp result to slide boundaries to prevent edge clipping.

        Args:
            shape: python-pptx Shape (non-placeholder).
            slide_width_emu: Slide width in EMU.

        Returns:
            True if the shape was mirrored, False otherwise.
        """
        try:
            left = shape.left
            width = shape.width
            if left is None or width is None:
                return False

            new_left = mirror_x(left, width, slide_width_emu)

            # Round 3: Tighter clamping — prevent edge clipping
            # Allow minimal bleed (-100000 EMU ≈ 0.11") for intentional edge shapes
            MIN_BLEED = -100_000
            MAX_BLEED = 100_000
            if new_left < MIN_BLEED:
                new_left = MIN_BLEED
            if width is not None and new_left + width > slide_width_emu + MAX_BLEED:
                new_left = slide_width_emu + MAX_BLEED - width

            if not bounds_check_emu(new_left, slide_width_emu):
                logger.debug('_mirror_freeform_shape OOB: %s new_left=%d', shape.name, new_left)
                return False

            if abs(new_left - left) < _POSITION_TOLERANCE_EMU:
                return False  # Negligible change — likely centred shape

            shape.left = new_left

            # Fix 4: Ensure no flipH/flipV was introduced on the shape
            self._ensure_no_content_flip(shape)

            return True
        except Exception as exc:
            logger.debug('_mirror_freeform_shape: %s', exc)
            return False

    def _ensure_no_content_flip(self, shape) -> None:
        """
        Fix 4: Handle flipH/flipV on shapes after mirroring.

        For CONNECTORS and LINES: TOGGLE flipH (mirroring reverses the
        direction of the line slope — if it went top-left to bottom-right,
        after mirror it should go top-right to bottom-left).

        For all OTHER shapes (images, text boxes): REMOVE flipH/flipV
        to avoid content flipping (maps, photos must not be mirrored).

        flipV is always removed (vertical orientation is preserved in RTL).
        """
        try:
            sp_el = shape._element
            tag = sp_el.tag.split('}')[-1] if '}' in sp_el.tag else sp_el.tag

            # --- Connector/Line detection (P0 fix — 3-model consensus) ---
            # 1. Primary: XML tag (cxnSp = connector shape)
            is_connector = (tag == 'cxnSp')

            # 2. Secondary: For sp elements, check preset geometry for line presets
            LINE_PRESETS = {'line', 'straightConnector1', 'bentConnector2',
                           'bentConnector3', 'bentConnector4', 'bentConnector5',
                           'curvedConnector2', 'curvedConnector3',
                           'curvedConnector4', 'curvedConnector5'}
            if not is_connector and tag == 'sp':
                try:
                    prst_geom = sp_el.find(f'.//{{{A_NS}}}prstGeom')
                    if prst_geom is not None:
                        prst_val = prst_geom.get('prst', '')
                        if prst_val in LINE_PRESETS or 'Connector' in prst_val:
                            is_connector = True
                except Exception:
                    pass

            # 3. Tertiary: Aspect ratio fallback for shapes without text
            if not is_connector:
                try:
                    w = int(shape.width or 0)
                    h = int(shape.height or 0)
                    has_text = getattr(shape, 'has_text_frame', False)
                    if h > 0 and w > 0 and not has_text:
                        aspect = max(w, h) / min(w, h)
                        # Tightened threshold: aspect > 20, minor < 50000 EMU
                        if aspect > 20 and min(w, h) < 50000:
                            is_connector = True
                    elif (h == 0 or w == 0) and not has_text:
                        # Perfectly straight line (zero in one dimension)
                        is_connector = True
                except (TypeError, ValueError, ZeroDivisionError):
                    pass

            for xfrm in sp_el.iter(f'{{{A_NS}}}xfrm'):
                if is_connector:
                    # TOGGLE flipH for connectors: mirroring reverses the slope
                    current_flipH = (xfrm.get('flipH', '0') or '0').lower()
                    if current_flipH in ('1', 'true'):
                        # Was flipped, mirror un-flips it
                        if 'flipH' in xfrm.attrib:
                            del xfrm.attrib['flipH']
                    else:
                        # Was NOT flipped, mirror introduces flip
                        xfrm.set('flipH', '1')
                    # Preserve flipV for connectors (Opus recommendation)
                    # Horizontal mirror doesn't affect vertical flip
                else:
                    # REMOVE flipH/flipV for content shapes
                    if xfrm.get('flipH'):
                        del xfrm.attrib['flipH']
                    # Remove flipV for content shapes only
                    if xfrm.get('flipV'):
                        del xfrm.attrib['flipV']
        except Exception:
            pass

    def _should_mirror_shape(self, shape, layout_type: str) -> bool:
        """
        Determine if a non-placeholder shape should be mirrored.

        Decision rules:
        - Full-width background shapes (width > 90% of slide): preserve (no mirror).
        - Logo images (small picture, < 20% width): mirror — handled by master,
          but any slide-level logos are also mirrored here.
        - 'blank' layout: mirror all freeform shapes.
        - 'secHead' layout: preserve all freeform shapes (decorative intent).
        - 'title' layout: preserve freeform shapes (brand elements stay put).
        - All other layouts: mirror content shapes, preserve full-width decoratives.

        Args:
            shape: python-pptx Shape.
            layout_type: ST_SlideLayoutType string.

        Returns:
            True if the shape should be mirrored.
        """
        try:
            slide_width = self._slide_width

            # Layout-specific rules (checked first — before positional checks)
            # secHead/title layouts are typically centered/symmetrical by design;
            # never mirror their shapes.
            if layout_type in ('secHead', 'title'):
                return False

            # Full-width background shapes: never mirror
            width = getattr(shape, 'width', None)
            if width is not None and width > slide_width * 0.90:
                return False

            # Footer-zone shapes — ALWAYS mirror regardless of layout
            top = getattr(shape, 'top', None)
            height = getattr(shape, 'height', None)
            if top is not None and height is not None:
                bottom = top + height
                slide_height = self._slide_height
                if bottom > slide_height * 0.88:
                    return True  # Footer zone — always mirror

            # By default, mirror content-sized shapes
            return True

        except Exception:
            return True  # Default to mirror on error

    # ─────────────────────────────────────────────────────────────────────
    # Text translation and RTL formatting
    # ─────────────────────────────────────────────────────────────────────

    def _fuzzy_lookup_translation(
        self, text: str, translations: Dict[str, str]
    ) -> Optional[str]:
        """
        Round 4: Flexible translation lookup with fallbacks.

        Order of attempts:
        1. Exact match on full text
        2. Exact match on stripped text
        3. Case-insensitive exact match
        4. Longest prefix match (>80% of text length must match)

        Returns:
            Arabic translation string, or None.
        """
        if not text or not text.strip():
            return None

        stripped = text.strip()

        # 1. Exact match
        result = translations.get(text) or translations.get(stripped)
        if result:
            return result

        # 2. Case-insensitive match (O(1) via pre-built index)
        text_lower = stripped.lower()
        result = self._translations_lower.get(text_lower)
        if result:
            return result

        # 3. Longest prefix match (for long text that may differ in suffix)
        #    Only triggered for text >40 chars where exact/case-insensitive failed
        if len(stripped) > 40:
            prefix_40 = stripped[:40]
            best_match = None
            best_match_len = 0
            for key, val in translations.items():
                if len(key) < 40:
                    continue
                # Quick prefix check before detailed comparison
                if key[:40] != prefix_40:
                    continue
                # Count matching prefix length
                common_len = min(len(key), len(stripped))
                match_len = 40
                for i in range(40, common_len):
                    if key[i] == stripped[i]:
                        match_len += 1
                    else:
                        break
                # Accept if >80% of the shorter text matches
                min_len = min(len(key), len(stripped))
                if match_len > min_len * 0.80:
                    if match_len > best_match_len:
                        best_match_len = match_len
                        best_match = val

            if best_match:
                return best_match

        return None

    def _apply_translation(self, shape, translations: Dict[str, str]) -> int:
        """
        Replace English text with Arabic translation in a shape's text frame,
        preserving run-level formatting (bold, italic, font size, color).

        After text replacement, sets RTL direction and appropriate alignment
        on every paragraph that contains Arabic text.

        Algorithm:
        1. For each paragraph, gather its runs' texts.
        2. Build the full paragraph text and look it up in the translation map.
        3. If a match is found: replace the first run's text with the Arabic
           translation, clear subsequent runs in the paragraph (to avoid
           duplicating content while preserving first-run formatting).
        4. Set rtl='1' on the paragraph's pPr.
        5. Determine alignment:
           - ctrTitle → 'ctr' (brand intent: centred titles stay centred)
           - footer/sldNum/dt → 'l' (footers left)
           - all other Arabic paragraphs → 'r'

        Round 4: Uses fuzzy matching for case-insensitive and prefix-based lookups.

        Args:
            shape: python-pptx Shape with a text frame.
            translations: Dict mapping English text → Arabic text.

        Returns:
            Count of paragraphs modified.
        """
        if not translations:
            return 0

        changes = 0
        ph_info = get_placeholder_info(shape)
        ph_type = ph_info[0] if ph_info else None

        try:
            tf = shape.text_frame
            for para in tf.paragraphs:
                para_text = para.text
                if not para_text or not para_text.strip():
                    continue

                # Look up translation (Round 4: with fuzzy matching)
                arabic_text = self._fuzzy_lookup_translation(
                    para_text, translations
                )

                if arabic_text:
                    # Replace text in runs — preserve first run's formatting
                    runs = para.runs
                    if runs:
                        runs[0].text = arabic_text
                        # Clear subsequent runs
                        for run in runs[1:]:
                            run.text = ''
                    else:
                        # No runs — paragraph has field/special elements
                        # Create a new run with the translation
                        p_elem = para._p
                        r_elem = etree.SubElement(p_elem, f'{{{A_NS}}}r')
                        t_elem = etree.SubElement(r_elem, f'{{{A_NS}}}t')
                        t_elem.text = arabic_text

                    # Set Arabic language on runs
                    for r_elem in para._p.findall(f'{{{A_NS}}}r'):
                        rPr = r_elem.find(f'{{{A_NS}}}rPr')
                        if rPr is not None:
                            rPr.set('lang', 'ar-SA')

                # Determine if paragraph now has Arabic content
                final_text = para.text or ''
                if not has_arabic(final_text) and not has_arabic(para_text):
                    continue

                # Apply RTL paragraph properties
                pPr = ensure_pPr(para._p)
                # Only set rtl='1' if content is actually Arabic
                if has_arabic(final_text):
                    pPr.set('rtl', '1')
                else:
                    pPr.set('rtl', '0')

                # Determine alignment
                alignment = self._compute_paragraph_alignment(final_text, ph_type)
                pPr.set('algn', alignment)

                # Set rtlCol on bodyPr
                try:
                    body_pr = tf._txBody.find(f'{{{A_NS}}}bodyPr')
                    if body_pr is not None:
                        body_pr.set('rtlCol', '1')
                except Exception:
                    pass

                changes += 1

        except Exception as exc:
            logger.warning('_apply_translation on shape "%s": %s',
                           getattr(shape, 'name', '?'), exc)

        return changes

    def _compute_paragraph_alignment(
        self, text: str, ph_type: Optional[str]
    ) -> str:
        """
        Compute the correct OOXML alignment value for a paragraph in RTL context.

        Rules (in priority order):
        1. Footer-type placeholders (ftr, sldNum, dt) → always 'l' (left).
        2. ctrTitle placeholder → 'ctr' (centred title is brand intent).
        3. Predominantly Arabic text → 'r' (right-align).
        4. Mixed/Latin text in Arabic context → 'r' (keep visual consistency).
        5. Pure LTR/numeric text → 'l'.

        Args:
            text: The current paragraph text.
            ph_type: Placeholder type string or None for freeforms.

        Returns:
            OOXML algn attribute value: 'l', 'r', or 'ctr'.
        """
        # Footer-type placeholders: check both frozenset and substring matching
        # because ph_type strings may include enum values like 'slide_number (13)'
        _footer_substrings = ('ftr', 'footer', 'sldnum', 'slide_number', 'date_time', 'date')
        if ph_type in _FOOTER_PH_TYPES or (ph_type and any(f in ph_type for f in _footer_substrings)):
            return 'l'

        # Title placeholders: center-align (preserve original design intent)
        # Match 'title (1)', 'center_title (3)', 'vertical_title (5)', 'ctrTitle'
        # but NOT 'subtitle (4)' — subtitles are body-like
        if ph_type and ('subtitle' not in ph_type.lower()) and ('title' in ph_type.lower()):
            return 'ctr'

        ratios = compute_script_ratio(text)
        arabic_ratio = ratios['arabic']
        latin_ratio = ratios['latin']

        if arabic_ratio > 0.70:
            return 'r'
        elif latin_ratio > 0.70 and not has_arabic(text):
            # Pure Latin text should remain left-aligned
            return 'l'
        else:
            # Mixed content or ambiguous — default to right in Arabic context
            return 'r'


    # ─────────────────────────────────────────────────────────────────────
    # Unconditional RTL alignment pass (Fix 1)
    # ─────────────────────────────────────────────────────────────────────

    def _set_rtl_alignment_unconditional(self, shape) -> int:
        """
        Set paragraph alignment for RTL layout on EVERY paragraph in a shape.

        CRITICAL RULE (Round 2 fix):
        - rtl='1' is ONLY set on paragraphs containing Arabic script characters.
          Setting rtl='1' on English text causes the OOXML bidi algorithm to
          reorder characters, move periods to line starts, and corrupt text.
        - For English-only paragraphs: ONLY set algn='r' (right-align).
          This achieves the visual RTL layout without corrupting text content.

        Called on every text-bearing shape on every slide AFTER
        _apply_translation().
        """
        changes = 0
        ph_info = get_placeholder_info(shape)
        ph_type = ph_info[0] if ph_info else None

        try:
            tf = shape.text_frame
            for para in tf.paragraphs:
                text = para.text or ''
                if not text.strip():
                    continue

                pPr = ensure_pPr(para._p)

                # ONLY set rtl='1' if paragraph actually contains Arabic text.
                # For English-only text, rtl='1' causes period displacement and
                # word reordering — a confirmed P0 bug from VQA Round 1.
                if has_arabic(text):
                    pPr.set('rtl', '1')
                else:
                    # Explicitly set rtl='0' to override any inherited rtl='1'
                    # from master/layout defaults
                    pPr.set('rtl', '0')

                # Compute alignment based on placeholder type and content
                alignment = self._compute_paragraph_alignment(text, ph_type)
                pPr.set('algn', alignment)

                changes += 1

            # Set rtlCol on bodyPr only if shape has Arabic content
            try:
                shape_text = tf.text or ''
                body_pr = tf._txBody.find(f'{{{A_NS}}}bodyPr')
                if body_pr is not None:
                    if has_arabic(shape_text):
                        body_pr.set('rtlCol', '1')
                    else:
                        # Don't set rtlCol for English-only shapes
                        pass
            except Exception:
                pass

        except Exception as exc:
            logger.debug('_set_rtl_alignment_unconditional on "%s": %s',
                         getattr(shape, 'name', '?'), exc)

        return changes

    # ─────────────────────────────────────────────────────────────────────
    # Table RTL transformation
    # ─────────────────────────────────────────────────────────────────────

    def _transform_table_rtl(self, shape) -> int:
        """
        Transform a table shape for RTL reading direction.

        Operations:
        1. Translate cell text using the translation map (Round 3 fix).
        2. Reverse the column widths in <a:tblGrid>.
        3. Reverse cell order within each row (deepcopy to avoid reference issues).
        4. Set rtl='1' on <a:tblPr> (table-level RTL flag).
        5. Set RTL paragraph properties on all cell text.

        Args:
            shape: python-pptx Shape with has_table == True.

        Returns:
            Count of changes made.
        """
        changes = 0
        try:
            table = shape.table
            num_cols = len(table.columns)
            if num_cols <= 1:
                return 0

            tbl_elem = table._tbl
            a_ns = A_NS

            # Round 3 Step 1: Translate cell text BEFORE reversing columns
            if self.translations:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            changes += self._translate_cell_text(cell.text_frame)

            # Step 2: Reverse column widths in <a:tblGrid>
            tbl_grid = tbl_elem.find(f'{{{a_ns}}}tblGrid')
            if tbl_grid is not None:
                grid_cols = tbl_grid.findall(f'{{{a_ns}}}gridCol')
                if len(grid_cols) == num_cols:
                    widths = [gc.get('w', '0') for gc in grid_cols]
                    for i, gc in enumerate(grid_cols):
                        gc.set('w', widths[num_cols - 1 - i])
                    changes += 1

            # Step 3: Reverse cell contents within each row
            for row in table.rows:
                tr_elem = row._tr
                tc_elems = tr_elem.findall(f'{{{a_ns}}}tc')

                # Deep copy, remove originals, re-insert in reversed order
                cell_copies = [deepcopy(tc) for tc in tc_elems]
                for tc in tc_elems:
                    tr_elem.remove(tc)
                for tc in reversed(cell_copies):
                    tr_elem.append(tc)
                changes += 1

            # Step 4: Set table-level RTL flag
            tbl_pr = tbl_elem.find(f'{{{a_ns}}}tblPr')
            if tbl_pr is None:
                tbl_pr = etree.SubElement(tbl_elem, f'{{{a_ns}}}tblPr')
                tbl_elem.insert(0, tbl_pr)
            tbl_pr.set('rtl', '1')
            changes += 1

            # Step 5: Set RTL properties on cell text
            for row in table.rows:
                for col_idx, cell in enumerate(row.cells):
                    if cell.text_frame:
                        self._set_cell_rtl_properties(
                            cell.text_frame, col_idx, num_cols
                        )
                        changes += 1

        except Exception as exc:
            logger.warning('_transform_table_rtl on "%s": %s',
                           getattr(shape, 'name', '?'), exc)

        # v1.1.3: Table transform telemetry
        try:
            logger.info(
                'Table RTL on "%s": %d cols, %d rows, tblPr.rtl=%s, changes=%d',
                getattr(shape, 'name', '?'), num_cols, len(table.rows),
                tbl_pr.get('rtl', 'not set') if tbl_pr is not None else 'no tblPr',
                changes,
            )
        except Exception:
            pass

        return changes

    def _translate_cell_text(self, text_frame) -> int:
        """
        Round 3: Translate text within a table cell's text frame.

        Same algorithm as _apply_translation but operates on a cell text frame
        instead of a shape's text frame.

        Returns:
            Count of paragraphs modified.
        """
        changes = 0
        if not self.translations:
            return 0

        try:
            for para in text_frame.paragraphs:
                para_text = para.text
                if not para_text or not para_text.strip():
                    continue

                arabic_text = self._fuzzy_lookup_translation(
                    para_text, self.translations
                )

                if arabic_text:
                    runs = para.runs
                    if runs:
                        runs[0].text = arabic_text
                        for run in runs[1:]:
                            run.text = ''
                    else:
                        p_elem = para._p
                        r_elem = etree.SubElement(p_elem, f'{{{A_NS}}}r')
                        t_elem = etree.SubElement(r_elem, f'{{{A_NS}}}t')
                        t_elem.text = arabic_text

                    # Set Arabic language on runs
                    for r_elem in para._p.findall(f'{{{A_NS}}}r'):
                        rPr = r_elem.find(f'{{{A_NS}}}rPr')
                        if rPr is not None:
                            rPr.set('lang', 'ar-SA')

                    changes += 1

        except Exception as exc:
            logger.debug('_translate_cell_text: %s', exc)

        return changes

    def _set_cell_rtl_properties(self, text_frame, col_idx: int, num_cols: int) -> None:
        """
        Apply RTL/alignment properties to a table cell text frame.

        Round 2 fix: Only set rtl='1' on cells containing Arabic text.
        English cells get rtl='0' to prevent period displacement.

        Numeric cells (mostly digits/currency) → left-align (numbers read LTR).
        All other cells → right-align.
        """
        cell_text = text_frame.text or ''

        # Detect numeric cell (>80% numeric/currency characters)
        numeric_chars = sum(
            1 for c in cell_text.strip()
            if c.isdigit() or c in '$€£¥%.,+-() '
        )
        total_nonws = len(cell_text.strip())
        is_numeric = total_nonws > 0 and (numeric_chars / total_nonws) > 0.80

        try:
            body_pr = text_frame._txBody.find(f'{{{A_NS}}}bodyPr')
            if body_pr is not None:
                body_pr.set('rtlCol', '1')
        except Exception:
            pass

        for para in text_frame.paragraphs:
            para_text = (para.text or '').strip()
            if not para_text:
                continue
            pPr = ensure_pPr(para._p)
            # Only set rtl='1' for Arabic content
            if has_arabic(para_text):
                pPr.set('rtl', '1')
            else:
                pPr.set('rtl', '0')
            if is_numeric:
                pPr.set('algn', 'l')
            else:
                pPr.set('algn', 'r')

    # ─────────────────────────────────────────────────────────────────────
    # Chart RTL transformation
    # ─────────────────────────────────────────────────────────────────────

    def _transform_chart_rtl(self, shape) -> int:
        """
        Apply RTL transformations to a chart shape.

        Operations (frontier-model-validated architecture):
        1. Detect chart types present in the plotArea.
        2. Reverse category axis AND dateAx direction for RTL reading.
        3. Collect axis IDs referenced by horizontal bar charts.
        4. Only reverse valAx orientation for axes used by horizontal bars.
        5. Mirror valAx position: only l↔r (NOT t↔b for horizontal axes).
        6. Remove crossesAt when setting crosses (mutually exclusive per spec).
        7. Mirror legend position.
        8. Translate chart labels.
        9. Set RTL on chart title (conditional alignment based on has_arabic).
        10. Remove manual legend layout to let PowerPoint reflow entries.

        Args:
            shape: python-pptx Shape with has_chart == True.

        Returns:
            Count of changes applied.
        """
        changes = 0
        try:
            chart = shape.chart
            chart_part = chart._part
            chart_elem = chart_part._element

            c_ns = 'http://schemas.openxmlformats.org/drawingml/2006/chart'

            # ── Detect chart types present ─────────────────────────────────
            axis_chart_types = set()  # chart types that use axes
            for tag in ('barChart', 'bar3DChart', 'lineChart', 'line3DChart',
                        'areaChart', 'area3DChart', 'scatterChart', 'radarChart',
                        'stockChart', 'surfaceChart', 'surface3DChart', 'bubbleChart'):
                if chart_elem.find(f'.//{{{c_ns}}}{tag}') is not None:
                    axis_chart_types.add(tag)

            # If only pie/doughnut charts, skip all axis manipulation
            has_axis_charts = bool(axis_chart_types)

            # Collect axis IDs referenced by horizontal bar charts
            horiz_bar_val_ax_ids = set()
            for bar_tag in ('barChart', 'bar3DChart'):
                for bar_chart in chart_elem.iter(f'{{{c_ns}}}{bar_tag}'):
                    bar_dir = bar_chart.find(f'{{{c_ns}}}barDir')
                    if bar_dir is not None and bar_dir.get('val') == 'bar':
                        for ax_id in bar_chart.iter(f'{{{c_ns}}}axId'):
                            horiz_bar_val_ax_ids.add(ax_id.get('val'))

            # Collect all valAx IDs for scatter chart handling
            # Scatter charts have 2 valAx and no catAx — identify the horizontal one
            scatter_horiz_val_ax_ids = set()
            for scatter in chart_elem.iter(f'{{{c_ns}}}scatterChart'):
                scatter_ax_ids = [ax_id.get('val') for ax_id in scatter.iter(f'{{{c_ns}}}axId')]
                for val_ax in chart_elem.iter(f'{{{c_ns}}}valAx'):
                    ax_id_elem = val_ax.find(f'{{{c_ns}}}axId')
                    if ax_id_elem is None:
                        continue
                    ax_id_val = ax_id_elem.get('val')
                    if ax_id_val not in scatter_ax_ids:
                        continue
                    # Horizontal valAx in scatter typically has axPos b or t
                    ax_pos = val_ax.find(f'{{{c_ns}}}axPos')
                    if ax_pos is not None and ax_pos.get('val') in ('b', 't'):
                        scatter_horiz_val_ax_ids.add(ax_id_val)

            # ── Step 1: Reverse category axis AND dateAx for RTL ───────────
            if has_axis_charts:
                for ax_tag in ('catAx', 'dateAx'):
                    for cat_ax in chart_elem.iter(f'{{{c_ns}}}{ax_tag}'):
                        # Remove crossesAt first (mutually exclusive with crosses)
                        crosses_at = cat_ax.find(f'{{{c_ns}}}crossesAt')
                        if crosses_at is not None:
                            cat_ax.remove(crosses_at)

                        crosses = cat_ax.find(f'{{{c_ns}}}crosses')
                        if crosses is None:
                            crosses = etree.SubElement(cat_ax, f'{{{c_ns}}}crosses')
                        crosses.set('val', 'max')

                        # Set axis orientation to maxMin for RTL reading
                        scaling = cat_ax.find(f'{{{c_ns}}}scaling')
                        if scaling is None:
                            scaling = etree.SubElement(cat_ax, f'{{{c_ns}}}scaling')
                            cat_ax.insert(0, scaling)
                        orientation = scaling.find(f'{{{c_ns}}}orientation')
                        if orientation is None:
                            orientation = etree.SubElement(scaling, f'{{{c_ns}}}orientation')
                        orientation.set('val', 'maxMin')  # Deterministic, not toggle
                        changes += 1

            # ── Step 2: Handle value axes ──────────────────────────────────
            if has_axis_charts:
                for val_ax in chart_elem.iter(f'{{{c_ns}}}valAx'):
                    ax_id_elem = val_ax.find(f'{{{c_ns}}}axId')
                    ax_id_val = ax_id_elem.get('val') if ax_id_elem is not None else None

                    # 2a: axPos mirroring — ONLY l↔r, NOT t↔b
                    # t↔b mirroring moves horizontal bar axes to the top, which is wrong
                    ax_pos = val_ax.find(f'{{{c_ns}}}axPos')
                    if ax_pos is not None:
                        pos = ax_pos.get('val', 'l')
                        if pos == 'l':
                            ax_pos.set('val', 'r')
                            changes += 1
                        elif pos == 'r':
                            ax_pos.set('val', 'l')
                            changes += 1
                        # t and b: leave unchanged (no vertical mirroring)

                    # 2b: valAx orientation reversal — ONLY for horizontal bar axes
                    # and horizontal scatter axes. NOT for line/area/column Y-axes.
                    is_horiz_bar = ax_id_val in horiz_bar_val_ax_ids
                    is_scatter_horiz = ax_id_val in scatter_horiz_val_ax_ids

                    if is_horiz_bar or is_scatter_horiz:
                        scaling = val_ax.find(f'{{{c_ns}}}scaling')
                        if scaling is None:
                            scaling = etree.SubElement(val_ax, f'{{{c_ns}}}scaling')
                            val_ax.insert(0, scaling)
                        orientation = scaling.find(f'{{{c_ns}}}orientation')
                        if orientation is None:
                            orientation = etree.SubElement(scaling, f'{{{c_ns}}}orientation')
                        orientation.set('val', 'maxMin')  # Deterministic set, not toggle

                        # Remove crossesAt (mutually exclusive with crosses)
                        crosses_at = val_ax.find(f'{{{c_ns}}}crossesAt')
                        if crosses_at is not None:
                            val_ax.remove(crosses_at)
                        changes += 1

            # ── Step 3: Handle serAx (3D series axis) ──────────────────────
            for ser_ax in chart_elem.iter(f'{{{c_ns}}}serAx'):
                ax_pos = ser_ax.find(f'{{{c_ns}}}axPos')
                if ax_pos is not None:
                    pos = ax_pos.get('val', 'b')
                    if pos == 'l':
                        ax_pos.set('val', 'r')
                        changes += 1
                    elif pos == 'r':
                        ax_pos.set('val', 'l')
                        changes += 1

            # ── Step 4: Mirror legend position ─────────────────────────────
            legend = chart_elem.find(f'.//{{{c_ns}}}legend')
            if legend is not None:
                leg_pos = legend.find(f'{{{c_ns}}}legendPos')
                if leg_pos is None:
                    # Default per OOXML is 'r', create and set to 'l' for RTL
                    leg_pos = etree.SubElement(legend, f'{{{c_ns}}}legendPos')
                    leg_pos.set('val', 'l')
                    changes += 1
                else:
                    pos_val = leg_pos.get('val', 'r')
                    mirror_map = {'r': 'l', 'l': 'r', 'tr': 'tl', 'tl': 'tr'}
                    new_pos = mirror_map.get(pos_val, pos_val)
                    leg_pos.set('val', new_pos)
                    if new_pos != pos_val:
                        changes += 1

                # Remove manual legend layout so PowerPoint can reflow entries
                legend_layout = legend.find(f'{{{c_ns}}}layout')
                if legend_layout is not None:
                    legend.remove(legend_layout)
                    changes += 1

            # ── Step 5: Translate chart category labels and series names ───
            if self.translations:
                changes += self._translate_chart_labels(chart_elem, c_ns)

            # ── Step 6: Set RTL on chart title text (conditional alignment) ─
            try:
                if chart.has_title and chart.chart_title.has_text_frame:
                    for para in chart.chart_title.text_frame.paragraphs:
                        para_text = (para.text or '').strip()
                        if para_text:
                            pPr = ensure_pPr(para._p)
                            if has_arabic(para_text):
                                pPr.set('rtl', '1')
                                pPr.set('algn', 'r')
                            else:
                                pPr.set('rtl', '0')
                                pPr.set('algn', 'l')
                            changes += 1
            except Exception:
                pass

        except Exception as exc:
            logger.warning('_transform_chart_rtl on "%s": %s',
                           getattr(shape, 'name', '?'), exc)

        return changes

    # ─────────────────────────────────────────────────────────────────────
    # Round 3: Chart label translation
    # ─────────────────────────────────────────────────────────────────────

    # Known Google Translate month name errors: words that are ambiguous
    # in English but get wrong Arabic translations when translated without context.
    _MONTH_CORRECTIONS = {
        'يمشي': 'مارس',       # "March" (verb walk → month)
        'يمكن': 'مايو',       # "May" (verb can → month)
        'يُولي و': 'يوليو',   # Garbled "July"
        'يمك': 'مايو',        # Truncated "May"
    }

    def _is_pie_or_doughnut_chart(self, chart_elem, c_ns: str) -> bool:
        """Check if chart contains ONLY pie/doughnut chart types (no axis charts)."""
        pie_types = ('pieChart', 'pie3DChart', 'doughnutChart', 'ofPieChart')
        axis_types = ('barChart', 'bar3DChart', 'lineChart', 'line3DChart',
                      'areaChart', 'area3DChart', 'scatterChart', 'radarChart',
                      'stockChart', 'surfaceChart', 'surface3DChart', 'bubbleChart')
        has_pie = any(
            chart_elem.find(f'.//{{{c_ns}}}{t}') is not None for t in pie_types
        )
        has_axis = any(
            chart_elem.find(f'.//{{{c_ns}}}{t}') is not None for t in axis_types
        )
        return has_pie and not has_axis

    def _pie_has_inside_labels(self, chart_elem, c_ns: str) -> bool:
        """
        Detect if pie/doughnut chart uses in-segment data labels
        (dLblPos = 'ctr', 'inEnd', 'inBase', or 'bestFit').

        These positions render text INSIDE the pie slice where space is
        extremely constrained. Arabic text is typically 20-40% wider than
        English and will truncate/overflow in these positions.
        """
        inside_positions = {'ctr', 'inEnd', 'inBase', 'bestFit'}

        # Check global data labels setting (c:dLbls at chart level)
        for dLbls in chart_elem.iter(f'{{{c_ns}}}dLbls'):
            pos = dLbls.find(f'{{{c_ns}}}dLblPos')
            if pos is not None and pos.get('val') in inside_positions:
                return True
            # Also check showCatName — if category names shown AND no explicit
            # dLblPos, pie defaults to 'bestFit' which is inside
            show_cat = dLbls.find(f'{{{c_ns}}}showCatName')
            if show_cat is not None and show_cat.get('val') == '1':
                # For pie/doughnut, if showCatName=1 and no outEnd position,
                # labels render inside by default
                if pos is None or pos.get('val') != 'outEnd':
                    return True

        # Check individual data label overrides
        for dLbl in chart_elem.iter(f'{{{c_ns}}}dLbl'):
            pos = dLbl.find(f'{{{c_ns}}}dLblPos')
            if pos is not None and pos.get('val') in inside_positions:
                return True

        return False

    def _translate_chart_labels(self, chart_elem, c_ns: str) -> int:
        """
        Round 3: Translate chart category labels and series names.

        Searches for <c:v> text elements within <c:cat> (category axis)
        and <c:tx> (series names) and replaces with Arabic translations.

        Fix 5 (frontier-validated): For pie/doughnut charts with in-segment
        labels, SKIP translating category cache values to prevent truncation.
        Arabic text is 20-40% wider than English and overflows in constrained
        pie segments. Legend and series names are still translated.

        Also applies month name corrections for known Google Translate errors.

        Returns:
            Count of labels translated.
        """
        changes = 0

        # Fix 5: Determine if this is a pie/doughnut with inside labels
        is_pie = self._is_pie_or_doughnut_chart(chart_elem, c_ns)
        skip_cat_translation = is_pie and self._pie_has_inside_labels(chart_elem, c_ns)

        if skip_cat_translation:
            logger.info('Fix 5: Skipping category label translation for '
                        'pie/doughnut chart with in-segment labels')

        # Translate category axis labels (skip for pie/doughnut with inside labels)
        if not skip_cat_translation:
            for cat in chart_elem.iter(f'{{{c_ns}}}cat'):
                for pt in cat.iter(f'{{{c_ns}}}pt'):
                    v = pt.find(f'{{{c_ns}}}v')
                    if v is not None and v.text:
                        original = v.text.strip()
                        arabic = (
                            self.translations.get(original)
                            or self.translations.get(original.title())
                            or self.translations.get(original.upper())
                            or self.translations.get(original.lower())
                        )
                        if arabic:
                            # Apply month corrections
                            corrected = self._MONTH_CORRECTIONS.get(arabic.strip(), arabic)
                            v.text = corrected
                            changes += 1

        # Translate series names (always — these appear in legend, not in segments)
        for ser in chart_elem.iter(f'{{{c_ns}}}ser'):
            for tx in ser.findall(f'{{{c_ns}}}tx'):
                for v in tx.iter(f'{{{c_ns}}}v'):
                    if v.text:
                        original = v.text.strip()
                        arabic = self.translations.get(original)
                        if arabic:
                            v.text = arabic
                            changes += 1

        return changes

    # ─────────────────────────────────────────────────────────────────────
    # Fix 8: Text box width validation
    # ─────────────────────────────────────────────────────────────────────

    def _validate_textbox_width(self, shape) -> None:
        """
        Fix 8: Ensure text box width is sufficient for short text content
        (e.g., page numbers). Prevents two-digit numbers from wrapping
        onto separate lines after mirroring narrows the text box.

        RTL-aware: expands to the LEFT (preserving right anchor) since
        neighboring shapes may be to the right after mirroring.
        """
        try:
            if not getattr(shape, 'has_text_frame', False):
                return
            text = shape.text_frame.text or ''
            stripped = text.strip()
            if not stripped or len(stripped) > 5:
                return  # Only check short text (page numbers, labels)

            # Minimum width: 0.5 inches = 457200 EMU
            min_width = 457200
            if shape.width is not None and shape.width < min_width:
                delta = min_width - shape.width
                shape.width = min_width
                # RTL-aware: shift left to preserve right anchor
                # Clamp symmetrically (P0 fix — 3-model consensus)
                if shape.left is not None:
                    shift = min(max(shape.left, 0), delta)
                    shape.left = shape.left - shift
        except Exception:
            pass

    # ─────────────────────────────────────────────────────────────────────
    # Fix 11: wrap="none" text boxes — widen for Arabic text
    # ─────────────────────────────────────────────────────────────────────

    def _fix_wrap_none_for_arabic(self, shape) -> int:
        """
        Fix 11: Text boxes with wrap="none" may be too narrow for Arabic text.

        Arabic text at the same font size is typically 20-40% wider than English.
        When a text box uses wrap="none" (auto-width in PowerPoint), the box width
        was sized for the English text. After Arabic translation, the text may
        overflow or break mid-word.

        Fix: For text boxes with wrap="none" that contain Arabic text:
        1. Change wrap to "square" (allow wrapping within bounds)
        2. Set normAutofit so text scales down if needed
        3. Expand cx if the text box is in a panel region with room

        Returns:
            Count of changes made.
        """
        changes = 0
        try:
            if not getattr(shape, 'has_text_frame', False):
                return 0
            tf = shape.text_frame
            body_pr = tf._txBody.find(f'{{{A_NS}}}bodyPr')
            if body_pr is None:
                return 0

            wrap_val = body_pr.get('wrap', 'square')
            if wrap_val != 'none':
                return 0

            # Check if shape contains Arabic text
            shape_text = tf.text or ''
            if not has_arabic(shape_text):
                return 0

            # Get current geometry
            sp_el = shape._element
            xfrm = sp_el.find(f'.//{{{A_NS}}}xfrm')
            if xfrm is None:
                return 0
            off = xfrm.find(f'{{{A_NS}}}off')
            ext = xfrm.find(f'{{{A_NS}}}ext')
            if off is None or ext is None:
                return 0

            try:
                x = int(off.get('x', 0))
                cx = int(ext.get('cx', 0))
            except (ValueError, TypeError):
                return 0

            # Change wrap to "square" for wrapping support
            body_pr.set('wrap', 'square')
            changes += 1

            # Remove spAutoFit (which auto-expands box) and add normAutofit
            # (which shrinks text to fit the box instead)
            for child_tag in ('spAutoFit', 'noAutofit'):
                for child in body_pr.findall(f'{{{A_NS}}}{child_tag}'):
                    body_pr.remove(child)

            # Only add normAutofit if not already present
            existing_autofit = body_pr.find(f'{{{A_NS}}}normAutofit')
            if existing_autofit is None:
                autofit = etree.SubElement(body_pr, f'{{{A_NS}}}normAutofit')
                autofit.set('fontScale', '80000')  # Allow 80% shrink
                changes += 1

            # ── RTL-Aware Width Expansion ──────────────────────────────
            # In RTL context, text is right-aligned. When we expand a text
            # box's width, we must preserve the RIGHT edge (anchor) and
            # grow to the LEFT. Otherwise, the expanded box swallows
            # neighboring shapes that are now to its right after mirroring.
            #
            # Rule: new_x = old_x - delta_width (shift left by the expansion amount)
            #        new_cx = old_cx + delta_width
            # This keeps right_edge = x + cx constant.
            # ─────────────────────────────────────────────────────────────

            # Estimate needed width: Arabic text is typically 20-40% wider.
            # Use a conservative 1.35x expansion factor, capped to avoid
            # excessive growth.
            ARABIC_EXPANSION_FACTOR = 1.35
            MAX_EXPANSION_FACTOR = 2.0

            estimated_cx = int(cx * ARABIC_EXPANSION_FACTOR)
            max_cx = int(cx * MAX_EXPANSION_FACTOR)

            # Cap expansion to 85% of slide width (P0 fix — 3-model consensus)
            # half_width was too aggressive for single-column layouts
            max_slide_cap = int(self._slide_width * 0.85)
            cap_cx = min(max_cx, max_slide_cap)
            new_cx = min(estimated_cx, cap_cx)

            if new_cx > cx:
                delta = new_cx - cx
                new_x = x - delta  # Shift LEFT to preserve right anchor
                # Clamp: don't let new_x go below 0 (slide left edge)
                if new_x < 0:
                    # Reduce expansion to what fits
                    delta = x  # expand only as much as x allows
                    new_x = 0
                    new_cx = cx + delta

                if new_cx > cx:  # still expanding after clamp
                    ext.set('cx', str(new_cx))
                    off.set('x', str(new_x))
                    changes += 1
                    logger.debug(
                        'Fix 11: RTL-aware expand textbox "%s" '
                        'cx %d->%d, x %d->%d (right edge preserved at %d)',
                        getattr(shape, 'name', '?'), cx, new_cx, x, new_x,
                        new_x + new_cx
                    )

        except Exception as exc:
            logger.debug('_fix_wrap_none_for_arabic: %s', exc)

        return changes

    # ─────────────────────────────────────────────────────────────────────
    # Fix 12: Title-body vertical overlap resolution for RTL
    # ─────────────────────────────────────────────────────────────────────

    def _fix_title_body_overlap(self, shapes: list, slide_number: int) -> int:
        """
        Fix 12: Detect and fix vertical overlap between title and body text boxes.

        In LTR, overlapping title and body boxes can coexist because the text
        flows in different horizontal zones (title centered/right, body left).
        In RTL, BOTH texts flow from the right edge, causing visual collision
        in the vertical overlap zone.

        Algorithm:
        1. Identify "title-like" shapes (large font, short text, near top of slide)
        2. Identify "body-like" shapes (smaller font, longer text)
        3. If a body shape's top edge is within a title shape's vertical extent,
           move the body shape down so it clears the title + a small gap.

        Returns:
            Count of shapes repositioned.
        """
        changes = 0
        GAP_EMU = 91440  # 0.1 inch gap

        # Collect text shapes with their geometry
        text_shapes = []
        for shape in shapes:
            try:
                if not getattr(shape, 'has_text_frame', False):
                    continue
                if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
                    continue
                text = (shape.text_frame.text or '').strip()
                if not text:
                    continue

                # Get max font size in this shape
                max_font_size = 0
                sp_el = shape._element
                for rPr in sp_el.iter(f'{{{A_NS}}}rPr'):
                    sz_str = rPr.get('sz')
                    if sz_str:
                        try:
                            max_font_size = max(max_font_size, int(sz_str))
                        except ValueError:
                            pass
                # Also check defRPr
                for def_rPr in sp_el.iter(f'{{{A_NS}}}defRPr'):
                    sz_str = def_rPr.get('sz')
                    if sz_str:
                        try:
                            max_font_size = max(max_font_size, int(sz_str))
                        except ValueError:
                            pass

                text_shapes.append({
                    'shape': shape,
                    'top': int(shape.top),
                    'bottom': int(shape.top) + int(shape.height),
                    'left': int(shape.left),
                    'right': int(shape.left) + int(shape.width),
                    'height': int(shape.height),
                    'font_size': max_font_size,
                    'text_len': len(text),
                })
            except Exception:
                continue

        # Identify title-like shapes: large font (>4000 = 40pt), in top 30% of slide
        top_30pct = self._slide_height * 0.30
        titles = [s for s in text_shapes
                  if s['font_size'] >= 4000 and s['top'] < top_30pct]

        if not titles:
            return 0

        # For each title, check if any body-like shape overlaps vertically
        for title in titles:
            for body in text_shapes:
                if body is title:
                    continue
                # Body must have smaller font than title
                if body['font_size'] >= title['font_size']:
                    continue
                # Body's top must be within title's vertical extent
                if body['top'] < title['top']:
                    continue  # Body is above title — not our concern
                if body['top'] >= title['bottom']:
                    continue  # Body is already below title — no overlap

                # OVERLAP DETECTED: body['top'] is between title['top'] and title['bottom']
                # Check if both shapes have significant horizontal overlap
                h_overlap = (min(title['right'], body['right']) -
                           max(title['left'], body['left']))
                if h_overlap <= 0:
                    continue  # No horizontal overlap — texts in different columns

                # Move body down to clear the title
                new_top = title['bottom'] + GAP_EMU
                old_top = body['top']
                try:
                    sp_el = body['shape']._element
                    xfrm = sp_el.find(f'.//{{{A_NS}}}xfrm')
                    if xfrm is not None:
                        off = xfrm.find(f'{{{A_NS}}}off')
                        if off is not None:
                            off.set('y', str(new_top))
                            body['top'] = new_top
                            body['bottom'] = new_top + body['height']
                            changes += 1
                            logger.debug(
                                'Fix 12 slide %d: moved body "%s" y %d -> %d to clear title',
                                slide_number,
                                getattr(body['shape'], 'name', '?'),
                                old_top, new_top
                            )
                except Exception as exc:
                    logger.debug('Fix 12: %s', exc)

        return changes

    # ─────────────────────────────────────────────────────────────────────
    # Fix 22: Post-translation Arabic autofit
    # ─────────────────────────────────────────────────────────────────────

    def _apply_arabic_autofit(self, shapes: list, slide_number: int) -> int:
        """
        Fix 22: Apply normAutofit to all text frames containing Arabic text.

        Arabic text is typically 20-40% wider and often taller (due to diacritics)
        than the English source at the same font size. Without autofit, text
        overflows the placeholder bounding box and collides with adjacent shapes.

        Two-pronged approach for maximum compatibility:
        1. Set <a:normAutofit fontScale="..."/> on bodyPr (PowerPoint respects this)
        2. Directly reduce font sizes on runs (LibreOffice ignores normAutofit in
           headless PDF export, so we must also reduce actual font sizes)

        Only applies to shapes that:
        - Have a text frame with Arabic content
        - Don't already have normAutofit set
        - Have text that plausibly overflows (estimated from character count vs box width)

        Returns:
            Count of shapes modified.
        """
        changes = 0

        for shape in shapes:
            try:
                if not getattr(shape, 'has_text_frame', False):
                    continue
                tf = shape.text_frame
                shape_text = tf.text or ''
                if not has_arabic(shape_text):
                    continue

                body_pr = tf._txBody.find(f'{{{A_NS}}}bodyPr')
                if body_pr is None:
                    continue

                # Skip if normAutofit is already set
                existing = body_pr.find(f'{{{A_NS}}}normAutofit')
                if existing is not None:
                    continue

                # Estimate whether text overflows its box.
                # Heuristic: Arabic text at the same font size needs ~30% more
                # horizontal space. Count characters and compare to box width.
                box_width_emu = None
                box_height_emu = None
                try:
                    if shape.width is not None:
                        box_width_emu = int(shape.width)
                    if shape.height is not None:
                        box_height_emu = int(shape.height)
                except (TypeError, ValueError):
                    pass

                # Get max font size in this shape (hundredths of a point)
                max_font_hundredths = 0
                for rPr in shape._element.iter(f'{{{A_NS}}}rPr'):
                    sz_str = rPr.get('sz')
                    if sz_str:
                        try:
                            max_font_hundredths = max(max_font_hundredths, int(sz_str))
                        except ValueError:
                            pass
                for def_rPr in shape._element.iter(f'{{{A_NS}}}defRPr'):
                    sz_str = def_rPr.get('sz')
                    if sz_str:
                        try:
                            max_font_hundredths = max(max_font_hundredths, int(sz_str))
                        except ValueError:
                            pass

                # Rough overflow estimate:
                # A character at N pt takes approximately N * 635 EMU width (monospace approx).
                # Arabic is ~1.3x wider. If estimated text width > box width, it overflows.
                char_count = len(shape_text.replace('\n', ''))
                if max_font_hundredths > 0 and box_width_emu and box_width_emu > 0:
                    font_pt = max_font_hundredths / 100
                    # Rough: each Arabic char takes ~0.6 * font_size_pt * 12700 EMU
                    estimated_char_width_emu = 0.6 * font_pt * 12700
                    # Count lines (by newlines) and estimate per-line width
                    lines = shape_text.split('\n')
                    max_line_chars = max(len(l.strip()) for l in lines) if lines else char_count
                    estimated_line_width = max_line_chars * estimated_char_width_emu * 1.3  # Arabic expansion
                    overflow_ratio = estimated_line_width / box_width_emu

                    # Also estimate height overflow
                    line_count = len(lines)
                    estimated_line_height_emu = font_pt * 12700 * 1.5  # 1.5x line spacing
                    estimated_total_height = line_count * estimated_line_height_emu
                    height_overflow = (estimated_total_height / box_height_emu) if box_height_emu else 0

                    # Only apply autofit if text plausibly overflows
                    if overflow_ratio < 0.9 and height_overflow < 0.9:
                        continue

                    # Calculate appropriate font scale
                    max_overflow = max(overflow_ratio, height_overflow)
                    # fontScale: 100000 = 100%. Minimum 50000 = 50%.
                    font_scale = max(int(100000 / max_overflow), 50000)
                    font_scale = min(font_scale, 100000)  # Don't exceed 100%
                else:
                    # Can't estimate — apply conservative autofit (90% scale)
                    font_scale = 90000

                # Remove any existing autofit variants
                for child_tag in ('spAutoFit', 'noAutofit'):
                    for child in body_pr.findall(f'{{{A_NS}}}{child_tag}'):
                        body_pr.remove(child)

                # Add normAutofit
                autofit_el = etree.SubElement(body_pr, f'{{{A_NS}}}normAutofit')
                autofit_el.set('fontScale', str(font_scale))

                # Also directly scale font sizes on runs for LibreOffice compatibility
                font_scale_ratio = font_scale / 100000
                for para in tf.paragraphs:
                    for run in para.runs:
                        if run.font.size is not None:
                            original_size = int(run.font.size)
                            new_size = max(int(original_size * font_scale_ratio), 8 * 12700)  # min 8pt
                            run.font.size = new_size
                        else:
                            # Read default font size from defRPr or endParaRPr
                            rPr = run._r.find(f'{{{A_NS}}}rPr')
                            if rPr is None:
                                rPr = etree.SubElement(run._r, f'{{{A_NS}}}rPr')
                                run._r.insert(0, rPr)
                            default_sz = None
                            for src in [para._p.find(f'.//{{{A_NS}}}defRPr'),
                                        para._p.find(f'{{{A_NS}}}endParaRPr')]:
                                if src is not None and src.get('sz'):
                                    try:
                                        default_sz = int(src.get('sz'))
                                    except ValueError:
                                        pass
                                    break
                            if default_sz:
                                new_sz = max(int(default_sz * font_scale_ratio), 800)
                                rPr.set('sz', str(new_sz))

                changes += 1
                logger.debug(
                    'Fix 22 slide %d: autofit on "%s" fontScale=%d%%',
                    slide_number, getattr(shape, 'name', '?'), font_scale // 1000,
                )

            except Exception as exc:
                logger.debug('Fix 22 slide %d: %s', slide_number, exc)

        return changes

    # ─────────────────────────────────────────────────────────────────────
    # Fix 13: Mirror slide-level logo images
    # ─────────────────────────────────────────────────────────────────────

    def _mirror_slide_level_logos(self, shapes: list) -> int:
        """
        Fix 13: Mirror small logo/brand images that appear directly on slides
        (not inherited from masters/layouts).

        The MasterLayoutTransformer mirrors logos on masters and layouts, but
        some decks place logo images directly on content slides. These need
        to be mirrored at the slide level too.

        Detection criteria (same as master-level):
        - Is a picture element (<p:pic>)
        - Width < 20% of slide width
        - Not a placeholder
        - Has an actual image (blipFill with rEmbed)

        Returns:
            Count of logos mirrored.
        """
        mirrored = 0
        for shape in shapes:
            try:
                sp_el = shape._element
                tag = sp_el.tag

                # Must be a picture element
                if not (tag.endswith('}pic') or tag == 'pic'):
                    continue

                # Must not be a placeholder
                nv_pr = sp_el.find(f'.//{{{P_NS}}}ph')
                if nv_pr is not None:
                    continue

                # Must have actual image
                blip_fill = sp_el.find(f'{{{P_NS}}}blipFill')
                if blip_fill is None:
                    blip_fill = sp_el.find(f'{{{A_NS}}}blipFill')
                if blip_fill is None:
                    continue
                blip = blip_fill.find(f'{{{A_NS}}}blip') if blip_fill is not None else None
                if blip is None or not blip.get(f'{{{R_NS}}}embed'):
                    continue

                # Width check
                width = getattr(shape, 'width', None)
                left = getattr(shape, 'left', None)
                if width is None or left is None:
                    continue
                if width >= self._slide_width * _LOGO_MAX_WIDTH_FRACTION:
                    continue

                # Already handled by _mirror_freeform_shape in _should_mirror_shape?
                # That function mirrors ALL content shapes. But logos are images,
                # not text shapes, so they go through the same path. However,
                # the _should_mirror_shape might skip some logos due to layout rules.
                # This explicit mirror ensures logos are always mirrored.
                # Skip if already near the mirrored position (tolerance check)
                new_left = mirror_x(left, width, self._slide_width)
                if abs(new_left - left) < _POSITION_TOLERANCE_EMU:
                    continue

                # Bounds check
                if not bounds_check_emu(new_left, self._slide_width):
                    continue

                shape.left = new_left
                mirrored += 1
                logger.debug('Fix 13: slide-level logo mirrored: %s %d -> %d',
                             getattr(shape, 'name', '?'), left, new_left)

            except Exception:
                continue

        return mirrored

    # ─────────────────────────────────────────────────────────────────────
    # Fix 14: Right-anchor cover title text for photo-background covers
    # ─────────────────────────────────────────────────────────────────────

    def _fix_cover_title_anchor(self, shapes: list, slide_number: int) -> int:
        """
        Fix 14: On photo-background cover slides, the title text box is often
        left-anchored (LTR convention). For RTL, it should be right-anchored.

        Detection: A slide is a photo-background cover when:
        - It has a large image (>60% of slide area) acting as background
        - It has a title-like text shape (large font, short text)
        - The title is in the left half of the slide

        Fix: Move the title text box to mirror position (right half) and
        set right-alignment on the text.

        Affected decks: R6_04, R6_13, R6_14
        """
        changes = 0
        try:
            half_width = self._slide_width // 2
            slide_area = self._slide_width * self._slide_height

            # Step 1: Detect if this is a photo-background slide
            # Only count actual images (pic elements or shapes with blipFill),
            # NOT solid/gradient filled rectangles which are decorative panels.
            has_large_bg_image = False
            for shape in shapes:
                sp_el = shape._element
                tag = sp_el.tag
                is_pic = tag.endswith('}pic')
                has_blip = sp_el.find(f'.//{{{A_NS}}}blipFill') is not None
                if not (is_pic or has_blip):
                    continue  # Skip non-image shapes (solid/gradient rects)
                w = getattr(shape, 'width', None)
                h = getattr(shape, 'height', None)
                if w is None or h is None:
                    continue
                # Large image covering >50% of slide area
                if int(w) * int(h) > slide_area * 0.50:
                    has_large_bg_image = True
                    break

            if not has_large_bg_image:
                return 0

            # Step 2: Find title-like text shapes in the LEFT half
            for shape in shapes:
                try:
                    if not getattr(shape, 'has_text_frame', False):
                        continue
                    left = getattr(shape, 'left', None)
                    width = getattr(shape, 'width', None)
                    if left is None or width is None:
                        continue

                    text = (shape.text_frame.text or '').strip()
                    if not text or not has_arabic(text):
                        continue

                    # Get font size
                    max_font = 0
                    sp_el = shape._element
                    for rPr in sp_el.iter(f'{{{A_NS}}}rPr'):
                        sz_str = rPr.get('sz')
                        if sz_str:
                            try:
                                max_font = max(max_font, int(sz_str))
                            except ValueError:
                                pass
                    for def_rPr in sp_el.iter(f'{{{A_NS}}}defRPr'):
                        sz_str = def_rPr.get('sz')
                        if sz_str:
                            try:
                                max_font = max(max_font, int(sz_str))
                            except ValueError:
                                pass

                    # Title-like: font >= 24pt (2400 hundredths), text < 200 chars
                    if max_font < 2400 or len(text) > 200:
                        continue

                    # Shape center is in left half?
                    shape_center_x = int(left) + int(width) // 2
                    if shape_center_x >= half_width:
                        continue  # Already on right side

                    # Mirror the title text box to the right half
                    new_left = mirror_x(int(left), int(width), self._slide_width)
                    new_left = clamp_emu(new_left, self._slide_width)

                    # Verify the new position is actually in the right half
                    new_center = new_left + int(width) // 2
                    if new_center <= half_width:
                        continue  # Mirror didn't move it to right half

                    shape.left = new_left
                    changes += 1
                    logger.debug(
                        'Fix 14 slide %d: cover title "%s" moved to right half %d -> %d',
                        slide_number, text[:30], int(left), new_left
                    )
                except Exception:
                    continue

        except Exception as exc:
            logger.debug('_fix_cover_title_anchor: %s', exc)

        return changes

    # ─────────────────────────────────────────────────────────────────────
    # Pre-mirror: Split-panel swap BEFORE individual mirroring
    # ─────────────────────────────────────────────────────────────────────

    def _pre_mirror_split_panel_swap(
        self, shapes: list, slide_number: int, handled_ids: set
    ) -> int:
        """
        Detect split-panel layouts (image on one side, content on other)
        BEFORE individual shape mirroring. If detected, swap ALL shapes
        in each panel as a unit, preserving internal spatial relationships.

        This replaces the old post-mirror _mirror_split_panel_layout approach
        which scrambled positions before detection could run.

        Args:
            shapes: all shapes on the slide
            slide_number: for logging
            handled_ids: set to populate with id() of shapes that were moved

        Returns:
            Count of changes made.
        """
        changes = 0
        try:
            half_width = self._slide_width // 2

            # Step 1: Find large panel-defining shapes (>35% width, >50% height)
            left_anchor = None  # large shape in left half
            right_anchor = None  # large shape in right half

            for shape in shapes:
                try:
                    w = getattr(shape, 'width', None)
                    h = getattr(shape, 'height', None)
                    l = getattr(shape, 'left', None)
                    if any(v is None for v in (w, h, l)):
                        continue
                    w, h, l = int(w), int(h), int(l)
                    if w < self._slide_width * 0.35 or h < self._slide_height * 0.50:
                        continue
                    center_x = l + w // 2
                    if center_x < half_width:
                        if left_anchor is None or int(left_anchor.width) < w:
                            left_anchor = shape
                    else:
                        if right_anchor is None or int(right_anchor.width) < w:
                            right_anchor = shape
                except Exception:
                    continue

            # Need exactly one anchor on each side
            if left_anchor is None or right_anchor is None:
                return 0

            # Check that one is an image and one is not
            def _is_image_shape(s):
                tag = s._element.tag
                if tag.endswith('}pic'):
                    return True
                return s._element.find(f'.//{{{A_NS}}}blipFill') is not None

            left_is_img = _is_image_shape(left_anchor)
            right_is_img = _is_image_shape(right_anchor)
            if left_is_img == right_is_img:
                return 0  # Both images or both non-images — not a split panel

            # Step 2: Classify ALL shapes into left-panel vs right-panel
            left_shapes = []  # shapes whose center is in left half
            right_shapes = []  # shapes whose center is in right half
            full_width_shapes = []  # shapes spanning >85% slide width

            for shape in shapes:
                try:
                    if getattr(shape, 'is_placeholder', False):
                        continue  # Placeholders handled separately
                    l = getattr(shape, 'left', None)
                    w = getattr(shape, 'width', None)
                    if l is None or w is None:
                        continue
                    l, w = int(l), int(w)
                    if w >= self._slide_width * 0.85:
                        full_width_shapes.append(shape)
                        continue
                    center_x = l + w // 2
                    if center_x < half_width:
                        left_shapes.append(shape)
                    else:
                        right_shapes.append(shape)
                except Exception:
                    continue

            if not left_shapes or not right_shapes:
                return 0

            # Step 3: Compute panel bounding boxes
            def _bbox(shape_list):
                min_x = min(int(s.left) for s in shape_list)
                max_r = max(int(s.left) + int(s.width) for s in shape_list)
                return min_x, max_r

            left_min_x, left_max_x = _bbox(left_shapes)
            right_min_x, right_max_x = _bbox(right_shapes)

            # Step 4: Swap panels — shift all left shapes to the right zone
            # and all right shapes to the left zone
            # Delta = how far to shift each panel
            # Left panel moves right by: right_min_x - left_min_x
            # Right panel moves left by: left_min_x - right_min_x
            shift = right_min_x - left_min_x

            for shape in left_shapes:
                try:
                    old_left = int(shape.left)
                    new_left = old_left + shift
                    new_left = max(0, min(new_left, self._slide_width - int(shape.width)))
                    shape.left = new_left
                    handled_ids.add(id(shape))
                    changes += 1
                except Exception:
                    continue

            for shape in right_shapes:
                try:
                    old_left = int(shape.left)
                    new_left = old_left - shift
                    new_left = max(0, min(new_left, self._slide_width - int(shape.width)))
                    shape.left = new_left
                    handled_ids.add(id(shape))
                    changes += 1
                except Exception:
                    continue

            if changes > 0:
                logger.debug(
                    'Pre-mirror panel swap slide %d: %d left shapes shifted +%d, '
                    '%d right shapes shifted -%d',
                    slide_number, len(left_shapes), shift,
                    len(right_shapes), shift
                )

        except Exception as exc:
            logger.debug('_pre_mirror_split_panel_swap: %s', exc)

        return changes

    def _exempt_map_overlay_shapes(
        self, shapes: list, slide_number: int, handled_ids: set
    ) -> int:
        """
        Detect slides with a large geographic map image and exempt small
        overlay shapes (bubbles, labels) from position mirroring.

        Geographic positions are absolute — mirroring them would place
        Americas data over Asia and vice versa.

        Detection: A slide has a large image (>40% slide area, width > 60%)
        and multiple small shapes overlaid on it (within the image bounds).

        Returns count of changes (always 0; this only marks shapes as handled).
        """
        try:
            slide_area = self._slide_width * self._slide_height

            # Find large background image
            bg_image = None
            bg_left = bg_top = bg_right = bg_bottom = 0
            for shape in shapes:
                sp_el = shape._element
                tag = sp_el.tag
                is_pic = tag.endswith('}pic')
                has_blip = sp_el.find(f'.//{{{A_NS}}}blipFill') is not None
                if not (is_pic or has_blip):
                    continue
                w = getattr(shape, 'width', None)
                h = getattr(shape, 'height', None)
                l = getattr(shape, 'left', None)
                t = getattr(shape, 'top', None)
                if any(v is None for v in (w, h, l, t)):
                    continue
                w, h, l, t = int(w), int(h), int(l), int(t)
                # Large image covering >40% area and >60% width
                if w * h > slide_area * 0.40 and w > self._slide_width * 0.60:
                    bg_image = shape
                    bg_left, bg_top = l, t
                    bg_right, bg_bottom = l + w, t + h
                    break

            if bg_image is None:
                return 0

            # Count small shapes overlaid on the image
            overlay_shapes = []
            for shape in shapes:
                if shape is bg_image:
                    continue
                if getattr(shape, 'is_placeholder', False):
                    continue
                l = getattr(shape, 'left', None)
                t = getattr(shape, 'top', None)
                w = getattr(shape, 'width', None)
                h = getattr(shape, 'height', None)
                if any(v is None for v in (l, t, w, h)):
                    continue
                l, t, w, h = int(l), int(t), int(w), int(h)
                # Shape must be small (<30% slide width) and within image bounds
                if w > self._slide_width * 0.30:
                    continue
                shape_cx = l + w // 2
                shape_cy = t + h // 2
                if bg_left <= shape_cx <= bg_right and bg_top <= shape_cy <= bg_bottom:
                    overlay_shapes.append(shape)

            # Need at least 3 overlay shapes to classify as map overlay
            if len(overlay_shapes) >= 3:
                for shape in overlay_shapes:
                    handled_ids.add(id(shape))
                logger.debug(
                    'Map overlay slide %d: exempted %d shapes from mirroring',
                    slide_number, len(overlay_shapes)
                )

        except Exception as exc:
            logger.debug('_exempt_map_overlay_shapes: %s', exc)

        return 0

    def _mirror_split_panel_layout(self, shapes: list, slide_number: int) -> int:
        """
        Fix 15: For split-panel slides (image-left / text-right or vice versa),
        swap the two halves so the reading flow is correct for RTL.

        Detection: A split-panel layout has:
        - Exactly 2 major non-overlapping regions occupying most of the slide
        - One region is predominantly image, the other predominantly text
        - The two regions are side-by-side horizontally

        Fix: Swap the X positions of the left and right panels.

        Affected decks: R6_15, R6_16, R6_18 (all covers with photo-left/text-right)
        """
        changes = 0
        try:
            half_width = self._slide_width // 2
            slide_height_80 = self._slide_height * 0.80

            # Classify shapes into left-panel and right-panel candidates
            # A panel shape is large (>35% slide width, >60% slide height)
            left_panels = []   # shapes predominantly in left half
            right_panels = []  # shapes predominantly in right half

            for shape in shapes:
                try:
                    w = getattr(shape, 'width', None)
                    h = getattr(shape, 'height', None)
                    l = getattr(shape, 'left', None)
                    t = getattr(shape, 'top', None)
                    if any(v is None for v in (w, h, l, t)):
                        continue
                    w, h, l, t = int(w), int(h), int(l), int(t)

                    # Must be a panel-sized shape (>35% slide width, >50% slide height)
                    if w < self._slide_width * 0.35 or h < self._slide_height * 0.50:
                        continue

                    # Determine which half it's in
                    center_x = l + w // 2
                    if center_x < half_width:
                        left_panels.append(shape)
                    else:
                        right_panels.append(shape)
                except Exception:
                    continue

            # Need exactly 1 panel on each side for a clear split layout
            if len(left_panels) != 1 or len(right_panels) != 1:
                return 0

            left_shape = left_panels[0]
            right_shape = right_panels[0]

            # Check that they're different types (one image, one text/container)
            left_is_image = left_shape._element.tag.endswith('}pic')
            right_is_image = right_shape._element.tag.endswith('}pic')

            # Also check if shape has blipFill (embedded image in a shape)
            if not left_is_image:
                left_is_image = left_shape._element.find(f'.//{{{A_NS}}}blipFill') is not None
            if not right_is_image:
                right_is_image = right_shape._element.find(f'.//{{{A_NS}}}blipFill') is not None

            # Only mirror if it's a clear image+text split (not two text panels)
            if left_is_image == right_is_image:
                return 0  # Both are images or both are text — don't swap

            # Perform the swap: exchange X positions and widths
            ll, lw = int(left_shape.left), int(left_shape.width)
            rl, rw = int(right_shape.left), int(right_shape.width)

            # Swap: left panel gets right panel's X, right panel gets left panel's X
            # Simple approach: mirror both shapes
            new_ll = mirror_x(ll, lw, self._slide_width)
            new_rl = mirror_x(rl, rw, self._slide_width)

            left_shape.left = new_ll
            right_shape.left = new_rl

            changes += 2
            logger.debug(
                'Fix 15 slide %d: split-panel swapped — left panel %d->%d, right panel %d->%d',
                slide_number, ll, new_ll, rl, new_rl
            )

        except Exception as exc:
            logger.debug('_mirror_split_panel_layout: %s', exc)

        return changes

    # ─────────────────────────────────────────────────────────────────────
    # Fix 16: Reverse timeline alternation pattern
    # ─────────────────────────────────────────────────────────────────────

    def _reverse_timeline_alternation(self, shapes: list, slide_number: int) -> int:
        """
        Fix 16: Timeline slides alternate text labels left/right of a vertical axis.
        For RTL, the alternation pattern must be reversed: first item's label
        should appear on the right side of the axis.

        Detection: A timeline is identified when:
        - There's a vertical line/connector in the center ~33-66% of slide width
        - Multiple small text shapes are arranged on alternating sides of it
        - Shapes have similar vertical spacing (timeline rhythm)

        Fix: For each pair of text shapes at the same Y level on opposite sides
        of the central axis, swap their X positions.

        Affected decks: R6_10 s15
        """
        changes = 0
        try:
            center_zone_left = self._slide_width * 0.33
            center_zone_right = self._slide_width * 0.67

            # Step 1: Find vertical axis (line/connector near center)
            has_central_axis = False
            for shape in shapes:
                try:
                    sp_el = shape._element
                    tag = sp_el.tag
                    # Connectors and thin vertical shapes
                    is_connector = tag.endswith('}cxnSp')
                    l = getattr(shape, 'left', None)
                    w = getattr(shape, 'width', None)
                    h = getattr(shape, 'height', None)
                    if l is None or w is None or h is None:
                        continue
                    l, w, h = int(l), int(w), int(h)
                    center = l + w // 2
                    # Vertical line: narrow (<5% width) and tall (>40% height)
                    is_vertical_line = (w < self._slide_width * 0.05 and
                                       h > self._slide_height * 0.40)
                    if (is_connector or is_vertical_line) and center_zone_left < center < center_zone_right:
                        has_central_axis = True
                        break
                except Exception:
                    continue

            if not has_central_axis:
                return 0

            # Step 2: Collect text shapes on left and right of center
            left_texts = []  # (shape, center_x, center_y)
            right_texts = []
            slide_center_x = self._slide_width // 2

            for shape in shapes:
                try:
                    if not getattr(shape, 'has_text_frame', False):
                        continue
                    text = (shape.text_frame.text or '').strip()
                    if not text:
                        continue
                    l = getattr(shape, 'left', None)
                    w = getattr(shape, 'width', None)
                    t = getattr(shape, 'top', None)
                    h = getattr(shape, 'height', None)
                    if any(v is None for v in (l, w, t, h)):
                        continue
                    l, w, t, h = int(l), int(w), int(t), int(h)
                    center_x = l + w // 2
                    center_y = t + h // 2

                    # Skip shapes in the center zone (axis labels, titles)
                    if center_zone_left < center_x < center_zone_right:
                        continue

                    # Skip very large shapes (titles spanning the whole width)
                    if w > self._slide_width * 0.50:
                        continue

                    if center_x < slide_center_x:
                        left_texts.append((shape, center_x, center_y))
                    else:
                        right_texts.append((shape, center_x, center_y))
                except Exception:
                    continue

            # Need shapes on both sides for a timeline
            if not left_texts or not right_texts:
                return 0

            # Minimum 2 total shapes to qualify as timeline
            if len(left_texts) + len(right_texts) < 3:
                return 0

            # Step 3: Pair shapes at similar Y levels and swap X positions
            Y_TOLERANCE = self._slide_height * 0.05  # 5% of slide height

            paired_swaps = []
            used_right = set()
            for ls, lcx, lcy in left_texts:
                best_match = None
                best_dist = float('inf')
                for ri, (rs, rcx, rcy) in enumerate(right_texts):
                    if ri in used_right:
                        continue
                    dist = abs(lcy - rcy)
                    if dist < Y_TOLERANCE and dist < best_dist:
                        best_dist = dist
                        best_match = (ri, rs, rcx, rcy)
                if best_match:
                    ri, rs, rcx, rcy = best_match
                    used_right.add(ri)
                    paired_swaps.append((ls, rs))

            # Swap X positions for each pair
            for left_shape, right_shape in paired_swaps:
                ll = int(left_shape.left)
                lw = int(left_shape.width)
                rl = int(right_shape.left)
                rw = int(right_shape.width)

                # Swap: left gets right's X, right gets left's X
                # Account for different widths by keeping the gap from center
                left_shape.left = rl + rw - lw  # Align right edge where right shape's right edge was
                right_shape.left = ll  # Align left edge where left shape's left edge was
                changes += 2

            if changes:
                logger.debug(
                    'Fix 16 slide %d: reversed timeline alternation — %d pairs swapped',
                    slide_number, len(paired_swaps)
                )

        except Exception as exc:
            logger.debug('_reverse_timeline_alternation: %s', exc)

        return changes

    # ─────────────────────────────────────────────────────────────────────
    # Fix 17: Reverse logo row ordering
    # ─────────────────────────────────────────────────────────────────────

    def _reverse_logo_row_order(self, shapes: list, slide_number: int) -> int:
        """
        Fix 17: Reverse the horizontal order of logos in a trust/partner logo row.

        Detection: A logo row is a group of 3+ small images arranged horizontally
        at approximately the same Y position, typically in the lower portion
        of the slide.

        Fix: Reverse the X positions of the logos in the row (first becomes last,
        last becomes first), preserving the gaps between them.

        Affected decks: R6_18 s1 (client logos)
        """
        changes = 0
        try:
            # Collect small image shapes (potential logos)
            logo_candidates = []  # (shape, left, top, width, height)
            for shape in shapes:
                try:
                    sp_el = shape._element
                    tag = sp_el.tag
                    is_pic = tag.endswith('}pic')
                    # Also check for shapes with blipFill (images in shape containers)
                    has_blip = sp_el.find(f'.//{{{A_NS}}}blipFill') is not None
                    if not (is_pic or has_blip):
                        continue

                    w = getattr(shape, 'width', None)
                    h = getattr(shape, 'height', None)
                    l = getattr(shape, 'left', None)
                    t = getattr(shape, 'top', None)
                    if any(v is None for v in (w, h, l, t)):
                        continue
                    w, h, l, t = int(w), int(h), int(l), int(t)

                    # Logo size: small relative to slide (<25% width, <25% height)
                    if w > self._slide_width * 0.25 or h > self._slide_height * 0.25:
                        continue

                    # Skip very tiny shapes (icons, bullets)
                    if w < self._slide_width * 0.03:
                        continue

                    logo_candidates.append((shape, l, t, w, h))
                except Exception:
                    continue

            if len(logo_candidates) < 3:
                return 0

            # Group logos by similar Y position (within 5% of slide height)
            Y_TOL = self._slide_height * 0.05
            logo_rows = []
            used = set()

            for i, (s1, l1, t1, w1, h1) in enumerate(logo_candidates):
                if i in used:
                    continue
                row = [(s1, l1, t1, w1, h1)]
                used.add(i)
                center_y1 = t1 + h1 // 2
                for j, (s2, l2, t2, w2, h2) in enumerate(logo_candidates):
                    if j in used:
                        continue
                    center_y2 = t2 + h2 // 2
                    if abs(center_y1 - center_y2) < Y_TOL:
                        row.append((s2, l2, t2, w2, h2))
                        used.add(j)
                if len(row) >= 3:
                    logo_rows.append(row)

            # Reverse each logo row's X positions
            for row in logo_rows:
                # Sort by current X position (left to right)
                row_sorted = sorted(row, key=lambda r: r[1])
                # Get the X positions in order
                x_positions = [r[1] for r in row_sorted]
                # Reverse: assign the reversed X positions
                n = len(row_sorted)
                for idx, (shape, orig_l, orig_t, orig_w, orig_h) in enumerate(row_sorted):
                    # New position: take the X from the mirror index
                    mirror_idx = n - 1 - idx
                    target_x = x_positions[mirror_idx]
                    # Adjust for width difference between this logo and the one at mirror position
                    target_w = row_sorted[mirror_idx][3]
                    # Center the logo at the target position's center
                    target_center = target_x + target_w // 2
                    new_left = target_center - orig_w // 2
                    if new_left != orig_l:
                        shape.left = new_left
                        changes += 1

                if changes:
                    logger.debug(
                        'Fix 17 slide %d: reversed logo row of %d logos',
                        slide_number, n
                    )

        except Exception as exc:
            logger.debug('_reverse_logo_row_order: %s', exc)

        return changes

    # ─────────────────────────────────────────────────────────────────────
    # Fix 18: Center text in circular/bounded container shapes
    # ─────────────────────────────────────────────────────────────────────

    def _center_text_in_container_shapes(self, shapes: list, slide_number: int) -> int:
        """
        Fix 18: When Arabic text is placed inside a circle, oval, or rounded
        container shape, it should be centered rather than right-aligned.

        Detection:
        - Shape has a preset geometry of circle, ellipse, roundRect, etc.
        - Shape has a text frame with Arabic text
        - Shape's aspect ratio is roughly 1:1 (circular) or shape is small

        Fix: Set center alignment on all paragraphs within the shape.

        Affected decks: R6_09 s1 (title in circle)
        """
        changes = 0
        _CONTAINER_PRESETS = frozenset({
            'ellipse', 'roundRect', 'snipRndRect', 'flowChartConnector',
            'flowChartTerminator', 'actionButtonBlank', 'pie', 'donut',
            'blockArc', 'cloud', 'octagon', 'hexagon', 'diamond',
        })

        try:
            for shape in shapes:
                try:
                    if not getattr(shape, 'has_text_frame', False):
                        continue
                    text = (shape.text_frame.text or '').strip()
                    if not text or not has_arabic(text):
                        continue

                    sp_el = shape._element
                    prst_geom = sp_el.find(f'.//{{{A_NS}}}prstGeom')
                    if prst_geom is None:
                        continue
                    prst = prst_geom.get('prst', '')

                    # Check if it's a container shape
                    is_container = prst in _CONTAINER_PRESETS

                    # Also detect circular shapes by aspect ratio
                    if not is_container:
                        w = getattr(shape, 'width', None)
                        h = getattr(shape, 'height', None)
                        if w and h and int(w) > 0 and int(h) > 0:
                            ratio = int(w) / int(h)
                            # Nearly square (circle-like) and small
                            if 0.8 <= ratio <= 1.2 and int(w) < self._slide_width * 0.35:
                                is_container = True

                    if not is_container:
                        continue

                    # Center-align all paragraphs in this shape
                    for para in shape.text_frame.paragraphs:
                        pPr = ensure_pPr(para._p)
                        pPr.set('algn', 'ctr')
                        changes += 1

                    logger.debug(
                        'Fix 18 slide %d: centered text in container shape "%s" (prst=%s)',
                        slide_number, text[:20], prst
                    )
                except Exception:
                    continue

        except Exception as exc:
            logger.debug('_center_text_in_container_shapes: %s', exc)

        return changes

    # ─────────────────────────────────────────────────────────────────────
    # Fix 19: Bidi base direction for mixed Arabic/English titles
    # ─────────────────────────────────────────────────────────────────────

    def _fix_bidi_base_direction(self, shapes: list, slide_number: int) -> int:
        """
        Fix 19: For text frames containing mixed Arabic and English text,
        ensure the text frame and all paragraphs have RTL base direction
        set explicitly so the Unicode bidi algorithm renders correctly.

        When Arabic text contains embedded English brand names, the bidi
        algorithm needs an explicit RTL base direction to avoid visual
        reordering issues.

        Fix: Set bidi-related attributes on the body and paragraph level:
        - bodyPr: rtlCol="1"
        - pPr: rtl="1"
        - Insert RLE (Right-to-Left Embedding U+202B) before mixed segments

        Affected decks: R6_11 s1 (mixed Arabic/English title)
        """
        changes = 0
        try:
            for shape in shapes:
                try:
                    if not getattr(shape, 'has_text_frame', False):
                        continue
                    text = (shape.text_frame.text or '').strip()
                    if not text:
                        continue

                    # Only process shapes with BOTH Arabic and Latin text
                    has_ar = has_arabic(text)
                    if not has_ar:
                        continue
                    # Check for Latin characters (A-Z, a-z)
                    has_latin = any(c.isascii() and c.isalpha() for c in text)
                    if not has_latin:
                        continue

                    # Get font size to focus on titles (>20pt = 2000)
                    max_font = 0
                    sp_el = shape._element
                    for rPr in sp_el.iter(f'{{{A_NS}}}rPr'):
                        sz_str = rPr.get('sz')
                        if sz_str:
                            try:
                                max_font = max(max_font, int(sz_str))
                            except ValueError:
                                pass

                    # Only apply to title-sized text (>= 16pt)
                    if max_font > 0 and max_font < 1600:
                        continue

                    # Ensure body-level RTL
                    body_pr = shape.text_frame._txBody.find(f'{{{A_NS}}}bodyPr')
                    if body_pr is not None:
                        if body_pr.get('rtlCol') != '1':
                            body_pr.set('rtlCol', '1')
                            changes += 1

                    # Ensure paragraph-level RTL and right alignment
                    for para in shape.text_frame.paragraphs:
                        pPr = ensure_pPr(para._p)
                        if pPr.get('rtl') != '1':
                            pPr.set('rtl', '1')
                            changes += 1
                        # Set alignment to right for mixed text
                        current_algn = pPr.get('algn', '')
                        if current_algn not in ('r', 'ctr'):
                            pPr.set('algn', 'r')
                            changes += 1

                    if changes:
                        logger.debug(
                            'Fix 19 slide %d: bidi direction set on mixed text "%s"',
                            slide_number, text[:40]
                        )
                except Exception:
                    continue

        except Exception as exc:
            logger.debug('_fix_bidi_base_direction: %s', exc)

        return changes

    # ─────────────────────────────────────────────────────────────────────
    # Arabic text auto-fit: scale font when placeholder was shrunk
    # ─────────────────────────────────────────────────────────────────────

    def _enable_autofit(self, shape, shrink_ratio: float) -> None:
        """
        Enable text auto-fit when a placeholder was shrunk to avoid logo overlap.

        Two-pronged approach for maximum compatibility:
        1. Set <a:normAutofit fontScale="..."/> on bodyPr (for PowerPoint).
        2. Directly reduce font sizes on all runs (for LibreOffice, which
           ignores normAutofit during PDF export).

        The font size reduction accounts for Arabic text being typically
        20-30% wider than equivalent English content.

        Args:
            shape: python-pptx shape with a text frame.
            shrink_ratio: Ratio of new_width / original_width (0.0 to 1.0).
        """
        try:
            if not getattr(shape, 'has_text_frame', False):
                return
            tf = shape.text_frame
            body_pr = tf._txBody.find(f'{{{A_NS}}}bodyPr')
            if body_pr is None:
                return

            # Remove any existing autofit/spAutoFit/noAutofit children
            for child_tag in ('normAutofit', 'spAutoFit', 'noAutofit'):
                for child in body_pr.findall(f'{{{A_NS}}}{child_tag}'):
                    body_pr.remove(child)

            # Calculate font scale: use the shrink ratio with 10% extra headroom
            # for Arabic expansion. Minimum scale = 50% to keep text readable.
            font_scale_pct = max(shrink_ratio * 0.85, 0.50)  # 0.0 to 1.0
            font_scale_ooxml = max(int(font_scale_pct * 100000), 50000)

            autofit = etree.SubElement(body_pr, f'{{{A_NS}}}normAutofit')
            autofit.set('fontScale', str(font_scale_ooxml))

            # Also directly scale font sizes on runs for LibreOffice compat.
            # LibreOffice ignores normAutofit during headless PDF export.
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        original_size = int(run.font.size)
                        new_size = max(int(original_size * font_scale_pct), 8 * 12700)  # min 8pt
                        run.font.size = new_size
                    else:
                        # No explicit size — read from defRPr or use common title size
                        # For title placeholders, typical size is 36-44pt
                        # Set an explicit reduced size based on common defaults
                        rPr = run._r.find(f'{{{A_NS}}}rPr')
                        if rPr is None:
                            rPr = etree.SubElement(run._r, f'{{{A_NS}}}rPr')
                            run._r.insert(0, rPr)
                        # Check defRPr for the paragraph's default font size
                        def_rPr = para._p.find(f'.//{{{A_NS}}}defRPr')
                        default_sz = None
                        if def_rPr is not None:
                            sz_str = def_rPr.get('sz')
                            if sz_str:
                                default_sz = int(sz_str)
                        # Also check endParaRPr
                        if default_sz is None:
                            end_rPr = para._p.find(f'{{{A_NS}}}endParaRPr')
                            if end_rPr is not None:
                                sz_str = end_rPr.get('sz')
                                if sz_str:
                                    default_sz = int(sz_str)
                        if default_sz:
                            new_sz = max(int(default_sz * font_scale_pct), 800)
                            rPr.set('sz', str(new_sz))

            logger.debug('Auto-fit enabled on "%s": fontScale=%d%%, runs scaled',
                         getattr(shape, 'name', '?'), int(font_scale_pct * 100))

        except Exception as exc:
            logger.debug('_enable_autofit: %s', exc)

    # ─────────────────────────────────────────────────────────────────────
    # Fix 20: Slide-number badge repositioning for RTL
    # ─────────────────────────────────────────────────────────────────────

    def _reposition_slide_number_badge(self, shapes: list, slide_number: int) -> int:
        """
        Move the slide-number badge from top-right to top-left for RTL output.

        In LTR, slide number badges sit in the top-right corner. In RTL, this
        collides with the title text which is also right-aligned. Moving the
        badge to the top-left avoids the collision.

        Detection: A small shape (<8% slide width, <8% slide height) in the
        top-right zone containing only a 1-3 digit number.
        """
        changes = 0
        try:
            for shape in shapes:
                try:
                    l = getattr(shape, 'left', None)
                    t = getattr(shape, 'top', None)
                    w = getattr(shape, 'width', None)
                    h = getattr(shape, 'height', None)
                    if any(v is None for v in (l, t, w, h)):
                        continue
                    l, t, w, h = int(l), int(t), int(w), int(h)

                    # Badge is small
                    if w > self._slide_width * 0.08 or h > self._slide_height * 0.08:
                        continue
                    # In top zone
                    if t > self._slide_height * 0.15:
                        continue
                    # Near right edge
                    right_edge = l + w
                    if right_edge < self._slide_width * 0.85:
                        continue

                    # Contains only a short number
                    if not getattr(shape, 'has_text_frame', False):
                        continue
                    text = (shape.text_frame.text or '').strip()
                    if not text or len(text) > 4:
                        continue
                    if not text.replace(' ', '').isdigit():
                        continue

                    # Move to top-left: use same margin from left as original had from right
                    margin = self._slide_width - right_edge
                    shape.left = margin
                    changes += 1
                    logger.debug(
                        'Fix 20 slide %d: moved badge "%s" from x=%d to x=%d',
                        slide_number, text, l, margin
                    )
                except Exception:
                    continue
        except Exception as exc:
            logger.debug('_reposition_slide_number_badge: %s', exc)
        return changes

    # ─────────────────────────────────────────────────────────────────────
    # Fix 9: Collision detection
    # ─────────────────────────────────────────────────────────────────────

    def _detect_collisions(self, shapes: List, slide_number: int) -> None:
        """
        Fix 9: Detect overlapping shapes after transformation and log warnings.
        Helps diagnose logo-heading collisions and footer overlaps.
        Limited to first 30 non-background shapes for performance.
        """
        # Performance guard: skip collision detection on complex slides
        if len(shapes) > 50:
            return

        rects = []
        for shape in shapes:
            try:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                if any(v is None for v in (left, top, width, height)):
                    continue
                # Skip full-slide background shapes
                if width > self._slide_width * 0.90 and height > self._slide_height * 0.90:
                    continue
                rects.append((
                    getattr(shape, 'name', '?'),
                    left, top, left + width, top + height
                ))
                if len(rects) >= 30:
                    break  # Cap for performance
            except Exception:
                continue

        collisions_found = 0
        for i, (name1, l1, t1, r1, b1) in enumerate(rects):
            if collisions_found >= 3:
                break  # Cap collision reports per slide
            for j, (name2, l2, t2, r2, b2) in enumerate(rects):
                if i >= j:
                    continue
                # Check overlap
                if l1 < r2 and r1 > l2 and t1 < b2 and b1 > t2:
                    overlap_w = min(r1, r2) - max(l1, l2)
                    overlap_h = min(b1, b2) - max(t1, t2)
                    overlap_area = overlap_w * overlap_h
                    # Only log very significant overlaps (> 2 sq inches)
                    if overlap_area > 2_000_000_000_000:  # ~2 sq inches
                        logger.debug(
                            'Slide %d: collision — "%s" overlaps "%s"',
                            slide_number, name1, name2
                        )
                        collisions_found += 1

    # ─────────────────────────────────────────────────────────────────────
    # Fix 10: Chevron/arrow direction reversal
    # ─────────────────────────────────────────────────────────────────────

    # Map of directional preset geometry types that need reversal
    _DIRECTIONAL_PRESETS = {
        # Type-swap pairs: swap the preset type name
        'rightArrow': 'leftArrow',
        'leftArrow': 'rightArrow',
        'rightArrowCallout': 'leftArrowCallout',
        'leftArrowCallout': 'rightArrowCallout',
        'curvedRightArrow': 'curvedLeftArrow',
        'curvedLeftArrow': 'curvedRightArrow',
        'leftRightArrow': 'leftRightArrow',  # symmetric — no change
        'upDownArrow': 'upDownArrow',  # symmetric — no change
        # FlipH-based: same type but flip horizontally
        'chevron': '_flipH',
        'homePlate': '_flipH',
        'notchedRightArrow': '_flipH',
        'stripedRightArrow': '_flipH',
        'bentArrow': '_flipH',
        'bentUpArrow': '_flipH',
        'circularArrow': '_flipH',
        'pentagon': '_flipH',
    }

    def _reverse_directional_shape(self, shape) -> int:
        """
        Fix 10: Reverse the direction of arrow/chevron shapes for RTL.

        For shapes with directional preset geometry:
        - Some get their preset type swapped (rightArrow → leftArrow)
        - Others get flipH toggled (chevron, homePlate, etc.)

        Returns:
            1 if the shape was modified, 0 otherwise.
        """
        try:
            sp_el = shape._element
            # Get the preset geometry type
            prst_geom = sp_el.find(f'.//{{{A_NS}}}prstGeom')
            if prst_geom is None:
                return 0

            prst = prst_geom.get('prst', '')
            action = self._DIRECTIONAL_PRESETS.get(prst)

            if action is None:
                return 0

            if action == '_flipH':
                # Toggle flipH on the shape's xfrm
                xfrm = sp_el.find(f'.//{{{A_NS}}}xfrm')
                if xfrm is not None:
                    current = xfrm.get('flipH', '0')
                    xfrm.set('flipH', '0' if current == '1' else '1')
                    return 1
            elif action == prst:
                # Symmetric shape — no change needed
                return 0
            else:
                # Swap to mirrored preset type
                prst_geom.set('prst', action)
                return 1

        except Exception as exc:
            logger.debug('_reverse_directional_shape: %s', exc)
        return 0

    def _reverse_connector_direction(self, shape) -> int:
        """Reverse connector arrow direction for RTL by swapping headEnd/tailEnd markers."""
        try:
            sp_el = shape._element
            if not (sp_el.tag.endswith('}cxnSp') or sp_el.tag == 'cxnSp'):
                return 0

            a_ns = A_NS
            ln = sp_el.find(f'.//{{{a_ns}}}ln')
            if ln is None:
                return 0

            head = ln.find(f'{{{a_ns}}}headEnd')
            tail = ln.find(f'{{{a_ns}}}tailEnd')

            if head is None and tail is None:
                return 0

            head_copy = deepcopy(head) if head is not None else None
            tail_copy = deepcopy(tail) if tail is not None else None

            if head is not None:
                ln.remove(head)
            if tail is not None:
                ln.remove(tail)

            # Swap: old tail becomes new head, old head becomes new tail
            if tail_copy is not None:
                tail_copy.tag = f'{{{a_ns}}}headEnd'
                ln.append(tail_copy)
            if head_copy is not None:
                head_copy.tag = f'{{{a_ns}}}tailEnd'
                ln.append(head_copy)

            logger.debug("Reversed connector arrowheads on %s", getattr(shape, 'name', '?'))
            return 1
        except (AttributeError, KeyError, ValueError, TypeError) as exc:
            logger.debug('_reverse_connector_direction: %s on %s', exc, getattr(shape, 'name', '?'))
            return 0

    def _reverse_line_arrowheads(self, shape) -> int:
        """Swap head/tail arrowhead markers on line shapes for RTL flow."""
        try:
            sp_el = shape._element
            a_ns = A_NS
            ln = sp_el.find(f'.//{{{a_ns}}}ln')
            if ln is None:
                return 0

            head_end = ln.find(f'{{{a_ns}}}headEnd')
            tail_end = ln.find(f'{{{a_ns}}}tailEnd')

            if head_end is None and tail_end is None:
                return 0

            head_attribs = dict(head_end.attrib) if head_end is not None else None
            tail_attribs = dict(tail_end.attrib) if tail_end is not None else None

            if head_end is not None:
                ln.remove(head_end)
            if tail_end is not None:
                ln.remove(tail_end)

            if tail_attribs:
                new_head = etree.SubElement(ln, f'{{{a_ns}}}headEnd')
                for k, v in tail_attribs.items():
                    new_head.set(k, v)
            if head_attribs:
                new_tail = etree.SubElement(ln, f'{{{a_ns}}}tailEnd')
                for k, v in head_attribs.items():
                    new_tail.set(k, v)

            return 1
        except (AttributeError, KeyError, ValueError, TypeError) as exc:
            logger.debug('_reverse_line_arrowheads: %s on %s', exc, getattr(shape, 'name', '?'))
            return 0


import logging
import re
import time
from pathlib import Path
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Callable, Any
from pptx import Presentation

logger = logging.getLogger(__name__)

# Broad Arabic detection covering all Arabic Unicode blocks:
# Basic Arabic (0600-06FF), Arabic Supplement (0750-077F),
# Arabic Extended-A (08A0-08FF), Presentation Forms-A (FB50-FDFF),
# Presentation Forms-B (FE70-FEFF)
_ARABIC_RE = re.compile(r"[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF\uFB50-\uFDFF\uFE70-\uFEFF]")

# ── Module imports — each isolated so one missing module doesn't break all ──

HAS_PROPERTY_RESOLVER = False
try:
    from slidearabi.property_resolver import PropertyResolver
    HAS_PROPERTY_RESOLVER = True
except ImportError as e:
    logger.warning(f"PropertyResolver not available: {e}")

HAS_LAYOUT_ANALYZER = False
try:
    from slidearabi.layout_analyzer import LayoutAnalyzer
    HAS_LAYOUT_ANALYZER = True
except ImportError as e:
    logger.warning(f"LayoutAnalyzer not available: {e}")

HAS_TEMPLATE_REGISTRY = False
try:
    from slidearabi.template_registry import TemplateRegistry
    HAS_TEMPLATE_REGISTRY = True
except ImportError as e:
    logger.warning(f"TemplateRegistry not available: {e}")

HAS_RTL_TRANSFORMS = False
try:
    from slidearabi.rtl_transforms import MasterLayoutTransformer, SlideContentTransformer
    HAS_RTL_TRANSFORMS = True
except ImportError as e:
    logger.warning(f"RTL transforms not available: {e}")

HAS_TYPOGRAPHY = False
try:
    from slidearabi.typography import TypographyNormalizer
    HAS_TYPOGRAPHY = True
except ImportError as e:
    logger.warning(f"TypographyNormalizer not available: {e}")

HAS_STRUCTURAL_VALIDATOR = False
try:
    from slidearabi.structural_validator import StructuralValidator
    HAS_STRUCTURAL_VALIDATOR = True
except ImportError as e:
    logger.warning(f"StructuralValidator not available: {e}")

try:
    from slidearabi.models import ResolvedPresentation, ValidationReport  # noqa: F401
except ImportError as e:
    logger.warning(f"Models not available (non-critical): {e}")


# ── PipelineConfig / PipelineResult — always available ──

@dataclass
class PipelineConfig:
    input_path: str
    output_path: str
    translate_fn: Optional[Callable[[List[str]], Dict[str, str]]] = None
    skip_translation: bool = False
    max_font_reduction_pct: float = 20.0
    log_level: str = 'INFO'
    enable_telemetry: bool = False

@dataclass
class PipelineResult:
    success: bool
    output_path: Optional[str]
    phase_reports: Dict[str, Any]
    validation_report: Optional[Any]
    total_duration_ms: float
    error: Optional[str] = None

# Phase 6 (VQA) and LLM translation — optional but wired when available
try:
    from slidearabi.visual_qa import run_vqa, VQAReport
    HAS_VQA = True
except ImportError:
    HAS_VQA = False

try:
    from slidearabi.llm_translator import DualLLMTranslator, TranslatorConfig
    HAS_LLM_TRANSLATOR = True
except ImportError:
    HAS_LLM_TRANSLATOR = False


class SlideArabiPipeline:
    """
    Template-first, deterministic RTL transformation pipeline.
    
    Phases execute as a DAG — each phase runs exactly once and produces
    immutable output consumed by subsequent phases. No fix loops.
    
    Phase 0: Parse & Resolve → ResolvedPresentation
    Phase 1: Translate (AI — GPT-5.2 primary, Claude 4.6 QA) → translation_map
    Phase 2: Master & Layout Transformation (deterministic) → TransformedPresentation
    Phase 3: Slide Content Transformation (deterministic) → TransformedPresentation
    Phase 4: Typography Normalization (deterministic) → FinalPresentation
    Phase 5: Structural Validation (read-only) → ValidationReport
    Phase 6: Visual QA (Gemini 3.1 Pro + Claude 4.6 dual-pass) → VQAReport
    """
    
    def __init__(self, config: 'PipelineConfig'):
        self.config = config
        self._phase_reports = {}
        
        # Configure logging
        numeric_level = getattr(logging, config.log_level.upper(), logging.INFO)
        logger.setLevel(numeric_level)
        
    def run(self) -> 'PipelineResult':
        """
        Execute the full pipeline. Returns PipelineResult.
        
        This is synchronous. Each phase completes before the next begins.
        Phases are NOT async because there's no I/O wait between them
        (except Phase 1 translation which calls an external API).
        """
        pipeline_start = time.monotonic()
        logger.info(f"Starting SlideArabi pipeline for {self.config.input_path}")
        
        try:
            # 1. Load presentation
            logger.info("Loading presentation...")
            try:
                prs = Presentation(self.config.input_path)
                self._prs = prs  # Store for supplementary text extraction
            except Exception as e:
                raise ValueError(f"Failed to load presentation: {e}")
            
            # Phase 0: Parse & Resolve
            resolved_prs = self._phase_0_resolve(prs)
            
            # Phase 1: Translate
            translation_map = self._phase_1_translate(resolved_prs)
            
            # Phase 2: Transform Masters & Layouts
            p2_report = self._phase_2_transform_masters_layouts(prs, resolved_prs)
            
            # Phase 3: Transform Slide Content
            p3_report = self._phase_3_transform_slides(prs, resolved_prs, translation_map)

            # ── PER-SLIDE COVERAGE CHECK (v1.1.3) ────────────────────────
            if not self.config.skip_translation and translation_map:
                untranslated = self._validate_slide_coverage(prs)
                if untranslated:
                    logger.error(
                        "QUALITY ALERT: Slides %s appear untranslated after Phase 3",
                        untranslated,
                    )
                    self._phase_reports['slide_coverage'] = {
                        "untranslated_slides": untranslated,
                    }

            # ── POST-TRANSFORM VERIFICATION ──────────────────────────────
            # Belt-and-suspenders: sample the in-memory prs for Arabic text
            # before saving. If translation_map had entries but no Arabic
            # text ended up in the presentation, something went wrong in
            # Phase 3's apply logic.
            # v1.1.1: Enhanced post-transform check — verify Arabic RATIO, not just existence.
            # A single cached Arabic string shouldn't mask 8 untranslated English slides.
            if not self.config.skip_translation and translation_map:
                arabic_text_count = 0
                total_text_count = 0
                sample_texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                text = para.text.strip()
                                if text and len(text) > 3:
                                    total_text_count += 1
                                    if _ARABIC_RE.search(text):
                                        arabic_text_count += 1
                                    elif len(sample_texts) < 10:
                                        sample_texts.append(text[:80])
                
                if total_text_count == 0:
                    logger.warning("Post-transform: no text found in output (empty presentation?)")
                elif arabic_text_count == 0:
                    raise RuntimeError(
                        "FATAL: Post-transform verification failed. "
                        f"Translation map has {len(translation_map)} entries but "
                        "NO Arabic text found in output presentation. "
                        f"Sample non-Arabic texts: {sample_texts[:5]}"
                    )
                else:
                    arabic_ratio = arabic_text_count / total_text_count
                    logger.info(
                        "Post-transform verification: %d/%d text frames Arabic (%.1f%%)",
                        arabic_text_count, total_text_count, arabic_ratio * 100
                    )
                    if arabic_ratio < 0.3:
                        # Less than 30% Arabic is a warning — translation partially failed
                        logger.warning(
                            "Post-transform WARNING: Only %.1f%% Arabic coverage. "
                            "Expected >50%%. Non-Arabic samples: %s",
                            arabic_ratio * 100, sample_texts[:5]
                        )
            
            # Phase 4: Typography Normalization
            p4_report = self._phase_4_typography(prs)
            
            # Phase 5: Validate
            val_report = self._phase_5_validate(prs, resolved_prs)
            
            # Save presentation
            logger.info(f"Saving transformed presentation to {self.config.output_path}...")
            try:
                # Ensure directory exists
                Path(self.config.output_path).parent.mkdir(parents=True, exist_ok=True)
                prs.save(self.config.output_path)
            except Exception as e:
                raise IOError(f"Failed to save presentation: {e}")
            
            # Phase 6: Visual QA (dual-pass: Gemini + Claude)
            vqa_report = self._phase_6_vqa()
            
            total_duration = (time.monotonic() - pipeline_start) * 1000
            logger.info(f"Pipeline completed successfully in {total_duration:.0f}ms")
            
            return PipelineResult(
                success=True,
                output_path=self.config.output_path,
                phase_reports=self._phase_reports,
                validation_report=val_report,
                total_duration_ms=total_duration
            )
            
        except Exception as e:
            logger.error(f"Pipeline failed: {e}", exc_info=True)
            total_duration = (time.monotonic() - pipeline_start) * 1000
            
            return PipelineResult(
                success=False,
                output_path=None,
                phase_reports=self._phase_reports,
                validation_report=None,
                total_duration_ms=total_duration,
                error=str(e)
            )
        
    def _phase_0_resolve(self, prs: Presentation) -> Any:
        """Phase 0: Parse the presentation and resolve all inherited properties.
        
        FAIL-CLOSED: If PropertyResolver is unavailable and translation is
        enabled, this raises immediately. A missing resolver means zero text
        extraction, zero translations, and an English-copy output — which is
        a silent failure the user would see as "not working".
        """
        start_time = time.monotonic()
        logger.info("Phase 0: Resolving properties...")
        
        if not HAS_PROPERTY_RESOLVER:
            if not self.config.skip_translation:
                raise RuntimeError(
                    "FATAL: PropertyResolver module not available but translation is enabled. "
                    "Cannot extract text without PropertyResolver. "
                    "Check that slidearabi.property_resolver imports successfully on this environment."
                )
            logger.warning("Phase 0 skipped: PropertyResolver not available (translation also skipped)")
            self._log_phase('phase_0_resolve', 0, {"status": "module_unavailable"})
            return None

        resolver = PropertyResolver(prs)
        resolved_prs = resolver.resolve_presentation()
        
        if not resolved_prs or not getattr(resolved_prs, 'slides', None):
            raise RuntimeError(
                f"FATAL: PropertyResolver returned empty result. "
                f"resolved_prs={resolved_prs}, "
                f"slides={getattr(resolved_prs, 'slides', 'N/A')}"
            )
        
        duration = (time.monotonic() - start_time) * 1000
        self._log_phase('phase_0_resolve', duration, {
            "status": "success",
            "slides_resolved": len(prs.slides),
        })
        return resolved_prs
            
    def _phase_1_translate(self, resolved: Any) -> Dict[str, str]:
        """Phase 1: Extract text and call translation function."""
        start_time = time.monotonic()
        logger.info("Phase 1: Translation...")
        
        if self.config.skip_translation:
            logger.info("Translation skipped per config.")
            self._log_phase('phase_1_translate', 0, {"status": "skipped"})
            return {}
            
        if not self.config.translate_fn:
            logger.warning("No translation function provided.")
            self._log_phase('phase_1_translate', 0, {"status": "no_function"})
            return {}
            
        texts_to_translate = self._extract_texts(resolved)
        logger.info(f"Extracted {len(texts_to_translate)} strings for translation.")
        
        # ── ZERO-TOLERANCE GATE: refuse to continue with zero texts ──
        if not texts_to_translate:
            raise RuntimeError(
                "FATAL: _extract_texts returned 0 strings. "
                f"resolved is {'None' if resolved is None else 'valid (' + str(len(getattr(resolved, 'slides', []))) + ' slides)'}. "
                f"HAS_PROPERTY_RESOLVER={HAS_PROPERTY_RESOLVER}. "
                "Refusing to proceed — output would be an untranslated English copy."
            )
        
        try:
            translation_map = self.config.translate_fn(texts_to_translate)
            
            # ── ZERO-TOLERANCE GATE: refuse empty translation map ──
            if not translation_map:
                raise RuntimeError(
                    f"FATAL: translate_fn returned 0 translations for "
                    f"{len(texts_to_translate)} input strings. "
                    "API call may have failed silently. "
                    "Refusing to save untranslated output."
                )
            
            coverage_pct = (len(translation_map) / len(texts_to_translate)) * 100
            logger.info(
                "Translation coverage: %d/%d strings (%.1f%%)",
                len(translation_map), len(texts_to_translate), coverage_pct,
            )
            if coverage_pct < 10:
                raise RuntimeError(
                    f"FATAL: Translation coverage critically low ({coverage_pct:.1f}%). "
                    f"Only {len(translation_map)}/{len(texts_to_translate)} strings translated. "
                    "Output would be mostly English."
                )
            
            # ── ZERO-TOLERANCE GATE: validate that translation values are actually Arabic ──
            # This catches the case where translation_map has entries but they
            # contain English text instead of Arabic (e.g., from a fallback path).
            arabic_value_count = 0
            sampled_count = 0
            sample_non_arabic = []
            for eng, ar in translation_map.items():
                eng_stripped = eng.strip()
                # Skip very short strings, numbers-only, abbreviations
                if (len(eng_stripped) <= 3
                        or re.match(r'^[\d\s\-+.,/%$\u20ac\u00a3]+$', eng_stripped)):
                    continue
                sampled_count += 1
                if _ARABIC_RE.search(ar):
                    arabic_value_count += 1
                elif len(sample_non_arabic) < 5:
                    sample_non_arabic.append(
                        f"EN: {eng_stripped[:60]} -> GOT: {ar[:60]}"
                    )
            
            if sampled_count > 0:
                arabic_ratio = arabic_value_count / sampled_count
                logger.info(
                    "Translation Arabic content ratio: %d/%d (%.1f%%)",
                    arabic_value_count, sampled_count, arabic_ratio * 100,
                )
                if arabic_ratio < 0.5:
                    raise RuntimeError(
                        f"FATAL: Translation content check failed. "
                        f"Only {arabic_value_count}/{sampled_count} translated strings "
                        f"({arabic_ratio*100:.1f}%) contain Arabic characters. "
                        f"This indicates the translation API returned English text. "
                        f"Samples without Arabic: {sample_non_arabic}"
                    )
            
            duration = (time.monotonic() - start_time) * 1000
            self._log_phase('phase_1_translate', duration, {
                "status": "success", 
                "strings_extracted": len(texts_to_translate),
                "strings_translated": len(translation_map),
                "coverage_pct": round(coverage_pct, 1),
            })
            return translation_map
            
        except RuntimeError:
            raise  # Re-raise our own zero-tolerance errors
        except Exception as e:
            logger.error(f"Translation failed: {e}")
            raise RuntimeError(f"Phase 1 (Translation) failed: {e}")
        
    def _phase_2_transform_masters_layouts(self, prs: Presentation, resolved: Any) -> Any:
        """Phase 2: Transform masters and layouts deterministically."""
        start_time = time.monotonic()
        logger.info("Phase 2: Transforming masters and layouts...")
        
        if not (HAS_TEMPLATE_REGISTRY and HAS_RTL_TRANSFORMS):
            logger.warning("Phase 2 skipped: TemplateRegistry or RTL transforms not available")
            self._log_phase('phase_2_transform_masters', 0, {"status": "module_unavailable"})
            return None

        registry = TemplateRegistry(prs.slide_width, prs.slide_height)
        transformer = MasterLayoutTransformer(prs, registry)
        
        master_report = transformer.transform_all_masters()
        layout_report = transformer.transform_all_layouts()
        
        duration = (time.monotonic() - start_time) * 1000
        self._log_phase('phase_2_transform_masters', duration, {
            "status": "success",
            "master_changes": master_report.total_changes,
            "layout_changes": layout_report.total_changes,
        })
        return layout_report
        
    def _phase_3_transform_slides(self, prs: Presentation, resolved: Any, translations: Dict[str, str]) -> Any:
        """Phase 3: Transform slide content deterministically."""
        start_time = time.monotonic()
        logger.info("Phase 3: Transforming slide content...")
        
        if not HAS_RTL_TRANSFORMS:
            logger.warning("Phase 3 skipped: RTL transforms not available")
            self._log_phase('phase_3_transform_slides', 0, {"status": "module_unavailable"})
            return None

        # Layout analysis — get per-slide layout classifications
        layout_classifications = {}
        if HAS_LAYOUT_ANALYZER:
            analyzer = LayoutAnalyzer(prs)
            raw_classifications = analyzer.analyze_all()
            # Convert LayoutClassification objects to simple type strings
            layout_classifications = {
                slide_num: cls.resolved_type
                for slide_num, cls in raw_classifications.items()
            }
        else:
            logger.warning("LayoutAnalyzer not available — proceeding without layout hints")

        # Build TemplateRegistry if available (for rule-based transforms)
        registry = None
        if HAS_TEMPLATE_REGISTRY:
            registry = TemplateRegistry(prs.slide_width, prs.slide_height)

        transformer = SlideContentTransformer(
            presentation=prs,
            template_registry=registry,
            layout_classifications=layout_classifications,
            translations=translations,
        )
        report = transformer.transform_all_slides()
        
        duration = (time.monotonic() - start_time) * 1000
        self._log_phase('phase_3_transform_slides', duration, {
            "status": "success",
            "total_changes": report.total_changes,
        })
        return report
        
    def _phase_4_typography(self, prs: Presentation) -> Any:
        """Phase 4: Normalize typography for Arabic."""
        start_time = time.monotonic()
        logger.info("Phase 4: Normalizing typography...")
        
        if not HAS_TYPOGRAPHY:
            logger.warning("Phase 4 skipped: TypographyNormalizer not available")
            self._log_phase('phase_4_typography', 0, {"status": "module_unavailable"})
            return None

        normalizer = TypographyNormalizer(prs)
        report = normalizer.normalize_all()
        
        duration = (time.monotonic() - start_time) * 1000
        self._log_phase('phase_4_typography', duration, {
            "status": "success",
            "total_changes": report.total_changes,
        })
        return report
        
    def _phase_5_validate(self, prs: Presentation, resolved: Any) -> Any:
        """Phase 5: Read-only structural validation."""
        start_time = time.monotonic()
        logger.info("Phase 5: Structural validation...")
        
        if not HAS_STRUCTURAL_VALIDATOR:
            logger.warning("Phase 5 skipped: StructuralValidator not available")
            self._log_phase('phase_5_validate', 0, {"status": "module_unavailable"})
            return None

        validator = StructuralValidator(prs, resolved)
        report = validator.validate()
        
        duration = (time.monotonic() - start_time) * 1000
        
        # NOTE: StructuralValidator.validate() returns structural_validator.ValidationReport
        # which has .passed (bool), .errors (int), .warnings (int) — NOT models.ValidationReport.
        log_data = {
            "status": "success",
            "passed": report.passed,
            "errors": getattr(report, 'errors', 0),
            "warnings": getattr(report, 'warnings', 0),
        }
        self._log_phase('phase_5_validate', duration, log_data)
        
        error_count = getattr(report, 'errors', 0)
        if not report.passed:
            logger.warning(f"Validation failed with {error_count} errors.")
            # We don't fail the pipeline, just report it
            
        return report
        
    def _walk_all_shapes_pptx(self, shapes):
        """Recursively yield all shapes including group children from a python-pptx shapes collection."""
        for shape in shapes:
            yield shape
            if hasattr(shape, 'shapes'):  # GroupShape
                yield from self._walk_all_shapes_pptx(shape.shapes)

    def _extract_texts(self, resolved: Any) -> List[str]:
        """Extract all translatable text strings from the resolved presentation.

        Uses two passes:
        1. Walk the resolved model (shapes → paragraphs → runs)
        2. Supplementary pass over raw python-pptx to capture table cell text
           and group shape children that the resolved model may miss.
        """
        texts = []
        if not resolved:
            return texts

        try:
            # Pass 1: Resolved model walk
            for slide in resolved.slides:
                for shape in slide.shapes:
                    for para in shape.paragraphs:
                        # Combine runs into a single translatable string per paragraph
                        para_text = "".join(run.text for run in para.runs if run.text).strip()
                        if para_text:
                            texts.append(para_text)

            # Deduplicate while preserving order
            seen = set()
            unique_texts = []
            for t in texts:
                if t not in seen:
                    seen.add(t)
                    unique_texts.append(t)

            # Pass 2: Supplementary extraction from raw python-pptx
            # Captures table cell text and group children missed by the resolver
            prs = getattr(self, '_prs', None)
            if prs is not None:
                table_count = 0
                group_count = 0
                for slide in prs.slides:
                    for shape in self._walk_all_shapes_pptx(slide.shapes):
                        # Text frames (catches group children)
                        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                text = para.text.strip()
                                if text and text not in seen:
                                    seen.add(text)
                                    unique_texts.append(text)
                                    group_count += 1
                        # Table cells
                        if hasattr(shape, 'has_table') and shape.has_table:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    if cell.text_frame:
                                        for para in cell.text_frame.paragraphs:
                                            text = para.text.strip()
                                            if text and text not in seen:
                                                seen.add(text)
                                                unique_texts.append(text)
                                                table_count += 1
                supplementary_total = table_count + group_count
                if supplementary_total > 0:
                    logger.info(
                        "Supplementary extraction found %d additional strings "
                        "(tables: %d, groups: %d)",
                        supplementary_total, table_count, group_count,
                    )

            return unique_texts

        except AttributeError as e:
            # Log the actual error so the zero-tolerance gate's error is debuggable
            logger.error(
                "_extract_texts failed due to unexpected structure in resolved presentation: %s",
                e,
                exc_info=True,
            )
            return []

    def _validate_slide_coverage(self, prs) -> List[int]:
        """Check each slide for Arabic content. Returns list of untranslated slide numbers."""
        untranslated_slides = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_text_count = 0
            arabic_count = 0
            for shape in slide.shapes:
                if getattr(shape, 'has_text_frame', False) and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = (para.text or '').strip()
                        if text and len(text) > 3:  # Skip very short strings
                            slide_text_count += 1
                            if _ARABIC_RE.search(text):
                                arabic_count += 1
                # Also check table cells
                if getattr(shape, 'has_table', False) and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                for para in cell.text_frame.paragraphs:
                                    text = (para.text or '').strip()
                                    if text and len(text) > 3:
                                        slide_text_count += 1
                                        if _ARABIC_RE.search(text):
                                            arabic_count += 1

            logger.info("Slide %d: %d/%d text elements contain Arabic", i, arabic_count, slide_text_count)

            if slide_text_count >= 3 and arabic_count == 0:
                logger.warning(
                    "Slide %d: 0/%d text elements contain Arabic — likely untranslated",
                    i, slide_text_count,
                )
                untranslated_slides.append(i)

        return untranslated_slides

    def _phase_6_vqa(
        self,
    ) -> Any:
        """Phase 6: Dual-pass Visual QA (Gemini 3.1 Pro + Claude Sonnet 4.6).
        
        Compares the original and converted PPTX visually:
        - Pass 1: Gemini evaluates slide composites in parallel
        - Pass 2: Claude adjudicates Gemini findings
        - Layer 3: Deterministic reconciliation
        - Remediates FAIL slides (font reduction, alignment, etc.)
        - Logs all issues to JSONL for root-cause analysis
        """
        start_time = time.monotonic()
        logger.info("Phase 6: Visual QA...")
        
        if not HAS_VQA:
            logger.warning("Phase 6 skipped: visual_qa module not available")
            self._log_phase('phase_6_vqa', 0, {"status": "module_unavailable"})
            return None
            
        if not self.config.input_path or not self.config.output_path:
            logger.warning("Phase 6 skipped: missing input/output paths")
            self._log_phase('phase_6_vqa', 0, {"status": "missing_paths"})
            return None
        
        # Set up issue log path alongside the output
        output_dir = Path(self.config.output_path).parent
        deck_stem = Path(self.config.output_path).stem
        issue_log_path = str(output_dir / f"{deck_stem}_vqa_issues.jsonl")
        
        try:
            report = run_vqa(
                original_pptx=self.config.input_path,
                converted_pptx=self.config.output_path,
                issue_log_path=issue_log_path,
                deck_name=deck_stem,
                enable_dual_pass=True,
            )
            
            duration = (time.monotonic() - start_time) * 1000
            self._log_phase('phase_6_vqa', duration, {
                "status": "success",
                "overall_rating": report.overall_rating.value,
                "pass_rate": round(report.pass_rate, 1),
                "slides_reviewed": report.slides_reviewed,
                "fail_count": report.fail_count,
                "remediation_attempted": report.remediation_attempted,
                "remediation_successful": report.remediation_successful,
                "issues_logged": report.issues_logged,
            })
            
            logger.info(report.summary())
            return report
            
        except Exception as e:
            logger.warning(f"Phase 6 VQA failed (non-fatal): {e}")
            duration = (time.monotonic() - start_time) * 1000
            self._log_phase('phase_6_vqa', duration, {
                "status": "error",
                "error": str(e),
            })
            return None

    def _log_phase(self, phase_name: str, duration_ms: float, report: Any):
        """Log phase completion with timing and summary."""
        self._phase_reports[phase_name] = {
            "duration_ms": duration_ms,
            "report": report
        }
        logger.debug(f"{phase_name} completed in {duration_ms:.1f}ms: {report}")

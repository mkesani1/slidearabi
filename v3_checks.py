"""
v3_checks.py — V3 XML Structural Checks & Auto-Fixes for SlideArabi

Sprint 1: Table checks (#1, #2, #9, #10, #11) + Fixes (#1, #2)
Sprint 2: Adds icon, page number, shape position, paragraph RTL, overlap
Sprint 3: Adds circular centering, master mirror, directional

All checks produce V3Defect objects. Gated by v3_config flags.
"""

from __future__ import annotations

import logging
import shutil
from copy import deepcopy
from dataclasses import field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from lxml import etree

logger = logging.getLogger(__name__)

# XML namespaces
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
SLIDEARABI_NS = 'https://slidearabi.ai/ns/transform'

# Idempotency marker attribute (shared between checks and fixes)
V3_PHYS_REV_MARKER = f'{{{SLIDEARABI_NS}}}v3PhysRev'

# Arabic character detection
import re
_ARABIC_RE = re.compile(r'[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF\uFB50-\uFDFF\uFE70-\uFEFF]')


def _has_arabic(text: str) -> bool:
    return bool(_ARABIC_RE.search(text))


def _get_cell_text(tc_elem) -> str:
    """Extract all text from a table cell element."""
    texts = []
    for t_elem in tc_elem.iter(f'{{{A_NS}}}t'):
        if t_elem.text:
            texts.append(t_elem.text)
    return ' '.join(texts).strip()


def _find_tables(slide_element) -> list:
    """Find all <a:tbl> elements in a slide."""
    return slide_element.findall(f'.//{{{A_NS}}}tbl')


def _get_table_rows(tbl) -> list:
    """Get all <a:tr> elements from a table."""
    return tbl.findall(f'{{{A_NS}}}tr')


def _get_row_cells(tr) -> list:
    """Get all <a:tc> elements from a table row."""
    return tr.findall(f'{{{A_NS}}}tc')


def _get_gridcols(tbl) -> list:
    """Get gridCol elements from table grid."""
    grid = tbl.find(f'{{{A_NS}}}tblGrid')
    if grid is None:
        return []
    return grid.findall(f'{{{A_NS}}}gridCol')


def _get_tbl_pr(tbl):
    """Get tblPr element from table."""
    return tbl.find(f'{{{A_NS}}}tblPr')


def _enumerate_slide_pics(element) -> list:
    """Enumerate all pic elements on a slide with deduplication.
    
    Shared between checks and fixes to ensure consistent indexing.
    Collects <p:pic> elements first, then adds any <a:blipFill> parents
    that aren't already in the set (prevents double-counting).
    """
    pics_set = set()
    pics = []
    for pic in element.findall(f'.//{{{P_NS}}}pic'):
        pics_set.add(id(pic))
        pics.append(pic)
    for parent in element.findall(f'.//{{{A_NS}}}blipFill/..'):
        if id(parent) not in pics_set:
            pics_set.add(id(parent))
            pics.append(parent)
    return pics


# ─────────────────────────────────────────────────────────────────────────────
# V3 XML CHECKS
# ─────────────────────────────────────────────────────────────────────────────

class V3XMLChecker:
    """V3 XML structural checks — produces V3Defect instances."""

    def __init__(self, slide_width: int = 9144000, slide_height: int = 6858000):
        self.slide_width = slide_width
        self.slide_height = slide_height

    def check_slide(
        self,
        slide_num: int,
        conv_slide_element,
        orig_slide_element=None,
    ) -> list:
        """Run all V3 checks on a single slide. Returns list of V3Defect."""
        from vqa_types import V3Defect
        defects: List[V3Defect] = []

        # Table checks (require both original and converted)
        if orig_slide_element is not None:
            defects.extend(self._check_table_column_order(
                slide_num, orig_slide_element, conv_slide_element))
            defects.extend(self._check_table_gridcol_reversal(
                slide_num, orig_slide_element, conv_slide_element))

        # Table checks (converted only)
        defects.extend(self._check_table_cell_alignment(slide_num, conv_slide_element))
        defects.extend(self._check_merged_cell_integrity(slide_num, conv_slide_element))
        defects.extend(self._check_table_grid_structural(slide_num, conv_slide_element))

        # Sprint 2 checks
        if orig_slide_element is not None:
            defects.extend(self._check_icon_table_correspondence(
                slide_num, orig_slide_element, conv_slide_element))
            defects.extend(self._check_shape_position_mirroring(
                slide_num, orig_slide_element, conv_slide_element))

        # Converted-only Sprint 2 checks
        defects.extend(self._check_page_number_duplication(slide_num, conv_slide_element))
        defects.extend(self._check_paragraph_rtl(slide_num, conv_slide_element))

        # Sprint 3 checks
        defects.extend(self._check_circular_text_centering(slide_num, conv_slide_element))
        defects.extend(self._check_directional_symbol_orientation(slide_num, conv_slide_element))
        if orig_slide_element is not None:
            defects.extend(self._check_master_element_mirroring(
                slide_num, orig_slide_element, conv_slide_element))

        return defects

    # ── Check #1: Table Column Order Reversal (CRITICAL) ──

    def _check_table_column_order(
        self, slide_num: int, orig_element, conv_element
    ) -> list:
        """Compare original vs converted table to detect unreversed columns."""
        from vqa_types import V3Defect, Severity, DefectStatus

        defects = []
        orig_tables = _find_tables(orig_element)
        conv_tables = _find_tables(conv_element)

        for tbl_idx, (orig_tbl, conv_tbl) in enumerate(
            zip(orig_tables, conv_tables)
        ):
            orig_rows = _get_table_rows(orig_tbl)
            conv_rows = _get_table_rows(conv_tbl)

            if not orig_rows or not conv_rows:
                continue

            # Use first row with meaningful text content as anchor
            anchor_orig = None
            anchor_conv = None
            for orig_row, conv_row in zip(orig_rows, conv_rows):
                orig_cells = _get_row_cells(orig_row)
                conv_cells = _get_row_cells(conv_row)
                if len(orig_cells) < 2 or len(conv_cells) < 2:
                    continue

                orig_texts = [_get_cell_text(c) for c in orig_cells]
                conv_texts = [_get_cell_text(c) for c in conv_cells]

                # Skip rows where all cells are empty or identical
                non_empty = [t for t in orig_texts if t.strip()]
                if len(non_empty) < 2:
                    continue

                # Check if texts have enough unique content to compare
                unique_orig = set(orig_texts)
                if len(unique_orig) < 2:
                    continue  # All same text — can't determine order

                anchor_orig = orig_texts
                anchor_conv = conv_texts
                break

            if anchor_orig is None or anchor_conv is None:
                continue

            # Check: if original order matches converted (not reversed), flag it
            # After proper RTL conversion, column order should be reversed
            if anchor_orig == anchor_conv:
                # Columns NOT reversed — this is a defect
                # Check idempotency marker (matches Fix #1 marker)
                conv_tbl_pr = _get_tbl_pr(conv_tbl)
                if conv_tbl_pr is not None and conv_tbl_pr.get(V3_PHYS_REV_MARKER) == '1':
                    continue  # Already processed by V3 fix

                defects.append(V3Defect(
                    code="TABLE_COLUMNS_NOT_REVERSED",
                    category="table",
                    severity=Severity.CRITICAL,
                    slide_idx=slide_num,
                    object_id=str(tbl_idx),
                    evidence={
                        'original_order': anchor_orig,
                        'converted_order': anchor_conv,
                        'row_count': len(orig_rows),
                        'col_count': len(anchor_orig),
                    },
                    fixable=True,
                    autofix_action="reverse_table_columns",
                    description=f"Table {tbl_idx+1} columns not reversed for RTL "
                                f"(slide {slide_num})",
                ))

        return defects

    # ── Check #2: Table Cell Alignment + RTL (HIGH) ──

    def _check_table_cell_alignment(self, slide_num: int, conv_element) -> list:
        """Check that Arabic table cells have rtl='1' and algn='r'."""
        from vqa_types import V3Defect, Severity

        defects = []
        tables = _find_tables(conv_element)

        for tbl_idx, tbl in enumerate(tables):
            for row in _get_table_rows(tbl):
                for cell_idx, cell in enumerate(_get_row_cells(row)):
                    cell_text = _get_cell_text(cell)
                    if not _has_arabic(cell_text):
                        continue

                    # Check each paragraph in the cell
                    for para in cell.iter(f'{{{A_NS}}}p'):
                        pPr = para.find(f'{{{A_NS}}}pPr')
                        
                        rtl_ok = False
                        algn_ok = False
                        
                        if pPr is not None:
                            rtl_ok = pPr.get('rtl') == '1'
                            algn_ok = pPr.get('algn') == 'r'

                        if not rtl_ok or not algn_ok:
                            issues = []
                            if not rtl_ok:
                                issues.append('rtl missing')
                            if not algn_ok:
                                issues.append('algn not right')

                            defects.append(V3Defect(
                                code="TABLE_CELL_RTL_MISSING",
                                category="alignment",
                                severity=Severity.HIGH,
                                slide_idx=slide_num,
                                object_id=f"{tbl_idx}:{cell_idx}",
                                evidence={
                                    'cell_text': cell_text[:50],
                                    'issues': issues,
                                },
                                fixable=True,
                                autofix_action="set_para_rtl",
                                description=f"Arabic cell missing RTL/alignment "
                                            f"(slide {slide_num}, table {tbl_idx+1})",
                            ))
                            break  # One defect per cell is enough

        return defects

    # ── Check #9: Table gridCol Width Reversal (HIGH) ──

    def _check_table_gridcol_reversal(
        self, slide_num: int, orig_element, conv_element
    ) -> list:
        """Check if gridCol widths are reversed between original and converted."""
        from vqa_types import V3Defect, Severity

        defects = []
        orig_tables = _find_tables(orig_element)
        conv_tables = _find_tables(conv_element)

        for tbl_idx, (orig_tbl, conv_tbl) in enumerate(
            zip(orig_tables, conv_tables)
        ):
            orig_cols = _get_gridcols(orig_tbl)
            conv_cols = _get_gridcols(conv_tbl)

            if len(orig_cols) < 2 or len(conv_cols) < 2:
                continue
            if len(orig_cols) != len(conv_cols):
                continue

            orig_widths = [c.get('w', '0') for c in orig_cols]
            conv_widths = [c.get('w', '0') for c in conv_cols]

            # If all widths are equal, no need to check reversal
            if len(set(orig_widths)) <= 1:
                continue

            # After RTL conversion, gridCol widths should be reversed
            if orig_widths == conv_widths:
                # Check idempotency (matches Fix #1 marker)
                conv_tbl_pr = _get_tbl_pr(conv_tbl)
                if conv_tbl_pr is not None and conv_tbl_pr.get(V3_PHYS_REV_MARKER) == '1':
                    continue

                defects.append(V3Defect(
                    code="TABLE_GRIDCOL_NOT_REVERSED",
                    category="table",
                    severity=Severity.HIGH,
                    slide_idx=slide_num,
                    object_id=str(tbl_idx),
                    evidence={
                        'original_widths': orig_widths,
                        'converted_widths': conv_widths,
                    },
                    fixable=True,
                    autofix_action="reverse_table_columns",
                    description=f"Table {tbl_idx+1} gridCol widths not reversed "
                                f"(slide {slide_num})",
                ))

        return defects

    # ── Check #10: Merged Cell Integrity (HIGH, flag only) ──

    def _check_merged_cell_integrity(self, slide_num: int, conv_element) -> list:
        """Detect merged cells — flag only, blocks auto-reversal."""
        from vqa_types import V3Defect, Severity

        defects = []
        tables = _find_tables(conv_element)

        for tbl_idx, tbl in enumerate(tables):
            has_merged = False
            for row in _get_table_rows(tbl):
                for cell in _get_row_cells(row):
                    grid_span = cell.get('gridSpan', '1')
                    row_span = cell.get('rowSpan', '1')
                    try:
                        if int(grid_span) > 1 or int(row_span) > 1:
                            has_merged = True
                            break
                    except ValueError:
                        pass
                if has_merged:
                    break

            if has_merged:
                defects.append(V3Defect(
                    code="TABLE_MERGED_CELL_INTEGRITY",
                    category="table",
                    severity=Severity.HIGH,
                    slide_idx=slide_num,
                    object_id=str(tbl_idx),
                    evidence={'has_merged_cells': True},
                    fixable=False,  # Flag only — blocks reversal
                    description=f"Table {tbl_idx+1} has merged cells, "
                                f"auto-reversal blocked (slide {slide_num})",
                ))

        return defects

    # ─────────────────────────────────────────────────────────────────────────
    # SPRINT 2 CHECKS: Icon (#3), Page Number (#4), Shape Position (#5),
    #                  Paragraph RTL (#6)
    # ─────────────────────────────────────────────────────────────────────────

    # ── Check #3: Icon/Image in Wrong Table Cell (HIGH) ──

    def _check_icon_table_correspondence(self, slide_num: int, orig_element, conv_element) -> list:
        """Compare icon/image positions relative to table cells between orig and conv.
        
        Icons in tables should follow the cell reversal — an icon in column 0
        of the original should be in the last column of the converted.
        """
        from vqa_types import V3Defect, Severity

        defects = []
        # Find pic elements (deduplicated via shared helper)
        orig_pics = _enumerate_slide_pics(orig_element)
        conv_pics = _enumerate_slide_pics(conv_element)

        if not orig_pics or not conv_pics:
            return defects

        # For each image, check if its x-position has been mirrored
        for pic_idx, (orig_pic, conv_pic) in enumerate(zip(orig_pics, conv_pics)):
            orig_off = orig_pic.find(f'.//{{{A_NS}}}off')
            conv_off = conv_pic.find(f'.//{{{A_NS}}}off')
            orig_ext = orig_pic.find(f'.//{{{A_NS}}}ext')

            if orig_off is None or conv_off is None or orig_ext is None:
                continue

            try:
                orig_x = int(orig_off.get('x', '0'))
                conv_x = int(conv_off.get('x', '0'))
                shape_w = int(orig_ext.get('cx', '0'))
            except (ValueError, TypeError):
                continue

            if shape_w == 0:
                continue

            expected_x = self.slide_width - orig_x - shape_w
            tolerance = max(int(self.slide_width * 0.02), 12000)  # 2% or 12000 EMU

            if abs(conv_x - expected_x) > tolerance:
                defects.append(V3Defect(
                    code="ICON_IN_WRONG_TABLE_CELL",
                    category="icon",
                    severity=Severity.HIGH,
                    slide_idx=slide_num,
                    object_id=str(pic_idx),
                    evidence={
                        'original_x': orig_x,
                        'converted_x': conv_x,
                        'expected_x': expected_x,
                        'shape_width': shape_w,
                        'delta': abs(conv_x - expected_x),
                    },
                    fixable=True,
                    autofix_action="reposition_icon",
                    description=f"Icon {pic_idx} not mirrored correctly "
                                f"(slide {slide_num}, delta={abs(conv_x - expected_x)} EMU)",
                ))

        return defects

    # ── Check #4: Page Number Duplication (MEDIUM-HIGH) ──

    def _check_page_number_duplication(self, slide_num: int, conv_element) -> list:
        """Detect duplicate page numbers caused by field element duplication.
        
        Looks for:
        1. Multiple <a:fld> with type containing 'slidenum' in same shape
        2. Doubled-string patterns like '1515' in slide number text
        3. Duplicate <a:r> runs adjacent to <a:fld>
        """
        from vqa_types import V3Defect, Severity

        defects = []

        # Search for all shapes with text bodies
        for sp_idx, sp in enumerate(conv_element.iter(f'{{{P_NS}}}sp')):
            txBody = sp.find(f'.//{{{P_NS}}}txBody')
            if txBody is None:
                txBody = sp.find(f'.//{{{A_NS}}}txBody')
            if txBody is None:
                continue

            for para in txBody.iter(f'{{{A_NS}}}p'):
                # Count slidenum fields in this paragraph
                fields = [f for f in para.iter(f'{{{A_NS}}}fld')
                          if 'slidenum' in (f.get('type', '') or '').lower()]

                if len(fields) >= 2:
                    defects.append(V3Defect(
                        code="PAGE_NUMBER_DUPLICATED",
                        category="numbering",
                        severity=Severity.HIGH,
                        slide_idx=slide_num,
                        object_id=str(sp_idx),
                        evidence={
                            'field_count': len(fields),
                        },
                        fixable=True,
                        autofix_action="dedup_page_number",
                        description=f"Duplicate page number fields "
                                    f"(slide {slide_num}, {len(fields)} fields)",
                    ))
                    break  # One defect per shape

                # Check for doubled-string pattern in text runs
                # Only check shapes that contain at least one slidenum field
                # (avoids false positives on e.g. "1010" as static content)
                has_any_field = any(
                    'slidenum' in (f.get('type', '') or '').lower()
                    for f in sp.iter(f'{{{A_NS}}}fld')
                )
                if not has_any_field:
                    continue

                all_text = ''.join(
                    t.text or ''
                    for t in para.iter(f'{{{A_NS}}}t')
                ).strip()

                if all_text and len(all_text) >= 4 and len(all_text) % 2 == 0:
                    # Require >= 4 chars to avoid false positive on "11" (page 11)
                    half = len(all_text) // 2
                    if all_text[:half] == all_text[half:] and all_text[:half].isdigit():
                        defects.append(V3Defect(
                            code="PAGE_NUMBER_DOUBLED_STRING",
                            category="numbering",
                            severity=Severity.MEDIUM,
                            slide_idx=slide_num,
                            object_id=str(sp_idx),
                            evidence={
                                'doubled_text': all_text,
                            },
                            fixable=True,
                            autofix_action="dedup_page_number",
                            description=f"Doubled page number string '{all_text}' "
                                        f"(slide {slide_num})",
                        ))
                        break

        return defects

    # ── Check #5: Shape Position Mirroring Verification (HIGH) ──

    def _check_shape_position_mirroring(
        self, slide_num: int, orig_element, conv_element
    ) -> list:
        """Verify that non-centered shapes have mirrored x-positions.
        
        Expected: conv_x = slide_width - orig_x - shape_width
        Tolerance: 2% of slide_width or 12000 EMU minimum.
        
        Skips: background shapes, centered shapes, shapes with SLIDEARABI_NS markers.
        """
        from vqa_types import V3Defect, Severity

        defects = []
        orig_shapes = list(orig_element.iter(f'{{{P_NS}}}sp'))
        conv_shapes = list(conv_element.iter(f'{{{P_NS}}}sp'))

        for shape_idx, (orig_sp, conv_sp) in enumerate(zip(orig_shapes, conv_shapes)):
            # Skip shapes with SLIDEARABI_NS markers (intentionally not mirrored)
            nvSpPr = conv_sp.find(f'{{{P_NS}}}nvSpPr')
            if nvSpPr is not None:
                cNvPr = nvSpPr.find(f'{{{P_NS}}}cNvPr')
                if cNvPr is not None:
                    if any(SLIDEARABI_NS in a for a in cNvPr.attrib):
                        continue  # Skip marked shapes

            # Get position and size
            orig_off = orig_sp.find(f'.//{{{A_NS}}}off')
            conv_off = conv_sp.find(f'.//{{{A_NS}}}off')
            orig_ext = orig_sp.find(f'.//{{{A_NS}}}ext')

            if orig_off is None or conv_off is None or orig_ext is None:
                continue

            try:
                orig_x = int(orig_off.get('x', '0'))
                orig_y = int(orig_off.get('y', '0'))
                conv_x = int(conv_off.get('x', '0'))
                shape_w = int(orig_ext.get('cx', '0'))
                shape_h = int(orig_ext.get('cy', '0'))
            except (ValueError, TypeError):
                continue

            # Skip very small shapes (decorative dots, bullets)
            if shape_w < 50000 and shape_h < 50000:
                continue

            # Skip shapes that span most of the slide width (backgrounds, full-width bars)
            if shape_w > self.slide_width * 0.85:
                continue

            # Skip centered shapes (within 5% tolerance of center)
            center_x = (self.slide_width - shape_w) // 2
            if abs(orig_x - center_x) < self.slide_width * 0.05:
                continue

            expected_x = self.slide_width - orig_x - shape_w
            if expected_x < 0:
                continue  # Shape extends beyond slide boundary, skip
            tolerance = max(int(self.slide_width * 0.02), 12000)

            if abs(conv_x - expected_x) > tolerance:
                defects.append(V3Defect(
                    code="SHAPE_NOT_MIRRORED_POSITION",
                    category="mirroring",
                    severity=Severity.HIGH,
                    slide_idx=slide_num,
                    object_id=str(shape_idx),
                    evidence={
                        'original_x': orig_x,
                        'converted_x': conv_x,
                        'expected_x': expected_x,
                        'shape_width': shape_w,
                        'delta': abs(conv_x - expected_x),
                    },
                    fixable=True,
                    autofix_action="mirror_shape_position",
                    description=f"Shape {shape_idx} not mirrored "
                                f"(slide {slide_num}, delta={abs(conv_x - expected_x)} EMU)",
                ))

        return defects

    # ── Check #6: Non-Table Paragraph RTL Missing (MEDIUM-HIGH) ──

    def _check_paragraph_rtl(self, slide_num: int, conv_element) -> list:
        """Check that Arabic paragraphs outside tables have rtl='1'.
        
        Skips table cells (covered by Check #2) and empty paragraphs.
        """
        from vqa_types import V3Defect, Severity

        defects = []
        # Collect all <a:p> that are NOT inside tables
        # Strategy: find all <a:p> and exclude those under <a:tbl>
        table_paras = set()
        for tbl in conv_element.iter(f'{{{A_NS}}}tbl'):
            for p in tbl.iter(f'{{{A_NS}}}p'):
                table_paras.add(id(p))

        flagged_shapes = set()  # One defect per shape
        for sp_idx, sp in enumerate(conv_element.iter(f'{{{P_NS}}}sp')):
            for para in sp.iter(f'{{{A_NS}}}p'):
                if id(para) in table_paras:
                    continue

                # Get paragraph text
                para_text = ''.join(
                    t.text or '' for t in para.iter(f'{{{A_NS}}}t')
                ).strip()

                if not para_text or not _has_arabic(para_text):
                    continue

                pPr = para.find(f'{{{A_NS}}}pPr')
                rtl_ok = pPr is not None and pPr.get('rtl') == '1'

                if not rtl_ok and sp_idx not in flagged_shapes:
                    flagged_shapes.add(sp_idx)
                    defects.append(V3Defect(
                        code="PARAGRAPH_RTL_MISSING",
                        category="alignment",
                        severity=Severity.MEDIUM,
                        slide_idx=slide_num,
                        object_id=str(sp_idx),
                        evidence={
                            'sample_text': para_text[:50],
                        },
                        fixable=True,
                        autofix_action="set_paragraph_rtl",
                        description=f"Arabic paragraph missing rtl='1' "
                                    f"(slide {slide_num}, shape {sp_idx})",
                    ))

        return defects

    # ─────────────────────────────────────────────────────────────────────────
    # SPRINT 3 CHECKS: Circular centering (#7), Master mirror (#8),
    #                  Directional orientation (#12)
    # ─────────────────────────────────────────────────────────────────────────

    # Preset geometry types for circular shapes
    CIRCULAR_PRESETS = frozenset([
        'ellipse', 'pie', 'donut', 'arc', 'blockArc',
        'flowChartConnector',  # circle connector in flowcharts
        'actionButtonBlank',  # often used as circle buttons
    ])

    # Preset geometry types for directional shapes (rightward by default)
    DIRECTIONAL_PRESETS = frozenset([
        'rightArrow', 'stripedRightArrow', 'notchedRightArrow',
        'bentArrow', 'uturnArrow', 'curvedRightArrow',
        'chevron', 'homePlate', 'rightArrowCallout',
        'circularArrow',
        'rtTriangle',  # OOXML name for right triangle
    ])

    # ── Check #7: Text Centering in Circular/Venn Shapes (MEDIUM) ──

    def _check_circular_text_centering(self, slide_num: int, conv_element) -> list:
        """Verify text in circular/near-square shapes has centered alignment.
        
        Checks prstGeom-based detection + aspect-ratio fallback (width within
        20% of height). Requires algn='ctr' and anchor='ctr' on text body.
        """
        from vqa_types import V3Defect, Severity

        defects = []

        for sp_idx, sp in enumerate(conv_element.iter(f'{{{P_NS}}}sp')):
            # Get preset geometry
            spPr = sp.find(f'{{{P_NS}}}spPr')
            if spPr is None:
                spPr = sp.find(f'{{{A_NS}}}spPr')
            if spPr is None:
                continue

            prstGeom = spPr.find(f'{{{A_NS}}}prstGeom')
            prst = prstGeom.get('prst', '') if prstGeom is not None else ''

            # Get dimensions for aspect ratio check
            ext = spPr.find(f'.//{{{A_NS}}}ext')
            is_circular = prst in self.CIRCULAR_PRESETS

            if not is_circular and ext is not None:
                try:
                    cx = int(ext.get('cx', '0'))
                    cy = int(ext.get('cy', '0'))
                    if cx > 0 and cy > 0:
                        ratio = max(cx, cy) / min(cx, cy)
                        if ratio <= 1.2:  # Near-square = likely circular
                            is_circular = True
                except (ValueError, TypeError):
                    pass

            if not is_circular:
                continue

            # Check for text content
            txBody = sp.find(f'{{{P_NS}}}txBody')
            if txBody is None:
                txBody = sp.find(f'{{{A_NS}}}txBody')
            if txBody is None:
                continue

            has_text = False
            for para in txBody.iter(f'{{{A_NS}}}p'):
                para_text = ''.join(
                    t.text or '' for t in para.iter(f'{{{A_NS}}}t')
                ).strip()
                if para_text:
                    has_text = True
                    break

            if not has_text:
                continue

            # Check centering: algn='ctr' on paragraphs and anchor='ctr' on bodyPr
            bodyPr = txBody.find(f'{{{A_NS}}}bodyPr')
            anchor_ok = bodyPr is not None and bodyPr.get('anchor') == 'ctr'

            algn_ok = True
            for para in txBody.iter(f'{{{A_NS}}}p'):
                para_text = ''.join(
                    t.text or '' for t in para.iter(f'{{{A_NS}}}t')
                ).strip()
                if not para_text:
                    continue
                pPr = para.find(f'{{{A_NS}}}pPr')
                if pPr is None or pPr.get('algn') != 'ctr':
                    algn_ok = False
                    break

            if not algn_ok or not anchor_ok:
                issues = []
                if not algn_ok:
                    issues.append('algn not ctr')
                if not anchor_ok:
                    issues.append('anchor not ctr')

                defects.append(V3Defect(
                    code="TEXT_NOT_CENTERED_IN_SHAPE",
                    category="alignment",
                    severity=Severity.MEDIUM,
                    slide_idx=slide_num,
                    object_id=str(sp_idx),
                    evidence={
                        'preset': prst,
                        'issues': issues,
                    },
                    fixable=True,
                    autofix_action="center_text_circular",
                    description=f"Text not centered in circular shape "
                                f"(slide {slide_num}, shape {sp_idx}, "
                                f"prst={prst or 'aspect-ratio'})",
                ))

        return defects

    # ── Check #8: Master/Layout Decorative Element Mirroring (MEDIUM) ──

    def _check_master_element_mirroring(
        self, slide_num: int, orig_element, conv_element
    ) -> list:
        """Check master/layout decorative shapes are mirrored.
        
        Looks at shape names containing common master patterns (line, rect,
        background, decorator) and verifies x-position mirroring.
        Skips centered and full-width shapes.
        """
        from vqa_types import V3Defect, Severity

        defects = []
        MASTER_PATTERNS = ('line', 'rect', 'background', 'decorator',
                          'border', 'stripe', 'bar', 'accent', 'freeform')

        orig_shapes = list(orig_element.iter(f'{{{P_NS}}}sp'))
        conv_shapes = list(conv_element.iter(f'{{{P_NS}}}sp'))

        for shape_idx, (orig_sp, conv_sp) in enumerate(zip(orig_shapes, conv_shapes)):
            # Check if this looks like a master/decorative element
            nvSpPr = conv_sp.find(f'{{{P_NS}}}nvSpPr')
            if nvSpPr is None:
                continue
            cNvPr = nvSpPr.find(f'{{{P_NS}}}cNvPr')
            if cNvPr is None:
                continue

            # Guard: verify shape names match between orig and conv
            # (prevents mispairing if shapes were added/removed)
            orig_nvSpPr = orig_sp.find(f'{{{P_NS}}}nvSpPr')
            if orig_nvSpPr is not None:
                orig_cNvPr = orig_nvSpPr.find(f'{{{P_NS}}}cNvPr')
                if orig_cNvPr is not None:
                    orig_name = (orig_cNvPr.get('name', '') or '').lower()
                    conv_name_check = (cNvPr.get('name', '') or '').lower()
                    if orig_name != conv_name_check:
                        continue  # Shape ordering diverged, skip

            name = (cNvPr.get('name', '') or '').lower()

            # Only check shapes that match master patterns
            is_master = any(pat in name for pat in MASTER_PATTERNS)
            if not is_master:
                continue

            # Skip shapes with SLIDEARABI_NS markers
            if any(SLIDEARABI_NS in a for a in cNvPr.attrib):
                continue

            # Get positions
            orig_off = orig_sp.find(f'.//{{{A_NS}}}off')
            conv_off = conv_sp.find(f'.//{{{A_NS}}}off')
            orig_ext = orig_sp.find(f'.//{{{A_NS}}}ext')

            if orig_off is None or conv_off is None or orig_ext is None:
                continue

            try:
                orig_x = int(orig_off.get('x', '0'))
                conv_x = int(conv_off.get('x', '0'))
                shape_w = int(orig_ext.get('cx', '0'))
            except (ValueError, TypeError):
                continue

            if shape_w == 0:
                continue

            # Skip full-width
            if shape_w > self.slide_width * 0.85:
                continue

            # Skip centered
            center_x = (self.slide_width - shape_w) // 2
            if abs(orig_x - center_x) < self.slide_width * 0.05:
                continue

            expected_x = self.slide_width - orig_x - shape_w
            if expected_x < 0:
                continue

            tolerance = max(int(self.slide_width * 0.02), 12000)

            if abs(conv_x - expected_x) > tolerance:
                defects.append(V3Defect(
                    code="MASTER_ELEMENT_NOT_MIRRORED",
                    category="mirroring",
                    severity=Severity.MEDIUM,
                    slide_idx=slide_num,
                    object_id=str(shape_idx),
                    evidence={
                        'shape_name': name,
                        'original_x': orig_x,
                        'converted_x': conv_x,
                        'expected_x': expected_x,
                        'shape_width': shape_w,
                        'delta': abs(conv_x - expected_x),
                    },
                    fixable=True,
                    autofix_action="mirror_shape_position",
                    description=f"Master element '{name}' not mirrored "
                                f"(slide {slide_num}, delta={abs(conv_x - expected_x)} EMU)",
                ))

        return defects

    # ── Check #12: Directional Symbol Orientation (HIGH) ──

    def _check_directional_symbol_orientation(
        self, slide_num: int, conv_element
    ) -> list:
        """Check that directional preset shapes (arrows, chevrons) have flipH='1'.
        
        In RTL conversion, rightward-pointing shapes must be horizontally flipped
        to point leftward. Detection is based on prstGeom, not shape name keywords.
        """
        from vqa_types import V3Defect, Severity

        defects = []

        for sp_idx, sp in enumerate(conv_element.iter(f'{{{P_NS}}}sp')):
            spPr = sp.find(f'{{{P_NS}}}spPr')
            if spPr is None:
                spPr = sp.find(f'{{{A_NS}}}spPr')
            if spPr is None:
                continue

            prstGeom = spPr.find(f'{{{A_NS}}}prstGeom')
            if prstGeom is None:
                continue

            prst = prstGeom.get('prst', '')
            if prst not in self.DIRECTIONAL_PRESETS:
                continue

            # Check for flipH on xfrm
            xfrm = spPr.find(f'{{{A_NS}}}xfrm')
            has_flipH = xfrm is not None and xfrm.get('flipH') == '1'

            if not has_flipH:
                defects.append(V3Defect(
                    code="DIRECTIONAL_SHAPE_NOT_FLIPPED",
                    category="mirroring",
                    severity=Severity.HIGH,
                    slide_idx=slide_num,
                    object_id=str(sp_idx),
                    evidence={
                        'preset': prst,
                        'current_flipH': xfrm.get('flipH', 'missing') if xfrm is not None else 'no_xfrm',
                    },
                    fixable=True,
                    autofix_action="flip_directional_shape",
                    description=f"Directional shape '{prst}' not flipped for RTL "
                                f"(slide {slide_num}, shape {sp_idx})",
                ))

        return defects

    # ─────────────────────────────────────────────────────────────────────────
    # END SPRINT 3 CHECKS
    # ─────────────────────────────────────────────────────────────────────────

    # ── Check #11: Table Grid Structural Integrity (CRITICAL, flag only) ──

    def _check_table_grid_structural(self, slide_num: int, conv_element) -> list:
        """Verify gridCol count matches row cell count."""
        from vqa_types import V3Defect, Severity

        defects = []
        tables = _find_tables(conv_element)

        for tbl_idx, tbl in enumerate(tables):
            gridcols = _get_gridcols(tbl)
            num_grid_cols = len(gridcols)

            if num_grid_cols == 0:
                continue

            for row_idx, row in enumerate(_get_table_rows(tbl)):
                cells = _get_row_cells(row)
                # Account for gridSpan in cell count
                effective_cols = 0
                for cell in cells:
                    try:
                        span = int(cell.get('gridSpan', '1'))
                    except (ValueError, TypeError):
                        span = 1
                    effective_cols += span

                if effective_cols != num_grid_cols:
                    defects.append(V3Defect(
                        code="TABLE_GRID_STRUCTURAL_ERROR",
                        category="table",
                        severity=Severity.CRITICAL,
                        slide_idx=slide_num,
                        object_id=f"{tbl_idx}:row{row_idx}",
                        evidence={
                            'grid_cols': num_grid_cols,
                            'effective_row_cols': effective_cols,
                            'row_index': row_idx,
                        },
                        fixable=False,  # Can't auto-fix structural corruption
                        description=f"Table {tbl_idx+1} row {row_idx} has "
                                    f"{effective_cols} effective cols but grid "
                                    f"has {num_grid_cols} (slide {slide_num})",
                    ))

        return defects


# ─────────────────────────────────────────────────────────────────────────────
# V3 AUTO-FIXES
# ─────────────────────────────────────────────────────────────────────────────

class V3AutoFixer:
    """V3 auto-fix engine — deterministic fixes with rollback safety."""

    MARKER_NS = SLIDEARABI_NS

    def __init__(self, slide_width: int = 9144000, slide_height: int = 6858000):
        self.slide_width = slide_width
        self.slide_height = slide_height
        self._dispatch = {
            'reverse_table_columns': self._fix_reverse_table_columns,
            'set_para_rtl': self._fix_set_para_rtl,
            'reposition_icon': self._fix_reposition_icon,
            'dedup_page_number': self._fix_dedup_page_number,
            'mirror_shape_position': self._fix_mirror_shape_position,
            'set_paragraph_rtl': self._fix_set_paragraph_rtl,
            # Sprint 3 fixes
            'center_text_circular': self._fix_center_text_circular,
            'flip_directional_shape': self._fix_flip_directional_shape,
        }

    def apply_fix(self, slide_element, defect) -> bool:
        """Apply a single fix. Returns True if successful."""
        action = defect.autofix_action
        if action not in self._dispatch:
            logger.warning(f"Unknown fix action: {action}")
            return False
        try:
            return self._dispatch[action](slide_element, defect)
        except Exception as e:
            logger.error(f"Fix {action} failed on slide {defect.slide_idx}: {e}")
            return False

    # ── Fix #1: Reverse Table Columns ──

    def _fix_reverse_table_columns(self, slide_element, defect) -> bool:
        """Reverse <a:tc> order in each row + gridCol widths.
        
        SAFETY: Aborts if any cell has gridSpan > 1 or rowSpan > 1.
        """
        tables = _find_tables(slide_element)
        tbl_idx = int(defect.object_id.split(':')[0]) if ':' in str(defect.object_id) else int(defect.object_id)

        if tbl_idx >= len(tables):
            return False

        tbl = tables[tbl_idx]

        # Safety: check for merged cells
        for row in _get_table_rows(tbl):
            for cell in _get_row_cells(row):
                try:
                    if int(cell.get('gridSpan', '1')) > 1:
                        logger.warning(f"Merged cells found, aborting reversal on slide {defect.slide_idx}")
                        return False
                    if int(cell.get('rowSpan', '1')) > 1:
                        logger.warning(f"Row-span found, aborting reversal on slide {defect.slide_idx}")
                        return False
                except ValueError:
                    pass

        # Check idempotency marker (shared constant with checks)
        tbl_pr = _get_tbl_pr(tbl)
        if tbl_pr is not None and tbl_pr.get(V3_PHYS_REV_MARKER) == '1':
            logger.debug(f"Table already reversed (v3 marker), skipping")
            return False

        # Reverse gridCol widths
        grid = tbl.find(f'{{{A_NS}}}tblGrid')
        if grid is not None:
            cols = grid.findall(f'{{{A_NS}}}gridCol')
            if len(cols) > 1:
                widths = [c.get('w') for c in cols]
                widths.reverse()
                for col, w in zip(cols, widths):
                    if w is not None:
                        col.set('w', w)

        # Reverse cell order in each row
        for row in _get_table_rows(tbl):
            cells = _get_row_cells(row)
            if len(cells) <= 1:
                continue
            # Remove all cells, re-add in reverse
            for cell in cells:
                row.remove(cell)
            for cell in reversed(cells):
                row.append(cell)

        # Set rtl='0' to prevent double-reversal
        if tbl_pr is None:
            tbl_pr = etree.Element(f'{{{A_NS}}}tblPr')
            tbl.insert(0, tbl_pr)  # tblPr must be first child per OOXML schema
        tbl_pr.set('rtl', '0')

        # Set idempotency marker
        tbl_pr.set(V3_PHYS_REV_MARKER, '1')

        # Swap firstCol/lastCol toggles
        first_col = tbl_pr.get('firstCol')
        last_col = tbl_pr.get('lastCol')
        if first_col or last_col:
            tbl_pr.set('firstCol', last_col or '0')
            tbl_pr.set('lastCol', first_col or '0')

        logger.info(f"V3 Fix: reversed table columns on slide {defect.slide_idx}, "
                     f"table {tbl_idx}")
        return True

    # ── Fix #2: Set Paragraph RTL + Alignment ──

    def _fix_set_para_rtl(self, slide_element, defect) -> bool:
        """Set rtl='1' and algn='r' on Arabic paragraph in table cell."""
        tables = _find_tables(slide_element)

        # Parse object_id "tbl_idx:cell_idx"
        parts = str(defect.object_id).split(':')
        if len(parts) < 2:
            return False
        tbl_idx = int(parts[0])
        # cell_idx is informational — we fix all Arabic cells in the table

        if tbl_idx >= len(tables):
            return False

        tbl = tables[tbl_idx]
        fixed = False

        for row in _get_table_rows(tbl):
            for cell in _get_row_cells(row):
                cell_text = _get_cell_text(cell)
                if not _has_arabic(cell_text):
                    continue

                for para in cell.iter(f'{{{A_NS}}}p'):
                    pPr = para.find(f'{{{A_NS}}}pPr')
                    if pPr is None:
                        pPr = etree.SubElement(para, f'{{{A_NS}}}pPr')
                        # Move to beginning of paragraph
                        para.remove(pPr)
                        para.insert(0, pPr)

                    changed = False
                    if pPr.get('rtl') != '1':
                        pPr.set('rtl', '1')
                        changed = True
                    if pPr.get('algn') != 'r':
                        pPr.set('algn', 'r')
                        changed = True

                    if changed:
                        fixed = True

        if fixed:
            logger.info(f"V3 Fix: set RTL/alignment on slide {defect.slide_idx}, "
                         f"table {tbl_idx}")
        return fixed

    # ─────────────────────────────────────────────────────────────────────────
    # SPRINT 2 FIXES
    # ─────────────────────────────────────────────────────────────────────────

    # ── Fix #3: Reposition Icon to Mirrored Position ──

    def _fix_reposition_icon(self, slide_element, defect) -> bool:
        """Move icon/image to its expected mirrored x-position."""
        pics = _enumerate_slide_pics(slide_element)

        try:
            pic_idx = int(defect.object_id)
        except (ValueError, TypeError):
            return False

        if pic_idx >= len(pics):
            return False

        pic = pics[pic_idx]
        off = pic.find(f'.//{{{A_NS}}}off')
        if off is None:
            return False

        expected_x = defect.evidence.get('expected_x')
        if expected_x is None:
            return False

        off.set('x', str(expected_x))
        logger.info(f"V3 Fix: repositioned icon {pic_idx} to x={expected_x} "
                    f"(slide {defect.slide_idx})")
        return True

    # ── Fix #4: Deduplicate Page Number ──

    def _fix_dedup_page_number(self, slide_element, defect) -> bool:
        """Remove duplicate page number elements.
        
        Handles two cases:
        - PAGE_NUMBER_DUPLICATED: duplicate <a:fld> elements → keep first, remove rest
        - PAGE_NUMBER_DOUBLED_STRING: doubled text like '1515' → truncate to first half
        """
        shapes = list(slide_element.iter(f'{{{P_NS}}}sp'))

        try:
            sp_idx = int(defect.object_id)
        except (ValueError, TypeError):
            return False

        if sp_idx >= len(shapes):
            return False

        sp = shapes[sp_idx]
        fixed = False

        for para in sp.iter(f'{{{A_NS}}}p'):
            # Case 1: Duplicate <a:fld> elements
            fields = [f for f in para.iter(f'{{{A_NS}}}fld')
                      if 'slidenum' in (f.get('type', '') or '').lower()]

            if len(fields) >= 2:
                for extra in fields[1:]:
                    parent = extra.getparent()
                    if parent is not None:
                        parent.remove(extra)
                        fixed = True
                continue

            # Case 2: Doubled-string in text runs
            if defect.code == 'PAGE_NUMBER_DOUBLED_STRING':
                doubled_text = defect.evidence.get('doubled_text', '')
                if doubled_text and len(doubled_text) >= 2:
                    half = doubled_text[:len(doubled_text) // 2]
                    for t_elem in para.iter(f'{{{A_NS}}}t'):
                        if t_elem.text and doubled_text in t_elem.text:
                            t_elem.text = t_elem.text.replace(doubled_text, half)
                            fixed = True

        if fixed:
            logger.info(f"V3 Fix: deduplicated page number on slide {defect.slide_idx}")
        return fixed

    # ── Fix #5: Mirror Shape Position ──

    def _fix_mirror_shape_position(self, slide_element, defect) -> bool:
        """Move shape to its expected mirrored x-position."""
        shapes = list(slide_element.iter(f'{{{P_NS}}}sp'))

        try:
            shape_idx = int(defect.object_id)
        except (ValueError, TypeError):
            return False

        if shape_idx >= len(shapes):
            return False

        sp = shapes[shape_idx]
        off = sp.find(f'.//{{{A_NS}}}off')
        if off is None:
            return False

        expected_x = defect.evidence.get('expected_x')
        if expected_x is None:
            return False

        off.set('x', str(expected_x))
        logger.info(f"V3 Fix: mirrored shape {shape_idx} to x={expected_x} "
                    f"(slide {defect.slide_idx})")
        return True

    # ── Fix #6: Set RTL on Non-Table Paragraphs ──

    def _fix_set_paragraph_rtl(self, slide_element, defect) -> bool:
        """Set rtl='1' on Arabic paragraphs in non-table shapes."""
        shapes = list(slide_element.iter(f'{{{P_NS}}}sp'))

        try:
            sp_idx = int(defect.object_id)
        except (ValueError, TypeError):
            return False

        if sp_idx >= len(shapes):
            return False

        sp = shapes[sp_idx]
        fixed = False

        # Collect table paragraphs to exclude
        table_paras = set()
        for tbl in slide_element.iter(f'{{{A_NS}}}tbl'):
            for p in tbl.iter(f'{{{A_NS}}}p'):
                table_paras.add(id(p))

        for para in sp.iter(f'{{{A_NS}}}p'):
            if id(para) in table_paras:
                continue

            para_text = ''.join(
                t.text or '' for t in para.iter(f'{{{A_NS}}}t')
            ).strip()

            if not para_text or not _has_arabic(para_text):
                continue

            pPr = para.find(f'{{{A_NS}}}pPr')
            if pPr is None:
                pPr = etree.Element(f'{{{A_NS}}}pPr')
                para.insert(0, pPr)

            if pPr.get('rtl') != '1':
                pPr.set('rtl', '1')
                fixed = True

            # Set alignment to right for Arabic text (unless explicitly centered)
            current_algn = pPr.get('algn', '')
            if current_algn != 'ctr' and current_algn != 'r':
                pPr.set('algn', 'r')
                fixed = True

        if fixed:
            logger.info(f"V3 Fix: set paragraph RTL on shape {sp_idx} "
                        f"(slide {defect.slide_idx})")
        return fixed

    # ─────────────────────────────────────────────────────────────────────────
    # SPRINT 3 FIXES
    # ─────────────────────────────────────────────────────────────────────────

    # ── Fix #7: Center Text in Circular Shapes ──

    def _fix_center_text_circular(self, slide_element, defect) -> bool:
        """Set algn='ctr' on paragraphs and anchor='ctr' on bodyPr."""
        shapes = list(slide_element.iter(f'{{{P_NS}}}sp'))

        try:
            sp_idx = int(defect.object_id)
        except (ValueError, TypeError):
            return False

        if sp_idx >= len(shapes):
            return False

        sp = shapes[sp_idx]
        txBody = sp.find(f'{{{P_NS}}}txBody')
        if txBody is None:
            txBody = sp.find(f'{{{A_NS}}}txBody')
        if txBody is None:
            return False

        fixed = False

        # Set anchor='ctr' on bodyPr
        bodyPr = txBody.find(f'{{{A_NS}}}bodyPr')
        if bodyPr is None:
            bodyPr = etree.SubElement(txBody, f'{{{A_NS}}}bodyPr')
            txBody.remove(bodyPr)
            txBody.insert(0, bodyPr)
        if bodyPr.get('anchor') != 'ctr':
            bodyPr.set('anchor', 'ctr')
            fixed = True

        # Set algn='ctr' on all text-bearing paragraphs
        for para in txBody.iter(f'{{{A_NS}}}p'):
            para_text = ''.join(
                t.text or '' for t in para.iter(f'{{{A_NS}}}t')
            ).strip()
            if not para_text:
                continue
            pPr = para.find(f'{{{A_NS}}}pPr')
            if pPr is None:
                pPr = etree.Element(f'{{{A_NS}}}pPr')
                para.insert(0, pPr)
            if pPr.get('algn') != 'ctr':
                pPr.set('algn', 'ctr')
                fixed = True

        if fixed:
            logger.info(f"V3 Fix: centered text in circular shape {sp_idx} "
                        f"(slide {defect.slide_idx})")
        return fixed

    # ── Fix #12: Flip Directional Shape for RTL ──

    def _fix_flip_directional_shape(self, slide_element, defect) -> bool:
        """Set flipH='1' on directional shape's xfrm."""
        shapes = list(slide_element.iter(f'{{{P_NS}}}sp'))

        try:
            sp_idx = int(defect.object_id)
        except (ValueError, TypeError):
            return False

        if sp_idx >= len(shapes):
            return False

        sp = shapes[sp_idx]
        spPr = sp.find(f'{{{P_NS}}}spPr')
        if spPr is None:
            spPr = sp.find(f'{{{A_NS}}}spPr')
        if spPr is None:
            return False

        xfrm = spPr.find(f'{{{A_NS}}}xfrm')
        if xfrm is None:
            xfrm = etree.SubElement(spPr, f'{{{A_NS}}}xfrm')
            spPr.remove(xfrm)
            spPr.insert(0, xfrm)

        if xfrm.get('flipH') == '1':
            return False  # Already flipped

        xfrm.set('flipH', '1')
        logger.info(f"V3 Fix: flipped directional shape {sp_idx} "
                    f"(slide {defect.slide_idx}, prst={defect.evidence.get('preset', '?')})")
        return True


# ─────────────────────────────────────────────────────────────────────────────
# V3 PIPELINE ENTRY POINTS
# ─────────────────────────────────────────────────────────────────────────────

def run_v3_xml_checks(
    orig_pptx_path: str,
    conv_pptx_path: str,
    slide_width_emu: int = 9144000,
    slide_height_emu: int = 6858000,
) -> Tuple[list, Dict[str, Any]]:
    """
    Run all V3 XML structural checks on all slides.
    
    Returns:
        (defects, metadata) where defects is List[V3Defect]
    """
    from pptx import Presentation as PptxPresentation

    checker = V3XMLChecker(slide_width_emu, slide_height_emu)
    all_defects = []

    conv_prs = PptxPresentation(conv_pptx_path)
    orig_prs = PptxPresentation(orig_pptx_path)

    for slide_idx in range(len(conv_prs.slides)):
        slide_num = slide_idx + 1
        conv_slide = conv_prs.slides[slide_idx]
        
        orig_slide_elem = None
        if slide_idx < len(orig_prs.slides):
            orig_slide_elem = orig_prs.slides[slide_idx]._element

        slide_defects = checker.check_slide(
            slide_num,
            conv_slide._element,
            orig_slide_elem,
        )
        all_defects.extend(slide_defects)

    from vqa_types import Severity
    metadata = {
        'slides_checked': len(conv_prs.slides),
        'total_defects': len(all_defects),
        'critical_count': sum(1 for d in all_defects if d.severity == Severity.CRITICAL),
        'high_count': sum(1 for d in all_defects if d.severity == Severity.HIGH),
        'fixable_count': sum(1 for d in all_defects if d.fixable),
    }

    logger.info(f"V3 XML checks: {metadata['total_defects']} defects "
                f"({metadata['critical_count']} critical, "
                f"{metadata['fixable_count']} fixable) "
                f"across {metadata['slides_checked']} slides")

    return all_defects, metadata


def safe_apply_fixes(
    pptx_path: str,
    defects: list,
    slide_width: int = 9144000,
    slide_height: int = 6858000,
) -> Tuple[list, list]:
    """
    Apply V3 fixes with per-fix validation and rollback.
    
    Returns:
        (applied, failed) — lists of V3Defect objects
    """
    from pptx import Presentation as PptxPresentation
    from vqa_types import DefectStatus
    import v3_config

    if not v3_config.ENABLE_V3_XML_AUTOFIX:
        logger.info("V3 XML autofix disabled, skipping")
        return [], list(defects)

    fixer = V3AutoFixer(slide_width, slide_height)
    applied = []
    failed = []

    # Group fixable defects by slide
    fixable = [d for d in defects if d.fixable and d.status.value == 'open']
    if not fixable:
        return [], []

    # Create backup
    backup_path = str(pptx_path) + '.v3_backup'
    shutil.copy2(pptx_path, backup_path)

    try:
        prs = PptxPresentation(pptx_path)

        # Group by slide
        by_slide: Dict[int, list] = {}
        for d in fixable:
            by_slide.setdefault(d.slide_idx, []).append(d)

        for slide_idx_0, slide in enumerate(prs.slides):
            slide_num = slide_idx_0 + 1
            slide_defects = by_slide.get(slide_num, [])
            
            for defect in slide_defects:
                # Check table autofix flag for table fixes
                if defect.autofix_action == 'reverse_table_columns':
                    if not v3_config.ENABLE_V3_TABLE_FIX:
                        logger.info(f"Table autofix disabled, skipping {defect.code}")
                        defect.status = DefectStatus.UNRESOLVED
                        failed.append(defect)
                        continue

                success = fixer.apply_fix(slide._element, defect)
                if success:
                    defect.status = DefectStatus.FIXED
                    applied.append(defect)
                else:
                    defect.status = DefectStatus.UNRESOLVED
                    failed.append(defect)

        # Save and validate
        prs.save(pptx_path)

        # Validation: try opening the saved file
        try:
            PptxPresentation(pptx_path)
            logger.info(f"V3 fixes applied: {len(applied)} succeeded, "
                        f"{len(failed)} failed")
            # Clean up backup on success
            try:
                Path(backup_path).unlink(missing_ok=True)
            except Exception:
                pass
        except Exception as e:
            logger.error(f"PPTX validation failed after fixes, rolling back: {e}")
            shutil.copy2(backup_path, pptx_path)
            # Mark all as failed
            for d in applied:
                d.status = DefectStatus.UNRESOLVED
            failed.extend(applied)
            applied = []

    except Exception as e:
        logger.error(f"Fix application failed, rolling back: {e}")
        shutil.copy2(backup_path, pptx_path)
        for d in fixable:
            d.status = DefectStatus.UNRESOLVED
        failed = list(fixable)
        applied = []

    return applied, failed


def compute_gate_decision(
    unresolved_defects: list,
    total_slides: int,
) -> 'VQAGateResult':
    """Compute quality gate decision from unresolved defects."""
    from vqa_types import (
        VQAGateResult, Severity, MUST_NOT_SHIP_CODES,
        HIGH_SEVERITY_THRESHOLD, HIGH_SLIDE_RATIO_THRESHOLD,
    )

    gate = VQAGateResult()

    critical = [d for d in unresolved_defects if d.severity == Severity.CRITICAL]
    high = [d for d in unresolved_defects if d.severity == Severity.HIGH]
    must_not_ship = [d for d in unresolved_defects if d.code in MUST_NOT_SHIP_CODES]

    gate.critical_remaining = len(critical)
    gate.high_remaining = len(high)
    gate.blocking_issues = [d.to_dict() for d in critical[:5]]
    gate.warning_issues = [d.to_dict() for d in high[:10]]

    # Gate logic (Sonnet's nuanced thresholds)
    if critical or must_not_ship:
        gate.status = "failed_qa"
    elif len(high) >= HIGH_SEVERITY_THRESHOLD:
        high_slides = len(set(d.slide_idx for d in high))
        if total_slides > 0 and (high_slides / total_slides) >= HIGH_SLIDE_RATIO_THRESHOLD:
            gate.status = "failed_qa"
        else:
            gate.status = "completed_with_warnings"
    elif high:
        gate.status = "completed_with_warnings"
    else:
        gate.status = "completed"

    return gate

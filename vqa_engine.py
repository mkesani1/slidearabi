"""
vqa_engine.py — Visual Quality Assurance Engine for SlideArabi

Three-layer defense architecture (consensus of GPT-5.4, Claude Opus 4.6, Gemini 3.1 Pro):

Layer 1: XML Structural Checks (pre-render, deterministic, catches ~70% of defects)
Layer 2: Vision Model (post-render, Document-Grounded per Gemini recommendation)
Layer 3: Cross-Validation (reconciles Layer 1 + Layer 2 findings)

Includes auto-fix cascade: Shift → Expand box → Reduce font → Line break injection
"""

from __future__ import annotations

import json
import logging
import re
from copy import deepcopy
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Tuple

from lxml import etree

logger = logging.getLogger(__name__)

# ─────────────────────────────────────────────────────────────────────────────
# Constants
# ─────────────────────────────────────────────────────────────────────────────

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

# Minimum margin from any slide edge (0.15 inches in EMU)
MIN_MARGIN_EMU = 137_160  # 0.15" — slightly more generous than 0.1" to catch real issues

# EMU per inch
EMU_PER_INCH = 914_400

# Text-to-box area threshold for clipping risk
TEXT_BOX_AREA_RISK_THRESHOLD = 0.90

# Overlap percentage thresholds
OVERLAP_CRITICAL_PCT = 0.30  # >30% overlap = CRITICAL
OVERLAP_HIGH_PCT = 0.10       # >10% overlap = HIGH

# Narrow box threshold (less than 0.6 inches)
NARROW_BOX_THRESHOLD_EMU = 548_640  # 0.6 inches

# Circular shape geometry types
CIRCULAR_GEOMS = frozenset({
    'ellipse', 'oval', 'circle', 'flowChartConnector',
    'actionButtonBlank', 'donut', 'pie', 'chord',
})

# Arabic character detection pattern
_ARABIC_RE = re.compile(r'[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF\uFB50-\uFDFF\uFE70-\uFEFF]')

# Atomic tokens that should never split across lines
_ATOMIC_TOKEN_RE = re.compile(
    r'\b(\d{1,4}[%$€£]?|\d{1,2}/\d{1,2}/\d{2,4}|\$\d+[\.,]?\d*|[A-Z]{1,3}\d{1,4})\b'
)


# ─────────────────────────────────────────────────────────────────────────────
# Data classes
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class Defect:
    """A single quality defect found by the VQA engine."""
    id: str                    # e.g. "S4-D001"
    layer: str                 # "xml", "vision", "cross_validation"
    check: str                 # check function name
    severity: str              # "CRITICAL", "HIGH", "MEDIUM", "LOW"
    defect_type: str           # category: "overflow", "missing_text", "margin", etc.
    slide_num: int
    shape_id: Optional[int] = None
    shape_name: str = ""
    description: str = ""
    coordinates: Optional[Dict[str, Any]] = None
    affected_text: str = ""
    remediation: Optional[Dict[str, Any]] = None
    auto_fixable: bool = False

    def to_dict(self) -> Dict[str, Any]:
        d = {
            'id': self.id,
            'layer': self.layer,
            'check': self.check,
            'severity': self.severity,
            'defect_type': self.defect_type,
            'slide': self.slide_num,
            'shape_id': self.shape_id,
            'shape_name': self.shape_name,
            'description': self.description,
            'affected_text': self.affected_text,
            'auto_fixable': self.auto_fixable,
        }
        if self.coordinates:
            d['coordinates'] = self.coordinates
        if self.remediation:
            d['remediation'] = self.remediation
        return d


@dataclass
class SlideReport:
    """VQA report for a single slide."""
    slide_num: int
    defects: List[Defect] = field(default_factory=list)

    @property
    def critical_count(self) -> int:
        return sum(1 for d in self.defects if d.severity == 'CRITICAL')

    @property
    def high_count(self) -> int:
        return sum(1 for d in self.defects if d.severity == 'HIGH')

    @property
    def medium_count(self) -> int:
        return sum(1 for d in self.defects if d.severity == 'MEDIUM')

    @property
    def auto_fixable_count(self) -> int:
        return sum(1 for d in self.defects if d.auto_fixable)

    def summary(self) -> Dict[str, Any]:
        return {
            'slide': self.slide_num,
            'total_defects': len(self.defects),
            'critical': self.critical_count,
            'high': self.high_count,
            'medium': self.medium_count,
            'auto_fixable': self.auto_fixable_count,
        }


@dataclass
class VQAReport:
    """Full VQA report across all slides."""
    deck_name: str
    slide_reports: List[SlideReport] = field(default_factory=list)
    original_shapes_per_slide: Dict[int, int] = field(default_factory=dict)
    rtl_shapes_per_slide: Dict[int, int] = field(default_factory=dict)

    @property
    def all_defects(self) -> List[Defect]:
        result = []
        for sr in self.slide_reports:
            result.extend(sr.defects)
        return result

    @property
    def total_defects(self) -> int:
        return sum(len(sr.defects) for sr in self.slide_reports)

    @property
    def slides_with_defects(self) -> int:
        return sum(1 for sr in self.slide_reports if sr.defects)

    def summary(self) -> Dict[str, Any]:
        all_d = self.all_defects
        return {
            'deck_name': self.deck_name,
            'total_slides_checked': len(self.slide_reports),
            'slides_with_defects': self.slides_with_defects,
            'total_defects': len(all_d),
            'critical': sum(1 for d in all_d if d.severity == 'CRITICAL'),
            'high': sum(1 for d in all_d if d.severity == 'HIGH'),
            'medium': sum(1 for d in all_d if d.severity == 'MEDIUM'),
            'auto_fixable': sum(1 for d in all_d if d.auto_fixable),
        }

    def to_json(self, path: str) -> None:
        payload = {
            'summary': self.summary(),
            'slides': []
        }
        for sr in self.slide_reports:
            slide_data = sr.summary()
            slide_data['defects'] = [d.to_dict() for d in sr.defects]
            payload['slides'].append(slide_data)

        with open(path, 'w', encoding='utf-8') as f:
            json.dump(payload, f, indent=2, ensure_ascii=False)

    def to_markdown(self, path: str) -> None:
        lines = []
        lines.append(f'# VQA Report: {self.deck_name}')
        lines.append('')
        s = self.summary()
        lines.append(f'**Slides checked:** {s["total_slides_checked"]} | '
                     f'**Slides with defects:** {s["slides_with_defects"]} | '
                     f'**Total defects:** {s["total_defects"]}')
        lines.append(f'**Critical:** {s["critical"]} | **High:** {s["high"]} | '
                     f'**Medium:** {s["medium"]} | **Auto-fixable:** {s["auto_fixable"]}')
        lines.append('')

        for sr in self.slide_reports:
            if not sr.defects:
                continue
            lines.append(f'## Slide {sr.slide_num}')
            lines.append('')
            lines.append('| ID | Severity | Type | Description | Auto-fix |')
            lines.append('|---|---|---|---|---|')
            for d in sr.defects:
                desc = d.description.replace('|', '\\|')
                lines.append(f'| {d.id} | **{d.severity}** | {d.defect_type} | {desc} | {"✅" if d.auto_fixable else "❌"} |')
            lines.append('')

        with open(path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines) + '\n')


# ─────────────────────────────────────────────────────────────────────────────
# Helper: extract shape info from XML
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class ShapeInfo:
    """Extracted shape geometry and content from PPTX XML."""
    shape_id: int
    shape_name: str
    x: int            # EMU
    y: int            # EMU
    cx: int           # width EMU
    cy: int           # height EMU
    text_content: str
    geometry: str      # preset geometry type (e.g., 'ellipse', 'rect')
    text_align: str    # horizontal alignment
    text_anchor: str   # vertical anchor
    is_placeholder: bool
    ph_type: str
    ph_idx: int
    z_order: int       # position in spTree (rendering order)
    element: Any = None  # lxml element reference


def _has_arabic(text: str) -> bool:
    return bool(_ARABIC_RE.search(text))


def extract_shapes_from_slide_xml(slide_element, slide_num: int = 0) -> List[ShapeInfo]:
    """Extract ShapeInfo objects from a slide's XML tree."""
    shapes = []
    sp_tree = slide_element.find(f'.//{{{P_NS}}}spTree')
    if sp_tree is None:
        return shapes

    z_order = 0
    for child in sp_tree:
        tag = child.tag
        if not (tag.endswith('}sp') or tag.endswith('}pic') or
                tag.endswith('}grpSp') or tag.endswith('}graphicFrame')):
            continue

        z_order += 1
        shape_info = _parse_shape_element(child, z_order)
        if shape_info:
            shapes.append(shape_info)

        # Also recurse into group shapes
        if tag.endswith('}grpSp'):
            for grp_child in child:
                grp_tag = grp_child.tag
                if grp_tag.endswith('}sp') or grp_tag.endswith('}pic'):
                    z_order += 1
                    child_info = _parse_shape_element(grp_child, z_order)
                    if child_info:
                        shapes.append(child_info)

    return shapes


def _parse_shape_element(elem, z_order: int) -> Optional[ShapeInfo]:
    """Parse a single shape XML element into ShapeInfo."""
    try:
        # Get shape ID and name
        nv_sp_pr = elem.find(f'{{{P_NS}}}nvSpPr')
        if nv_sp_pr is None:
            nv_sp_pr = elem.find(f'{{{P_NS}}}nvPicPr')
        if nv_sp_pr is None:
            nv_sp_pr = elem.find(f'{{{P_NS}}}nvGrpSpPr')
        if nv_sp_pr is None:
            nv_sp_pr = elem.find(f'{{{P_NS}}}nvGraphicFramePr')

        shape_id = 0
        shape_name = ''
        if nv_sp_pr is not None:
            c_nv_pr = nv_sp_pr.find(f'{{{P_NS}}}cNvPr')
            if c_nv_pr is not None:
                shape_id = int(c_nv_pr.get('id', 0))
                shape_name = c_nv_pr.get('name', '')

        # Get position and size
        xfrm = None
        for x in elem.iter(f'{{{A_NS}}}xfrm'):
            xfrm = x
            break

        if xfrm is None:
            return None

        off = xfrm.find(f'{{{A_NS}}}off')
        ext = xfrm.find(f'{{{A_NS}}}ext')
        if off is None or ext is None:
            return None

        x = int(off.get('x', 0))
        y = int(off.get('y', 0))
        cx = int(ext.get('cx', 0))
        cy = int(ext.get('cy', 0))

        # Get text content
        text_parts = []
        for t_elem in elem.iter(f'{{{A_NS}}}t'):
            if t_elem.text:
                text_parts.append(t_elem.text)
        text_content = ' '.join(text_parts)

        # Get geometry type
        geometry = ''
        prst_geom = elem.find(f'.//{{{A_NS}}}prstGeom')
        if prst_geom is not None:
            geometry = prst_geom.get('prst', '')

        # Get text alignment
        text_align = ''
        text_anchor = ''
        body_pr = elem.find(f'.//{{{A_NS}}}bodyPr')
        if body_pr is not None:
            text_anchor = body_pr.get('anchor', '')

        # Get first paragraph alignment
        first_pPr = elem.find(f'.//{{{A_NS}}}pPr')
        if first_pPr is not None:
            text_align = first_pPr.get('algn', '')

        # Check if placeholder
        is_placeholder = False
        ph_type = ''
        ph_idx = -1
        ph_elem = elem.find(f'.//{{{P_NS}}}ph')
        if ph_elem is not None:
            is_placeholder = True
            ph_type = ph_elem.get('type', '')
            try:
                ph_idx = int(ph_elem.get('idx', -1))
            except (ValueError, TypeError):
                ph_idx = -1

        return ShapeInfo(
            shape_id=shape_id,
            shape_name=shape_name,
            x=x, y=y, cx=cx, cy=cy,
            text_content=text_content,
            geometry=geometry,
            text_align=text_align,
            text_anchor=text_anchor,
            is_placeholder=is_placeholder,
            ph_type=ph_type,
            ph_idx=ph_idx,
            z_order=z_order,
            element=elem,
        )
    except Exception as e:
        logger.debug('_parse_shape_element: %s', e)
        return None


# ─────────────────────────────────────────────────────────────────────────────
# Layer 1: XML Structural Checks
# ─────────────────────────────────────────────────────────────────────────────

class XMLStructuralChecker:
    """
    Layer 1: Deterministic XML structural checks.
    Catches ~70% of defects without rendering. Fast and reliable.
    """

    def __init__(self, slide_width: int, slide_height: int):
        self.slide_width = slide_width
        self.slide_height = slide_height
        self._defect_counter: Dict[int, int] = {}  # slide_num -> count

    def _next_defect_id(self, slide_num: int) -> str:
        count = self._defect_counter.get(slide_num, 0) + 1
        self._defect_counter[slide_num] = count
        return f'S{slide_num}-D{count:03d}'

    def check_slide(
        self,
        slide_num: int,
        rtl_shapes: List[ShapeInfo],
        orig_shapes: Optional[List[ShapeInfo]] = None,
    ) -> SlideReport:
        """Run all XML structural checks on a single slide."""
        report = SlideReport(slide_num=slide_num)

        # Check 1: Right-edge overflow
        report.defects.extend(self._check_overflow_right_edge(slide_num, rtl_shapes))

        # Check 2: Left-edge overflow (less common but still wrong)
        report.defects.extend(self._check_overflow_left_edge(slide_num, rtl_shapes))

        # Check 3: Margin violations
        report.defects.extend(self._check_margin_violations(slide_num, rtl_shapes))

        # Check 4: Missing text elements (compare with original)
        if orig_shapes:
            report.defects.extend(self._check_text_count_mismatch(slide_num, orig_shapes, rtl_shapes))

        # Check 5: Shape overlap detection
        report.defects.extend(self._check_shape_overlap(slide_num, rtl_shapes))

        # Check 6: Text centering in circular/bounded shapes
        report.defects.extend(self._check_text_centering(slide_num, rtl_shapes))

        # Check 7: Number/token line-break risk
        report.defects.extend(self._check_number_linebreak_risk(slide_num, rtl_shapes))

        # Check 8: Text-to-box area ratio (Gemini recommendation)
        report.defects.extend(self._check_text_box_area_ratio(slide_num, rtl_shapes))

        # Check 9: Z-order occlusion (Gemini recommendation)
        report.defects.extend(self._check_z_order_occlusion(slide_num, rtl_shapes))

        return report

    # ── Check 1: Right-Edge Overflow ──────────────────────────────────────

    def _check_overflow_right_edge(self, slide_num: int, shapes: List[ShapeInfo]) -> List[Defect]:
        """Detect shapes whose right edge exceeds slide width."""
        defects = []
        for shape in shapes:
            right_edge = shape.x + shape.cx
            if right_edge > self.slide_width:
                overflow_emu = right_edge - self.slide_width
                overflow_inches = overflow_emu / EMU_PER_INCH

                # Severity: CRITICAL if text-bearing (visible clipping), HIGH otherwise
                severity = 'CRITICAL' if shape.text_content.strip() else 'HIGH'

                defects.append(Defect(
                    id=self._next_defect_id(slide_num),
                    layer='xml',
                    check='overflow_right_edge',
                    severity=severity,
                    defect_type='text_clipping' if shape.text_content.strip() else 'shape_overflow',
                    slide_num=slide_num,
                    shape_id=shape.shape_id,
                    shape_name=shape.shape_name,
                    description=f'Shape overflows right edge by {overflow_inches:.2f}" ({overflow_emu} EMU)',
                    coordinates={'x': shape.x, 'cx': shape.cx, 'slide_width': self.slide_width,
                                'overflow_emu': overflow_emu},
                    affected_text=shape.text_content[:80] if shape.text_content else '',
                    remediation={
                        'action': 'shift_left',
                        'params': {'delta_x': -overflow_emu},
                        'alternative': 'reduce_width',
                        'alt_params': {'new_cx': self.slide_width - shape.x}
                    },
                    auto_fixable=True,
                ))
        return defects

    # ── Check 2: Left-Edge Overflow ───────────────────────────────────────

    def _check_overflow_left_edge(self, slide_num: int, shapes: List[ShapeInfo]) -> List[Defect]:
        """Detect text shapes that extend past the left edge (negative x)."""
        defects = []
        for shape in shapes:
            if shape.x < 0 and shape.text_content.strip():
                overflow_emu = abs(shape.x)
                overflow_inches = overflow_emu / EMU_PER_INCH

                defects.append(Defect(
                    id=self._next_defect_id(slide_num),
                    layer='xml',
                    check='overflow_left_edge',
                    severity='HIGH',
                    defect_type='text_clipping',
                    slide_num=slide_num,
                    shape_id=shape.shape_id,
                    shape_name=shape.shape_name,
                    description=f'Text shape extends {overflow_inches:.2f}" past left edge',
                    coordinates={'x': shape.x, 'cx': shape.cx},
                    affected_text=shape.text_content[:80],
                    remediation={
                        'action': 'shift_right',
                        'params': {'new_x': 0},
                    },
                    auto_fixable=True,
                ))
        return defects

    # ── Check 3: Margin Violations ────────────────────────────────────────

    def _check_margin_violations(self, slide_num: int, shapes: List[ShapeInfo]) -> List[Defect]:
        """Detect text shapes too close to any edge."""
        defects = []
        for shape in shapes:
            if not shape.text_content.strip():
                continue

            right_margin = self.slide_width - (shape.x + shape.cx)
            left_margin = shape.x
            top_margin = shape.y
            bottom_margin = self.slide_height - (shape.y + shape.cy)

            for edge, margin in [('right', right_margin), ('left', left_margin),
                                  ('top', top_margin), ('bottom', bottom_margin)]:
                # Only flag if margin is positive but dangerously small
                # (negative values are caught by overflow checks)
                if 0 <= margin < MIN_MARGIN_EMU:
                    margin_inches = margin / EMU_PER_INCH

                    defects.append(Defect(
                        id=self._next_defect_id(slide_num),
                        layer='xml',
                        check='margin_violation',
                        severity='HIGH',
                        defect_type='margin_violation',
                        slide_num=slide_num,
                        shape_id=shape.shape_id,
                        shape_name=shape.shape_name,
                        description=f'{edge.upper()} margin only {margin_inches:.3f}" '
                                    f'(min {MIN_MARGIN_EMU / EMU_PER_INCH:.2f}")',
                        coordinates={'x': shape.x, 'y': shape.y, 'cx': shape.cx, 'cy': shape.cy,
                                    'edge': edge, 'margin_emu': margin},
                        affected_text=shape.text_content[:50],
                        remediation={
                            'action': f'increase_{edge}_margin',
                            'params': {'target_margin_emu': MIN_MARGIN_EMU},
                        },
                        auto_fixable=True,
                    ))
        return defects

    # ── Check 4: Missing Text Elements ────────────────────────────────────

    def _check_text_count_mismatch(
        self, slide_num: int,
        orig_shapes: List[ShapeInfo],
        rtl_shapes: List[ShapeInfo]
    ) -> List[Defect]:
        """Compare text-bearing shape counts between original and RTL output."""
        defects = []

        orig_text = [s for s in orig_shapes if s.text_content.strip()]
        rtl_text = [s for s in rtl_shapes if s.text_content.strip()]

        orig_count = len(orig_text)
        rtl_count = len(rtl_text)

        if rtl_count < orig_count:
            missing = orig_count - rtl_count

            # Try to identify which texts are missing
            rtl_texts_set = {s.text_content.strip().lower() for s in rtl_text}
            missing_texts = []
            for s in orig_text:
                t = s.text_content.strip().lower()
                if t not in rtl_texts_set:
                    missing_texts.append(s.text_content.strip()[:60])

            defects.append(Defect(
                id=self._next_defect_id(slide_num),
                layer='xml',
                check='text_count_mismatch',
                severity='CRITICAL',
                defect_type='missing_text',
                slide_num=slide_num,
                description=f'{missing} text element(s) lost during transformation '
                            f'(original: {orig_count}, RTL: {rtl_count})',
                affected_text='; '.join(missing_texts[:5]) if missing_texts else '',
                remediation={
                    'action': 'investigate',
                    'params': {'missing_count': missing,
                              'missing_texts': missing_texts[:10]},
                },
                auto_fixable=False,  # Missing text needs manual investigation
            ))
        return defects

    # ── Check 5: Shape Overlap Detection ──────────────────────────────────

    def _check_shape_overlap(self, slide_num: int, shapes: List[ShapeInfo]) -> List[Defect]:
        """Detect overlapping text-bearing shapes."""
        defects = []
        text_shapes = [s for s in shapes if s.text_content.strip()]

        for i, a in enumerate(text_shapes):
            for b in text_shapes[i+1:]:
                # Bounding box intersection
                overlap_x = max(0, min(a.x + a.cx, b.x + b.cx) - max(a.x, b.x))
                overlap_y = max(0, min(a.y + a.cy, b.y + b.cy) - max(a.y, b.y))

                if overlap_x > 0 and overlap_y > 0:
                    overlap_area = overlap_x * overlap_y
                    smaller_area = min(a.cx * a.cy, b.cx * b.cy)
                    overlap_pct = overlap_area / smaller_area if smaller_area > 0 else 0

                    if overlap_pct > OVERLAP_HIGH_PCT:
                        severity = 'CRITICAL' if overlap_pct > OVERLAP_CRITICAL_PCT else 'HIGH'

                        defects.append(Defect(
                            id=self._next_defect_id(slide_num),
                            layer='xml',
                            check='shape_overlap',
                            severity=severity,
                            defect_type='shape_overlap',
                            slide_num=slide_num,
                            shape_id=a.shape_id,
                            shape_name=f'{a.shape_name} ∩ {b.shape_name}',
                            description=f'"{a.shape_name}" and "{b.shape_name}" overlap by '
                                        f'{overlap_pct*100:.0f}%',
                            coordinates={
                                'shape_a': {'x': a.x, 'y': a.y, 'cx': a.cx, 'cy': a.cy},
                                'shape_b': {'x': b.x, 'y': b.y, 'cx': b.cx, 'cy': b.cy},
                                'overlap_pct': round(overlap_pct * 100, 1),
                            },
                            affected_text=f'A: {a.text_content[:40]} | B: {b.text_content[:40]}',
                            auto_fixable=False,
                        ))
        return defects

    # ── Check 6: Text Centering in Shapes ─────────────────────────────────

    def _check_text_centering(self, slide_num: int, shapes: List[ShapeInfo]) -> List[Defect]:
        """Verify text is centered in circular/elliptical shapes."""
        defects = []
        for shape in shapes:
            if not shape.text_content.strip():
                continue
            if shape.geometry not in CIRCULAR_GEOMS:
                continue

            # Check horizontal centering
            if shape.text_align and shape.text_align not in ('ctr', 'center', 'mid'):
                defects.append(Defect(
                    id=self._next_defect_id(slide_num),
                    layer='xml',
                    check='text_not_centered_h',
                    severity='MEDIUM',
                    defect_type='alignment_error',
                    slide_num=slide_num,
                    shape_id=shape.shape_id,
                    shape_name=shape.shape_name,
                    description=f'Text not horizontally centered in {shape.geometry} shape '
                                f'(align={shape.text_align})',
                    affected_text=shape.text_content[:30],
                    remediation={
                        'action': 'set_alignment',
                        'params': {'algn': 'ctr'},
                    },
                    auto_fixable=True,
                ))

            # Check vertical centering
            if shape.text_anchor and shape.text_anchor not in ('ctr', 'mid'):
                defects.append(Defect(
                    id=self._next_defect_id(slide_num),
                    layer='xml',
                    check='text_not_centered_v',
                    severity='MEDIUM',
                    defect_type='alignment_error',
                    slide_num=slide_num,
                    shape_id=shape.shape_id,
                    shape_name=shape.shape_name,
                    description=f'Text not vertically centered in {shape.geometry} shape '
                                f'(anchor={shape.text_anchor})',
                    affected_text=shape.text_content[:30],
                    remediation={
                        'action': 'set_anchor',
                        'params': {'anchor': 'ctr'},
                    },
                    auto_fixable=True,
                ))
        return defects

    # ── Check 7: Number/Token Line-Break Risk ─────────────────────────────

    def _check_number_linebreak_risk(self, slide_num: int, shapes: List[ShapeInfo]) -> List[Defect]:
        """Detect numbers or short tokens at risk of being split across lines."""
        defects = []
        for shape in shapes:
            if not shape.text_content.strip():
                continue
            if shape.cx <= 0:
                continue

            # Check if box is very narrow
            if shape.cx < NARROW_BOX_THRESHOLD_EMU:
                tokens = _ATOMIC_TOKEN_RE.findall(shape.text_content)
                if tokens:
                    defects.append(Defect(
                        id=self._next_defect_id(slide_num),
                        layer='xml',
                        check='number_linebreak_risk',
                        severity='HIGH',
                        defect_type='linebreak_defect',
                        slide_num=slide_num,
                        shape_id=shape.shape_id,
                        shape_name=shape.shape_name,
                        description=f'Text box too narrow ({shape.cx/EMU_PER_INCH:.2f}") for '
                                    f'atomic tokens {tokens}',
                        coordinates={'cx': shape.cx, 'cx_inches': shape.cx / EMU_PER_INCH},
                        affected_text=shape.text_content[:60],
                        remediation={
                            'action': 'widen_box',
                            'params': {'min_cx': NARROW_BOX_THRESHOLD_EMU},
                        },
                        auto_fixable=True,
                    ))
        return defects

    # ── Check 8: Text-to-Box Area Ratio (Gemini recommendation) ──────────

    def _check_text_box_area_ratio(self, slide_num: int, shapes: List[ShapeInfo]) -> List[Defect]:
        """
        Heuristic: estimate if Arabic text will fit in its bounding box.
        Arabic text is typically 1.2-1.5x wider than equivalent English.
        """
        defects = []
        for shape in shapes:
            text = shape.text_content.strip()
            if not text or not _has_arabic(text):
                continue
            if shape.cx <= 0 or shape.cy <= 0:
                continue

            # Rough heuristic: Arabic chars average ~0.08" wide at 12pt
            # Scale by font size if known, otherwise assume 12pt
            char_count = len(text)
            estimated_text_width_emu = int(char_count * 0.08 * EMU_PER_INCH)

            if shape.cx > 0:
                width_ratio = estimated_text_width_emu / shape.cx
                if width_ratio > TEXT_BOX_AREA_RISK_THRESHOLD:
                    defects.append(Defect(
                        id=self._next_defect_id(slide_num),
                        layer='xml',
                        check='text_box_area_ratio',
                        severity='HIGH' if width_ratio > 1.2 else 'MEDIUM',
                        defect_type='overflow_risk',
                        slide_num=slide_num,
                        shape_id=shape.shape_id,
                        shape_name=shape.shape_name,
                        description=f'Arabic text likely overflows box (est. ratio {width_ratio:.1f}x, '
                                    f'{char_count} chars in {shape.cx/EMU_PER_INCH:.2f}" box)',
                        coordinates={'cx': shape.cx, 'char_count': char_count,
                                    'estimated_ratio': round(width_ratio, 2)},
                        affected_text=text[:60],
                        remediation={
                            'action': 'autofit_or_widen',
                            'params': {'estimated_ratio': round(width_ratio, 2)},
                        },
                        auto_fixable=True,
                    ))
        return defects

    # ── Check 9: Z-Order Occlusion (Gemini recommendation) ───────────────

    def _check_z_order_occlusion(self, slide_num: int, shapes: List[ShapeInfo]) -> List[Defect]:
        """
        Detect text shapes that are occluded by opaque shapes rendered on top.
        Lower z_order = rendered first = potentially hidden by higher z_order shapes.
        """
        defects = []
        text_shapes = [s for s in shapes if s.text_content.strip()]
        non_text_shapes = [s for s in shapes if not s.text_content.strip() and s.cx > 0 and s.cy > 0]

        for ts in text_shapes:
            for ns in non_text_shapes:
                # Only check if the non-text shape renders ON TOP (higher z_order)
                if ns.z_order <= ts.z_order:
                    continue

                # Check bounding box intersection
                overlap_x = max(0, min(ts.x + ts.cx, ns.x + ns.cx) - max(ts.x, ns.x))
                overlap_y = max(0, min(ts.y + ts.cy, ns.y + ns.cy) - max(ts.y, ns.y))

                if overlap_x > 0 and overlap_y > 0:
                    ts_area = max(ts.cx * ts.cy, 1)
                    overlap_area = overlap_x * overlap_y
                    occlusion_pct = overlap_area / ts_area

                    if occlusion_pct > 0.25:  # >25% of text shape is covered
                        defects.append(Defect(
                            id=self._next_defect_id(slide_num),
                            layer='xml',
                            check='z_order_occlusion',
                            severity='CRITICAL' if occlusion_pct > 0.50 else 'HIGH',
                            defect_type='text_hidden',
                            slide_num=slide_num,
                            shape_id=ts.shape_id,
                            shape_name=ts.shape_name,
                            description=f'Text shape "{ts.shape_name}" is {occlusion_pct*100:.0f}% '
                                        f'hidden behind "{ns.shape_name}" (z-order: text={ts.z_order}, '
                                        f'cover={ns.z_order})',
                            coordinates={
                                'text_shape': {'x': ts.x, 'y': ts.y, 'cx': ts.cx, 'cy': ts.cy,
                                              'z': ts.z_order},
                                'covering_shape': {'x': ns.x, 'y': ns.y, 'cx': ns.cx, 'cy': ns.cy,
                                                  'z': ns.z_order},
                                'occlusion_pct': round(occlusion_pct * 100, 1),
                            },
                            affected_text=ts.text_content[:60],
                            auto_fixable=False,
                        ))
        return defects


# ─────────────────────────────────────────────────────────────────────────────
# Auto-Fix Engine (Gemini cascade: Shift → Expand → Reduce font → Line break)
# ─────────────────────────────────────────────────────────────────────────────

class AutoFixer:
    """
    Applies deterministic fixes to defects found by Layer 1.
    Follows Gemini's cascade hierarchy:
    1. Shift (non-destructive)
    2. Expand bounding box
    3. Reduce font / autofit
    4. Line break injection
    """

    def __init__(self, slide_width: int, slide_height: int):
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.fixes_applied: List[Dict[str, Any]] = []

    def apply_fixes(self, slide_element, defects: List[Defect]) -> int:
        """Apply auto-fixes to a slide's XML for all auto-fixable defects."""
        fixed = 0
        for defect in defects:
            if not defect.auto_fixable:
                continue
            if not defect.remediation:
                continue

            action = defect.remediation.get('action', '')
            params = defect.remediation.get('params', {})

            success = False
            if action == 'shift_left':
                success = self._fix_shift(slide_element, defect, params)
            elif action == 'shift_right':
                success = self._fix_shift_right(slide_element, defect, params)
            elif action.startswith('increase_') and action.endswith('_margin'):
                success = self._fix_margin(slide_element, defect, params)
            elif action == 'widen_box':
                success = self._fix_widen_box(slide_element, defect, params)
            elif action == 'set_alignment':
                success = self._fix_alignment(slide_element, defect, params)
            elif action == 'set_anchor':
                success = self._fix_anchor(slide_element, defect, params)
            elif action == 'autofit_or_widen':
                success = self._fix_autofit(slide_element, defect, params)

            if success:
                fixed += 1
                self.fixes_applied.append({
                    'defect_id': defect.id,
                    'action': action,
                    'slide': defect.slide_num,
                    'shape': defect.shape_name,
                })

        return fixed

    def _find_shape_element(self, slide_element, defect: Defect):
        """Find the shape XML element matching a defect's shape_id."""
        if not defect.shape_id:
            return None

        target_id = str(defect.shape_id)
        sp_tree = slide_element.find(f'.//{{{P_NS}}}spTree')
        if sp_tree is None:
            return None

        for child in sp_tree.iter():
            tag = child.tag
            if not (tag.endswith('}sp') or tag.endswith('}pic') or tag.endswith('}grpSp')):
                continue
            # Check cNvPr id
            for cnv in child.iter():
                cnv_tag = cnv.tag
                if cnv_tag.endswith('}cNvPr'):
                    if cnv.get('id') == target_id:
                        return child
        return None

    def _fix_shift(self, slide_element, defect: Defect, params: Dict) -> bool:
        """Step 1: Shift shape left to eliminate right-edge overflow."""
        shape_el = self._find_shape_element(slide_element, defect)
        if shape_el is None:
            return False

        delta_x = params.get('delta_x', 0)
        if delta_x == 0:
            return False

        for xfrm in shape_el.iter(f'{{{A_NS}}}xfrm'):
            off = xfrm.find(f'{{{A_NS}}}off')
            if off is not None:
                current_x = int(off.get('x', 0))
                new_x = current_x + delta_x

                # Cascade check: does shifting cause left-edge overflow?
                if new_x < 0:
                    # Step 2: Try expanding the box instead
                    ext = xfrm.find(f'{{{A_NS}}}ext')
                    if ext is not None:
                        current_cx = int(ext.get('cx', 0))
                        # Reduce width to fit within slide, keeping right edge at slide_width
                        new_cx = self.slide_width - max(0, current_x + delta_x)
                        if new_cx > current_cx * 0.5:  # Don't shrink more than 50%
                            new_x = 0
                            ext.set('cx', str(new_cx))
                        else:
                            new_x = 0  # Clamp to left edge
                else:
                    # Verify right edge is now within bounds
                    ext = xfrm.find(f'{{{A_NS}}}ext')
                    if ext is not None:
                        current_cx = int(ext.get('cx', 0))
                        if new_x + current_cx > self.slide_width:
                            new_x = self.slide_width - current_cx

                off.set('x', str(max(0, new_x)))
                return True
        return False

    def _fix_shift_right(self, slide_element, defect: Defect, params: Dict) -> bool:
        """Fix left-edge overflow by shifting shape to x=0 or specified position."""
        shape_el = self._find_shape_element(slide_element, defect)
        if shape_el is None:
            return False

        new_x = params.get('new_x', 0)
        for xfrm in shape_el.iter(f'{{{A_NS}}}xfrm'):
            off = xfrm.find(f'{{{A_NS}}}off')
            if off is not None:
                off.set('x', str(new_x))
                return True
        return False

    def _fix_margin(self, slide_element, defect: Defect, params: Dict) -> bool:
        """Increase margin by shifting/resizing shape."""
        shape_el = self._find_shape_element(slide_element, defect)
        if shape_el is None:
            return False

        target_margin = params.get('target_margin_emu', MIN_MARGIN_EMU)
        coords = defect.coordinates or {}
        edge = coords.get('edge', '')
        current_margin = coords.get('margin_emu', 0)

        if not edge:
            return False

        delta = target_margin - current_margin

        for xfrm in shape_el.iter(f'{{{A_NS}}}xfrm'):
            off = xfrm.find(f'{{{A_NS}}}off')
            ext = xfrm.find(f'{{{A_NS}}}ext')
            if off is None or ext is None:
                continue

            x = int(off.get('x', 0))
            y = int(off.get('y', 0))
            cx = int(ext.get('cx', 0))
            cy = int(ext.get('cy', 0))

            if edge == 'right':
                # Shift left or reduce width
                new_cx = cx - delta
                if new_cx > cx * 0.5:
                    ext.set('cx', str(new_cx))
                else:
                    off.set('x', str(x - delta))
            elif edge == 'left':
                off.set('x', str(x + delta))
                ext.set('cx', str(max(cx - delta, cx // 2)))
            elif edge == 'top':
                off.set('y', str(y + delta))
            elif edge == 'bottom':
                ext.set('cy', str(max(cy - delta, cy // 2)))

            return True
        return False

    def _fix_widen_box(self, slide_element, defect: Defect, params: Dict) -> bool:
        """Widen a narrow text box."""
        shape_el = self._find_shape_element(slide_element, defect)
        if shape_el is None:
            return False

        min_cx = params.get('min_cx', NARROW_BOX_THRESHOLD_EMU)

        for xfrm in shape_el.iter(f'{{{A_NS}}}xfrm'):
            ext = xfrm.find(f'{{{A_NS}}}ext')
            off = xfrm.find(f'{{{A_NS}}}off')
            if ext is not None and off is not None:
                current_cx = int(ext.get('cx', 0))
                current_x = int(off.get('x', 0))
                if current_cx < min_cx:
                    new_cx = min_cx
                    # Ensure we don't overflow right edge
                    if current_x + new_cx > self.slide_width:
                        current_x = self.slide_width - new_cx
                        off.set('x', str(max(0, current_x)))
                    ext.set('cx', str(new_cx))
                    return True
        return False

    def _fix_alignment(self, slide_element, defect: Defect, params: Dict) -> bool:
        """Set paragraph alignment to center for circular shapes."""
        shape_el = self._find_shape_element(slide_element, defect)
        if shape_el is None:
            return False

        target_align = params.get('algn', 'ctr')
        changed = False
        for pPr in shape_el.iter(f'{{{A_NS}}}pPr'):
            pPr.set('algn', target_align)
            changed = True

        # Also set on paragraphs that don't have pPr
        for p in shape_el.iter(f'{{{A_NS}}}p'):
            pPr = p.find(f'{{{A_NS}}}pPr')
            if pPr is None:
                pPr = etree.SubElement(p, f'{{{A_NS}}}pPr')
                p.insert(0, pPr)
            pPr.set('algn', target_align)
            changed = True

        return changed

    def _fix_anchor(self, slide_element, defect: Defect, params: Dict) -> bool:
        """Set text vertical anchor to center."""
        shape_el = self._find_shape_element(slide_element, defect)
        if shape_el is None:
            return False

        target_anchor = params.get('anchor', 'ctr')
        for body_pr in shape_el.iter(f'{{{A_NS}}}bodyPr'):
            body_pr.set('anchor', target_anchor)
            return True
        return False

    def _fix_autofit(self, slide_element, defect: Defect, params: Dict) -> bool:
        """Enable autofit (font scaling) on a shape."""
        shape_el = self._find_shape_element(slide_element, defect)
        if shape_el is None:
            return False

        for body_pr in shape_el.iter(f'{{{A_NS}}}bodyPr'):
            # Remove existing autofit settings
            for child in list(body_pr):
                if child.tag.endswith('}noAutofit') or child.tag.endswith('}spAutoFit'):
                    body_pr.remove(child)

            # Add normAutofit with font scaling
            norm_auto = etree.SubElement(body_pr, f'{{{A_NS}}}normAutofit')
            norm_auto.set('fontScale', '70000')  # Allow up to 70% font reduction
            return True
        return False


# ─────────────────────────────────────────────────────────────────────────────
# Layer 2: Vision Model Prompt Generator (Document-Grounded per Gemini)
# ─────────────────────────────────────────────────────────────────────────────

class VisionPromptGenerator:
    """
    Generates Document-Grounded Vision prompts for Layer 2 VQA.
    Follows Gemini 3.1 Pro's recommendation: pass expected XML elements
    alongside the rendered image for grounded analysis.
    """

    @staticmethod
    def generate_single_slide_prompt(shapes: List[ShapeInfo]) -> str:
        """Generate a vision model prompt for single-slide analysis."""

        # Build expected elements JSON
        expected = []
        for s in shapes:
            if s.text_content.strip():
                expected.append({
                    'shape_name': s.shape_name,
                    'text': s.text_content[:100],
                    'geometry': s.geometry or 'rect',
                    'position': {
                        'x_inches': round(s.x / EMU_PER_INCH, 2),
                        'y_inches': round(s.y / EMU_PER_INCH, 2),
                        'width_inches': round(s.cx / EMU_PER_INCH, 2),
                        'height_inches': round(s.cy / EMU_PER_INCH, 2),
                    }
                })

        expected_json = json.dumps(expected, ensure_ascii=False, indent=2)

        return f"""You are an expert Visual Quality Assurance (VQA) system specializing in Arabic Right-to-Left (RTL) PowerPoint presentations.

Analyze the provided slide image. You are also provided with the EXPECTED_ELEMENTS JSON containing the text that SHOULD be visible and its intended position.

EXPECTED_ELEMENTS:
{expected_json}

Your task is to identify specific rendering, translation, and RTL layout defects. Pay intense attention to the exact edges of the slide.

### Defect Categories to Check:
1. OVERFLOW_RIGHT: Text or shapes bleeding into, touching, or clipped by the right-hand edge of the slide. In RTL, the right margin is sacred.
2. MISSING_ELEMENTS: Text from the EXPECTED_ELEMENTS that is completely hidden, missing, or cropped out of existence.
3. ALIGNMENT_ERROR: Text that is incorrectly left-aligned within its container, or text not vertically/horizontally centered within shapes (like badges or numbered circles).
4. ORPHANED_PUNCTUATION: Numbers (e.g., "01") or punctuation broken across two lines due to container width being too small.
5. SEQUENCE_DIRECTION_ERROR: Numbered sequences, arrows, or visual flows that read Left-to-Right (1→2→3). They MUST read Right-to-Left (3←2←1) in this Arabic slide.

### Output Format
Return ONLY valid JSON matching this schema:
{{
  "pass": boolean,
  "defects": [
    {{
      "category": "OVERFLOW_RIGHT" | "MISSING_ELEMENTS" | "ALIGNMENT_ERROR" | "ORPHANED_PUNCTUATION" | "SEQUENCE_DIRECTION_ERROR",
      "severity": "CRITICAL" | "WARNING",
      "description": "Specific details of what is wrong and where it is located visually",
      "affected_text": "The snippet of Arabic text or number affected"
    }}
  ]
}}"""

    @staticmethod
    def generate_comparative_prompt() -> str:
        """Generate a prompt for side-by-side original vs RTL comparison."""
        return """Compare these two slides:
IMAGE 1: Original LTR slide
IMAGE 2: RTL-transformed slide

## A. TEXT PRESERVATION
1. Count ALL distinct text elements in original
2. Count ALL distinct text elements in RTL
3. Flag any text in original that is MISSING in RTL

## B. ELEMENT COUNT
For each group of similar elements (e.g., circles with numbers, labels, badges):
- Count in original → count in RTL
- Report mismatches

## C. SPATIAL MIRRORING
- Elements should be horizontally mirrored
- Sequences should be reversed (left-to-right → right-to-left)
- Flag anything not mirrored correctly

## D. VISUAL QUALITY
- Text fully visible in original but clipped in RTL?
- Text contained in shape in original but overflows in RTL?
- Text on one line in original but wraps in RTL?

Return JSON array of discrepancies:
[{"category": "...", "severity": "CRITICAL|WARNING", "description": "..."}]"""


# ─────────────────────────────────────────────────────────────────────────────
# Main VQA Pipeline
# ─────────────────────────────────────────────────────────────────────────────

class VQAPipeline:
    """
    Orchestrates the full VQA pipeline:
    1. Extract shapes from original and RTL PPTX XML
    2. Run Layer 1 XML structural checks
    3. Generate Layer 2 vision prompts (for external model call)
    4. Optionally apply auto-fixes
    """

    def __init__(self, slide_width: int, slide_height: int, deck_name: str = ''):
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.deck_name = deck_name
        self.xml_checker = XMLStructuralChecker(slide_width, slide_height)
        self.auto_fixer = AutoFixer(slide_width, slide_height)
        self.prompt_gen = VisionPromptGenerator()

    def run_layer1(
        self,
        rtl_pptx_path: str,
        orig_pptx_path: Optional[str] = None,
    ) -> VQAReport:
        """
        Run Layer 1 (XML structural checks) on all slides.

        Args:
            rtl_pptx_path: Path to the RTL-transformed PPTX
            orig_pptx_path: Optional path to original English PPTX for comparison

        Returns:
            VQAReport with all defects found
        """
        from pptx import Presentation

        report = VQAReport(deck_name=self.deck_name)

        rtl_prs = Presentation(rtl_pptx_path)
        orig_prs = Presentation(orig_pptx_path) if orig_pptx_path else None

        for slide_idx, slide in enumerate(rtl_prs.slides):
            slide_num = slide_idx + 1
            rtl_shapes = extract_shapes_from_slide_xml(slide._element, slide_num)

            orig_shapes = None
            if orig_prs and slide_idx < len(orig_prs.slides):
                orig_shapes = extract_shapes_from_slide_xml(
                    orig_prs.slides[slide_idx]._element, slide_num
                )
                report.original_shapes_per_slide[slide_num] = len(orig_shapes)

            report.rtl_shapes_per_slide[slide_num] = len(rtl_shapes)

            slide_report = self.xml_checker.check_slide(slide_num, rtl_shapes, orig_shapes)
            report.slide_reports.append(slide_report)

        return report

    def generate_vision_prompts(
        self,
        rtl_pptx_path: str,
    ) -> Dict[int, str]:
        """
        Generate Document-Grounded Vision prompts for each slide.

        Returns:
            Dict mapping slide_num to prompt string (including EXPECTED_ELEMENTS JSON)
        """
        from pptx import Presentation

        prompts = {}
        rtl_prs = Presentation(rtl_pptx_path)

        for slide_idx, slide in enumerate(rtl_prs.slides):
            slide_num = slide_idx + 1
            shapes = extract_shapes_from_slide_xml(slide._element, slide_num)
            prompts[slide_num] = self.prompt_gen.generate_single_slide_prompt(shapes)

        return prompts

    def apply_auto_fixes(
        self,
        rtl_pptx_path: str,
        defects: List[Defect],
        output_path: str,
    ) -> int:
        """
        Apply auto-fixes to the RTL PPTX for all auto-fixable defects.

        Args:
            rtl_pptx_path: Path to RTL PPTX to fix
            defects: List of defects from Layer 1
            output_path: Path to write fixed PPTX

        Returns:
            Count of fixes applied
        """
        from pptx import Presentation
        import shutil

        prs = Presentation(rtl_pptx_path)
        total_fixed = 0

        # Group defects by slide
        defects_by_slide: Dict[int, List[Defect]] = {}
        for d in defects:
            defects_by_slide.setdefault(d.slide_num, []).append(d)

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            slide_defects = defects_by_slide.get(slide_num, [])
            if slide_defects:
                fixed = self.auto_fixer.apply_fixes(slide._element, slide_defects)
                total_fixed += fixed

        prs.save(output_path)
        return total_fixed

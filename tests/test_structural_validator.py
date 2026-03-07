import unittest
from unittest.mock import MagicMock
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from slidearabi.structural_validator import StructuralValidator, IssueType

class TestStructuralValidator(unittest.TestCase):
    
    def setUp(self):
        self.mock_prs = MagicMock()
        self.mock_prs.slide_width = 10000000
        self.mock_prs.slide_height = 7500000
        self.mock_prs.slides = []
        self.mock_prs.slide_masters = []
        
        self.validator = StructuralValidator(self.mock_prs)
        
    def test_rtl_direction_check(self):
        # Create a shape with Arabic text but missing RTL attribute
        mock_shape = MagicMock()
        mock_shape.shape_id = 1
        mock_shape.name = "TextBox 1"
        mock_shape.has_text_frame = True
        
        mock_paragraph = MagicMock()
        mock_paragraph.text = "هذا نص عربي"  # Arabic text
        mock_paragraph._pPr = MagicMock()
        mock_paragraph._pPr.get.return_value = None  # Missing rtl='1'
        
        mock_shape.text_frame.paragraphs = [mock_paragraph]
        
        issues = self.validator._check_rtl_direction(mock_shape, 1)
        
        self.assertEqual(len(issues), 1)
        self.assertEqual(issues[0].issue_type, IssueType.RTL_MISSING)
        self.assertEqual(issues[0].severity, 'error')
        
    def test_alignment_check(self):
        # Arabic text should not be left-aligned
        mock_shape = MagicMock()
        mock_shape.shape_id = 1
        mock_shape.name = "TextBox 1"
        mock_shape.is_placeholder = False
        
        mock_paragraph = MagicMock()
        mock_paragraph.text = "هذا نص عربي"
        mock_paragraph.alignment = PP_ALIGN.LEFT
        
        mock_shape.text_frame.paragraphs = [mock_paragraph]
        
        issues = self.validator._check_alignment(mock_shape, 1)
        
        self.assertEqual(len(issues), 1)
        self.assertEqual(issues[0].issue_type, IssueType.ALIGNMENT_WRONG)
        self.assertEqual(issues[0].severity, 'error')
        
    def test_shape_bounds_check(self):
        # Shape outside slide boundaries
        mock_shape = MagicMock()
        mock_shape.shape_id = 1
        mock_shape.name = "TextBox 1"
        mock_shape.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_shape.has_text_frame = True
        
        # Position past the right edge
        mock_shape.left = self.validator.slide_width + 500000
        mock_shape.top = 100000
        mock_shape.width = 1000000
        mock_shape.height = 500000
        
        issues = self.validator._check_shape_bounds(mock_shape, 1)
        
        self.assertEqual(len(issues), 1)
        self.assertEqual(issues[0].issue_type, IssueType.SHAPE_OUT_OF_BOUNDS)
        self.assertEqual(issues[0].severity, 'error')
        
    def test_overlap_check(self):
        mock_slide = MagicMock()
        
        # Two text boxes that overlap significantly
        s1 = MagicMock()
        s1.shape_id = 1
        s1.name = "T1"
        s1.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        s1.has_text_frame = True
        s1.left, s1.top, s1.width, s1.height = 100, 100, 100, 100  # 100x100 at 100,100
        
        s2 = MagicMock()
        s2.shape_id = 2
        s2.name = "T2"
        s2.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        s2.has_text_frame = True
        s2.left, s2.top, s2.width, s2.height = 150, 150, 100, 100  # Overlaps s1 by 50x50 (25% area)
        
        mock_slide.shapes = [s1, s2]
        
        issues = self.validator._check_shape_overlaps(mock_slide, 1)
        
        self.assertEqual(len(issues), 1)
        self.assertEqual(issues[0].issue_type, IssueType.SHAPE_OVERLAP)
        self.assertEqual(issues[0].severity, 'error')

if __name__ == '__main__':
    unittest.main()

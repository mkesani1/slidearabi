"""
Tests for SlideArabi PropertyResolver and Models.

These tests construct mock OOXML XML elements using lxml to simulate
the 7-level inheritance chain without needing real PPTX files.

Test categories:
1. Model construction and immutability
2. Font size resolution through all 7 levels
3. Font name resolution including theme font references
4. Alignment resolution
5. Bold/italic resolution
6. Placeholder matching (by idx and by type)
7. Master txStyles resolution
8. RTL resolution
9. Edge cases and defaults
"""

from __future__ import annotations

import pytest
from lxml import etree
from unittest.mock import MagicMock, PropertyMock, patch
from dataclasses import FrozenInstanceError

from slidearabi.models import (
    A_NS,
    DEFAULT_ALIGNMENT,
    DEFAULT_BOLD,
    DEFAULT_FONT_NAME,
    DEFAULT_FONT_SIZE_PT,
    DEFAULT_ITALIC,
    DEFAULT_RTL,
    P_NS,
    R_NS,
    ResolvedLayout,
    ResolvedMaster,
    ResolvedParagraph,
    ResolvedPresentation,
    ResolvedRun,
    ResolvedShape,
    ResolvedSlide,
    TransformAction,
    TransformPlan,
    ValidationIssue,
    ValidationReport,
)
from slidearabi.property_resolver import PropertyResolver, _qn, _get_attr_int, _get_attr_bool


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# XML BUILDER HELPERS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


def make_rPr(**attrs) -> etree._Element:
    """Build an a:rPr element with given attributes."""
    rPr = etree.Element(f'{{{A_NS}}}rPr')
    for k, v in attrs.items():
        rPr.set(k, str(v))
    return rPr


def make_run(text: str = 'Hello', sz: int = None, b: str = None, i: str = None,
             font_name: str = None) -> etree._Element:
    """Build an a:r element with optional rPr attributes."""
    r = etree.SubElement(etree.Element('dummy'), f'{{{A_NS}}}r')
    r = etree.Element(f'{{{A_NS}}}r')

    if sz is not None or b is not None or i is not None or font_name is not None:
        rPr = etree.SubElement(r, f'{{{A_NS}}}rPr')
        if sz is not None:
            rPr.set('sz', str(sz))
        if b is not None:
            rPr.set('b', str(b))
        if i is not None:
            rPr.set('i', str(i))
        if font_name is not None:
            latin = etree.SubElement(rPr, f'{{{A_NS}}}latin')
            latin.set('typeface', font_name)

    t = etree.SubElement(r, f'{{{A_NS}}}t')
    t.text = text
    return r


def make_paragraph(algn: str = None, lvl: int = None,
                   defRPr_sz: int = None, defRPr_b: str = None,
                   defRPr_i: str = None, defRPr_font: str = None,
                   rtl: str = None,
                   runs: list = None) -> etree._Element:
    """Build an a:p element with optional pPr and runs."""
    p = etree.Element(f'{{{A_NS}}}p')

    need_pPr = algn is not None or lvl is not None or defRPr_sz is not None \
               or defRPr_b is not None or defRPr_i is not None \
               or defRPr_font is not None or rtl is not None
    if need_pPr:
        pPr = etree.SubElement(p, f'{{{A_NS}}}pPr')
        if algn is not None:
            pPr.set('algn', algn)
        if lvl is not None:
            pPr.set('lvl', str(lvl))
        if rtl is not None:
            pPr.set('rtl', str(rtl))
        if defRPr_sz is not None or defRPr_b is not None or defRPr_i is not None or defRPr_font is not None:
            defRPr = etree.SubElement(pPr, f'{{{A_NS}}}defRPr')
            if defRPr_sz is not None:
                defRPr.set('sz', str(defRPr_sz))
            if defRPr_b is not None:
                defRPr.set('b', str(defRPr_b))
            if defRPr_i is not None:
                defRPr.set('i', str(defRPr_i))
            if defRPr_font is not None:
                latin = etree.SubElement(defRPr, f'{{{A_NS}}}latin')
                latin.set('typeface', defRPr_font)

    if runs:
        for r in runs:
            p.append(r)

    return p


def make_lst_style(level_props: dict = None) -> etree._Element:
    """Build an a:lstStyle element.

    Args:
        level_props: dict of level_num → dict of properties.
            Example: {1: {'sz': '2400', 'b': '1', 'algn': 'ctr', 'font': 'Arial'}}
    """
    lstStyle = etree.Element(f'{{{A_NS}}}lstStyle')
    if level_props:
        for level_num, props in level_props.items():
            lvl_pPr = etree.SubElement(lstStyle, f'{{{A_NS}}}lvl{level_num}pPr')
            if 'algn' in props:
                lvl_pPr.set('algn', props['algn'])
            if 'rtl' in props:
                lvl_pPr.set('rtl', props['rtl'])
            if any(k in props for k in ('sz', 'b', 'i', 'font')):
                defRPr = etree.SubElement(lvl_pPr, f'{{{A_NS}}}defRPr')
                if 'sz' in props:
                    defRPr.set('sz', str(props['sz']))
                if 'b' in props:
                    defRPr.set('b', str(props['b']))
                if 'i' in props:
                    defRPr.set('i', str(props['i']))
                if 'font' in props:
                    latin = etree.SubElement(defRPr, f'{{{A_NS}}}latin')
                    latin.set('typeface', props['font'])
    return lstStyle


def make_txBody(lst_style_props: dict = None) -> etree._Element:
    """Build an a:txBody element with optional lstStyle."""
    txBody = etree.Element(f'{{{A_NS}}}txBody')
    etree.SubElement(txBody, f'{{{A_NS}}}bodyPr')
    if lst_style_props:
        lstStyle = make_lst_style(lst_style_props)
        txBody.append(lstStyle)
    else:
        etree.SubElement(txBody, f'{{{A_NS}}}lstStyle')
    return txBody


def make_sp_element(ph_type: str = None, ph_idx: int = None,
                    lst_style_props: dict = None,
                    has_xfrm: bool = False) -> etree._Element:
    """Build a p:sp element (shape) with optional placeholder and lstStyle."""
    sp = etree.Element(f'{{{P_NS}}}sp')

    # nvSpPr
    nvSpPr = etree.SubElement(sp, f'{{{P_NS}}}nvSpPr')
    cNvPr = etree.SubElement(nvSpPr, f'{{{P_NS}}}cNvPr')
    cNvPr.set('id', '1')
    cNvPr.set('name', 'Shape 1')
    cNvSpPr = etree.SubElement(nvSpPr, f'{{{P_NS}}}cNvSpPr')
    nvPr = etree.SubElement(nvSpPr, f'{{{P_NS}}}nvPr')

    if ph_type is not None or ph_idx is not None:
        ph = etree.SubElement(nvPr, f'{{{P_NS}}}ph')
        if ph_type is not None:
            ph.set('type', ph_type)
        if ph_idx is not None:
            ph.set('idx', str(ph_idx))

    # spPr
    spPr = etree.SubElement(sp, f'{{{P_NS}}}spPr')
    if has_xfrm:
        xfrm = etree.SubElement(spPr, f'{{{A_NS}}}xfrm')
        off = etree.SubElement(xfrm, f'{{{A_NS}}}off')
        off.set('x', '100000')
        off.set('y', '200000')
        ext = etree.SubElement(xfrm, f'{{{A_NS}}}ext')
        ext.set('cx', '3000000')
        ext.set('cy', '1500000')

    # txBody
    txBody = make_txBody(lst_style_props)
    sp.append(txBody)

    return sp


def make_master_txstyles(
    title_props: dict = None,
    body_props: dict = None,
    other_props: dict = None,
) -> etree._Element:
    """Build a p:txStyles element for a slide master.

    Args:
        title_props, body_props, other_props: dict of level_num → properties dict.
    """
    txStyles = etree.Element(f'{{{P_NS}}}txStyles')

    for style_name, props in [
        ('titleStyle', title_props),
        ('bodyStyle', body_props),
        ('otherStyle', other_props),
    ]:
        if props:
            style_elem = etree.SubElement(txStyles, f'{{{P_NS}}}{style_name}')
            for level_num, level_props in props.items():
                lvl_pPr = etree.SubElement(style_elem, f'{{{A_NS}}}lvl{level_num}pPr')
                if 'algn' in level_props:
                    lvl_pPr.set('algn', level_props['algn'])
                if 'rtl' in level_props:
                    lvl_pPr.set('rtl', level_props['rtl'])
                if any(k in level_props for k in ('sz', 'b', 'i', 'font')):
                    defRPr = etree.SubElement(lvl_pPr, f'{{{A_NS}}}defRPr')
                    if 'sz' in level_props:
                        defRPr.set('sz', str(level_props['sz']))
                    if 'b' in level_props:
                        defRPr.set('b', str(level_props['b']))
                    if 'i' in level_props:
                        defRPr.set('i', str(level_props['i']))
                    if 'font' in level_props:
                        latin = etree.SubElement(defRPr, f'{{{A_NS}}}latin')
                        latin.set('typeface', level_props['font'])

    return txStyles


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# MOCK OBJECT HELPERS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class MockRun:
    """Mock of a python-pptx Run object."""
    def __init__(self, r_elem: etree._Element, text: str = 'Hello'):
        self._r = r_elem
        self.text = text


class MockParagraph:
    """Mock of a python-pptx Paragraph object."""
    def __init__(self, p_elem: etree._Element, runs: list = None):
        self._p = p_elem
        self._runs = runs or []
        self.text = ''.join(r.text for r in self._runs)

    @property
    def runs(self):
        return self._runs

    @property
    def level(self):
        pPr = self._p.find(f'{{{A_NS}}}pPr')
        if pPr is not None:
            lvl = pPr.get('lvl')
            if lvl is not None:
                return int(lvl)
        return 0


class MockTextFrame:
    """Mock of a python-pptx TextFrame object."""
    def __init__(self, txBody_elem: etree._Element, paragraphs: list = None):
        self._txBody = txBody_elem
        self._paragraphs = paragraphs or []

    @property
    def paragraphs(self):
        return self._paragraphs


class MockPlaceholderFormat:
    """Mock of python-pptx PlaceholderFormat."""
    def __init__(self, idx=None, ph_type=None):
        self.idx = idx
        self.type = ph_type
        self._element = None


class MockShape:
    """Mock of a python-pptx Shape object."""
    def __init__(self, sp_elem: etree._Element, text_frame: MockTextFrame = None,
                 is_placeholder: bool = False, ph_idx: int = None,
                 ph_type=None, shape_id: int = 1, name: str = 'Shape 1',
                 left: int = 0, top: int = 0, width: int = 1000000,
                 height: int = 500000, rotation: float = 0.0):
        self._element = sp_elem
        self._text_frame = text_frame
        self.is_placeholder = is_placeholder
        self.shape_id = shape_id
        self.name = name
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.rotation = rotation
        self.has_text_frame = text_frame is not None
        self.shape_type = None

        if is_placeholder:
            self.placeholder_format = MockPlaceholderFormat(idx=ph_idx, ph_type=ph_type)
        else:
            self.placeholder_format = None

    @property
    def text_frame(self):
        return self._text_frame


class MockPlaceholderShape:
    """Mock of a placeholder shape on a layout/master."""
    def __init__(self, sp_elem):
        self._element = sp_elem


class MockLayoutOrMaster:
    """Mock of a python-pptx SlideLayout or SlideMaster with placeholders."""
    def __init__(self, placeholders: list = None, shapes: list = None,
                 name: str = 'Mock', element: etree._Element = None):
        self._placeholders = placeholders or []
        self._shapes = shapes or placeholders or []
        self.name = name
        self._element = element or etree.Element(f'{{{P_NS}}}sldLayout')

    @property
    def placeholders(self):
        return self._placeholders

    @property
    def shapes(self):
        return self._shapes

    @property
    def slide_layouts(self):
        return []


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TEST: Models
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestResolvedRun:
    """Tests for the ResolvedRun frozen dataclass."""

    def test_construction(self):
        run = ResolvedRun(
            text='Hello',
            effective_font_size_pt=18.0,
            effective_font_name='Calibri',
            effective_bold=False,
            effective_italic=False,
            effective_color=None,
            effective_underline=False,
            source_font_size_level='run',
        )
        assert run.text == 'Hello'
        assert run.effective_font_size_pt == 18.0
        assert run.effective_font_name == 'Calibri'
        assert run.effective_bold is False
        assert run.source_font_size_level == 'run'

    def test_frozen(self):
        run = ResolvedRun(
            text='Hello', effective_font_size_pt=18.0,
            effective_font_name='Calibri', effective_bold=False,
            effective_italic=False, effective_color=None,
            effective_underline=False, source_font_size_level='run',
        )
        with pytest.raises(FrozenInstanceError):
            run.text = 'World'

    def test_color_none_allowed(self):
        run = ResolvedRun(
            text='X', effective_font_size_pt=12.0,
            effective_font_name='Arial', effective_bold=True,
            effective_italic=True, effective_color=None,
            effective_underline=True, source_font_size_level='default',
        )
        assert run.effective_color is None

    def test_color_hex(self):
        run = ResolvedRun(
            text='X', effective_font_size_pt=12.0,
            effective_font_name='Arial', effective_bold=False,
            effective_italic=False, effective_color='FF0000',
            effective_underline=False, source_font_size_level='run',
        )
        assert run.effective_color == 'FF0000'


class TestResolvedParagraph:
    def test_construction(self):
        run = ResolvedRun(
            text='Test', effective_font_size_pt=14.0,
            effective_font_name='Arial', effective_bold=True,
            effective_italic=False, effective_color=None,
            effective_underline=False, source_font_size_level='paragraph',
        )
        para = ResolvedParagraph(
            runs=(run,),
            effective_alignment='ctr',
            effective_rtl=False,
            effective_level=0,
            effective_bullet_type=None,
            effective_line_spacing=1.5,
            effective_space_before=6.0,
            effective_space_after=12.0,
        )
        assert len(para.runs) == 1
        assert para.effective_alignment == 'ctr'
        assert para.effective_line_spacing == 1.5

    def test_frozen(self):
        para = ResolvedParagraph(
            runs=(), effective_alignment='l',
            effective_rtl=False, effective_level=0,
            effective_bullet_type=None,
            effective_line_spacing=None,
            effective_space_before=None,
            effective_space_after=None,
        )
        with pytest.raises(FrozenInstanceError):
            para.effective_alignment = 'r'


class TestResolvedShape:
    def test_construction(self):
        shape = ResolvedShape(
            shape_id=1, shape_name='Title 1',
            shape_type='placeholder', placeholder_type='title',
            placeholder_idx=0, x_emu=100000, y_emu=200000,
            width_emu=8000000, height_emu=1200000,
            rotation_degrees=0.0, paragraphs=(),
            is_master_inherited=False, source_level='slide',
            has_local_position_override=True, has_text=True,
        )
        assert shape.shape_id == 1
        assert shape.is_placeholder is True
        assert shape.source_level == 'slide'

    def test_full_text(self):
        r1 = ResolvedRun(
            text='Hello ', effective_font_size_pt=18.0,
            effective_font_name='Calibri', effective_bold=False,
            effective_italic=False, effective_color=None,
            effective_underline=False, source_font_size_level='run',
        )
        r2 = ResolvedRun(
            text='World', effective_font_size_pt=18.0,
            effective_font_name='Calibri', effective_bold=False,
            effective_italic=False, effective_color=None,
            effective_underline=False, source_font_size_level='run',
        )
        p1 = ResolvedParagraph(
            runs=(r1, r2), effective_alignment='l',
            effective_rtl=False, effective_level=0,
            effective_bullet_type=None,
            effective_line_spacing=None,
            effective_space_before=None,
            effective_space_after=None,
        )
        shape = ResolvedShape(
            shape_id=1, shape_name='Text 1',
            shape_type='textbox', placeholder_type=None,
            placeholder_idx=None, x_emu=0, y_emu=0,
            width_emu=1000000, height_emu=500000,
            rotation_degrees=0.0, paragraphs=(p1,),
            is_master_inherited=False, source_level='slide',
            has_local_position_override=False, has_text=True,
        )
        assert shape.full_text == 'Hello World'


class TestTransformPlan:
    def test_add_actions(self):
        plan = TransformPlan()
        action = TransformAction(shape_id=1, action_type='mirror', params={'slide_width_emu': 12192000})
        plan.add_slide_action(1, action)
        assert plan.total_actions == 1
        assert len(plan.slide_actions[1]) == 1

    def test_invalid_action_type(self):
        with pytest.raises(ValueError, match="Invalid action_type"):
            TransformAction(shape_id=1, action_type='invalid_action')


class TestValidationReport:
    def test_counts(self):
        issues = (
            ValidationIssue(severity='error', slide_number=1, shape_id=1,
                          issue_type='rtl_missing', message='RTL not set'),
            ValidationIssue(severity='warning', slide_number=1, shape_id=2,
                          issue_type='font_small', message='Font too small'),
            ValidationIssue(severity='info', slide_number=2, shape_id=None,
                          issue_type='note', message='Info only'),
        )
        report = ValidationReport(issues=issues, total_shapes_checked=10, total_slides_checked=2)
        assert report.error_count == 1
        assert report.warning_count == 1
        assert report.info_count == 1
        assert report.has_errors is True
        assert report.passed is False

    def test_passed_with_warnings(self):
        issues = (
            ValidationIssue(severity='warning', slide_number=1, shape_id=1,
                          issue_type='minor', message='Minor issue'),
        )
        report = ValidationReport(issues=issues)
        assert report.passed is True  # Warnings don't fail


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TEST: XML Helpers
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestXmlHelpers:
    def test_qn(self):
        assert _qn('a:rPr') == f'{{{A_NS}}}rPr'
        assert _qn('p:sp') == f'{{{P_NS}}}sp'

    def test_get_attr_int(self):
        elem = etree.Element('test')
        elem.set('sz', '1800')
        assert _get_attr_int(elem, 'sz') == 1800
        assert _get_attr_int(elem, 'missing') is None
        assert _get_attr_int(None, 'sz') is None

    def test_get_attr_bool(self):
        elem = etree.Element('test')
        elem.set('b', '1')
        elem.set('i', 'true')
        elem.set('u', '0')
        elem.set('s', 'false')
        assert _get_attr_bool(elem, 'b') is True
        assert _get_attr_bool(elem, 'i') is True
        assert _get_attr_bool(elem, 'u') is False
        assert _get_attr_bool(elem, 's') is False
        assert _get_attr_bool(elem, 'missing') is None
        assert _get_attr_bool(None, 'b') is None


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TEST: Font Size Resolution Chain
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestFontSizeResolution:
    """Test the 7-level font size inheritance chain."""

    def setup_method(self):
        """Create a basic PropertyResolver with a mock presentation."""
        self.prs = MagicMock()
        self.prs.slide_width = 12192000
        self.prs.slide_height = 6858000
        self.prs.slide_masters = []
        self.prs.slides = []
        self.resolver = PropertyResolver(self.prs)

    def test_level1_run_explicit(self):
        """Level 1: Font size set directly on the run's rPr."""
        r = make_run('Hello', sz=2400)
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        size, source = self.resolver.resolve_font_size(
            run, para, shape, tf, None, None, None, 0
        )
        assert size == 24.0
        assert source == 'run'

    def test_level2_paragraph_defRPr(self):
        """Level 2: Font size from paragraph's defRPr."""
        r = make_run('Hello')  # No sz on run
        run = MockRun(r, 'Hello')
        p = make_paragraph(defRPr_sz=2000)
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        size, source = self.resolver.resolve_font_size(
            run, para, shape, tf, None, None, None, 0
        )
        assert size == 20.0
        assert source == 'paragraph'

    def test_level3_textframe_lstStyle(self):
        """Level 3: Font size from text frame's lstStyle at the correct level."""
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph()  # No defRPr
        para = MockParagraph(p)
        txBody = make_txBody(lst_style_props={1: {'sz': '1600'}})
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        size, source = self.resolver.resolve_font_size(
            run, para, shape, tf, None, None, None, 0
        )
        assert size == 16.0
        assert source == 'textframe'

    def test_level3_paragraph_level_matching(self):
        """Level 3: lstStyle should match paragraph level (lvl2pPr for level 1)."""
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph(lvl=1)  # Paragraph level 1
        para = MockParagraph(p)
        txBody = make_txBody(lst_style_props={
            1: {'sz': '2400'},  # For level 0 (lvl1pPr)
            2: {'sz': '2000'},  # For level 1 (lvl2pPr)
        })
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        size, source = self.resolver.resolve_font_size(
            run, para, shape, tf, None, None, None, 1  # level=1
        )
        assert size == 20.0
        assert source == 'textframe'

    def test_level5_layout_placeholder(self):
        """Level 5: Font size from layout placeholder's lstStyle."""
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()  # Empty lstStyle
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        # Layout placeholder with lstStyle
        layout_sp = make_sp_element(
            ph_type='body', ph_idx=1,
            lst_style_props={1: {'sz': '1400'}}
        )

        size, source = self.resolver.resolve_font_size(
            run, para, shape, tf, layout_sp, None, None, 0
        )
        assert size == 14.0
        assert source == 'layout'

    def test_level6_master_placeholder(self):
        """Level 6: Font size from master placeholder's lstStyle."""
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        # Master placeholder with lstStyle
        master_sp = make_sp_element(
            ph_type='body', ph_idx=1,
            lst_style_props={1: {'sz': '2800'}}
        )

        size, source = self.resolver.resolve_font_size(
            run, para, shape, tf, None, master_sp, None, 0
        )
        assert size == 28.0
        assert source == 'master'

    def test_level7_master_txstyles(self):
        """Level 7: Font size from master's txStyles (bodyStyle)."""
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        # Master txStyles
        txStyles = make_master_txstyles(
            body_props={1: {'sz': '3200'}}
        )
        body_style = txStyles.find(f'{{{P_NS}}}bodyStyle')

        size, source = self.resolver.resolve_font_size(
            run, para, shape, tf, None, None, body_style, 0
        )
        assert size == 32.0
        assert source == 'txstyles'

    def test_fallback_default(self):
        """Chain exhausted — should return the default (18pt)."""
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        size, source = self.resolver.resolve_font_size(
            run, para, shape, tf, None, None, None, 0
        )
        assert size == DEFAULT_FONT_SIZE_PT
        assert source == 'default'

    def test_priority_run_over_paragraph(self):
        """Level 1 (run) should beat level 2 (paragraph defRPr)."""
        r = make_run('Hello', sz=3600)
        run = MockRun(r, 'Hello')
        p = make_paragraph(defRPr_sz=2000)
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        size, source = self.resolver.resolve_font_size(
            run, para, shape, tf, None, None, None, 0
        )
        assert size == 36.0
        assert source == 'run'

    def test_priority_paragraph_over_lstStyle(self):
        """Level 2 (paragraph defRPr) should beat level 3 (lstStyle)."""
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph(defRPr_sz=2200)
        para = MockParagraph(p)
        txBody = make_txBody(lst_style_props={1: {'sz': '1400'}})
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        size, source = self.resolver.resolve_font_size(
            run, para, shape, tf, None, None, None, 0
        )
        assert size == 22.0
        assert source == 'paragraph'

    def test_priority_layout_over_master(self):
        """Level 5 (layout) should beat level 6 (master)."""
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        layout_sp = make_sp_element(lst_style_props={1: {'sz': '1600'}})
        master_sp = make_sp_element(lst_style_props={1: {'sz': '2400'}})

        size, source = self.resolver.resolve_font_size(
            run, para, shape, tf, layout_sp, master_sp, None, 0
        )
        assert size == 16.0
        assert source == 'layout'


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TEST: Font Name Resolution
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestFontNameResolution:
    def setup_method(self):
        self.prs = MagicMock()
        self.prs.slide_width = 12192000
        self.prs.slide_height = 6858000
        self.prs.slide_masters = []
        self.prs.slides = []
        self.resolver = PropertyResolver(self.prs)

    def test_run_explicit_font(self):
        """Font name set directly on run's rPr."""
        r = make_run('Hello', font_name='Arial')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        name = self.resolver.resolve_font_name(
            run, para, shape, tf, None, None, None, None, 0
        )
        assert name == 'Arial'

    def test_paragraph_defRPr_font(self):
        """Font name from paragraph's defRPr."""
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph(defRPr_font='Times New Roman')
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        name = self.resolver.resolve_font_name(
            run, para, shape, tf, None, None, None, None, 0
        )
        assert name == 'Times New Roman'

    def test_lstStyle_font(self):
        """Font name from text frame lstStyle."""
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody(lst_style_props={1: {'font': 'Segoe UI'}})
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        name = self.resolver.resolve_font_name(
            run, para, shape, tf, None, None, None, None, 0
        )
        assert name == 'Segoe UI'

    def test_theme_font_major_latin(self):
        """Theme font reference +mj-lt should resolve to major latin font."""
        r = make_run('Hello', font_name='+mj-lt')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        # Build a mock theme
        theme_elem = etree.Element(f'{{{A_NS}}}theme')
        themeElements = etree.SubElement(theme_elem, f'{{{A_NS}}}themeElements')
        fontScheme = etree.SubElement(themeElements, f'{{{A_NS}}}fontScheme')
        majorFont = etree.SubElement(fontScheme, f'{{{A_NS}}}majorFont')
        latin = etree.SubElement(majorFont, f'{{{A_NS}}}latin')
        latin.set('typeface', 'Cambria')
        minorFont = etree.SubElement(fontScheme, f'{{{A_NS}}}minorFont')
        latin2 = etree.SubElement(minorFont, f'{{{A_NS}}}latin')
        latin2.set('typeface', 'Calibri')

        master_obj = MagicMock()
        theme_mock = MagicMock()
        theme_mock._element = theme_elem
        master_obj.theme = theme_mock

        name = self.resolver.resolve_font_name(
            run, para, shape, tf, None, None, None, master_obj, 0
        )
        assert name == 'Cambria'

    def test_theme_font_minor_latin(self):
        """Theme font reference +mn-lt should resolve to minor latin font."""
        r = make_run('Hello', font_name='+mn-lt')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        theme_elem = etree.Element(f'{{{A_NS}}}theme')
        themeElements = etree.SubElement(theme_elem, f'{{{A_NS}}}themeElements')
        fontScheme = etree.SubElement(themeElements, f'{{{A_NS}}}fontScheme')
        majorFont = etree.SubElement(fontScheme, f'{{{A_NS}}}majorFont')
        latin = etree.SubElement(majorFont, f'{{{A_NS}}}latin')
        latin.set('typeface', 'Cambria')
        minorFont = etree.SubElement(fontScheme, f'{{{A_NS}}}minorFont')
        latin2 = etree.SubElement(minorFont, f'{{{A_NS}}}latin')
        latin2.set('typeface', 'Calibri Body')

        master_obj = MagicMock()
        theme_mock = MagicMock()
        theme_mock._element = theme_elem
        master_obj.theme = theme_mock

        name = self.resolver.resolve_font_name(
            run, para, shape, tf, None, None, None, master_obj, 0
        )
        assert name == 'Calibri Body'

    def test_fallback_default_font(self):
        """Chain exhausted — should return default 'Calibri'."""
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        name = self.resolver.resolve_font_name(
            run, para, shape, tf, None, None, None, None, 0
        )
        assert name == DEFAULT_FONT_NAME


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TEST: Alignment Resolution
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestAlignmentResolution:
    def setup_method(self):
        self.prs = MagicMock()
        self.prs.slide_width = 12192000
        self.prs.slide_height = 6858000
        self.prs.slide_masters = []
        self.prs.slides = []
        self.resolver = PropertyResolver(self.prs)

    def test_explicit_paragraph_alignment(self):
        """Alignment set directly on paragraph pPr."""
        p = make_paragraph(algn='ctr')
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        algn = self.resolver.resolve_alignment(para, shape, None, None, None, 0)
        assert algn == 'ctr'

    def test_lstStyle_alignment(self):
        """Alignment from text frame lstStyle."""
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody(lst_style_props={1: {'algn': 'r'}})
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        algn = self.resolver.resolve_alignment(para, shape, None, None, None, 0)
        assert algn == 'r'

    def test_layout_placeholder_alignment(self):
        """Alignment from layout placeholder lstStyle."""
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        layout_sp = make_sp_element(lst_style_props={1: {'algn': 'just'}})

        algn = self.resolver.resolve_alignment(para, shape, layout_sp, None, None, 0)
        assert algn == 'just'

    def test_master_txstyles_alignment(self):
        """Alignment from master txStyles."""
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        txStyles = make_master_txstyles(title_props={1: {'algn': 'ctr'}})
        title_style = txStyles.find(f'{{{P_NS}}}titleStyle')

        algn = self.resolver.resolve_alignment(para, shape, None, None, title_style, 0)
        assert algn == 'ctr'

    def test_fallback_left(self):
        """Chain exhausted — should return default 'l'."""
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        algn = self.resolver.resolve_alignment(para, shape, None, None, None, 0)
        assert algn == DEFAULT_ALIGNMENT


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TEST: Bold / Italic Resolution
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestBoldItalicResolution:
    def setup_method(self):
        self.prs = MagicMock()
        self.prs.slide_width = 12192000
        self.prs.slide_height = 6858000
        self.prs.slide_masters = []
        self.prs.slides = []
        self.resolver = PropertyResolver(self.prs)

    def test_run_explicit_bold(self):
        r = make_run('Hello', b='1')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        bold = self.resolver.resolve_bold(
            run, para, shape, tf, None, None, None, 0
        )
        assert bold is True

    def test_run_explicit_not_bold(self):
        r = make_run('Hello', b='0')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        bold = self.resolver.resolve_bold(
            run, para, shape, tf, None, None, None, 0
        )
        assert bold is False

    def test_paragraph_defRPr_bold(self):
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph(defRPr_b='1')
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        bold = self.resolver.resolve_bold(
            run, para, shape, tf, None, None, None, 0
        )
        assert bold is True

    def test_lstStyle_bold(self):
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody(lst_style_props={1: {'b': '1'}})
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        bold = self.resolver.resolve_bold(
            run, para, shape, tf, None, None, None, 0
        )
        assert bold is True

    def test_txstyles_bold(self):
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        txStyles = make_master_txstyles(title_props={1: {'b': '1'}})
        title_style = txStyles.find(f'{{{P_NS}}}titleStyle')

        bold = self.resolver.resolve_bold(
            run, para, shape, tf, None, None, title_style, 0
        )
        assert bold is True

    def test_fallback_bold_false(self):
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        bold = self.resolver.resolve_bold(
            run, para, shape, tf, None, None, None, 0
        )
        assert bold is DEFAULT_BOLD

    def test_run_explicit_italic(self):
        r = make_run('Hello', i='1')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody()
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        italic = self.resolver.resolve_italic(
            run, para, shape, tf, None, None, None, 0
        )
        assert italic is True

    def test_lstStyle_italic(self):
        r = make_run('Hello')
        run = MockRun(r, 'Hello')
        p = make_paragraph()
        para = MockParagraph(p)
        txBody = make_txBody(lst_style_props={1: {'i': '1'}})
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        italic = self.resolver.resolve_italic(
            run, para, shape, tf, None, None, None, 0
        )
        assert italic is True


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TEST: RTL Resolution
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestRTLResolution:
    def setup_method(self):
        self.prs = MagicMock()
        self.prs.slide_width = 12192000
        self.prs.slide_height = 6858000
        self.prs.slide_masters = []
        self.prs.slides = []
        self.resolver = PropertyResolver(self.prs)

    def test_explicit_rtl(self):
        p = make_paragraph(rtl='1')
        rtl = self.resolver._resolve_rtl(p, None, None, None, 0)
        assert rtl is True

    def test_explicit_ltr(self):
        p = make_paragraph(rtl='0')
        rtl = self.resolver._resolve_rtl(p, None, None, None, 0)
        assert rtl is False

    def test_layout_lstStyle_rtl(self):
        p = make_paragraph()  # No explicit RTL
        layout_sp = make_sp_element(lst_style_props={1: {'rtl': '1'}})
        rtl = self.resolver._resolve_rtl(p, layout_sp, None, None, 0)
        assert rtl is True

    def test_master_txstyles_rtl(self):
        p = make_paragraph()
        txStyles = make_master_txstyles(body_props={1: {'rtl': '1'}})
        body_style = txStyles.find(f'{{{P_NS}}}bodyStyle')
        rtl = self.resolver._resolve_rtl(p, None, None, body_style, 0)
        assert rtl is True

    def test_fallback_rtl_false(self):
        p = make_paragraph()
        rtl = self.resolver._resolve_rtl(p, None, None, None, 0)
        assert rtl is DEFAULT_RTL


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TEST: Placeholder Matching
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestPlaceholderMatching:
    def setup_method(self):
        self.prs = MagicMock()
        self.prs.slide_width = 12192000
        self.prs.slide_height = 6858000
        self.prs.slide_masters = []
        self.prs.slides = []
        self.resolver = PropertyResolver(self.prs)

    def test_match_by_idx(self):
        """Placeholder should match by idx first."""
        # Source shape with idx=1
        sp = make_sp_element(ph_type='body', ph_idx=1)
        shape = MockShape(
            sp, is_placeholder=True, ph_idx=1,
            ph_type=MagicMock(return_value='body')
        )
        shape.placeholder_format = MockPlaceholderFormat(idx=1, ph_type=None)

        # Layout has a placeholder with idx=1
        layout_ph_sp = make_sp_element(ph_type='body', ph_idx=1, lst_style_props={1: {'sz': '1400'}})
        layout_ph = MockPlaceholderShape(layout_ph_sp)

        # Layout also has a placeholder with idx=2
        layout_ph_sp2 = make_sp_element(ph_type='body', ph_idx=2, lst_style_props={1: {'sz': '1200'}})
        layout_ph2 = MockPlaceholderShape(layout_ph_sp2)

        mock_layout = MockLayoutOrMaster(placeholders=[layout_ph, layout_ph2])

        result = self.resolver._find_matching_placeholder(shape, mock_layout)
        assert result is not None
        # Verify it matched idx=1, not idx=2
        ph_el = result.find(f'.//{{{P_NS}}}ph')
        assert ph_el is not None
        assert ph_el.get('idx') == '1'

    def test_match_by_type_fallback(self):
        """When idx doesn't match, fall back to type matching."""
        sp = make_sp_element(ph_type='title', ph_idx=99)
        shape = MockShape(sp, is_placeholder=True, ph_idx=99)
        shape.placeholder_format = MockPlaceholderFormat(idx=99, ph_type=None)

        # Patch _get_placeholder_type_str to return 'title'
        original_method = self.resolver._get_placeholder_type_str
        self.resolver._get_placeholder_type_str = lambda pf: 'title'

        # Layout has no idx=99 but has type='title'
        layout_ph_sp = make_sp_element(ph_type='title', ph_idx=0)
        layout_ph = MockPlaceholderShape(layout_ph_sp)
        mock_layout = MockLayoutOrMaster(placeholders=[layout_ph])

        result = self.resolver._find_matching_placeholder(shape, mock_layout)
        assert result is not None

        # Restore
        self.resolver._get_placeholder_type_str = original_method

    def test_no_match(self):
        """When no placeholder matches, return None."""
        sp = make_sp_element(ph_type='chart', ph_idx=50)
        shape = MockShape(sp, is_placeholder=True, ph_idx=50)
        shape.placeholder_format = MockPlaceholderFormat(idx=50, ph_type=None)

        # Layout has no chart placeholder
        layout_ph_sp = make_sp_element(ph_type='body', ph_idx=1)
        layout_ph = MockPlaceholderShape(layout_ph_sp)
        mock_layout = MockLayoutOrMaster(placeholders=[layout_ph])

        result = self.resolver._find_matching_placeholder(shape, mock_layout)
        assert result is None


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TEST: Master txStyles Selection
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestMasterTxStyles:
    def setup_method(self):
        self.prs = MagicMock()
        self.prs.slide_width = 12192000
        self.prs.slide_height = 6858000
        self.prs.slide_masters = []
        self.prs.slides = []
        self.resolver = PropertyResolver(self.prs)

    def _make_master_with_txstyles(self):
        """Create a mock master with p:txStyles containing all three styles."""
        master_elem = etree.Element(f'{{{P_NS}}}sldMaster')
        txStyles = make_master_txstyles(
            title_props={1: {'sz': '4400', 'b': '1'}},
            body_props={1: {'sz': '2400', 'b': '0'}},
            other_props={1: {'sz': '1800', 'b': '0'}},
        )
        master_elem.append(txStyles)

        master = MagicMock()
        master._element = master_elem
        return master

    def test_title_placeholder_gets_titleStyle(self):
        """Title placeholder should use p:titleStyle."""
        master = self._make_master_with_txstyles()
        from pptx.enum.shapes import PP_PLACEHOLDER
        shape = MagicMock()
        shape.is_placeholder = True
        shape.placeholder_format = MockPlaceholderFormat(
            idx=0, ph_type=PP_PLACEHOLDER.TITLE
        )

        result = self.resolver._get_master_txstyle_for_shape(shape, master)
        assert result is not None
        assert result.tag == f'{{{P_NS}}}titleStyle'

    def test_body_placeholder_gets_bodyStyle(self):
        """Body placeholder should use p:bodyStyle."""
        master = self._make_master_with_txstyles()
        from pptx.enum.shapes import PP_PLACEHOLDER
        shape = MagicMock()
        shape.is_placeholder = True
        shape.placeholder_format = MockPlaceholderFormat(
            idx=1, ph_type=PP_PLACEHOLDER.BODY
        )

        result = self.resolver._get_master_txstyle_for_shape(shape, master)
        assert result is not None
        assert result.tag == f'{{{P_NS}}}bodyStyle'

    def test_non_placeholder_gets_otherStyle(self):
        """Non-placeholder shapes should use p:otherStyle."""
        master = self._make_master_with_txstyles()
        shape = MagicMock()
        shape.is_placeholder = False

        result = self.resolver._get_master_txstyle_for_shape(shape, master)
        assert result is not None
        assert result.tag == f'{{{P_NS}}}otherStyle'


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TEST: Edge Cases
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestEdgeCases:
    def setup_method(self):
        self.prs = MagicMock()
        self.prs.slide_width = 12192000
        self.prs.slide_height = 6858000
        self.prs.slide_masters = []
        self.prs.slides = []
        self.resolver = PropertyResolver(self.prs)

    def test_empty_paragraph_no_runs(self):
        """Paragraph with no runs should produce empty runs tuple."""
        p_elem = make_paragraph()
        para = MockParagraph(p_elem, runs=[])
        para.text = ''

        sp = make_sp_element()
        txBody = make_txBody()
        tf = MockTextFrame(txBody, paragraphs=[para])
        shape = MockShape(sp, text_frame=tf)

        resolved = self.resolver._resolve_paragraph(
            para, shape, tf, None, None, None, None
        )
        assert len(resolved.runs) == 0

    def test_has_local_position_override(self):
        """Shape with explicit xfrm should report local position override."""
        sp = make_sp_element(has_xfrm=True)
        result = self.resolver._has_local_position_override(sp)
        assert result is True

    def test_no_local_position_override(self):
        """Shape without explicit xfrm should not report override."""
        sp = make_sp_element(has_xfrm=False)
        result = self.resolver._has_local_position_override(sp)
        assert result is False

    def test_paragraph_level_clamped(self):
        """Paragraph level should be clamped to 0-8."""
        p = etree.Element(f'{{{A_NS}}}p')
        pPr = etree.SubElement(p, f'{{{A_NS}}}pPr')
        pPr.set('lvl', '15')  # Way too high
        level = self.resolver._get_paragraph_level(p)
        assert level == 8  # Clamped to max

    def test_paragraph_negative_level(self):
        """Negative level should be clamped to 0."""
        p = etree.Element(f'{{{A_NS}}}p')
        pPr = etree.SubElement(p, f'{{{A_NS}}}pPr')
        pPr.set('lvl', '-1')
        level = self.resolver._get_paragraph_level(p)
        assert level == 0

    def test_bullet_type_char(self):
        """Bullet character detection."""
        p = etree.Element(f'{{{A_NS}}}p')
        pPr = etree.SubElement(p, f'{{{A_NS}}}pPr')
        buChar = etree.SubElement(pPr, f'{{{A_NS}}}buChar')
        buChar.set('char', '→')

        bullet = self.resolver._resolve_bullet_type(p)
        assert bullet == 'char:→'

    def test_bullet_type_auto_num(self):
        """Auto-numbered bullet detection."""
        p = etree.Element(f'{{{A_NS}}}p')
        pPr = etree.SubElement(p, f'{{{A_NS}}}pPr')
        buAutoNum = etree.SubElement(pPr, f'{{{A_NS}}}buAutoNum')
        buAutoNum.set('type', 'arabicPeriod')

        bullet = self.resolver._resolve_bullet_type(p)
        assert bullet == 'auto:arabicPeriod'

    def test_bullet_none(self):
        """Explicit buNone should return None."""
        p = etree.Element(f'{{{A_NS}}}p')
        pPr = etree.SubElement(p, f'{{{A_NS}}}pPr')
        etree.SubElement(pPr, f'{{{A_NS}}}buNone')

        bullet = self.resolver._resolve_bullet_type(p)
        assert bullet is None

    def test_line_spacing_percentage(self):
        """Line spacing as percentage (150000 = 1.5x)."""
        p = etree.Element(f'{{{A_NS}}}p')
        pPr = etree.SubElement(p, f'{{{A_NS}}}pPr')
        lnSpc = etree.SubElement(pPr, f'{{{A_NS}}}lnSpc')
        spcPct = etree.SubElement(lnSpc, f'{{{A_NS}}}spcPct')
        spcPct.set('val', '150000')

        spacing = self.resolver._resolve_line_spacing(p)
        assert spacing == 1.5

    def test_line_spacing_points(self):
        """Line spacing as points (1200 = 12pt)."""
        p = etree.Element(f'{{{A_NS}}}p')
        pPr = etree.SubElement(p, f'{{{A_NS}}}pPr')
        lnSpc = etree.SubElement(pPr, f'{{{A_NS}}}lnSpc')
        spcPts = etree.SubElement(lnSpc, f'{{{A_NS}}}spcPts')
        spcPts.set('val', '1200')

        spacing = self.resolver._resolve_line_spacing(p)
        assert spacing == 12.0

    def test_space_before(self):
        """Space before in hundredths of a point."""
        p = etree.Element(f'{{{A_NS}}}p')
        pPr = etree.SubElement(p, f'{{{A_NS}}}pPr')
        spcBef = etree.SubElement(pPr, f'{{{A_NS}}}spcBef')
        spcPts = etree.SubElement(spcBef, f'{{{A_NS}}}spcPts')
        spcPts.set('val', '600')

        space = self.resolver._resolve_space_before(p)
        assert space == 6.0

    def test_underline_single(self):
        """Underline set to 'sng' should be True."""
        r = etree.Element(f'{{{A_NS}}}r')
        rPr = etree.SubElement(r, f'{{{A_NS}}}rPr')
        rPr.set('u', 'sng')
        run = MockRun(r)

        result = self.resolver._resolve_underline(run)
        assert result is True

    def test_underline_none(self):
        """Underline set to 'none' should be False."""
        r = etree.Element(f'{{{A_NS}}}r')
        rPr = etree.SubElement(r, f'{{{A_NS}}}rPr')
        rPr.set('u', 'none')
        run = MockRun(r)

        result = self.resolver._resolve_underline(run)
        assert result is False

    def test_color_srgb(self):
        """Color from solidFill/srgbClr."""
        r = etree.Element(f'{{{A_NS}}}r')
        rPr = etree.SubElement(r, f'{{{A_NS}}}rPr')
        solidFill = etree.SubElement(rPr, f'{{{A_NS}}}solidFill')
        srgb = etree.SubElement(solidFill, f'{{{A_NS}}}srgbClr')
        srgb.set('val', 'FF5500')
        run = MockRun(r)

        color = self.resolver._resolve_color(run)
        assert color == 'FF5500'


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# TEST: Full Chain Integration
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class TestFullChainIntegration:
    """Test that all levels interact correctly in a realistic scenario."""

    def setup_method(self):
        self.prs = MagicMock()
        self.prs.slide_width = 12192000
        self.prs.slide_height = 6858000
        self.prs.slide_masters = []
        self.prs.slides = []
        self.resolver = PropertyResolver(self.prs)

    def test_full_chain_font_size_each_level_wins(self):
        """Verify that removing a value at each level causes the next level to win."""
        # Build the full chain
        layout_sp = make_sp_element(lst_style_props={1: {'sz': '1400'}})
        master_sp = make_sp_element(lst_style_props={1: {'sz': '2800'}})
        txStyles = make_master_txstyles(body_props={1: {'sz': '3200'}})
        body_style = txStyles.find(f'{{{P_NS}}}bodyStyle')

        # Test: All levels present — run wins
        r = make_run('Hi', sz=2400)
        run = MockRun(r, 'Hi')
        p = make_paragraph(defRPr_sz=2000)
        para = MockParagraph(p)
        txBody = make_txBody(lst_style_props={1: {'sz': '1600'}})
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        size, src = self.resolver.resolve_font_size(
            run, para, shape, tf, layout_sp, master_sp, body_style, 0
        )
        assert size == 24.0
        assert src == 'run'

        # Remove run-level → paragraph wins
        r2 = make_run('Hi')
        run2 = MockRun(r2, 'Hi')
        size, src = self.resolver.resolve_font_size(
            run2, para, shape, tf, layout_sp, master_sp, body_style, 0
        )
        assert size == 20.0
        assert src == 'paragraph'

        # Remove paragraph-level → textframe wins
        p2 = make_paragraph()
        para2 = MockParagraph(p2)
        size, src = self.resolver.resolve_font_size(
            run2, para2, shape, tf, layout_sp, master_sp, body_style, 0
        )
        assert size == 16.0
        assert src == 'textframe'

        # Remove textframe → layout wins
        txBody_empty = make_txBody()
        tf_empty = MockTextFrame(txBody_empty)
        sp_empty = make_sp_element()
        shape_empty = MockShape(sp_empty, text_frame=tf_empty)
        size, src = self.resolver.resolve_font_size(
            run2, para2, shape_empty, tf_empty, layout_sp, master_sp, body_style, 0
        )
        assert size == 14.0
        assert src == 'layout'

        # Remove layout → master wins
        size, src = self.resolver.resolve_font_size(
            run2, para2, shape_empty, tf_empty, None, master_sp, body_style, 0
        )
        assert size == 28.0
        assert src == 'master'

        # Remove master → txstyles wins
        size, src = self.resolver.resolve_font_size(
            run2, para2, shape_empty, tf_empty, None, None, body_style, 0
        )
        assert size == 32.0
        assert src == 'txstyles'

        # Remove txstyles → default wins
        size, src = self.resolver.resolve_font_size(
            run2, para2, shape_empty, tf_empty, None, None, None, 0
        )
        assert size == DEFAULT_FONT_SIZE_PT
        assert src == 'default'

    def test_multi_level_paragraph_lststyle(self):
        """Different paragraph levels pick up different lstStyle entries."""
        txBody = make_txBody(lst_style_props={
            1: {'sz': '2400', 'b': '1'},
            2: {'sz': '2000', 'b': '0'},
            3: {'sz': '1800', 'b': '0'},
        })
        tf = MockTextFrame(txBody)
        sp = make_sp_element()
        shape = MockShape(sp, text_frame=tf)

        for level, expected_sz in [(0, 24.0), (1, 20.0), (2, 18.0)]:
            r = make_run('Hi')
            run = MockRun(r, 'Hi')
            p = make_paragraph(lvl=level)
            para = MockParagraph(p)

            size, src = self.resolver.resolve_font_size(
                run, para, shape, tf, None, None, None, level
            )
            assert size == expected_sz, f"Level {level}: expected {expected_sz}, got {size}"
            assert src == 'textframe'

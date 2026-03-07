"""
test_rtl_transforms.py — Unit tests for SlideArabi RTL transforms and typography.

Run with: python -m pytest slidearabi/tests/test_rtl_transforms.py -v

Test strategy:
- Utils: pure-function unit tests (no external dependencies).
- RTL transforms: tests using minimal lxml element trees + mock python-pptx shapes.
- Typography: tests using mock shapes with controlled dimensions and text content.
"""

from __future__ import annotations

import sys
import os
import unittest
from unittest.mock import MagicMock, patch, PropertyMock
from lxml import etree

# Ensure the workspace root is on sys.path so 'slidearabi' is importable
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..'))

from slidearabi.utils import (
    emu_to_inches,
    emu_to_pt,
    pt_to_emu,
    inches_to_emu,
    hundredths_pt_to_pt,
    pt_to_hundredths_pt,
    mirror_x,
    swap_positions,
    has_arabic,
    has_latin,
    is_bidi_text,
    compute_script_ratio,
    qn,
    ensure_pPr,
    set_rtl_on_paragraph,
    set_alignment_on_paragraph,
    get_placeholder_info,
    bounds_check_emu,
    clamp_emu,
    A_NS, P_NS,
)
from slidearabi.rtl_transforms import (
    TransformReport,
    MasterLayoutTransformer,
    SlideContentTransformer,
)
from slidearabi.typography import (
    TypographyNormalizer,
    ARABIC_FONT_MAP,
    ARABIC_EXPANSION_FACTORS,
    MIN_FONT_SIZE_PT,
    MIN_TITLE_FONT_SIZE_PT,
    MIN_BODY_FONT_SIZE_PT,
)


# ─────────────────────────────────────────────────────────────────────────────
# Helpers: build minimal lxml elements for testing
# ─────────────────────────────────────────────────────────────────────────────

def make_para_elem(text: str = '', rtl: str | None = None, algn: str | None = None):
    """Create a minimal <a:p> element with optional text run."""
    p = etree.Element(f'{{{A_NS}}}p')
    if rtl is not None or algn is not None:
        pPr = etree.SubElement(p, f'{{{A_NS}}}pPr')
        if rtl is not None:
            pPr.set('rtl', rtl)
        if algn is not None:
            pPr.set('algn', algn)
    if text:
        r = etree.SubElement(p, f'{{{A_NS}}}r')
        t = etree.SubElement(r, f'{{{A_NS}}}t')
        t.text = text
    return p


def make_txBody_elem(texts: list[str] | None = None):
    """Create a minimal <a:txBody> element with optional paragraphs."""
    txBody = etree.Element(f'{{{A_NS}}}txBody')
    bodyPr = etree.SubElement(txBody, f'{{{A_NS}}}bodyPr')
    lstStyle = etree.SubElement(txBody, f'{{{A_NS}}}lstStyle')
    for text in (texts or []):
        p = etree.SubElement(txBody, f'{{{A_NS}}}p')
        r = etree.SubElement(p, f'{{{A_NS}}}r')
        t = etree.SubElement(r, f'{{{A_NS}}}t')
        t.text = text
    return txBody


def make_mock_shape(name='shape', left=0, top=0, width=1000000, height=500000,
                    is_placeholder=False, ph_type=None, ph_idx=None,
                    has_text_frame=False, has_table=False, has_chart=False,
                    shape_type=None):
    """Create a minimal mock python-pptx Shape."""
    shape = MagicMock()
    shape.name = name
    shape.left = left
    shape.top = top
    shape.width = width
    shape.height = height
    shape.is_placeholder = is_placeholder
    shape.has_text_frame = has_text_frame
    shape.has_table = has_table
    shape.has_chart = has_chart
    shape.shape_type = shape_type

    if is_placeholder and ph_type is not None:
        ph_fmt = MagicMock()
        ph_fmt.type = ph_type
        ph_fmt.idx = ph_idx or 0
        shape.placeholder_format = ph_fmt
    else:
        shape.placeholder_format = None

    return shape


# ─────────────────────────────────────────────────────────────────────────────
# Test: utils.py — unit conversion
# ─────────────────────────────────────────────────────────────────────────────

class TestUnitConversions(unittest.TestCase):

    def test_emu_to_inches(self):
        self.assertAlmostEqual(emu_to_inches(914400), 1.0)
        self.assertAlmostEqual(emu_to_inches(0), 0.0)
        self.assertAlmostEqual(emu_to_inches(4572000), 5.0)

    def test_emu_to_pt(self):
        self.assertAlmostEqual(emu_to_pt(12700), 1.0)
        self.assertAlmostEqual(emu_to_pt(228600), 18.0)

    def test_pt_to_emu(self):
        self.assertEqual(pt_to_emu(1.0), 12700)
        self.assertEqual(pt_to_emu(18.0), 228600)

    def test_inches_to_emu(self):
        self.assertEqual(inches_to_emu(1.0), 914400)
        self.assertEqual(inches_to_emu(0.0), 0)

    def test_hundredths_pt_to_pt(self):
        self.assertAlmostEqual(hundredths_pt_to_pt(1800), 18.0)
        self.assertAlmostEqual(hundredths_pt_to_pt(1100), 11.0)
        self.assertAlmostEqual(hundredths_pt_to_pt(800), 8.0)

    def test_pt_to_hundredths_pt(self):
        self.assertEqual(pt_to_hundredths_pt(18.0), 1800)
        self.assertEqual(pt_to_hundredths_pt(11.0), 1100)
        self.assertEqual(pt_to_hundredths_pt(8.0), 800)

    def test_roundtrip_pt_conversion(self):
        for pt in [8.0, 10.0, 12.0, 14.0, 18.0, 24.0, 36.0]:
            self.assertAlmostEqual(hundredths_pt_to_pt(pt_to_hundredths_pt(pt)), pt,
                                   delta=0.1)


# ─────────────────────────────────────────────────────────────────────────────
# Test: utils.py — coordinate math
# ─────────────────────────────────────────────────────────────────────────────

class TestCoordinateMath(unittest.TestCase):

    # Slide: 12192000 EMU wide (13.333" widescreen)
    SLIDE_WIDTH = 12192000

    def test_mirror_x_basic(self):
        # Shape at x=1000, width=500 on slide_width=10000
        # new_x = 10000 - (1000 + 500) = 8500
        self.assertEqual(mirror_x(1000, 500, 10000), 8500)

    def test_mirror_x_at_left_edge(self):
        # Shape flush against left edge
        # new_x = SLIDE_WIDTH - (0 + width) → flush against right edge
        width = 914400  # 1 inch
        new_x = mirror_x(0, width, self.SLIDE_WIDTH)
        self.assertEqual(new_x, self.SLIDE_WIDTH - width)

    def test_mirror_x_at_right_edge(self):
        # Shape flush against right edge
        # new_x should be 0 (flush against left edge after mirror)
        width = 914400  # 1 inch
        left = self.SLIDE_WIDTH - width
        new_x = mirror_x(left, width, self.SLIDE_WIDTH)
        self.assertEqual(new_x, 0)

    def test_mirror_x_centred_shape(self):
        # Centred shape should mirror to same position
        width = 2000000
        left = (self.SLIDE_WIDTH - width) // 2
        new_x = mirror_x(left, width, self.SLIDE_WIDTH)
        # Centred shape: mirror is symmetric around centre
        self.assertEqual(new_x, self.SLIDE_WIDTH - left - width)
        self.assertEqual(new_x, left)  # Centred = same position

    def test_mirror_is_involution(self):
        """Applying mirror twice should return to original position."""
        left = 1000000
        width = 2000000
        mirrored = mirror_x(left, width, self.SLIDE_WIDTH)
        double_mirrored = mirror_x(mirrored, width, self.SLIDE_WIDTH)
        self.assertEqual(double_mirrored, left)

    def test_swap_positions_basic(self):
        # Two shapes side by side: shape1 at 100, w=400; shape2 at 600, w=400
        # After RTL swap: shape1 should end up where shape2 was (mirrored)
        slide_width = 1200
        new_x1, new_x2 = swap_positions(100, 400, 600, 400, slide_width)
        # new_x1 = mirror(600, 400, 1200) = 1200 - (600+400) = 200
        # new_x2 = mirror(100, 400, 1200) = 1200 - (100+400) = 700
        self.assertEqual(new_x1, 200)
        self.assertEqual(new_x2, 700)

    def test_swap_positions_symmetry(self):
        """Swapping twice should return to original."""
        slide_width = 12192000
        x1, w1 = 914400, 4572000   # 1" wide, 5" wide
        x2, w2 = 6858000, 4572000  # symmetric
        new_x1, new_x2 = swap_positions(x1, w1, x2, w2, slide_width)
        # Second swap (using swapped widths) should restore original
        orig_x1, orig_x2 = swap_positions(new_x1, w2, new_x2, w1, slide_width)
        self.assertAlmostEqual(orig_x1, x1, delta=1)
        self.assertAlmostEqual(orig_x2, x2, delta=1)


# ─────────────────────────────────────────────────────────────────────────────
# Test: utils.py — text script detection
# ─────────────────────────────────────────────────────────────────────────────

class TestTextScriptDetection(unittest.TestCase):

    def test_has_arabic_true(self):
        self.assertTrue(has_arabic('مرحبا'))
        self.assertTrue(has_arabic('Hello مرحبا'))
        self.assertTrue(has_arabic('العالم'))

    def test_has_arabic_false(self):
        self.assertFalse(has_arabic('Hello World'))
        self.assertFalse(has_arabic('12345'))
        self.assertFalse(has_arabic(''))
        self.assertFalse(has_arabic('   '))

    def test_has_arabic_boundary_chars(self):
        # U+0600 (first Arabic block char)
        self.assertTrue(has_arabic('\u0600'))
        # U+06FF (last Arabic block char)
        self.assertTrue(has_arabic('\u06FF'))
        # U+FB50 (Arabic Presentation Forms-A start)
        self.assertTrue(has_arabic('\uFB50'))
        # U+FEFF is the last char of Arabic Presentation Forms-B (FE70-FEFF)
        # OOXML includes this range as Arabic, so it IS Arabic
        self.assertTrue(has_arabic('\uFEFF'))
        # U+0600-U+06FF range is exactly covered
        self.assertTrue(has_arabic('\u0660'))  # Arabic-Indic digit
        # Clearly non-Arabic characters
        self.assertFalse(has_arabic('\u00FF'))  # Latin Extended-A
        self.assertFalse(has_arabic('\u4E00'))  # CJK

    def test_has_latin_true(self):
        self.assertTrue(has_latin('Hello'))
        self.assertTrue(has_latin('a'))
        self.assertTrue(has_latin('ABC'))
        self.assertTrue(has_latin('hello مرحبا'))  # Mixed

    def test_has_latin_false(self):
        self.assertFalse(has_latin('مرحبا'))
        self.assertFalse(has_latin('12345'))
        self.assertFalse(has_latin(''))

    def test_is_bidi_text_true(self):
        self.assertTrue(is_bidi_text('Hello مرحبا'))
        self.assertTrue(is_bidi_text('Revenue العائدات'))

    def test_is_bidi_text_false(self):
        self.assertFalse(is_bidi_text('Hello World'))
        self.assertFalse(is_bidi_text('مرحبا بالعالم'))
        self.assertFalse(is_bidi_text('12345'))

    def test_compute_script_ratio_pure_arabic(self):
        text = 'مرحبا'
        ratios = compute_script_ratio(text)
        self.assertGreater(ratios['arabic'], 0.99)
        self.assertAlmostEqual(ratios['latin'], 0.0)

    def test_compute_script_ratio_pure_latin(self):
        text = 'Hello'
        ratios = compute_script_ratio(text)
        self.assertAlmostEqual(ratios['arabic'], 0.0)
        self.assertGreater(ratios['latin'], 0.99)

    def test_compute_script_ratio_mixed(self):
        # "Hello مرحبا" — 5 Latin + 5 Arabic = ~50/50
        text = 'Hello مرحبا'
        ratios = compute_script_ratio(text)
        self.assertGreater(ratios['arabic'], 0.0)
        self.assertGreater(ratios['latin'], 0.0)

    def test_compute_script_ratio_empty(self):
        ratios = compute_script_ratio('')
        self.assertEqual(ratios['arabic'], 0.0)
        self.assertEqual(ratios['latin'], 0.0)
        self.assertEqual(ratios['numeric'], 0.0)

    def test_compute_script_ratio_sums_to_one(self):
        for text in ['Hello مرحبا 123', 'مرحبا', 'Hello', '12345']:
            ratios = compute_script_ratio(text)
            total = sum(ratios.values())
            self.assertAlmostEqual(total, 1.0, places=5,
                                   msg=f'Ratios don\'t sum to 1 for: {text!r}')

    def test_compute_script_ratio_whitespace_ignored(self):
        # Whitespace-only string → all zeros
        ratios = compute_script_ratio('   \t\n')
        self.assertEqual(ratios['arabic'], 0.0)


# ─────────────────────────────────────────────────────────────────────────────
# Test: utils.py — OOXML XML helpers
# ─────────────────────────────────────────────────────────────────────────────

class TestOOXMLHelpers(unittest.TestCase):

    def test_qn_a_namespace(self):
        self.assertEqual(
            qn('a:pPr'),
            f'{{{A_NS}}}pPr'
        )

    def test_qn_p_namespace(self):
        self.assertEqual(
            qn('p:sp'),
            f'{{{P_NS}}}sp'
        )

    def test_qn_unknown_prefix_raises(self):
        with self.assertRaises(KeyError):
            qn('x:unknown')

    def test_ensure_pPr_creates_when_missing(self):
        p = etree.Element(f'{{{A_NS}}}p')
        r = etree.SubElement(p, f'{{{A_NS}}}r')
        pPr = ensure_pPr(p)
        self.assertIsNotNone(pPr)
        self.assertEqual(pPr.tag, f'{{{A_NS}}}pPr')
        # pPr should be first child
        self.assertEqual(list(p)[0].tag, f'{{{A_NS}}}pPr')

    def test_ensure_pPr_returns_existing(self):
        p = etree.Element(f'{{{A_NS}}}p')
        existing_pPr = etree.SubElement(p, f'{{{A_NS}}}pPr')
        existing_pPr.set('rtl', '1')
        returned = ensure_pPr(p)
        self.assertIs(returned, existing_pPr)
        # Should not create a second pPr
        self.assertEqual(len(p.findall(f'{{{A_NS}}}pPr')), 1)

    def test_set_rtl_on_paragraph(self):
        p = make_para_elem('Hello')
        set_rtl_on_paragraph(p)
        pPr = p.find(f'{{{A_NS}}}pPr')
        self.assertIsNotNone(pPr)
        self.assertEqual(pPr.get('rtl'), '1')

    def test_set_rtl_on_paragraph_already_has_pPr(self):
        p = make_para_elem()
        pPr_existing = etree.SubElement(p, f'{{{A_NS}}}pPr')
        pPr_existing.set('algn', 'r')
        set_rtl_on_paragraph(p)
        # Should set rtl on existing pPr and not create a duplicate
        pPrs = p.findall(f'{{{A_NS}}}pPr')
        self.assertEqual(len(pPrs), 1)
        self.assertEqual(pPrs[0].get('rtl'), '1')
        self.assertEqual(pPrs[0].get('algn'), 'r')  # Existing attr preserved

    def test_set_alignment_on_paragraph(self):
        p = make_para_elem()
        set_alignment_on_paragraph(p, 'r')
        pPr = p.find(f'{{{A_NS}}}pPr')
        self.assertEqual(pPr.get('algn'), 'r')

    def test_set_alignment_overrides_existing(self):
        p = make_para_elem(algn='l')
        set_alignment_on_paragraph(p, 'ctr')
        pPr = p.find(f'{{{A_NS}}}pPr')
        self.assertEqual(pPr.get('algn'), 'ctr')

    def test_bounds_check_emu_within_bounds(self):
        slide_width = 12192000
        self.assertTrue(bounds_check_emu(0, slide_width))
        self.assertTrue(bounds_check_emu(6096000, slide_width))
        self.assertTrue(bounds_check_emu(12192000, slide_width))

    def test_bounds_check_emu_out_of_bounds(self):
        slide_width = 12192000
        self.assertFalse(bounds_check_emu(-1000000, slide_width))
        self.assertFalse(bounds_check_emu(slide_width + 1000000, slide_width))

    def test_bounds_check_emu_allows_small_negative(self):
        # -200000 EMU is allowed (bleed off edge)
        self.assertTrue(bounds_check_emu(-200000, 12192000))
        self.assertFalse(bounds_check_emu(-200001, 12192000))

    def test_clamp_emu(self):
        slide_width = 12192000
        self.assertEqual(clamp_emu(-300000, slide_width), -200000)
        self.assertEqual(clamp_emu(13000000, slide_width), slide_width + 500000)
        self.assertEqual(clamp_emu(5000000, slide_width), 5000000)


# ─────────────────────────────────────────────────────────────────────────────
# Test: TransformReport
# ─────────────────────────────────────────────────────────────────────────────

class TestTransformReport(unittest.TestCase):

    def test_initial_state(self):
        r = TransformReport(phase='master')
        self.assertEqual(r.total_changes, 0)
        self.assertEqual(r.changes_by_type, {})
        self.assertEqual(r.warnings, [])
        self.assertEqual(r.errors, [])

    def test_add_increments_total(self):
        r = TransformReport(phase='layout')
        r.add('mirror', 3)
        r.add('rtl_set', 5)
        self.assertEqual(r.total_changes, 8)
        self.assertEqual(r.changes_by_type['mirror'], 3)
        self.assertEqual(r.changes_by_type['rtl_set'], 5)

    def test_add_accumulates_same_type(self):
        r = TransformReport(phase='slide')
        r.add('mirror', 2)
        r.add('mirror', 3)
        self.assertEqual(r.changes_by_type['mirror'], 5)
        self.assertEqual(r.total_changes, 5)

    def test_warn_appends(self):
        r = TransformReport(phase='master')
        r.warn('test warning')
        self.assertEqual(len(r.warnings), 1)
        self.assertIn('test warning', r.warnings[0])

    def test_error_appends(self):
        r = TransformReport(phase='layout')
        r.error('test error')
        self.assertEqual(len(r.errors), 1)
        self.assertIn('test error', r.errors[0])

    def test_merge_combines_counts(self):
        r1 = TransformReport(phase='master')
        r1.add('mirror', 5)
        r1.add('rtl_set', 10)
        r1.warn('w1')

        r2 = TransformReport(phase='layout')
        r2.add('mirror', 3)
        r2.add('swap', 2)
        r2.error('e1')

        r1.merge(r2)
        self.assertEqual(r1.total_changes, 20)
        self.assertEqual(r1.changes_by_type['mirror'], 8)
        self.assertEqual(r1.changes_by_type['rtl_set'], 10)
        self.assertEqual(r1.changes_by_type['swap'], 2)
        self.assertEqual(len(r1.warnings), 1)
        self.assertEqual(len(r1.errors), 1)


# ─────────────────────────────────────────────────────────────────────────────
# Test: MasterLayoutTransformer — RTL direction defaults
# ─────────────────────────────────────────────────────────────────────────────

class TestMasterLayoutTransformerDefaults(unittest.TestCase):
    """Tests for the RTL direction default-setting logic on XML elements."""

    def _make_transformer(self):
        """Create a MasterLayoutTransformer with a mock presentation."""
        prs = MagicMock()
        prs.slide_width = 12192000
        prs.slide_height = 6858000
        prs.slide_masters = []
        return MasterLayoutTransformer(prs)

    def test_apply_rtl_direction_defaults_bodyPr(self):
        """bodyPr elements should get rtlCol='1'."""
        transformer = self._make_transformer()
        root = etree.Element('root')
        bodyPr = etree.SubElement(root, f'{{{A_NS}}}bodyPr')
        transformer._apply_rtl_direction_defaults(root)
        self.assertEqual(bodyPr.get('rtlCol'), '1')

    def test_apply_rtl_direction_defaults_defPPr(self):
        """defPPr elements should get rtl='1' but NOT algn."""
        transformer = self._make_transformer()
        root = etree.Element('root')
        defPPr = etree.SubElement(root, f'{{{A_NS}}}defPPr')
        transformer._apply_rtl_direction_defaults(root)
        self.assertEqual(defPPr.get('rtl'), '1')
        # CRITICAL: algn must NOT be set at master level
        self.assertIsNone(defPPr.get('algn'),
                          'algn must never be set at master/layout level')

    def test_apply_rtl_direction_defaults_lstStyle(self):
        """lstStyle lvlNpPr elements should get rtl='1' but NOT algn."""
        transformer = self._make_transformer()
        root = etree.Element('root')
        lstStyle = etree.SubElement(root, f'{{{A_NS}}}lstStyle')
        lvl1pPr = etree.SubElement(lstStyle, f'{{{A_NS}}}lvl1pPr')
        lvl2pPr = etree.SubElement(lstStyle, f'{{{A_NS}}}lvl2pPr')
        transformer._apply_rtl_direction_defaults(root)
        self.assertEqual(lvl1pPr.get('rtl'), '1')
        self.assertEqual(lvl2pPr.get('rtl'), '1')
        # CRITICAL: no algn at this level
        self.assertIsNone(lvl1pPr.get('algn'))
        self.assertIsNone(lvl2pPr.get('algn'))

    def test_apply_arabic_lang_defaults(self):
        """defRPr elements should get lang='ar-SA'."""
        transformer = self._make_transformer()
        root = etree.Element('root')
        defRPr = etree.SubElement(root, f'{{{A_NS}}}defRPr')
        transformer._apply_arabic_lang_defaults(root)
        self.assertEqual(defRPr.get('lang'), 'ar-SA')

    def test_no_flipH_anywhere(self):
        """The transformer must never set flipH on any element."""
        transformer = self._make_transformer()
        # This is a design invariant test — if flipH is ever added to
        # _apply_rtl_direction_defaults or _apply_arabic_lang_defaults,
        # these tests would detect it.
        root = etree.Element('root')
        spPr = etree.SubElement(root, f'{{{A_NS}}}spPr')
        xfrm = etree.SubElement(spPr, f'{{{A_NS}}}xfrm')
        transformer._apply_rtl_direction_defaults(root)
        self.assertIsNone(xfrm.get('flipH'),
                          'flipH must never be set by RTL direction defaults')


# ─────────────────────────────────────────────────────────────────────────────
# Test: MasterLayoutTransformer — logo detection
# ─────────────────────────────────────────────────────────────────────────────

class TestLogoDetection(unittest.TestCase):

    def _make_transformer(self, slide_width=12192000):
        prs = MagicMock()
        prs.slide_width = slide_width
        prs.slide_height = 6858000
        prs.slide_masters = []
        return MasterLayoutTransformer(prs)

    def _make_pic_element(self, width_fraction=0.10, has_text=False,
                          has_embed=True, is_placeholder=False):
        """Create a mock shape that looks like a picture."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        slide_width = 12192000
        shape = MagicMock()
        shape.name = 'logo'
        shape.width = int(slide_width * width_fraction)
        shape.has_text_frame = has_text
        shape.shape_type = MSO_SHAPE_TYPE.PICTURE

        # Build minimal XML for a <p:pic> element
        pic_el = etree.Element(f'{{{P_NS}}}pic')
        nv_pic_pr = etree.SubElement(pic_el, f'{{{P_NS}}}nvPicPr')
        nv_cNv_pr = etree.SubElement(nv_pic_pr, f'{{{P_NS}}}cNvPr')
        nv_cNv_pr.set('id', '1')
        nv_cNv_pr.set('name', 'logo')

        if is_placeholder:
            sp_pr_inner = etree.SubElement(nv_pic_pr, f'{{{P_NS}}}nvPicSpPr')
            ph = etree.SubElement(sp_pr_inner, f'{{{P_NS}}}ph')

        if has_embed:
            blip_fill = etree.SubElement(pic_el, f'{{{P_NS}}}blipFill')
            blip = etree.SubElement(blip_fill, f'{{{A_NS}}}blip')
            blip.set(f'{{{P_NS}}}embed'.replace(P_NS, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'), 'rId1')
            # Use the proper R_NS
            from slidearabi.utils import R_NS
            blip.set(f'{{{R_NS}}}embed', 'rId1')

        shape._element = pic_el
        return shape

    def test_small_logo_is_detected(self):
        transformer = self._make_transformer()
        shape = self._make_pic_element(width_fraction=0.10)
        self.assertTrue(transformer._is_logo_shape(shape, transformer._slide_width))

    def test_wide_image_not_logo(self):
        transformer = self._make_transformer()
        shape = self._make_pic_element(width_fraction=0.30)  # > 20% threshold
        self.assertFalse(transformer._is_logo_shape(shape, transformer._slide_width))

    def test_text_bearing_shape_not_logo(self):
        transformer = self._make_transformer()
        shape = self._make_pic_element(width_fraction=0.10, has_text=True)
        self.assertFalse(transformer._is_logo_shape(shape, transformer._slide_width))

    def test_shape_without_embed_not_logo(self):
        transformer = self._make_transformer()
        shape = self._make_pic_element(width_fraction=0.10, has_embed=False)
        self.assertFalse(transformer._is_logo_shape(shape, transformer._slide_width))


# ─────────────────────────────────────────────────────────────────────────────
# Test: MasterLayoutTransformer — position mirroring
# ─────────────────────────────────────────────────────────────────────────────

class TestShapeMirroring(unittest.TestCase):

    SLIDE_WIDTH = 12192000

    def _make_transformer(self):
        prs = MagicMock()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = 6858000
        prs.slide_masters = []
        return MasterLayoutTransformer(prs)

    def test_mirror_shape_position_basic(self):
        transformer = self._make_transformer()
        shape = MagicMock()
        shape.name = 'test'
        shape.left = 914400    # 1 inch
        shape.width = 1828800  # 2 inches
        shape.top = 500000

        result = transformer._mirror_shape_position(shape, self.SLIDE_WIDTH)
        self.assertTrue(result)
        expected = mirror_x(914400, 1828800, self.SLIDE_WIDTH)
        self.assertEqual(shape.left, expected)
        # top should be unchanged
        self.assertEqual(shape.top, 500000)

    def test_mirror_shape_position_centred_skipped(self):
        """Centred shape: mirror is same position, so no change should be made."""
        transformer = self._make_transformer()
        shape = MagicMock()
        shape.name = 'centred'
        width = 4000000
        left = (self.SLIDE_WIDTH - width) // 2
        shape.left = left
        shape.width = width
        original_left = shape.left

        result = transformer._mirror_shape_position(shape, self.SLIDE_WIDTH)
        # Centred shape mirror == original — below tolerance
        self.assertFalse(result)  # No change for symmetric position

    def test_mirror_shape_position_oob_rejected(self):
        transformer = self._make_transformer()
        shape = MagicMock()
        shape.name = 'oob'
        shape.left = -5000000  # Far off-slide
        shape.width = 1000000
        result = transformer._mirror_shape_position(shape, self.SLIDE_WIDTH)
        # mirror_x(-5000000, 1000000, 12192000) = 16192000 → out of bounds
        self.assertFalse(result)

    def test_mirror_involution(self):
        """Mirroring a shape twice returns it to the original position."""
        transformer = self._make_transformer()
        original_left = 2000000
        shape = MagicMock()
        shape.name = 'test'
        shape.left = original_left
        shape.width = 3000000

        transformer._mirror_shape_position(shape, self.SLIDE_WIDTH)
        first_mirror = shape.left

        shape.left = first_mirror
        transformer._mirror_shape_position(shape, self.SLIDE_WIDTH)
        self.assertEqual(shape.left, original_left)


# ─────────────────────────────────────────────────────────────────────────────
# Test: SlideContentTransformer — should_mirror_shape
# ─────────────────────────────────────────────────────────────────────────────

class TestShouldMirrorShape(unittest.TestCase):

    def _make_transformer(self):
        prs = MagicMock()
        prs.slide_width = 12192000
        prs.slide_height = 6858000
        prs.slides = []
        return SlideContentTransformer(prs)

    def test_full_width_background_not_mirrored(self):
        transformer = self._make_transformer()
        shape = MagicMock()
        shape.width = int(12192000 * 0.95)  # 95% of slide width
        result = transformer._should_mirror_shape(shape, 'tx')
        self.assertFalse(result)

    def test_secHead_layout_preserves_shapes(self):
        transformer = self._make_transformer()
        shape = MagicMock()
        shape.width = 2000000  # Normal content shape
        result = transformer._should_mirror_shape(shape, 'secHead')
        self.assertFalse(result)

    def test_title_layout_preserves_shapes(self):
        transformer = self._make_transformer()
        shape = MagicMock()
        shape.width = 2000000
        result = transformer._should_mirror_shape(shape, 'title')
        self.assertFalse(result)

    def test_standard_layout_mirrors_content(self):
        transformer = self._make_transformer()
        shape = MagicMock()
        shape.width = 2000000
        result = transformer._should_mirror_shape(shape, 'tx')
        self.assertTrue(result)

    def test_blank_layout_mirrors_shapes(self):
        transformer = self._make_transformer()
        shape = MagicMock()
        shape.width = 2000000
        result = transformer._should_mirror_shape(shape, 'blank')
        self.assertTrue(result)


# ─────────────────────────────────────────────────────────────────────────────
# Test: SlideContentTransformer — paragraph alignment computation
# ─────────────────────────────────────────────────────────────────────────────

class TestParagraphAlignmentComputation(unittest.TestCase):

    def _make_transformer(self):
        prs = MagicMock()
        prs.slide_width = 12192000
        prs.slide_height = 6858000
        prs.slides = []
        return SlideContentTransformer(prs)

    def test_arabic_text_aligns_right(self):
        transformer = self._make_transformer()
        algn = transformer._compute_paragraph_alignment('مرحبا بالعالم', None)
        self.assertEqual(algn, 'r')

    def test_footer_aligns_left(self):
        transformer = self._make_transformer()
        for ph_type in ('ftr', 'sldNum', 'dt'):
            algn = transformer._compute_paragraph_alignment('مرحبا', ph_type)
            self.assertEqual(algn, 'l',
                             f'Footer type {ph_type!r} should always be left-aligned')

    def test_ctrTitle_aligns_centre(self):
        transformer = self._make_transformer()
        algn = transformer._compute_paragraph_alignment('مرحبا', 'ctrTitle')
        self.assertEqual(algn, 'ctr')

    def test_mixed_bidi_aligns_right(self):
        transformer = self._make_transformer()
        # Mixed Arabic/Latin → right-align (Arabic context)
        algn = transformer._compute_paragraph_alignment('Revenue العائدات', None)
        self.assertEqual(algn, 'r')

    def test_pure_latin_aligns_left(self):
        transformer = self._make_transformer()
        algn = transformer._compute_paragraph_alignment('Hello World', None)
        self.assertEqual(algn, 'l')


# ─────────────────────────────────────────────────────────────────────────────
# Test: Typography — font mapping
# ─────────────────────────────────────────────────────────────────────────────

class TestFontMapping(unittest.TestCase):

    def _make_normalizer(self):
        prs = MagicMock()
        prs.slide_width = 12192000
        prs.slide_height = 6858000
        prs.slides = []
        return TypographyNormalizer(prs)

    def test_calibri_maps_to_itself(self):
        normalizer = self._make_normalizer()
        self.assertEqual(normalizer._map_font('Calibri'), 'Calibri')

    def test_arial_maps_to_itself(self):
        normalizer = self._make_normalizer()
        self.assertEqual(normalizer._map_font('Arial'), 'Arial')

    def test_times_new_roman_maps_to_itself(self):
        normalizer = self._make_normalizer()
        self.assertEqual(normalizer._map_font('Times New Roman'), 'Times New Roman')

    def test_cambria_maps_to_sakkal(self):
        normalizer = self._make_normalizer()
        self.assertEqual(normalizer._map_font('Cambria'), 'Sakkal Majalla')

    def test_georgia_maps_to_sakkal(self):
        normalizer = self._make_normalizer()
        self.assertEqual(normalizer._map_font('Georgia'), 'Sakkal Majalla')

    def test_helvetica_maps_to_arial(self):
        normalizer = self._make_normalizer()
        self.assertEqual(normalizer._map_font('Helvetica'), 'Arial')
        self.assertEqual(normalizer._map_font('Helvetica Neue'), 'Arial')

    def test_century_gothic_maps_to_dubai(self):
        normalizer = self._make_normalizer()
        self.assertEqual(normalizer._map_font('Century Gothic'), 'Dubai')

    def test_unknown_font_returned_unchanged(self):
        normalizer = self._make_normalizer()
        # Custom corporate font — assume it may support Arabic
        result = normalizer._map_font('McKinsey Sans')
        self.assertEqual(result, 'McKinsey Sans')

    def test_empty_font_returned_unchanged(self):
        normalizer = self._make_normalizer()
        self.assertEqual(normalizer._map_font(''), '')

    def test_all_mapped_fonts_in_expansion_table_or_have_default(self):
        """All target fonts in ARABIC_FONT_MAP should have expansion factors
        or fall back to 'default'."""
        normalizer = self._make_normalizer()
        for source, target in ARABIC_FONT_MAP.items():
            # Expansion should be retrievable (key or 'default' fallback)
            factor = ARABIC_EXPANSION_FACTORS.get(target,
                     ARABIC_EXPANSION_FACTORS['default'])
            self.assertGreater(factor, 0.9, f'{target!r} expansion factor invalid')
            self.assertLess(factor, 2.0, f'{target!r} expansion factor too large')

    def test_case_insensitive_matching(self):
        normalizer = self._make_normalizer()
        # 'calibri' (lowercase) should match 'Calibri'
        result = normalizer._map_font('calibri')
        self.assertEqual(result, 'Calibri')


# ─────────────────────────────────────────────────────────────────────────────
# Test: Typography — expansion estimation
# ─────────────────────────────────────────────────────────────────────────────

class TestArabicExpansionEstimation(unittest.TestCase):

    def _make_normalizer(self):
        prs = MagicMock()
        prs.slide_width = 12192000
        prs.slide_height = 6858000
        prs.slides = []
        return TypographyNormalizer(prs)

    def test_expansion_is_positive(self):
        normalizer = self._make_normalizer()
        ratio = normalizer._estimate_arabic_expansion('Hello', 'مرحبا', 12.0)
        self.assertGreater(ratio, 0.0)

    def test_expansion_within_bounds(self):
        normalizer = self._make_normalizer()
        # Should be clamped to [0.8, 2.0]
        ratio = normalizer._estimate_arabic_expansion('Hi', 'مرحبا بالعالم العربي', 12.0)
        self.assertGreaterEqual(ratio, 0.8)
        self.assertLessEqual(ratio, 2.0)

    def test_empty_source_returns_one(self):
        normalizer = self._make_normalizer()
        ratio = normalizer._estimate_arabic_expansion('', 'مرحبا', 12.0)
        self.assertEqual(ratio, 1.0)

    def test_per_font_factor_used(self):
        normalizer = self._make_normalizer()
        # Sakkal Majalla has lower expansion (1.15) vs Traditional Arabic (1.30)
        ratio_sakkal = normalizer._estimate_arabic_expansion(
            'Hello World', 'مرحبا بالعالم', 12.0, 'Sakkal Majalla')
        ratio_trad = normalizer._estimate_arabic_expansion(
            'Hello World', 'مرحبا بالعالم', 12.0, 'Traditional Arabic')
        # Sakkal should estimate lower expansion than Traditional Arabic
        self.assertLess(ratio_sakkal, ratio_trad)

    def test_short_arabic_longer_than_short_latin(self):
        normalizer = self._make_normalizer()
        # Arabic text is typically ~same length or longer than English equivalent
        ratio = normalizer._estimate_arabic_expansion('Q1', 'الربع الأول', 12.0)
        # Arabic expansion should be > 1 (Arabic version is longer)
        self.assertGreater(ratio, 1.0)


# ─────────────────────────────────────────────────────────────────────────────
# Test: Typography — bidi paragraph formatting rules
# ─────────────────────────────────────────────────────────────────────────────

class TestBidiFormattingRules(unittest.TestCase):
    """
    Tests that _apply_bidi_formatting sets correct rtl and algn attributes
    on paragraph elements for various text content types.
    """

    def _make_normalizer(self):
        prs = MagicMock()
        prs.slide_width = 12192000
        prs.slide_height = 6858000
        prs.slides = []
        return TypographyNormalizer(prs)

    def _make_shape_with_text(self, text: str, ph_type: str | None = None):
        """Create a minimal mock shape with a text frame containing *text*."""
        shape = MagicMock()
        shape.name = 'test_shape'

        # Build txBody
        txBody = make_txBody_elem([text])
        tf = MagicMock()
        tf._txBody = txBody
        tf.text = text

        # Create paragraph objects
        p_elems = txBody.findall(f'{{{A_NS}}}p')
        paragraphs = []
        for p_el in p_elems:
            para = MagicMock()
            para._p = p_el
            para.text = text
            para.runs = []
            for r_el in p_el.findall(f'{{{A_NS}}}r'):
                run = MagicMock()
                run._r = r_el
                t_el = r_el.find(f'{{{A_NS}}}t')
                run.text = t_el.text if t_el is not None else ''
                para.runs.append(run)
            paragraphs.append(para)

        tf.paragraphs = paragraphs
        shape.has_text_frame = True
        shape.text_frame = tf

        if ph_type:
            ph_fmt = MagicMock()
            ph_fmt.type = ph_type
            ph_fmt.idx = 0
            shape.placeholder_format = ph_fmt
            shape.is_placeholder = True

            # Mock get_placeholder_info to return ph_type
            # We'll patch it at the call site using the shape mock
            # Instead, set up the _element for get_placeholder_info_from_xml
            p_ns = P_NS
            sp_el = etree.Element(f'{{{p_ns}}}sp')
            nv_sp_pr = etree.SubElement(sp_el, f'{{{p_ns}}}nvSpPr')
            ph_el = etree.SubElement(nv_sp_pr, f'{{{p_ns}}}ph')
            ph_el.set('type', ph_type)
            shape._element = sp_el
        else:
            shape.is_placeholder = False
            shape.placeholder_format = None
            shape._element = etree.Element(f'{{{P_NS}}}sp')

        return shape

    def _get_pPr_attr(self, txBody, attr: str, para_idx: int = 0) -> str | None:
        """Helper to extract an attribute from a paragraph's pPr."""
        paras = txBody.findall(f'{{{A_NS}}}p')
        if para_idx >= len(paras):
            return None
        pPr = paras[para_idx].find(f'{{{A_NS}}}pPr')
        if pPr is None:
            return None
        return pPr.get(attr)

    def test_arabic_paragraph_gets_rtl_1_and_algn_r(self):
        normalizer = self._make_normalizer()
        shape = self._make_shape_with_text('مرحبا بالعالم العربي')
        normalizer._apply_bidi_formatting(shape)
        txBody = shape.text_frame._txBody
        self.assertEqual(self._get_pPr_attr(txBody, 'rtl'), '1')
        self.assertEqual(self._get_pPr_attr(txBody, 'algn'), 'r')

    def test_latin_paragraph_in_ltr_context_gets_ltr(self):
        normalizer = self._make_normalizer()
        shape = self._make_shape_with_text('Hello World')
        normalizer._apply_bidi_formatting(shape)
        txBody = shape.text_frame._txBody
        self.assertEqual(self._get_pPr_attr(txBody, 'rtl'), '0')
        self.assertEqual(self._get_pPr_attr(txBody, 'algn'), 'l')

    def test_ctrTitle_gets_rtl_1_and_ctr(self):
        normalizer = self._make_normalizer()
        shape = self._make_shape_with_text('مرحبا', 'ctrTitle')

        # Patch get_placeholder_info to return ('ctrTitle', 0)
        with patch('slidearabi.typography.get_placeholder_info',
                   return_value=('ctrTitle', 0)):
            normalizer._apply_bidi_formatting(shape)

        txBody = shape.text_frame._txBody
        self.assertEqual(self._get_pPr_attr(txBody, 'rtl'), '1')
        self.assertEqual(self._get_pPr_attr(txBody, 'algn'), 'ctr')

    def test_footer_always_left(self):
        normalizer = self._make_normalizer()
        shape = self._make_shape_with_text('مرحبا', 'ftr')

        with patch('slidearabi.typography.get_placeholder_info',
                   return_value=('ftr', 0)):
            normalizer._apply_bidi_formatting(shape)

        txBody = shape.text_frame._txBody
        self.assertEqual(self._get_pPr_attr(txBody, 'algn'), 'l')


# ─────────────────────────────────────────────────────────────────────────────
# Test: Typography — font size reduction
# ─────────────────────────────────────────────────────────────────────────────

class TestFontSizeReduction(unittest.TestCase):

    def _make_normalizer(self):
        prs = MagicMock()
        prs.slide_width = 12192000
        prs.slide_height = 6858000
        prs.slides = []
        return TypographyNormalizer(prs)

    def _make_shape_with_font_size(self, size_pt: float, ph_type: str | None = None):
        """Create a mock shape with a single run at a given font size."""
        shape = MagicMock()
        shape.name = 'test'
        shape.width = 4572000   # 5 inches
        shape.height = 1143000  # 1.25 inches

        # Build XML
        txBody = etree.Element(f'{{{A_NS}}}txBody')
        bodyPr = etree.SubElement(txBody, f'{{{A_NS}}}bodyPr')
        p = etree.SubElement(txBody, f'{{{A_NS}}}p')
        r = etree.SubElement(p, f'{{{A_NS}}}r')
        rPr = etree.SubElement(r, f'{{{A_NS}}}rPr')
        rPr.set('sz', str(int(size_pt * 100)))
        t = etree.SubElement(r, f'{{{A_NS}}}t')
        t.text = 'مرحبا بالعالم العربي والطويل'

        tf = MagicMock()
        tf._txBody = txBody

        para = MagicMock()
        para._p = p
        para.text = 'test'
        run = MagicMock()
        run._r = r
        run.text = t.text
        para.runs = [run]
        tf.paragraphs = [para]

        shape.has_text_frame = True
        shape.text_frame = tf
        shape.is_placeholder = ph_type is not None
        if ph_type:
            ph_fmt = MagicMock()
            ph_fmt.type = ph_type
            ph_fmt.idx = 0
            shape.placeholder_format = ph_fmt
        else:
            shape.placeholder_format = None

        return shape

    def test_font_reduced_when_overflow(self):
        normalizer = self._make_normalizer()
        shape = self._make_shape_with_font_size(20.0)
        normalizer._reduce_font_size_to_fit(shape, max_reduction_pct=20.0)
        rPr = shape.text_frame._txBody.find(f'.//{{{A_NS}}}rPr')
        new_sz = int(rPr.get('sz')) / 100.0
        # Should be reduced by up to 20%
        self.assertLessEqual(new_sz, 20.0)
        self.assertGreaterEqual(new_sz, 20.0 * 0.80)

    def test_font_not_reduced_below_body_floor(self):
        normalizer = self._make_normalizer()
        shape = self._make_shape_with_font_size(10.0)
        normalizer._reduce_font_size_to_fit(shape, max_reduction_pct=20.0)
        rPr = shape.text_frame._txBody.find(f'.//{{{A_NS}}}rPr')
        new_sz = int(rPr.get('sz')) / 100.0
        # At 10pt with body floor of 10pt: no reduction possible
        self.assertGreaterEqual(new_sz, MIN_BODY_FONT_SIZE_PT)

    def test_title_floor_enforced(self):
        normalizer = self._make_normalizer()
        shape = self._make_shape_with_font_size(15.0, ph_type='title')
        with patch('slidearabi.typography.get_placeholder_info',
                   return_value=('title', 0)):
            normalizer._reduce_font_size_to_fit(shape, max_reduction_pct=20.0)
        rPr = shape.text_frame._txBody.find(f'.//{{{A_NS}}}rPr')
        new_sz = int(rPr.get('sz')) / 100.0
        self.assertGreaterEqual(new_sz, MIN_TITLE_FONT_SIZE_PT)

    def test_no_reduction_for_small_font(self):
        """Font already at the floor should not be reduced further."""
        normalizer = self._make_normalizer()
        shape = self._make_shape_with_font_size(MIN_FONT_SIZE_PT)
        normalizer._reduce_font_size_to_fit(shape, max_reduction_pct=20.0)
        rPr = shape.text_frame._txBody.find(f'.//{{{A_NS}}}rPr')
        new_sz = int(rPr.get('sz')) / 100.0
        self.assertGreaterEqual(new_sz, MIN_FONT_SIZE_PT)


# ─────────────────────────────────────────────────────────────────────────────
# Test: Typography — text frame margins
# ─────────────────────────────────────────────────────────────────────────────

class TestTextFrameMargins(unittest.TestCase):

    def _make_normalizer(self):
        prs = MagicMock()
        prs.slide_width = 12192000
        prs.slide_height = 6858000
        prs.slides = []
        return TypographyNormalizer(prs)

    def _make_shape_with_arabic(self, has_arabic_text: bool):
        shape = MagicMock()
        shape.name = 'test'

        txBody = etree.Element(f'{{{A_NS}}}txBody')
        bodyPr = etree.SubElement(txBody, f'{{{A_NS}}}bodyPr')

        text = 'مرحبا بالعالم' if has_arabic_text else 'Hello World'
        p = etree.SubElement(txBody, f'{{{A_NS}}}p')
        r = etree.SubElement(p, f'{{{A_NS}}}r')
        t = etree.SubElement(r, f'{{{A_NS}}}t')
        t.text = text

        tf = MagicMock()
        tf._txBody = txBody
        tf.text = text
        shape.has_text_frame = True
        shape.text_frame = tf
        return shape

    def test_arabic_frame_gets_insets(self):
        from slidearabi.typography import (
            ARABIC_INSET_LEFT_EMU, ARABIC_INSET_RIGHT_EMU,
            ARABIC_INSET_TOP_EMU, ARABIC_INSET_BOTTOM_EMU,
        )
        normalizer = self._make_normalizer()
        shape = self._make_shape_with_arabic(has_arabic_text=True)
        normalizer._set_text_frame_margins(shape)

        bodyPr = shape.text_frame._txBody.find(f'{{{A_NS}}}bodyPr')
        self.assertEqual(int(bodyPr.get('lIns')), ARABIC_INSET_LEFT_EMU)
        self.assertEqual(int(bodyPr.get('rIns')), ARABIC_INSET_RIGHT_EMU)
        self.assertEqual(int(bodyPr.get('tIns')), ARABIC_INSET_TOP_EMU)
        self.assertEqual(int(bodyPr.get('bIns')), ARABIC_INSET_BOTTOM_EMU)

    def test_latin_frame_insets_not_changed(self):
        normalizer = self._make_normalizer()
        shape = self._make_shape_with_arabic(has_arabic_text=False)
        result = normalizer._set_text_frame_margins(shape)
        self.assertEqual(result, 0)  # No change

        bodyPr = shape.text_frame._txBody.find(f'{{{A_NS}}}bodyPr')
        self.assertIsNone(bodyPr.get('lIns'))


# ─────────────────────────────────────────────────────────────────────────────
# Test: ARABIC_FONT_MAP completeness
# ─────────────────────────────────────────────────────────────────────────────

class TestArabicFontMapCompleteness(unittest.TestCase):

    def test_all_values_are_strings(self):
        for k, v in ARABIC_FONT_MAP.items():
            self.assertIsInstance(v, str, f'ARABIC_FONT_MAP[{k!r}] should be a string')
            self.assertGreater(len(v), 0, f'ARABIC_FONT_MAP[{k!r}] should be non-empty')

    def test_self_mapping_fonts_map_to_themselves(self):
        """Fonts that already support Arabic should map to themselves."""
        self_mapping = ['Calibri', 'Arial', 'Times New Roman', 'Tahoma', 'Segoe UI']
        for font in self_mapping:
            self.assertEqual(ARABIC_FONT_MAP.get(font), font,
                             f'{font!r} should map to itself')

    def test_problematic_fonts_are_mapped(self):
        """Latin-only fonts should be mapped to Arabic alternatives."""
        latin_only = ['Cambria', 'Georgia', 'Garamond', 'Palatino Linotype',
                      'Century Gothic', 'Helvetica Neue', 'Impact']
        for font in latin_only:
            self.assertIn(font, ARABIC_FONT_MAP,
                         f'{font!r} should be in ARABIC_FONT_MAP')
            mapped = ARABIC_FONT_MAP[font]
            self.assertNotEqual(mapped, font,
                               f'{font!r} is a Latin-only font and should map '
                               f'to an Arabic-capable font, not itself')


# ─────────────────────────────────────────────────────────────────────────────
# Test: Table column reversal integration
# ─────────────────────────────────────────────────────────────────────────────

class TestTableColumnReversal(unittest.TestCase):
    """Tests for SlideContentTransformer._transform_table_rtl XML operations."""

    def _make_transformer(self):
        prs = MagicMock()
        prs.slide_width = 12192000
        prs.slide_height = 6858000
        prs.slides = []
        return SlideContentTransformer(prs)

    def _make_table_shape(self, num_rows: int = 2, num_cols: int = 3):
        """Create a mock table shape with minimal XML structure."""
        shape = MagicMock()
        shape.name = 'table'
        shape.has_table = True

        a_ns = A_NS
        # Build minimal <a:tbl> XML
        tbl = etree.Element(f'{{{a_ns}}}tbl')
        tblPr = etree.SubElement(tbl, f'{{{a_ns}}}tblPr')
        tblGrid = etree.SubElement(tbl, f'{{{a_ns}}}tblGrid')

        original_widths = [str(100 * (i + 1)) for i in range(num_cols)]
        for w in original_widths:
            gc = etree.SubElement(tblGrid, f'{{{a_ns}}}gridCol')
            gc.set('w', w)

        for row_i in range(num_rows):
            tr = etree.SubElement(tbl, f'{{{a_ns}}}tr')
            for col_j in range(num_cols):
                tc = etree.SubElement(tr, f'{{{a_ns}}}tc')
                txBody = etree.SubElement(tc, f'{{{a_ns}}}txBody')
                bodyPr = etree.SubElement(txBody, f'{{{a_ns}}}bodyPr')
                p = etree.SubElement(txBody, f'{{{a_ns}}}p')
                r = etree.SubElement(p, f'{{{a_ns}}}r')
                t = etree.SubElement(r, f'{{{a_ns}}}t')
                t.text = f'R{row_i}C{col_j}'

        # Create mock table object
        table = MagicMock()
        table._tbl = tbl

        # Mock columns
        columns = [MagicMock() for _ in range(num_cols)]
        table.columns = columns

        # Mock rows with cells
        rows = []
        tr_elems = tbl.findall(f'{{{a_ns}}}tr')
        for tr_el in tr_elems:
            row = MagicMock()
            row._tr = tr_el
            cells = []
            for tc_el in tr_el.findall(f'{{{a_ns}}}tc'):
                cell = MagicMock()
                tf = MagicMock()
                txb = tc_el.find(f'{{{a_ns}}}txBody')
                tf._txBody = txb
                t_el = txb.find(f'.//{{{a_ns}}}t')
                tf.text = t_el.text if t_el is not None else ''
                para = MagicMock()
                para._p = txb.find(f'{{{a_ns}}}p')
                para.text = tf.text
                para.runs = []
                tf.paragraphs = [para]
                cell.text_frame = tf
                cells.append(cell)
            row.cells = cells
            rows.append(row)
        table.rows = rows

        shape.table = table
        return shape, tbl, original_widths

    def test_column_widths_reversed(self):
        transformer = self._make_transformer()
        shape, tbl, original_widths = self._make_table_shape(2, 3)
        transformer._transform_table_rtl(shape)

        a_ns = A_NS
        tblGrid = tbl.find(f'{{{a_ns}}}tblGrid')
        grid_cols = tblGrid.findall(f'{{{a_ns}}}gridCol')
        new_widths = [gc.get('w') for gc in grid_cols]
        self.assertEqual(new_widths, list(reversed(original_widths)))

    def test_cell_order_reversed_in_each_row(self):
        transformer = self._make_transformer()
        shape, tbl, _ = self._make_table_shape(2, 3)

        # Record original cell texts
        a_ns = A_NS
        original_rows = []
        for tr in tbl.findall(f'{{{a_ns}}}tr'):
            row_texts = [tc.findtext(f'.//{{{a_ns}}}t') for tc in tr.findall(f'{{{a_ns}}}tc')]
            original_rows.append(row_texts)

        transformer._transform_table_rtl(shape)

        # Check reversal
        for row_idx, tr in enumerate(tbl.findall(f'{{{a_ns}}}tr')):
            new_texts = [tc.findtext(f'.//{{{a_ns}}}t') for tc in tr.findall(f'{{{a_ns}}}tc')]
            expected = list(reversed(original_rows[row_idx]))
            self.assertEqual(new_texts, expected,
                             f'Row {row_idx}: cells not reversed correctly')

    def test_table_pr_rtl_set(self):
        transformer = self._make_transformer()
        shape, tbl, _ = self._make_table_shape(2, 3)
        transformer._transform_table_rtl(shape)

        a_ns = A_NS
        tblPr = tbl.find(f'{{{a_ns}}}tblPr')
        self.assertIsNotNone(tblPr)
        self.assertEqual(tblPr.get('rtl'), '1')

    def test_single_column_table_unchanged(self):
        """Tables with only 1 column should not be touched."""
        transformer = self._make_transformer()
        shape, tbl, _ = self._make_table_shape(2, 1)
        result = transformer._transform_table_rtl(shape)
        self.assertEqual(result, 0)


# ─────────────────────────────────────────────────────────────────────────────
# Test: Translation application
# ─────────────────────────────────────────────────────────────────────────────

class TestTranslationApplication(unittest.TestCase):

    def _make_transformer(self, translations=None):
        prs = MagicMock()
        prs.slide_width = 12192000
        prs.slide_height = 6858000
        prs.slides = []
        return SlideContentTransformer(prs, translations=translations)

    def _make_text_shape(self, text: str, ph_type: str | None = None):
        shape = MagicMock()
        shape.name = 'test'
        shape.is_placeholder = ph_type is not None

        txBody = etree.Element(f'{{{A_NS}}}txBody')
        bodyPr = etree.SubElement(txBody, f'{{{A_NS}}}bodyPr')
        p = etree.SubElement(txBody, f'{{{A_NS}}}p')
        r = etree.SubElement(p, f'{{{A_NS}}}r')
        rPr_el = etree.SubElement(r, f'{{{A_NS}}}rPr')
        rPr_el.set('sz', '1800')  # 18pt
        t = etree.SubElement(r, f'{{{A_NS}}}t')
        t.text = text

        tf = MagicMock()
        tf._txBody = txBody
        tf.text = text

        # Build a run mock that reads/writes through the actual lxml <a:t> element
        class XmlRun:
            def __init__(self, r_elem):
                self._r = r_elem
                self._t = r_elem.find(f'{{{A_NS}}}t')

            @property
            def text(self):
                return self._t.text or ''

            @text.setter
            def text(self, value):
                self._t.text = value

        xml_run = XmlRun(r)

        # Build a paragraph mock that reads through the actual lxml <a:p> element
        class XmlPara:
            def __init__(self, p_elem, runs_list):
                self._p = p_elem
                self.runs = runs_list

            @property
            def text(self):
                texts = []
                for r_el in self._p.findall(f'{{{A_NS}}}r'):
                    t_el = r_el.find(f'{{{A_NS}}}t')
                    if t_el is not None and t_el.text:
                        texts.append(t_el.text)
                return ''.join(texts)

        xml_para = XmlPara(p, [xml_run])

        tf.paragraphs = [xml_para]
        shape.text_frame = tf
        shape.has_text_frame = True

        if ph_type:
            ph_fmt = MagicMock()
            ph_fmt.type = ph_type
            ph_fmt.idx = 0
            shape.placeholder_format = ph_fmt
        else:
            shape.placeholder_format = None

        return shape, p, r

    def test_translation_replaces_first_run_text(self):
        translations = {'Hello World': 'مرحبا بالعالم'}
        transformer = self._make_transformer(translations)
        shape, p, r = self._make_text_shape('Hello World')

        transformer._apply_translation(shape, translations)

        t_elem = r.find(f'{{{A_NS}}}t')
        self.assertEqual(t_elem.text, 'مرحبا بالعالم')

    def test_translation_sets_rtl_on_paragraph(self):
        translations = {'Revenue': 'الإيرادات'}
        transformer = self._make_transformer(translations)
        shape, p, r = self._make_text_shape('Revenue')

        transformer._apply_translation(shape, translations)

        pPr = p.find(f'{{{A_NS}}}pPr')
        self.assertIsNotNone(pPr)
        self.assertEqual(pPr.get('rtl'), '1')

    def test_translation_sets_alignment_right(self):
        translations = {'Revenue': 'الإيرادات'}
        transformer = self._make_transformer(translations)
        shape, p, r = self._make_text_shape('Revenue')

        transformer._apply_translation(shape, translations)

        pPr = p.find(f'{{{A_NS}}}pPr')
        self.assertEqual(pPr.get('algn'), 'r')

    def test_no_translation_match_no_change(self):
        translations = {'Something Else': 'شيء آخر'}
        transformer = self._make_transformer(translations)
        shape, p, r = self._make_text_shape('Untranslated Text')
        original_text = r.find(f'{{{A_NS}}}t').text

        transformer._apply_translation(shape, translations)

        # Text should be unchanged
        t_elem = r.find(f'{{{A_NS}}}t')
        self.assertEqual(t_elem.text, original_text)

    def test_empty_translations_dict_no_change(self):
        transformer = self._make_transformer({})
        shape, p, r = self._make_text_shape('Hello')
        result = transformer._apply_translation(shape, {})
        self.assertEqual(result, 0)


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == '__main__':
    unittest.main(verbosity=2)

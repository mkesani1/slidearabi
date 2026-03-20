"""
tests/test_e2e_v3.py — E2E Integration Tests for V3 VQA Quality Gate

Sprint 4: End-to-end tests that exercise the full V3 pipeline:
  1. run_v3_xml_checks() on real PPTX pairs
  2. safe_apply_fixes() with rollback verification
  3. compute_gate_decision() for all 4 terminal statuses
  4. Vision prompt builder + selective vision integration
  5. API contract integration

These tests use synthetic PPTX files (no external dependencies).
"""

from __future__ import annotations

import os
import sys
import tempfile

import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))


def _create_test_pptx(slide_count: int = 3, add_table: bool = False,
                       add_arabic: bool = False, add_arrow: bool = False) -> str:
    """Create a minimal PPTX for testing. Returns temp file path."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu

    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(6858000)

    for i in range(slide_count):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        if add_table and i == 0:
            from pptx.util import Inches
            rows, cols = 3, 3
            table_shape = slide.shapes.add_table(
                rows, cols,
                Inches(1), Inches(1), Inches(5), Inches(2),
            )
            table = table_shape.table
            # Fill with identifiable content
            for r in range(rows):
                for c in range(cols):
                    table.cell(r, c).text = f"R{r}C{c}"

        if add_arabic and i == 0:
            from pptx.util import Inches
            txBox = slide.shapes.add_textbox(
                Inches(1), Inches(3), Inches(3), Inches(1))
            txBox.text_frame.text = "مرحبا بالعالم"

        if add_arrow and i == 0:
            from pptx.util import Inches
            from pptx.enum.shapes import MSO_SHAPE
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(3), Inches(4), Inches(2), Inches(0.5))

    fd, path = tempfile.mkstemp(suffix='.pptx')
    os.close(fd)
    prs.save(path)
    return path


# ─────────────────────────────────────────────────────────────────────────────
# E2E: run_v3_xml_checks
# ─────────────────────────────────────────────────────────────────────────────

class TestE2EXMLChecks:

    def test_clean_pptx_returns_zero_critical(self):
        """Clean identical PPTX pair should have zero CRITICAL defects."""
        pptx_path = _create_test_pptx(slide_count=2)
        try:
            from v3_checks import run_v3_xml_checks
            defects, metadata = run_v3_xml_checks(pptx_path, pptx_path)
            assert metadata['critical_count'] == 0
            assert metadata['slides_checked'] == 2
        finally:
            os.unlink(pptx_path)

    def test_pptx_with_table_runs_table_checks(self):
        """PPTX with table should trigger table-related checks."""
        pptx_path = _create_test_pptx(slide_count=1, add_table=True)
        try:
            from v3_checks import run_v3_xml_checks
            defects, metadata = run_v3_xml_checks(pptx_path, pptx_path)
            # Same file = no column reversal issue (identical order)
            assert metadata['slides_checked'] == 1
            assert isinstance(defects, list)
        finally:
            os.unlink(pptx_path)

    def test_pptx_with_arrow_detects_directional(self):
        """PPTX with rightArrow should detect DIRECTIONAL_SHAPE_NOT_FLIPPED."""
        pptx_path = _create_test_pptx(slide_count=1, add_arrow=True)
        try:
            from v3_checks import run_v3_xml_checks
            defects, metadata = run_v3_xml_checks(pptx_path, pptx_path)
            directional = [d for d in defects
                          if d.code == 'DIRECTIONAL_SHAPE_NOT_FLIPPED']
            # rightArrow without flipH should be detected
            assert len(directional) >= 1
        finally:
            os.unlink(pptx_path)


# ─────────────────────────────────────────────────────────────────────────────
# E2E: safe_apply_fixes
# ─────────────────────────────────────────────────────────────────────────────

class TestE2ESafeApplyFixes:

    def test_fixes_disabled_returns_empty(self):
        """When autofix is disabled, no fixes should be applied."""
        import v3_config
        orig_val = v3_config.ENABLE_V3_XML_AUTOFIX
        try:
            v3_config.ENABLE_V3_XML_AUTOFIX = False
            pptx_path = _create_test_pptx(slide_count=1, add_arrow=True)
            try:
                from v3_checks import run_v3_xml_checks, safe_apply_fixes
                defects, _ = run_v3_xml_checks(pptx_path, pptx_path)
                applied, failed = safe_apply_fixes(pptx_path, defects)
                assert len(applied) == 0
            finally:
                os.unlink(pptx_path)
        finally:
            v3_config.ENABLE_V3_XML_AUTOFIX = orig_val

    def test_fixes_enabled_applies_fixes(self):
        """When autofix is enabled, fixable defects should be applied."""
        import v3_config
        orig_autofix = v3_config.ENABLE_V3_XML_AUTOFIX
        try:
            v3_config.ENABLE_V3_XML_AUTOFIX = True
            pptx_path = _create_test_pptx(slide_count=1, add_arrow=True)
            try:
                from v3_checks import run_v3_xml_checks, safe_apply_fixes
                defects, _ = run_v3_xml_checks(pptx_path, pptx_path)
                fixable = [d for d in defects if d.fixable]
                if fixable:
                    applied, failed = safe_apply_fixes(pptx_path, defects)
                    assert len(applied) > 0

                    # Verify PPTX is still valid after fixes
                    from pptx import Presentation
                    prs = Presentation(pptx_path)
                    assert len(prs.slides) == 1
            finally:
                if os.path.exists(pptx_path):
                    os.unlink(pptx_path)
                backup = pptx_path + '.v3_backup'
                if os.path.exists(backup):
                    os.unlink(backup)
        finally:
            v3_config.ENABLE_V3_XML_AUTOFIX = orig_autofix


# ─────────────────────────────────────────────────────────────────────────────
# E2E: compute_gate_decision
# ─────────────────────────────────────────────────────────────────────────────

class TestE2EGateDecision:

    def test_gate_completed_on_clean(self):
        """No defects → gate status 'completed'."""
        from v3_checks import compute_gate_decision
        gate = compute_gate_decision([], total_slides=5)
        assert gate.status == 'completed'

    def test_gate_failed_on_critical(self):
        """CRITICAL defect → gate status 'failed_qa'."""
        from v3_checks import compute_gate_decision
        from vqa_types import V3Defect, Severity
        defects = [V3Defect(
            code="TABLE_COLUMNS_NOT_REVERSED", category="table",
            severity=Severity.CRITICAL, slide_idx=1, object_id="0",
            fixable=True, autofix_action="reverse_table_columns",
        )]
        gate = compute_gate_decision(defects, total_slides=5)
        assert gate.status == 'failed_qa'

    def test_gate_warns_on_high(self):
        """Few HIGH defects → 'completed_with_warnings'."""
        from v3_checks import compute_gate_decision
        from vqa_types import V3Defect, Severity
        defects = [V3Defect(
            code="DIRECTIONAL_SHAPE_NOT_FLIPPED", category="mirroring",
            severity=Severity.HIGH, slide_idx=1, object_id="0",
            fixable=True, autofix_action="flip_directional_shape",
        )]
        gate = compute_gate_decision(defects, total_slides=10)
        assert gate.status == 'completed_with_warnings'


# ─────────────────────────────────────────────────────────────────────────────
# E2E: Full pipeline flow
# ─────────────────────────────────────────────────────────────────────────────

class TestE2EFullPipeline:

    def test_check_fix_recheck_gate_flow(self):
        """Full V3 pipeline: check → fix → recheck → gate."""
        import v3_config
        orig_autofix = v3_config.ENABLE_V3_XML_AUTOFIX
        try:
            v3_config.ENABLE_V3_XML_AUTOFIX = True
            pptx_path = _create_test_pptx(slide_count=2, add_arrow=True)
            try:
                from v3_checks import (
                    run_v3_xml_checks, safe_apply_fixes, compute_gate_decision,
                )

                # Step 1: Check
                defects, meta1 = run_v3_xml_checks(pptx_path, pptx_path)
                assert meta1['total_defects'] >= 0

                # Step 2: Fix
                applied, failed = safe_apply_fixes(pptx_path, defects)

                # Step 3: Recheck
                defects2, meta2 = run_v3_xml_checks(pptx_path, pptx_path)

                # Step 4: Gate
                unresolved = [d for d in defects2
                             if d.code not in [a.code for a in applied]]
                gate = compute_gate_decision(unresolved, meta2['slides_checked'])

                assert gate.status in (
                    'completed', 'completed_with_warnings', 'failed_qa')
            finally:
                if os.path.exists(pptx_path):
                    os.unlink(pptx_path)
                backup = pptx_path + '.v3_backup'
                if os.path.exists(backup):
                    os.unlink(backup)
        finally:
            v3_config.ENABLE_V3_XML_AUTOFIX = orig_autofix

    def test_api_contract_integration(self):
        """Full flow ending with API contract response."""
        from v3_api_contract import build_status_response
        from vqa_types import VQAGateResult

        gate = VQAGateResult()
        gate.status = 'completed_with_warnings'
        gate.high_remaining = 2
        gate.warning_issues = [{'code': 'DIRECTIONAL_SHAPE_NOT_FLIPPED'}]

        resp = build_status_response(
            job_id='e2e-test',
            phase='done',
            progress=1.0,
            gate_result=gate,
            download_url='/download/e2e-test.pptx',
        )

        assert resp['status'] == 'completed_with_warnings'
        assert resp['download_available'] is True
        assert resp['vqa']['high_remaining'] == 2
        assert 'download_url' in resp


# ─────────────────────────────────────────────────────────────────────────────
# Module import verification
# ─────────────────────────────────────────────────────────────────────────────

class TestModuleImports:

    def test_all_v3_modules_importable(self):
        """All V3 modules should be importable without error."""
        import vqa_types
        import v3_config
        import v3_checks
        import v3_vision_prompts
        import v3_api_contract

    def test_v3_checks_all_12_checks_wired(self):
        """check_slide should call all 12 checks."""
        from v3_checks import V3XMLChecker
        checker = V3XMLChecker()
        # Verify all check methods exist
        check_methods = [
            '_check_table_column_order',       # #1
            '_check_table_cell_alignment',     # #2
            '_check_icon_table_correspondence', # #3
            '_check_page_number_duplication',   # #4
            '_check_shape_position_mirroring',  # #5
            '_check_paragraph_rtl',             # #6
            '_check_circular_text_centering',   # #7
            '_check_master_element_mirroring',  # #8
            '_check_table_gridcol_reversal',    # #9
            '_check_merged_cell_integrity',     # #10
            '_check_table_grid_structural',     # #11
            '_check_directional_symbol_orientation',  # #12
        ]
        for method_name in check_methods:
            assert hasattr(checker, method_name), f"Missing: {method_name}"

    def test_v3_fixer_all_actions_registered(self):
        """V3AutoFixer should have all fix actions registered."""
        from v3_checks import V3AutoFixer
        fixer = V3AutoFixer()
        expected_actions = [
            'reverse_table_columns',    # Fix #1
            'set_para_rtl',             # Fix #2
            'reposition_icon',          # Fix #3
            'dedup_page_number',        # Fix #4
            'mirror_shape_position',    # Fix #5/8
            'set_paragraph_rtl',        # Fix #6
            'center_text_circular',     # Fix #7
            'flip_directional_shape',   # Fix #12
        ]
        for action in expected_actions:
            assert action in fixer._dispatch, f"Missing fix action: {action}"

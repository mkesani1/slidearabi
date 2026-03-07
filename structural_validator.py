import re
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Any
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

class IssueType:
    RTL_MISSING = 'rtl_missing'              # Arabic paragraph without rtl='1'
    ALIGNMENT_WRONG = 'alignment_wrong'       # Arabic text with wrong alignment
    SHAPE_OUT_OF_BOUNDS = 'shape_out_of_bounds'  # Shape extends past slide edge
    SHAPE_OVERLAP = 'shape_overlap'           # Significant overlap between shapes
    FONT_TOO_SMALL = 'font_too_small'        # Font size below minimum
    TABLE_NOT_RTL = 'table_not_rtl'          # Table without RTL processing
    CHART_NOT_RTL = 'chart_not_rtl'          # Chart without RTL axis reversal
    PLACEHOLDER_MISMATCH = 'placeholder_mismatch'  # Placeholder not matching layout
    MASTER_NO_RTL = 'master_no_rtl'          # Master missing RTL defaults
    POSITION_OVERRIDE = 'position_override'   # Unnecessary local position override


@dataclass
class ValidationIssue:
    severity: str  # 'error', 'warning', 'info'
    slide_number: int
    shape_id: Optional[int]
    shape_name: Optional[str]
    issue_type: str  # Category of issue
    message: str
    expected_value: Any = None
    actual_value: Any = None


@dataclass
class ValidationReport:
    total_issues: int = 0
    errors: int = 0
    warnings: int = 0
    info: int = 0
    issues: List[ValidationIssue] = field(default_factory=list)
    slides_checked: int = 0
    shapes_checked: int = 0
    pass_rate: float = 100.0  # Percentage of shapes that passed all checks
    
    @property
    def passed(self) -> bool:
        """Pipeline passes if there are zero errors (warnings are acceptable)."""
        return self.errors == 0


class StructuralValidator:
    """
    Read-only validator for Phase 5. Checks the transformed presentation
    against expected RTL properties. Reports issues but NEVER modifies anything.
    
    This replaces the VQA fix loops from v1.
    """
    
    def __init__(self, presentation, resolved_presentation=None):
        """
        Args:
            presentation: The transformed python-pptx Presentation object
            resolved_presentation: Optional ResolvedPresentation from Phase 0 for comparison
        """
        self.prs = presentation
        self.resolved_prs = resolved_presentation
        
        # Dimensions for bounds checking
        self.slide_width = presentation.slide_width
        self.slide_height = presentation.slide_height
        
        # Constants
        self.ARABIC_REGEX = re.compile(r'[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF]')
        self.MIN_TITLE_PT = 14
        self.MIN_BODY_PT = 10
        self.MIN_ABS_PT = 8
        self.OVERLAP_TOLERANCE = 0.10  # 10% area overlap allowed
        self.BOUNDS_TOLERANCE = 100000 # ~1cm tolerance for bounds
        
    def validate(self) -> ValidationReport:
        """Run all validation checks. Returns comprehensive report."""
        report = ValidationReport()
        issues = []
        
        shapes_checked = 0
        
        # 1. Check Master Defaults
        issues.extend(self._check_master_rtl_defaults())
        
        # 2. Check each slide
        for i, slide in enumerate(self.prs.slides):
            slide_number = i + 1
            report.slides_checked += 1
            
            # Slide-level checks
            issues.extend(self._check_placeholder_positions(slide, slide_number))
            issues.extend(self._check_shape_overlaps(slide, slide_number))
            
            # Shape-level checks
            for shape in slide.shapes:
                shapes_checked += 1
                
                # Check bounds
                issues.extend(self._check_shape_bounds(shape, slide_number))
                
                if shape.has_text_frame:
                    issues.extend(self._check_rtl_direction(shape, slide_number))
                    issues.extend(self._check_alignment(shape, slide_number))
                    issues.extend(self._check_font_sizes(shape, slide_number))
                    
                if shape.has_table:
                    issues.extend(self._check_table_column_order(shape, slide_number))
                    
                if shape.has_chart:
                    issues.extend(self._check_chart_rtl(shape, slide_number))
        
        report.shapes_checked = shapes_checked
        report.issues = issues
        report.total_issues = len(issues)
        report.errors = sum(1 for i in issues if i.severity == 'error')
        report.warnings = sum(1 for i in issues if i.severity == 'warning')
        report.info = sum(1 for i in issues if i.severity == 'info')
        
        if shapes_checked > 0:
            # Very rough heuristic: erroring shapes vs total
            # We don't track issues per shape exactly, but we can approximate
            erroring_shapes = len(set(i.shape_id for i in issues if i.severity == 'error' and i.shape_id is not None))
            report.pass_rate = max(0.0, 100.0 * (shapes_checked - erroring_shapes) / shapes_checked)
            
        return report
        
    def _is_mostly_arabic(self, text: str) -> bool:
        """Check if text is >50% Arabic characters."""
        if not text:
            return False
        
        chars = [c for c in text if c.strip()]
        if not chars:
            return False
            
        arabic_chars = sum(1 for c in chars if self.ARABIC_REGEX.search(c))
        return (arabic_chars / len(chars)) > 0.5
        
    def _has_arabic(self, text: str) -> bool:
        """Check if text contains any Arabic characters."""
        if not text:
            return False
        return bool(self.ARABIC_REGEX.search(text))

    def _check_rtl_direction(self, shape, slide_number: int) -> List[ValidationIssue]:
        """Check that all paragraphs containing Arabic text have rtl='1' set."""
        issues = []
        
        for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
            text = paragraph.text
            if not self._has_arabic(text):
                continue
                
            # Check pPr for rtl="1"
            pPr = paragraph._pPr
            rtl_attr = pPr.get('rtl') if pPr is not None else None
            
            mostly_arabic = self._is_mostly_arabic(text)
            
            if mostly_arabic and rtl_attr != '1':
                issues.append(ValidationIssue(
                    severity='error',
                    slide_number=slide_number,
                    shape_id=shape.shape_id,
                    shape_name=shape.name,
                    issue_type=IssueType.RTL_MISSING,
                    message=f"Paragraph {p_idx+1} is mostly Arabic but missing rtl='1'",
                    expected_value='1',
                    actual_value=rtl_attr
                ))
            elif not mostly_arabic and rtl_attr != '1':
                issues.append(ValidationIssue(
                    severity='warning',
                    slide_number=slide_number,
                    shape_id=shape.shape_id,
                    shape_name=shape.name,
                    issue_type=IssueType.RTL_MISSING,
                    message=f"Paragraph {p_idx+1} contains mixed bidi text but missing rtl='1'",
                    expected_value='1',
                    actual_value=rtl_attr
                ))
                
        return issues
        
    def _check_alignment(self, shape, slide_number: int) -> List[ValidationIssue]:
        """Check text alignment is appropriate for RTL."""
        issues = []
        
        is_title = shape.is_placeholder and shape.placeholder_format.type in (1, 3) # TITLE or CENTER_TITLE
        
        for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
            text = paragraph.text
            if not self._is_mostly_arabic(text):
                continue
                
            alignment = paragraph.alignment
            
            # Left alignment is always an error for Arabic
            if alignment == PP_ALIGN.LEFT:
                issues.append(ValidationIssue(
                    severity='error',
                    slide_number=slide_number,
                    shape_id=shape.shape_id,
                    shape_name=shape.name,
                    issue_type=IssueType.ALIGNMENT_WRONG,
                    message=f"Arabic text in paragraph {p_idx+1} is left-aligned",
                    expected_value='RIGHT or JUSTIFY',
                    actual_value='LEFT'
                ))
            # Center alignment is okay for titles, warning for body text
            elif alignment == PP_ALIGN.CENTER and not is_title:
                issues.append(ValidationIssue(
                    severity='warning',
                    slide_number=slide_number,
                    shape_id=shape.shape_id,
                    shape_name=shape.name,
                    issue_type=IssueType.ALIGNMENT_WRONG,
                    message=f"Arabic body text in paragraph {p_idx+1} is center-aligned",
                    expected_value='RIGHT or JUSTIFY',
                    actual_value='CENTER'
                ))
                
        return issues
        
    def _check_placeholder_positions(self, slide, slide_number: int) -> List[ValidationIssue]:
        """Check that placeholder shapes are in expected positions."""
        issues = []
        
        # This requires the layout shape for comparison
        layout = slide.slide_layout
        
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
                
            # Find matching placeholder on layout
            layout_shape = None
            for ls in layout.shapes:
                if ls.is_placeholder and ls.placeholder_format.idx == shape.placeholder_format.idx:
                    layout_shape = ls
                    break
                    
            if layout_shape:
                # Check if slide shape has local overrides (differs from layout)
                # We allow small differences due to rendering float rounding, but warn on significant overrides
                if abs(shape.left - layout_shape.left) > 10000 or abs(shape.top - layout_shape.top) > 10000:
                    issues.append(ValidationIssue(
                        severity='warning',
                        slide_number=slide_number,
                        shape_id=shape.shape_id,
                        shape_name=shape.name,
                        issue_type=IssueType.POSITION_OVERRIDE,
                        message="Placeholder has local position overrides that differ from layout",
                        expected_value=f"x={layout_shape.left}, y={layout_shape.top}",
                        actual_value=f"x={shape.left}, y={shape.top}"
                    ))
                    
        return issues
        
    def _check_shape_bounds(self, shape, slide_number: int) -> List[ValidationIssue]:
        """Check that no shapes extend outside slide boundaries."""
        issues = []
        
        # Decorative shapes are often designed to bleed off the edge
        if shape.shape_type in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM):
            return issues
            
        left_bound = -self.BOUNDS_TOLERANCE
        top_bound = -self.BOUNDS_TOLERANCE
        right_bound = self.slide_width + self.BOUNDS_TOLERANCE
        bottom_bound = self.slide_height + self.BOUNDS_TOLERANCE
        
        if shape.left < left_bound or shape.top < top_bound or \
           (shape.left + shape.width) > right_bound or \
           (shape.top + shape.height) > bottom_bound:
           
            # Only warn for non-text shapes, error for text/tables/charts
            severity = 'error' if shape.has_text_frame or shape.has_table or shape.has_chart else 'warning'
            
            issues.append(ValidationIssue(
                severity=severity,
                slide_number=slide_number,
                shape_id=shape.shape_id,
                shape_name=shape.name,
                issue_type=IssueType.SHAPE_OUT_OF_BOUNDS,
                message="Shape extends outside slide boundaries",
                actual_value=f"x={shape.left}, w={shape.width}"
            ))
            
        return issues
        
    def _get_intersection_area(self, r1, r2) -> float:
        """Calculate intersection area of two rectangles (left, top, right, bottom)."""
        x_left = max(r1[0], r2[0])
        y_top = max(r1[1], r2[1])
        x_right = min(r1[2], r2[2])
        y_bottom = min(r1[3], r2[3])

        if x_right < x_left or y_bottom < y_top:
            return 0.0

        return (x_right - x_left) * (y_bottom - y_top)

    def _check_shape_overlaps(self, slide, slide_number: int) -> List[ValidationIssue]:
        """Check for significant shape overlaps that indicate mirroring errors."""
        issues = []
        
        # Get bounding boxes for all non-decorative shapes
        shapes_data = []
        for shape in slide.shapes:
            # Skip lines and background/decorative elements
            if shape.shape_type in (MSO_SHAPE_TYPE.LINE, 1000):
                continue
                
            # Rect: (left, top, right, bottom, area, shape)
            rect = (
                shape.left, 
                shape.top, 
                shape.left + shape.width, 
                shape.top + shape.height,
                shape.width * shape.height,
                shape
            )
            # Only consider shapes with actual area
            if rect[4] > 0:
                shapes_data.append(rect)
                
        # Compare all pairs
        for i in range(len(shapes_data)):
            for j in range(i + 1, len(shapes_data)):
                s1_data = shapes_data[i]
                s2_data = shapes_data[j]
                
                # Check for z-order containment (e.g. text on top of a colored box)
                # If one shape is completely inside another, it's likely intentional
                intersection = self._get_intersection_area(s1_data[:4], s2_data[:4])
                
                if intersection > 0:
                    # Calculate overlap percentage relative to the smaller shape
                    min_area = min(s1_data[4], s2_data[4])
                    overlap_pct = intersection / min_area
                    
                    # If overlap is between 10% and 90%, it's likely a collision
                    # >90% is usually intentional layering (text over background box)
                    if self.OVERLAP_TOLERANCE < overlap_pct < 0.90:
                        s1, s2 = s1_data[5], s2_data[5]
                        
                        # Only error if both are content shapes (text, table, chart, picture)
                        s1_content = s1.has_text_frame or s1.has_table or s1.has_chart or s1.shape_type == MSO_SHAPE_TYPE.PICTURE
                        s2_content = s2.has_text_frame or s2.has_table or s2.has_chart or s2.shape_type == MSO_SHAPE_TYPE.PICTURE
                        
                        severity = 'error' if (s1_content and s2_content) else 'warning'
                        
                        issues.append(ValidationIssue(
                            severity=severity,
                            slide_number=slide_number,
                            shape_id=s1.shape_id,
                            shape_name=s1.name,
                            issue_type=IssueType.SHAPE_OVERLAP,
                            message=f"Shape overlaps significantly ({overlap_pct:.0%}) with shape '{s2.name}' (id:{s2.shape_id})"
                        ))
                        
        return issues
        
    def _check_table_column_order(self, shape, slide_number: int) -> List[ValidationIssue]:
        """Check that tables have been processed for RTL."""
        issues = []
        
        # Look for the rightToLeft attribute on the tblPr element
        tbl = shape.table._tbl
        tblPr = tbl.tblPr
        
        if tblPr is not None:
            rtl = tblPr.get('rtl')
            if rtl != '1' and rtl != 'true':
                issues.append(ValidationIssue(
                    severity='error',
                    slide_number=slide_number,
                    shape_id=shape.shape_id,
                    shape_name=shape.name,
                    issue_type=IssueType.TABLE_NOT_RTL,
                    message="Table is missing right-to-left band direction",
                    expected_value='1',
                    actual_value=rtl
                ))
        else:
            issues.append(ValidationIssue(
                severity='error',
                slide_number=slide_number,
                shape_id=shape.shape_id,
                shape_name=shape.name,
                issue_type=IssueType.TABLE_NOT_RTL,
                message="Table properties (tblPr) missing, RTL cannot be set"
            ))
            
        return issues
        
    def _check_chart_rtl(self, shape, slide_number: int) -> List[ValidationIssue]:
        """Check that charts have been transformed for RTL."""
        issues = []
        
        # A full implementation would check orientation attributes in chart XML
        # For this skeleton, we assume it checks axis reversal
        
        try:
            chart = shape.chart
            # If there's a category axis, it should be reversed for RTL
            if chart.has_category_axis:
                axis = chart.category_axis
                # To truly check, we'd need to inspect the c:scaling/c:orientation element
                # which isn't directly exposed as 'reverse_order' in older python-pptx
                # We'd use axis._element.xpath('.//c:orientation/@val')
                orientations = axis._element.xpath('.//c:scaling/c:orientation/@val')
                if orientations and orientations[0] != 'maxMin':
                    issues.append(ValidationIssue(
                        severity='error',
                        slide_number=slide_number,
                        shape_id=shape.shape_id,
                        shape_name=shape.name,
                        issue_type=IssueType.CHART_NOT_RTL,
                        message="Chart category axis is not reversed for RTL",
                        expected_value='maxMin',
                        actual_value=orientations[0]
                    ))
        except Exception as e:
            # Ignore charts that don't have standard axes (pie charts, etc)
            pass
            
        return issues
        
    def _check_font_sizes(self, shape, slide_number: int) -> List[ValidationIssue]:
        """Check that font sizes are within acceptable bounds."""
        issues = []
        
        is_title = shape.is_placeholder and shape.placeholder_format.type in (1, 3)
        min_pt = self.MIN_TITLE_PT if is_title else self.MIN_BODY_PT
        
        # We only check explicitly set sizes here. A full check would use ResolvedPresentation
        # to get inherited sizes.
        for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
            for r_idx, run in enumerate(paragraph.runs):
                if run.font.size is not None:
                    size_pt = run.font.size.pt
                    
                    if size_pt < self.MIN_ABS_PT:
                        issues.append(ValidationIssue(
                            severity='error',
                            slide_number=slide_number,
                            shape_id=shape.shape_id,
                            shape_name=shape.name,
                            issue_type=IssueType.FONT_TOO_SMALL,
                            message=f"Font size {size_pt}pt in paragraph {p_idx+1} is below absolute minimum ({self.MIN_ABS_PT}pt)",
                            expected_value=f">={self.MIN_ABS_PT}",
                            actual_value=size_pt
                        ))
                    elif size_pt < min_pt:
                        issues.append(ValidationIssue(
                            severity='warning',
                            slide_number=slide_number,
                            shape_id=shape.shape_id,
                            shape_name=shape.name,
                            issue_type=IssueType.FONT_TOO_SMALL,
                            message=f"Font size {size_pt}pt is smaller than recommended minimum for this shape type ({min_pt}pt)",
                            expected_value=f">={min_pt}",
                            actual_value=size_pt
                        ))
                        
        return issues
        
    def _check_master_rtl_defaults(self) -> List[ValidationIssue]:
        """Check that slide masters have RTL defaults set."""
        issues = []
        
        for i, master in enumerate(self.prs.slide_masters):
            # Check txStyles (titleStyle, bodyStyle, otherStyle)
            txStyles = master._element.xpath('./p:txStyles')
            if not txStyles:
                continue
                
            for style_name in ['titleStyle', 'bodyStyle', 'otherStyle']:
                styles = txStyles[0].xpath(f'./p:{style_name}')
                if not styles:
                    continue
                    
                # Check lvl1pPr RTL setting
                lvl1pPr = styles[0].xpath('./a:lvl1pPr')
                if lvl1pPr:
                    rtl = lvl1pPr[0].get('rtl')
                    if rtl != '1':
                        issues.append(ValidationIssue(
                            severity='warning',
                            slide_number=0, # Master slide
                            shape_id=None,
                            shape_name=f"Master {i+1}",
                            issue_type=IssueType.MASTER_NO_RTL,
                            message=f"Master {style_name} level 1 is missing default rtl='1'",
                            expected_value='1',
                            actual_value=rtl
                        ))
                        
        return issues

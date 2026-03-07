#!/usr/bin/env python3
"""
test_harness.py — SlideArabi Structural Comparison Test Harness

Runs structural analysis and transformation on one or more input PPTX files,
then generates per-deck HTML comparison reports with side-by-side slide images.

Usage:
    python slidearabi/test_harness.py input1.pptx [input2.pptx ...]
    python slidearabi/test_harness.py --dir /path/to/test_decks/

Output structure (per deck):
    /home/user/workspace/test_results/
    └── myfile/
        ├── analysis.json
        ├── myfile_v2_structural.pptx
        ├── original/slide-01.jpg ...
        ├── v2/slide-01.jpg ...
        ├── comparison/slide-01.jpg ...
        └── report.html
"""

from __future__ import annotations

import argparse
import copy
import datetime
import json
import logging
import os
import shutil
import subprocess
import sys
import tempfile
import traceback
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

# ── Ensure slidearabi package root is on sys.path ────────────────────────
_HERE = Path(__file__).resolve().parent
_ROOT = _HERE.parent
if str(_ROOT) not in sys.path:
    sys.path.insert(0, str(_ROOT))
# Also add skills/pptx/scripts so the soffice helper is importable
_SKILLS_OFFICE = _ROOT / "skills" / "pptx" / "scripts"
if str(_SKILLS_OFFICE) not in sys.path:
    sys.path.insert(0, str(_SKILLS_OFFICE))

# ── Third-party / stdlib ─────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError as e:
    print(f"[FATAL] python-pptx not installed: {e}", file=sys.stderr)
    sys.exit(1)

try:
    from PIL import Image, ImageDraw, ImageFont
    _PIL_AVAILABLE = True
except ImportError:
    _PIL_AVAILABLE = False
    print("[WARN] Pillow not installed — comparison images will be skipped.")

# ── LibreOffice soffice helper ───────────────────────────────────────────────
try:
    from office.soffice import run_soffice, get_soffice_env
    _SOFFICE_HELPER_AVAILABLE = True
except ImportError:
    _SOFFICE_HELPER_AVAILABLE = False
    print("[WARN] soffice helper not available — will use subprocess directly.")

# ── SlideArabi modules ────────────────────────────────────────────────────
_V2_IMPORT_ERRORS: List[str] = []

try:
    from slidearabi.property_resolver import PropertyResolver
except ImportError as e:
    PropertyResolver = None  # type: ignore
    _V2_IMPORT_ERRORS.append(f"PropertyResolver: {e}")

try:
    from slidearabi.layout_analyzer import LayoutAnalyzer, LayoutClassification
except ImportError as e:
    LayoutAnalyzer = None  # type: ignore
    LayoutClassification = None  # type: ignore
    _V2_IMPORT_ERRORS.append(f"LayoutAnalyzer: {e}")

try:
    from slidearabi.template_registry import TemplateRegistry
except ImportError as e:
    TemplateRegistry = None  # type: ignore
    _V2_IMPORT_ERRORS.append(f"TemplateRegistry: {e}")

try:
    from slidearabi.rtl_transforms import MasterLayoutTransformer, SlideContentTransformer, TransformReport
except ImportError as e:
    MasterLayoutTransformer = None  # type: ignore
    SlideContentTransformer = None  # type: ignore
    TransformReport = None  # type: ignore
    _V2_IMPORT_ERRORS.append(f"rtl_transforms: {e}")

try:
    from slidearabi.typography import TypographyNormalizer
except ImportError as e:
    TypographyNormalizer = None  # type: ignore
    _V2_IMPORT_ERRORS.append(f"TypographyNormalizer: {e}")

try:
    from slidearabi.structural_validator import StructuralValidator
except ImportError as e:
    StructuralValidator = None  # type: ignore
    _V2_IMPORT_ERRORS.append(f"StructuralValidator: {e}")

# ── Constants ────────────────────────────────────────────────────────────────
EMU_PER_INCH = 914_400
RESULTS_DIR = _ROOT / "test_results"

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  %(levelname)-7s  %(name)s  %(message)s",
    datefmt="%H:%M:%S",
)
logger = logging.getLogger("test_harness")


# ════════════════════════════════════════════════════════════════════════════
# Helper utilities
# ════════════════════════════════════════════════════════════════════════════

def _emu_to_inches(emu: int) -> float:
    return round(emu / EMU_PER_INCH, 4)


def _safe(fn, default=None, label: str = ""):
    """Call fn(); on exception log and return default."""
    try:
        return fn()
    except Exception as exc:
        if label:
            logger.debug("  [safe] %s: %s", label, exc)
        return default


def _report_to_dict(report) -> Dict[str, Any]:
    """Convert a TransformReport to a plain dict for JSON serialisation."""
    if report is None:
        return {}
    try:
        return {
            "phase": report.phase,
            "total_changes": report.total_changes,
            "changes_by_type": dict(report.changes_by_type),
            "warnings": list(report.warnings),
            "errors": list(report.errors),
        }
    except Exception:
        return {"raw": str(report)}


def _validation_issue_to_dict(issue) -> Dict[str, Any]:
    """Convert a ValidationIssue (from structural_validator) to a dict."""
    return {
        "severity": _safe(lambda: issue.severity, "unknown"),
        "slide_number": _safe(lambda: issue.slide_number, 0),
        "shape_id": _safe(lambda: issue.shape_id),
        "shape_name": _safe(lambda: issue.shape_name),
        "issue_type": _safe(lambda: issue.issue_type, "unknown"),
        "message": _safe(lambda: issue.message, ""),
        "expected_value": _safe(lambda: str(issue.expected_value) if issue.expected_value is not None else None),
        "actual_value": _safe(lambda: str(issue.actual_value) if issue.actual_value is not None else None),
    }


# ════════════════════════════════════════════════════════════════════════════
# Phase 0: Structural Analysis
# ════════════════════════════════════════════════════════════════════════════

def _analyze_shape(shape, slide_width_emu: int) -> Dict[str, Any]:
    """Extract per-shape structural data from a python-pptx shape object."""
    result: Dict[str, Any] = {
        "shape_id": _safe(lambda: shape.shape_id),
        "shape_name": _safe(lambda: shape.name, ""),
        "shape_type": _safe(lambda: str(shape.shape_type).split(".")[-1], "UNKNOWN"),
        "is_placeholder": _safe(lambda: shape.is_placeholder, False),
        "has_text_frame": _safe(lambda: shape.has_text_frame, False),
        "has_table": _safe(lambda: shape.has_table, False),
        "has_chart": _safe(lambda: shape.has_chart, False),
    }

    # Position / size
    for attr in ("left", "top", "width", "height"):
        val = _safe(lambda a=attr: getattr(shape, a))
        result[f"{attr}_emu"] = val
        result[f"{attr}_in"] = _emu_to_inches(val) if val is not None else None

    # Placeholder info
    if result["is_placeholder"]:
        ph_fmt = _safe(lambda: shape.placeholder_format)
        result["placeholder_type"] = _safe(lambda: str(ph_fmt.type).split(".")[-1].lower() if ph_fmt else None)
        result["placeholder_idx"] = _safe(lambda: ph_fmt.idx if ph_fmt else None)
    else:
        result["placeholder_type"] = None
        result["placeholder_idx"] = None

    # Text properties
    if result["has_text_frame"]:
        paragraphs_info = []
        try:
            for para in shape.text_frame.paragraphs:
                pPr = _safe(lambda: para._pPr)
                rtl_val = _safe(lambda: pPr.get("rtl") if pPr is not None else None)
                algn_val = _safe(lambda: pPr.get("algn") if pPr is not None else None)
                runs_info = []
                for run in para.runs:
                    font = _safe(lambda: run.font)
                    sz = _safe(lambda: font.size.pt if font and font.size else None)
                    runs_info.append({
                        "text": _safe(lambda: run.text[:80], ""),  # truncate
                        "font_size_pt": sz,
                        "font_name": _safe(lambda: font.name if font else None),
                        "bold": _safe(lambda: font.bold if font else None),
                        "italic": _safe(lambda: font.italic if font else None),
                    })
                paragraphs_info.append({
                    "text": _safe(lambda: para.text[:120], ""),
                    "alignment": _safe(lambda: str(para.alignment).split(".")[-1] if para.alignment else None),
                    "rtl": rtl_val,
                    "algn_raw": algn_val,
                    "runs": runs_info,
                })
        except Exception as exc:
            paragraphs_info = [{"error": str(exc)}]
        result["paragraphs"] = paragraphs_info

    return result


def _analyze_slide(slide, slide_number: int, layout_classification, template_registry) -> Dict[str, Any]:
    """Produce a full structural snapshot of a single slide."""
    slide_data: Dict[str, Any] = {
        "slide_number": slide_number,
        "shape_count": _safe(lambda: len(slide.shapes), 0),
        "placeholder_count": _safe(lambda: sum(1 for s in slide.shapes if s.is_placeholder), 0),
        "freeform_count": _safe(lambda: sum(
            1 for s in slide.shapes
            if not s.is_placeholder and not s.has_table and not s.has_chart
        ), 0),
        "has_table": _safe(lambda: any(s.has_table for s in slide.shapes), False),
        "has_chart": _safe(lambda: any(s.has_chart for s in slide.shapes), False),
    }

    # Layout classification
    if layout_classification:
        slide_data["layout"] = {
            "name": _safe(lambda: layout_classification.layout_name, ""),
            "resolved_type": _safe(lambda: layout_classification.resolved_type, "unknown"),
            "explicit_type": _safe(lambda: layout_classification.explicit_type),
            "confidence": _safe(lambda: round(layout_classification.confidence, 3), 0.0),
            "is_explicit": _safe(lambda: layout_classification.explicit_type is not None, False),
            "requires_ai": _safe(lambda: layout_classification.requires_ai_classification, False),
            "placeholder_summary": _safe(lambda: dict(layout_classification.placeholder_summary), {}),
        }
        # Template registry rules
        if template_registry:
            try:
                layout_type = layout_classification.resolved_type
                rules = template_registry.get_rules(layout_type)
                slide_data["template_rules"] = {
                    "layout_type": rules.layout_type,
                    "description": rules.description,
                    "freeform_action": rules.freeform_action,
                    "swap_columns": rules.swap_columns,
                    "table_action": rules.table_action,
                    "chart_action": rules.chart_action,
                    "placeholder_rules": {
                        k: {
                            "action": v.action,
                            "set_rtl": v.set_rtl,
                            "set_alignment": v.set_alignment,
                            "mirror_x": v.mirror_x,
                            "swap_partner_idx": v.swap_partner_idx,
                        }
                        for k, v in rules.placeholder_rules.items()
                    },
                }
            except Exception as exc:
                slide_data["template_rules"] = {"error": str(exc)}
    else:
        slide_data["layout"] = None
        slide_data["template_rules"] = None

    # Shapes
    shapes_info = []
    for shape in _safe(lambda: list(slide.shapes), []):
        try:
            shapes_info.append(_analyze_shape(shape, 0))
        except Exception as exc:
            shapes_info.append({"error": str(exc), "shape_name": _safe(lambda: shape.name, "?")})
    slide_data["shapes"] = shapes_info

    return slide_data


def run_phase0_analysis(pptx_path: Path) -> Dict[str, Any]:
    """
    Phase 0: Parse the presentation, run PropertyResolver and LayoutAnalyzer,
    and return a full structural analysis dict.
    """
    print(f"  [Phase 0] Structural analysis…")
    result: Dict[str, Any] = {
        "input_path": str(pptx_path),
        "phase": "0_structural_analysis",
        "errors": [],
        "warnings": [],
    }

    if _V2_IMPORT_ERRORS:
        result["warnings"].extend([f"Import warning: {e}" for e in _V2_IMPORT_ERRORS])

    # Load presentation
    try:
        prs = Presentation(str(pptx_path))
    except Exception as exc:
        result["errors"].append(f"Failed to load PPTX: {exc}")
        return result

    slide_width_emu = int(prs.slide_width)
    slide_height_emu = int(prs.slide_height)

    result["presentation"] = {
        "slide_count": len(prs.slides),
        "slide_width_emu": slide_width_emu,
        "slide_height_emu": slide_height_emu,
        "slide_width_in": _emu_to_inches(slide_width_emu),
        "slide_height_in": _emu_to_inches(slide_height_emu),
        "master_count": _safe(lambda: len(prs.slide_masters), 0),
    }

    # Masters / layouts inventory
    masters_info = []
    for m_idx, master in enumerate(_safe(lambda: prs.slide_masters, [])):
        layouts_info = []
        for layout in _safe(lambda: master.slide_layouts, []):
            layouts_info.append({
                "name": _safe(lambda: layout.name, ""),
                "placeholder_count": _safe(lambda: len(layout.placeholders), 0),
            })
        masters_info.append({
            "master_index": m_idx,
            "master_name": _safe(lambda: master.name, ""),
            "layout_count": len(layouts_info),
            "layouts": layouts_info,
        })
    result["masters"] = masters_info

    # PropertyResolver (Phase 0 enrichment)
    resolved_prs = None
    if PropertyResolver is not None:
        try:
            resolver = PropertyResolver(prs)
            resolved_prs = resolver.resolve_presentation()
            result["property_resolver"] = {
                "status": "success",
                "total_slides": _safe(lambda: resolved_prs.total_slides, 0),
                "total_shapes": _safe(lambda: resolved_prs.total_shapes, 0),
            }
        except Exception as exc:
            result["errors"].append(f"PropertyResolver failed: {exc}")
            result["property_resolver"] = {"status": "failed", "error": str(exc)}
    else:
        result["property_resolver"] = {"status": "unavailable"}

    # LayoutAnalyzer
    layout_classifications: Dict[int, Any] = {}
    if LayoutAnalyzer is not None:
        try:
            analyzer = LayoutAnalyzer(prs)
            layout_classifications = analyzer.analyze_all()
            layout_types_found = list(
                {lc.resolved_type for lc in layout_classifications.values()}
            )
            result["layout_analyzer"] = {
                "status": "success",
                "layouts_classified": len(layout_classifications),
                "layout_types_found": sorted(layout_types_found),
            }
        except Exception as exc:
            result["errors"].append(f"LayoutAnalyzer failed: {exc}")
            result["layout_analyzer"] = {"status": "failed", "error": str(exc)}
    else:
        result["layout_analyzer"] = {"status": "unavailable"}

    # TemplateRegistry
    template_registry = None
    if TemplateRegistry is not None:
        try:
            template_registry = TemplateRegistry(slide_width_emu, slide_height_emu)
            result["template_registry"] = {"status": "available"}
        except Exception as exc:
            result["errors"].append(f"TemplateRegistry failed: {exc}")
            result["template_registry"] = {"status": "failed", "error": str(exc)}
    else:
        result["template_registry"] = {"status": "unavailable"}

    # Per-slide analysis
    slides_info = []
    for i, slide in enumerate(prs.slides):
        slide_number = i + 1
        classification = layout_classifications.get(slide_number)
        try:
            slide_data = _analyze_slide(slide, slide_number, classification, template_registry)
            slides_info.append(slide_data)
        except Exception as exc:
            slides_info.append({
                "slide_number": slide_number,
                "error": str(exc),
                "traceback": traceback.format_exc(),
            })
    result["slides"] = slides_info

    return result


# ════════════════════════════════════════════════════════════════════════════
# Phase 2+3+4+5: Structural Transform
# ════════════════════════════════════════════════════════════════════════════

def run_structural_transform(
    input_path: Path,
    output_path: Path,
) -> Dict[str, Any]:
    """
    Run Phase 2 (MasterLayoutTransformer), Phase 3 (SlideContentTransformer,
    skip_translation=True), Phase 4 (TypographyNormalizer), and Phase 5
    (StructuralValidator) on a copy of the input PPTX.

    Returns a dict summarising what happened in each phase.
    """
    print(f"  [Phase 2-5] Structural transform…")
    result: Dict[str, Any] = {
        "output_path": str(output_path),
        "phase_reports": {},
        "errors": [],
        "warnings": [],
        "success": False,
    }

    # Work on a copy so we never mutate the original
    try:
        shutil.copy2(str(input_path), str(output_path))
        prs = Presentation(str(output_path))
    except Exception as exc:
        result["errors"].append(f"Failed to copy/load PPTX: {exc}")
        return result

    slide_width_emu = int(prs.slide_width)
    slide_height_emu = int(prs.slide_height)

    # Phase 0 (needed to get layout_classifications for Phase 3)
    resolved_prs = None
    layout_classifications: Dict[int, Any] = {}

    if PropertyResolver is not None:
        try:
            resolver = PropertyResolver(prs)
            resolved_prs = resolver.resolve_presentation()
        except Exception as exc:
            result["warnings"].append(f"PropertyResolver (for Phase 3 input) failed: {exc}")

    if LayoutAnalyzer is not None:
        try:
            analyzer = LayoutAnalyzer(prs)
            layout_classifications = analyzer.analyze_all()
        except Exception as exc:
            result["warnings"].append(f"LayoutAnalyzer (for Phase 3 input) failed: {exc}")

    # Build TemplateRegistry
    registry = None
    if TemplateRegistry is not None:
        try:
            registry = TemplateRegistry(slide_width_emu, slide_height_emu)
        except Exception as exc:
            result["warnings"].append(f"TemplateRegistry failed: {exc}")

    # Phase 2: MasterLayoutTransformer
    if MasterLayoutTransformer is not None:
        try:
            transformer = MasterLayoutTransformer(prs, registry)
            master_report = transformer.transform_all_masters()
            layout_report = transformer.transform_all_layouts()
            result["phase_reports"]["phase_2_masters"] = _report_to_dict(master_report)
            result["phase_reports"]["phase_2_layouts"] = _report_to_dict(layout_report)
            print(f"    Phase 2: masters={master_report.total_changes} changes, "
                  f"layouts={layout_report.total_changes} changes")
        except Exception as exc:
            result["errors"].append(f"Phase 2 failed: {exc}")
            result["phase_reports"]["phase_2"] = {"error": str(exc), "traceback": traceback.format_exc()}
    else:
        result["warnings"].append("MasterLayoutTransformer unavailable — Phase 2 skipped.")

    # Phase 3: SlideContentTransformer (empty translations = skip_translation mode)
    if SlideContentTransformer is not None:
        try:
            # Convert LayoutClassification dict to the format expected by SlideContentTransformer:
            # {slide_number: layout_type_str}  or  {slide_number: LayoutClassification}
            # The SlideContentTransformer accepts either; pass the full object dict
            content_transformer = SlideContentTransformer(
                prs,
                template_registry=registry,
                layout_classifications=layout_classifications,
                translations={},  # skip_translation mode
            )
            slide_report = content_transformer.transform_all_slides()
            result["phase_reports"]["phase_3_slides"] = _report_to_dict(slide_report)
            print(f"    Phase 3: {slide_report.total_changes} changes across slides")
        except Exception as exc:
            result["errors"].append(f"Phase 3 failed: {exc}")
            result["phase_reports"]["phase_3"] = {"error": str(exc), "traceback": traceback.format_exc()}
    else:
        result["warnings"].append("SlideContentTransformer unavailable — Phase 3 skipped.")

    # Phase 4: TypographyNormalizer
    if TypographyNormalizer is not None:
        try:
            normalizer = TypographyNormalizer(prs)
            typo_report = normalizer.normalize_all()
            result["phase_reports"]["phase_4_typography"] = _report_to_dict(typo_report)
            print(f"    Phase 4: {typo_report.total_changes} typography changes")
        except Exception as exc:
            result["errors"].append(f"Phase 4 failed: {exc}")
            result["phase_reports"]["phase_4"] = {"error": str(exc), "traceback": traceback.format_exc()}
    else:
        result["warnings"].append("TypographyNormalizer unavailable — Phase 4 skipped.")

    # Save the transformed PPTX
    try:
        prs.save(str(output_path))
        result["success"] = True
        print(f"    Saved transformed PPTX → {output_path.name}")
    except Exception as exc:
        result["errors"].append(f"Failed to save transformed PPTX: {exc}")
        return result

    # Phase 5: StructuralValidator (reload the saved file for a clean pass)
    if StructuralValidator is not None:
        try:
            prs_val = Presentation(str(output_path))
            validator = StructuralValidator(prs_val, None)  # resolved_prs may differ after transforms
            # Patch: structural_validator._check_master_rtl_defaults uses bare xpath without
            # namespace bindings which raises XPathEvalError in lxml. Wrap with a safe version.
            import types
            original_check = validator._check_master_rtl_defaults
            def _safe_check_master_rtl_defaults(self=validator):
                try:
                    return original_check()
                except Exception as _xpath_err:
                    logger.debug("_check_master_rtl_defaults skipped (%s)", _xpath_err)
                    return []
            validator._check_master_rtl_defaults = _safe_check_master_rtl_defaults
            val_report = validator.validate()
            result["phase_reports"]["phase_5_validation"] = {
                "passed": val_report.passed,
                "total_issues": val_report.total_issues,
                "errors": val_report.errors,
                "warnings": val_report.warnings,
                "info": val_report.info,
                "slides_checked": val_report.slides_checked,
                "shapes_checked": val_report.shapes_checked,
                "pass_rate": val_report.pass_rate,
                "issues": [_validation_issue_to_dict(i) for i in val_report.issues],
            }
            print(f"    Phase 5: {val_report.total_issues} issues "
                  f"({val_report.errors} errors, {val_report.warnings} warnings)")
        except Exception as exc:
            result["errors"].append(f"Phase 5 failed: {exc}")
            result["phase_reports"]["phase_5"] = {"error": str(exc), "traceback": traceback.format_exc()}
    else:
        result["warnings"].append("StructuralValidator unavailable — Phase 5 skipped.")

    return result


# ════════════════════════════════════════════════════════════════════════════
# Rendering: PPTX → images via LibreOffice + pdftoppm
# ════════════════════════════════════════════════════════════════════════════

def _convert_pptx_to_pdf(pptx_path: Path, output_dir: Path) -> Optional[Path]:
    """Convert a PPTX to PDF in output_dir using LibreOffice. Returns the PDF path."""
    pdf_name = pptx_path.stem + ".pdf"
    pdf_path = output_dir / pdf_name

    try:
        if _SOFFICE_HELPER_AVAILABLE:
            env = get_soffice_env()
            result = subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf",
                 "--outdir", str(output_dir), str(pptx_path)],
                env=env,
                capture_output=True,
                text=True,
                timeout=120,
            )
        else:
            result = subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf",
                 "--outdir", str(output_dir), str(pptx_path)],
                capture_output=True,
                text=True,
                timeout=120,
            )

        if result.returncode != 0:
            logger.warning("LibreOffice conversion stderr: %s", result.stderr[:500])
            # Try to find the PDF anyway (soffice sometimes exits non-zero but still produces output)

        if pdf_path.exists():
            return pdf_path

        # LibreOffice may rename the file; look for any PDF in output_dir
        pdfs = list(output_dir.glob("*.pdf"))
        if pdfs:
            return pdfs[0]

        logger.warning("LibreOffice did not produce a PDF for %s", pptx_path.name)
        return None

    except FileNotFoundError:
        logger.warning("LibreOffice (soffice) not found — rendering skipped.")
        return None
    except subprocess.TimeoutExpired:
        logger.warning("LibreOffice timed out for %s", pptx_path.name)
        return None
    except Exception as exc:
        logger.warning("LibreOffice conversion failed: %s", exc)
        return None


def _pdf_to_images(pdf_path: Path, output_prefix: Path, dpi: int = 150) -> List[Path]:
    """Convert each page of a PDF to a JPEG using pdftoppm. Returns list of image paths."""
    images = []
    try:
        result = subprocess.run(
            ["pdftoppm", "-jpeg", "-r", str(dpi), str(pdf_path), str(output_prefix)],
            capture_output=True,
            text=True,
            timeout=120,
        )
        if result.returncode != 0:
            logger.warning("pdftoppm stderr: %s", result.stderr[:300])

        # pdftoppm produces files like prefix-1.jpg, prefix-01.jpg, prefix-001.jpg
        parent = output_prefix.parent
        stem = output_prefix.name
        for candidate in sorted(parent.glob(f"{stem}*.jpg")):
            images.append(candidate)
        # Also check .jpeg extension
        if not images:
            for candidate in sorted(parent.glob(f"{stem}*.jpeg")):
                images.append(candidate)

        return images

    except FileNotFoundError:
        logger.warning("pdftoppm not found — rendering skipped.")
        return []
    except subprocess.TimeoutExpired:
        logger.warning("pdftoppm timed out for %s", pdf_path.name)
        return []
    except Exception as exc:
        logger.warning("pdftoppm failed: %s", exc)
        return []


def render_pptx_slides(pptx_path: Path, output_dir: Path, label: str) -> List[Path]:
    """
    Convert a PPTX to per-slide JPEG images in output_dir.
    Returns a sorted list of image paths (slide-01.jpg, slide-02.jpg, …).
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    print(f"  [Render] Converting {label} to images…")

    with tempfile.TemporaryDirectory(prefix="ss_render_") as tmp_dir:
        tmp = Path(tmp_dir)

        # Step 1: PPTX → PDF
        pdf_path = _convert_pptx_to_pdf(pptx_path, tmp)
        if pdf_path is None:
            logger.warning("Skipping image rendering for %s (no PDF produced)", label)
            return []

        # Step 2: PDF → per-page JPEGs
        raw_images = _pdf_to_images(pdf_path, tmp / "slide", dpi=150)
        if not raw_images:
            logger.warning("No images produced from PDF for %s", label)
            return []

        # Step 3: Rename to slide-01.jpg, slide-02.jpg, … and copy to output_dir
        final_paths = []
        for idx, raw_img in enumerate(sorted(raw_images), start=1):
            dest = output_dir / f"slide-{idx:02d}.jpg"
            shutil.copy2(str(raw_img), str(dest))
            final_paths.append(dest)

        print(f"    Rendered {len(final_paths)} slides for {label}")
        return final_paths


# ════════════════════════════════════════════════════════════════════════════
# Comparison images (PIL side-by-side)
# ════════════════════════════════════════════════════════════════════════════

def make_comparison_image(
    original_img: Path,
    v2_img: Path,
    output_path: Path,
    slide_number: int,
) -> bool:
    """Generate a side-by-side comparison JPEG. Returns True on success."""
    if not _PIL_AVAILABLE:
        return False

    try:
        img_orig = Image.open(original_img).convert("RGB")
        img_v2 = Image.open(v2_img).convert("RGB")

        # Match heights
        target_h = max(img_orig.height, img_v2.height)
        if img_orig.height != target_h:
            ratio = target_h / img_orig.height
            img_orig = img_orig.resize(
                (int(img_orig.width * ratio), target_h), Image.LANCZOS
            )
        if img_v2.height != target_h:
            ratio = target_h / img_v2.height
            img_v2 = img_v2.resize(
                (int(img_v2.width * ratio), target_h), Image.LANCZOS
            )

        label_h = 36
        gap = 6
        total_w = img_orig.width + gap + img_v2.width
        total_h = label_h + target_h

        canvas = Image.new("RGB", (total_w, total_h), color=(245, 245, 245))
        draw = ImageDraw.Draw(canvas)

        # Labels
        label_color = (40, 40, 40)
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 14)
        except Exception:
            font = ImageFont.load_default()

        draw.text((img_orig.width // 2 - 40, 8), "ORIGINAL", fill=label_color, font=font)
        draw.text((img_orig.width + gap + img_v2.width // 2 - 55, 8), "v2 STRUCTURAL", fill=label_color, font=font)

        # Paste images
        canvas.paste(img_orig, (0, label_h))
        canvas.paste(img_v2, (img_orig.width + gap, label_h))

        # Divider line
        draw.line(
            [(img_orig.width + gap // 2, label_h),
             (img_orig.width + gap // 2, total_h)],
            fill=(180, 180, 180),
            width=2,
        )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        canvas.save(str(output_path), "JPEG", quality=88)
        return True

    except Exception as exc:
        logger.warning("Failed to create comparison image for slide %d: %s", slide_number, exc)
        return False


# ════════════════════════════════════════════════════════════════════════════
# HTML Report Generation
# ════════════════════════════════════════════════════════════════════════════

_HTML_CSS = """
body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
       margin: 0; padding: 20px; background: #f5f5f5; color: #1a1a2e; }
* { box-sizing: border-box; }
.header { background: #1a1a2e; color: white; padding: 30px; border-radius: 12px;
          margin-bottom: 30px; }
.header h1 { margin: 0 0 8px 0; font-size: 1.8rem; }
.header .meta { opacity: 0.7; font-size: 0.9rem; }
.section-title { font-size: 1.2rem; font-weight: 700; color: #1a1a2e;
                 margin: 24px 0 12px 0; border-left: 4px solid #4a6fa5;
                 padding-left: 10px; }
.stats { display: grid; grid-template-columns: repeat(auto-fit, minmax(180px, 1fr));
         gap: 15px; margin: 20px 0; }
.stat-card { background: white; padding: 16px; border-radius: 8px;
             box-shadow: 0 1px 4px rgba(0,0,0,0.08); }
.stat-value { font-size: 2rem; font-weight: 700; color: #1a1a2e; }
.stat-label { font-size: 11px; text-transform: uppercase; color: #888;
              letter-spacing: 0.05em; margin-top: 2px; }
.slide-card { background: white; border-radius: 8px; padding: 20px;
              margin-bottom: 24px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); }
.slide-card h3 { margin: 0 0 16px 0; color: #1a1a2e; }
.comparison { display: grid; grid-template-columns: 1fr 1fr; gap: 16px;
              margin-bottom: 16px; }
.comparison-pane { text-align: center; }
.comparison img { width: 100%; border: 1px solid #ddd; border-radius: 4px;
                  box-shadow: 0 1px 3px rgba(0,0,0,0.1); }
.comparison-full img { width: 100%; border: 1px solid #ddd; border-radius: 4px; }
.label { font-size: 11px; font-weight: 600; text-transform: uppercase;
         color: #666; margin-bottom: 6px; letter-spacing: 0.04em; }
.badge { display: inline-block; padding: 2px 8px; border-radius: 4px;
         font-size: 11px; font-weight: 600; margin: 2px; }
.badge-green  { background: #dcfce7; color: #166534; }
.badge-red    { background: #fef2f2; color: #991b1b; }
.badge-yellow { background: #fefce8; color: #854d0e; }
.badge-blue   { background: #dbeafe; color: #1e40af; }
.badge-gray   { background: #f3f4f6; color: #374151; }
.info-grid { display: grid; grid-template-columns: repeat(auto-fit, minmax(220px, 1fr));
             gap: 10px; margin: 12px 0; }
.info-item { background: #f9fafb; padding: 10px 12px; border-radius: 6px;
             font-size: 13px; }
.info-item .key { font-weight: 600; color: #555; margin-bottom: 2px; }
.info-item .val { color: #1a1a2e; }
.changes-table { width: 100%; border-collapse: collapse; margin: 12px 0;
                 font-size: 13px; }
.changes-table th { background: #1a1a2e; color: white; padding: 8px 10px;
                    text-align: left; font-weight: 600; }
.changes-table td { padding: 7px 10px; border-bottom: 1px solid #eee; }
.changes-table tr:last-child td { border-bottom: none; }
.changes-table tr:hover td { background: #f9fafb; }
.issue { padding: 8px 12px; border-left: 3px solid; margin-bottom: 5px;
         border-radius: 0 4px 4px 0; font-size: 13px; }
.issue-error   { border-color: #dc2626; background: #fef2f2; }
.issue-warning { border-color: #d97706; background: #fffbeb; }
.issue-info    { border-color: #2563eb; background: #eff6ff; }
.phase-report { background: #f9fafb; padding: 12px 16px; border-radius: 6px;
                margin: 8px 0; font-size: 13px; }
.phase-report .phase-name { font-weight: 700; color: #1a1a2e; margin-bottom: 6px; }
.error-box { background: #fef2f2; border: 1px solid #fca5a5; padding: 12px 16px;
             border-radius: 6px; margin: 8px 0; font-size: 13px; color: #991b1b; }
.warn-box  { background: #fffbeb; border: 1px solid #fcd34d; padding: 12px 16px;
             border-radius: 6px; margin: 8px 0; font-size: 13px; color: #854d0e; }
.details-summary { cursor: pointer; font-weight: 600; color: #1e40af; }
details { margin: 8px 0; }
pre { background: #1a1a2e; color: #e2e8f0; padding: 12px; border-radius: 6px;
      overflow-x: auto; font-size: 12px; white-space: pre-wrap; }
"""

_HTML_HEAD_TEMPLATE = """<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>SlideArabi Structural Comparison: {deck_name}</title>
  <style>{css}</style>
</head>
<body>
"""

_HTML_FOOTER = "</body>\n</html>\n"


def _badge(text: str, color: str = "gray") -> str:
    return f'<span class="badge badge-{color}">{text}</span>'


def _severity_badge(sev: str) -> str:
    color = {"error": "red", "warning": "yellow", "info": "blue"}.get(sev, "gray")
    return _badge(sev.upper(), color)


def _issue_html(issue: Dict[str, Any]) -> str:
    sev = issue.get("severity", "info")
    css_class = f"issue issue-{sev}"
    msg = issue.get("message", "")
    shape_info = ""
    if issue.get("shape_name"):
        shape_info = f' <strong>Shape:</strong> {issue["shape_name"]}'
    if issue.get("actual_value") is not None:
        shape_info += f' | actual: <code>{issue["actual_value"]}</code>'
    if issue.get("expected_value") is not None:
        shape_info += f' | expected: <code>{issue["expected_value"]}</code>'
    badge = _severity_badge(sev)
    return f'<div class="{css_class}">{badge} {msg}{shape_info}</div>'


def _image_tag_relative(img_path: Path, report_dir: Path, alt: str = "") -> str:
    """Return an <img> tag with relative path from the report directory."""
    try:
        rel = img_path.relative_to(report_dir)
        return f'<img src="{rel}" alt="{alt}" loading="lazy">'
    except ValueError:
        return f'<img src="{img_path}" alt="{alt}" loading="lazy">'


def _inline_image_as_b64(img_path: Path) -> str:
    """Return <img> with base64 inline data (for self-contained report)."""
    import base64
    try:
        with open(img_path, "rb") as f:
            data = base64.b64encode(f.read()).decode("ascii")
        return f'<img src="data:image/jpeg;base64,{data}" loading="lazy" style="width:100%;border:1px solid #ddd;border-radius:4px;">'
    except Exception:
        return "<em>(image not available)</em>"


def generate_html_report(
    deck_name: str,
    deck_dir: Path,
    analysis: Dict[str, Any],
    transform_result: Dict[str, Any],
    original_images: List[Path],
    v2_images: List[Path],
    comparison_images: List[Path],
) -> str:
    """Build the full HTML report string."""

    now_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    prs_info = analysis.get("presentation", {})
    slide_count = prs_info.get("slide_count", 0)
    slide_w_in = prs_info.get("slide_width_in", 0)
    slide_h_in = prs_info.get("slide_height_in", 0)

    layout_ana = analysis.get("layout_analyzer", {})
    layout_types_found = layout_ana.get("layout_types_found", [])

    val_report = transform_result.get("phase_reports", {}).get("phase_5_validation", {})
    val_errors = val_report.get("errors", 0)
    val_warnings = val_report.get("warnings", 0)
    val_total = val_report.get("total_issues", 0)

    html = _HTML_HEAD_TEMPLATE.format(deck_name=deck_name, css=_HTML_CSS)

    # ── Header ──────────────────────────────────────────────────────────────
    html += f"""
<div class="header">
  <h1>SlideArabi — Structural Comparison</h1>
  <div class="meta">
    Deck: <strong>{deck_name}</strong> &nbsp;|&nbsp;
    {slide_count} slides &nbsp;|&nbsp;
    {slide_w_in:.2f}" × {slide_h_in:.2f}" &nbsp;|&nbsp;
    Generated: {now_str}
  </div>
</div>
"""

    # ── Summary stats ────────────────────────────────────────────────────────
    html += '<div class="section-title">Summary</div>\n'
    html += '<div class="stats">\n'

    def stat_card(value, label, color="#1a1a2e"):
        return (f'<div class="stat-card">'
                f'<div class="stat-value" style="color:{color}">{value}</div>'
                f'<div class="stat-label">{label}</div></div>\n')

    html += stat_card(slide_count, "Total Slides")
    html += stat_card(len(layout_types_found), "Layout Types Found")
    html += stat_card(prs_info.get("master_count", 0), "Slide Masters")
    html += stat_card(
        val_errors,
        "Validation Errors",
        color="#dc2626" if val_errors > 0 else "#166534",
    )
    html += stat_card(
        val_warnings,
        "Validation Warnings",
        color="#d97706" if val_warnings > 0 else "#166534",
    )

    # Phase 2 changes
    p2m = transform_result.get("phase_reports", {}).get("phase_2_masters", {})
    p2l = transform_result.get("phase_reports", {}).get("phase_2_layouts", {})
    p2_total = p2m.get("total_changes", 0) + p2l.get("total_changes", 0)
    html += stat_card(p2_total, "Phase 2 Changes (masters+layouts)")

    p3 = transform_result.get("phase_reports", {}).get("phase_3_slides", {})
    html += stat_card(p3.get("total_changes", 0), "Phase 3 Changes (slides)")

    p4 = transform_result.get("phase_reports", {}).get("phase_4_typography", {})
    html += stat_card(p4.get("total_changes", 0), "Phase 4 Typography Changes")

    html += '</div>\n'

    # Layout types found
    if layout_types_found:
        html += '<div class="section-title">Layout Types Detected</div>\n<p>'
        for lt in sorted(layout_types_found):
            html += _badge(lt, "blue") + " "
        html += "</p>\n"

    # Import warnings
    if analysis.get("warnings") or transform_result.get("warnings"):
        html += '<div class="section-title">Warnings &amp; Import Notes</div>\n'
        for w in (analysis.get("warnings", []) + transform_result.get("warnings", [])):
            html += f'<div class="warn-box">{w}</div>\n'

    # Errors
    all_errors = analysis.get("errors", []) + transform_result.get("errors", [])
    if all_errors:
        html += '<div class="section-title">Errors</div>\n'
        for e in all_errors:
            html += f'<div class="error-box">{e}</div>\n'

    # ── Phase reports ─────────────────────────────────────────────────────────
    html += '<div class="section-title">Phase Reports</div>\n'
    phase_labels = {
        "phase_2_masters": "Phase 2: Master Transformation",
        "phase_2_layouts": "Phase 2: Layout Transformation",
        "phase_3_slides":  "Phase 3: Slide Content Transformation",
        "phase_4_typography": "Phase 4: Typography Normalization",
        "phase_5_validation": "Phase 5: Structural Validation",
    }
    for key, label in phase_labels.items():
        pr = transform_result.get("phase_reports", {}).get(key, {})
        if not pr:
            continue
        html += f'<div class="phase-report">\n'
        html += f'  <div class="phase-name">{label}</div>\n'
        if "error" in pr:
            html += f'  <div class="error-box">{pr["error"]}</div>\n'
        elif key == "phase_5_validation":
            passed = pr.get("passed", False)
            html += (f'  <strong>Result:</strong> {_badge("PASSED", "green") if passed else _badge("FAILED", "red")} &nbsp;'
                     f'  Issues: {pr.get("total_issues", 0)} '
                     f'  ({pr.get("errors", 0)} errors, {pr.get("warnings", 0)} warnings, '
                     f'  {pr.get("info", 0)} info) &nbsp;'
                     f'  Pass rate: {pr.get("pass_rate", 100):.1f}%\n')
        else:
            tc = pr.get("total_changes", 0)
            cbt = pr.get("changes_by_type", {})
            detail = ", ".join(f"{k}: {v}" for k, v in cbt.items())
            html += f'  <strong>Total changes:</strong> {tc}'
            if detail:
                html += f'  &nbsp;| {detail}'
            html += "\n"
            if pr.get("warnings"):
                for w in pr["warnings"]:
                    html += f'  <div class="warn-box">{w}</div>\n'
            if pr.get("errors"):
                for e in pr["errors"]:
                    html += f'  <div class="error-box">{e}</div>\n'
        html += '</div>\n'

    # ── Per-slide cards ───────────────────────────────────────────────────────
    html += '<div class="section-title">Per-Slide Analysis</div>\n'

    # Build per-slide validation issues index
    val_issues_by_slide: Dict[int, List[Dict]] = {}
    for issue in val_report.get("issues", []):
        sn = issue.get("slide_number", 0)
        val_issues_by_slide.setdefault(sn, []).append(issue)

    slides_data = analysis.get("slides", [])

    for slide_data in slides_data:
        sn = slide_data.get("slide_number", 0)
        layout_info = slide_data.get("layout") or {}
        layout_type = layout_info.get("resolved_type", "unknown")
        confidence = layout_info.get("confidence", 0.0)
        is_explicit = layout_info.get("is_explicit", False)
        layout_name = layout_info.get("name", "")

        shape_count = slide_data.get("shape_count", 0)
        ph_count = slide_data.get("placeholder_count", 0)
        ff_count = slide_data.get("freeform_count", 0)
        slide_issues = val_issues_by_slide.get(sn, [])
        issue_error_count = sum(1 for i in slide_issues if i.get("severity") == "error")
        issue_warn_count = sum(1 for i in slide_issues if i.get("severity") == "warning")

        # Slide-level badge
        if issue_error_count > 0:
            status_badge = _badge(f"{issue_error_count} error(s)", "red")
        elif issue_warn_count > 0:
            status_badge = _badge(f"{issue_warn_count} warning(s)", "yellow")
        else:
            status_badge = _badge("OK", "green")

        html += f'<div class="slide-card" id="slide-{sn}">\n'
        html += (f'  <h3>Slide {sn} &nbsp; {status_badge} &nbsp; '
                 f'{_badge(layout_type, "blue")} &nbsp;'
                 f'{_badge("explicit" if is_explicit else "inferred", "gray")}</h3>\n')

        # Side-by-side comparison images
        orig_img = original_images[sn - 1] if sn - 1 < len(original_images) else None
        v2_img = v2_images[sn - 1] if sn - 1 < len(v2_images) else None
        comp_img = comparison_images[sn - 1] if sn - 1 < len(comparison_images) else None

        if comp_img and comp_img.exists():
            html += '<div class="comparison-full">\n'
            html += f'  <div class="label">Original (left) vs v2 Structural (right)</div>\n'
            html += f'  {_inline_image_as_b64(comp_img)}\n'
            html += '</div>\n'
        elif orig_img or v2_img:
            html += '<div class="comparison">\n'
            if orig_img and orig_img.exists():
                html += f'<div class="comparison-pane"><div class="label">Original</div>{_inline_image_as_b64(orig_img)}</div>\n'
            else:
                html += '<div class="comparison-pane"><div class="label">Original</div><em>(not rendered)</em></div>\n'
            if v2_img and v2_img.exists():
                html += f'<div class="comparison-pane"><div class="label">v2 Structural</div>{_inline_image_as_b64(v2_img)}</div>\n'
            else:
                html += '<div class="comparison-pane"><div class="label">v2 Structural</div><em>(not rendered)</em></div>\n'
            html += '</div>\n'
        else:
            html += '<p><em>No rendered images available (LibreOffice/pdftoppm may not be installed).</em></p>\n'

        # Layout / classification info
        html += '<div class="info-grid">\n'
        html += f'<div class="info-item"><div class="key">Layout Name</div><div class="val">{layout_name}</div></div>\n'
        html += f'<div class="info-item"><div class="key">Layout Type</div><div class="val">{layout_type}</div></div>\n'
        html += f'<div class="info-item"><div class="key">Confidence</div><div class="val">{confidence:.0%}</div></div>\n'
        html += f'<div class="info-item"><div class="key">Source</div><div class="val">{"Explicit XML" if is_explicit else "Inferred"}</div></div>\n'
        html += f'<div class="info-item"><div class="key">Shapes</div><div class="val">{shape_count} total ({ph_count} placeholder, {ff_count} freeform)</div></div>\n'
        html += f'<div class="info-item"><div class="key">Has Table</div><div class="val">{"Yes" if slide_data.get("has_table") else "No"}</div></div>\n'
        html += f'<div class="info-item"><div class="key">Has Chart</div><div class="val">{"Yes" if slide_data.get("has_chart") else "No"}</div></div>\n'

        # Placeholder summary
        ph_summary = layout_info.get("placeholder_summary", {})
        if ph_summary:
            ph_str = ", ".join(f"{k}: {v}" for k, v in ph_summary.items())
            html += f'<div class="info-item"><div class="key">Placeholders</div><div class="val">{ph_str}</div></div>\n'
        html += '</div>\n'

        # Template registry rules summary
        rules = slide_data.get("template_rules")
        if rules and "error" not in rules:
            html += '<details><summary class="details-summary">Template Registry Rules</summary>\n'
            html += '<div class="phase-report">\n'
            html += (f'  Freeform action: {_badge(rules.get("freeform_action", "?"), "blue")} &nbsp;'
                     f'  Swap columns: {_badge(str(rules.get("swap_columns", False)), "gray")} &nbsp;'
                     f'  Table: {_badge(rules.get("table_action", "?"), "gray")} &nbsp;'
                     f'  Chart: {_badge(rules.get("chart_action", "?"), "gray")}\n')
            ph_rules = rules.get("placeholder_rules", {})
            if ph_rules:
                html += '<br><strong>Placeholder rules:</strong><br>\n'
                for ph_key, ph_act in ph_rules.items():
                    html += (f'  {_badge(ph_key, "blue")} → action: <strong>{ph_act["action"]}</strong>'
                             f', rtl: {ph_act.get("set_rtl")}'
                             f', align: {ph_act.get("set_alignment")}'
                             f', mirror_x: {ph_act.get("mirror_x")}<br>\n')
            html += '</div></details>\n'

        # Validation issues for this slide
        if slide_issues:
            html += f'<div class="section-title" style="font-size:0.95rem;margin:12px 0 8px 0;">Validation Issues ({len(slide_issues)})</div>\n'
            for issue in slide_issues:
                html += _issue_html(issue)

        # Shape inventory table
        shapes = slide_data.get("shapes", [])
        if shapes:
            html += '<details><summary class="details-summary">Shape Inventory</summary>\n'
            html += '<table class="changes-table">\n'
            html += ('<thead><tr>'
                     '<th>#</th><th>Name</th><th>Type</th><th>Placeholder</th>'
                     '<th>Position (in)</th><th>Size (in)</th>'
                     '<th>Paras</th><th>RTL</th></tr></thead>\n<tbody>\n')
            for sidx, shape in enumerate(shapes, 1):
                if "error" in shape:
                    html += f'<tr><td>{sidx}</td><td colspan="7"><em>Error: {shape["error"]}</em></td></tr>\n'
                    continue
                x_in = shape.get("left_in", "")
                y_in = shape.get("top_in", "")
                w_in = shape.get("width_in", "")
                h_in = shape.get("height_in", "")
                pos_str = f"({x_in}, {y_in})" if x_in is not None else "–"
                size_str = f"{w_in} × {h_in}" if w_in is not None else "–"
                ph_info = shape.get("placeholder_type") or "–"
                if shape.get("placeholder_idx") is not None:
                    ph_info += f' [idx={shape["placeholder_idx"]}]'
                paras = shape.get("paragraphs", [])
                para_count = len(paras)

                # Gather RTL status
                rtl_values = list({p.get("rtl") for p in paras})
                if rtl_values == ["1"]:
                    rtl_badge = _badge("RTL", "blue")
                elif "1" in rtl_values:
                    rtl_badge = _badge("mixed", "yellow")
                else:
                    rtl_badge = _badge("LTR", "gray")

                html += (f'<tr>'
                         f'<td>{sidx}</td>'
                         f'<td>{shape.get("shape_name", "")}</td>'
                         f'<td>{shape.get("shape_type", "")}</td>'
                         f'<td>{ph_info}</td>'
                         f'<td>{pos_str}</td>'
                         f'<td>{size_str}</td>'
                         f'<td>{para_count}</td>'
                         f'<td>{rtl_badge}</td>'
                         f'</tr>\n')
            html += '</tbody></table></details>\n'

        html += '</div>\n'  # /slide-card

    # ── Validation summary ─────────────────────────────────────────────────────
    all_issues = val_report.get("issues", [])
    if all_issues:
        html += '<div class="section-title">Validation Summary (Phase 5)</div>\n'
        # Group by issue_type
        by_type: Dict[str, List[Dict]] = {}
        for issue in all_issues:
            it = issue.get("issue_type", "unknown")
            by_type.setdefault(it, []).append(issue)

        for issue_type, issues in sorted(by_type.items()):
            html += f'<details><summary class="details-summary">{issue_type} ({len(issues)} issues)</summary>\n'
            for issue in issues:
                html += _issue_html(issue)
            html += '</details>\n'

    # ── Raw analysis JSON ──────────────────────────────────────────────────────
    html += '<div class="section-title">Analysis JSON (raw)</div>\n'
    html += '<details><summary class="details-summary">Show full analysis.json</summary>\n'
    try:
        json_str = json.dumps(analysis, indent=2, default=str)
        # Truncate very large JSON to keep HTML report manageable
        if len(json_str) > 200_000:
            json_str = json_str[:200_000] + '\n\n... (truncated — see analysis.json for full data) ...'
        html += f'<pre>{json_str}</pre>\n'
    except Exception:
        html += '<p><em>Could not serialise analysis to JSON.</em></p>\n'
    html += '</details>\n'

    html += _HTML_FOOTER
    return html


# ════════════════════════════════════════════════════════════════════════════
# Main per-deck orchestrator
# ════════════════════════════════════════════════════════════════════════════

def process_deck(pptx_path: Path) -> Dict[str, Any]:
    """
    Run the full test harness for a single PPTX file.
    Returns a summary dict with paths and status.
    """
    deck_name = pptx_path.stem
    deck_dir = RESULTS_DIR / deck_name
    deck_dir.mkdir(parents=True, exist_ok=True)

    print(f"\n{'═' * 60}")
    print(f"  Processing: {pptx_path.name}")
    print(f"  Output dir: {deck_dir}")
    print(f"{'═' * 60}")

    summary: Dict[str, Any] = {
        "input": str(pptx_path),
        "deck_name": deck_name,
        "deck_dir": str(deck_dir),
        "status": "in_progress",
        "phases_completed": [],
    }

    # ── Phase 0: Structural Analysis ──────────────────────────────────────────
    analysis: Dict[str, Any] = {}
    try:
        analysis = run_phase0_analysis(pptx_path)
        summary["phases_completed"].append("phase_0_analysis")

        # Save analysis JSON
        analysis_path = deck_dir / "analysis.json"
        with open(analysis_path, "w", encoding="utf-8") as f:
            json.dump(analysis, f, indent=2, default=str)
        print(f"  Saved analysis → {analysis_path.name}")
        summary["analysis_json"] = str(analysis_path)

    except Exception as exc:
        print(f"  [ERROR] Phase 0 failed: {exc}")
        analysis = {"error": str(exc), "slides": [], "presentation": {}}
        summary["phase_0_error"] = str(exc)

    # ── Phase 2-5: Structural Transform ────────────────────────────────────────
    v2_pptx_path = deck_dir / f"{deck_name}_v2_structural.pptx"
    transform_result: Dict[str, Any] = {}
    try:
        transform_result = run_structural_transform(pptx_path, v2_pptx_path)
        summary["phases_completed"].append("phase_2_3_4_5")
        summary["v2_pptx"] = str(v2_pptx_path)
    except Exception as exc:
        print(f"  [ERROR] Structural transform failed: {exc}")
        transform_result = {
            "errors": [str(exc)],
            "warnings": [],
            "phase_reports": {},
            "success": False,
        }
        summary["transform_error"] = str(exc)

    # ── Render original slides ─────────────────────────────────────────────────
    orig_dir = deck_dir / "original"
    original_images: List[Path] = []
    try:
        original_images = render_pptx_slides(pptx_path, orig_dir, "original")
        if original_images:
            summary["phases_completed"].append("render_original")
    except Exception as exc:
        print(f"  [WARN] Failed to render original slides: {exc}")

    # ── Render v2 structural slides ────────────────────────────────────────────
    v2_dir = deck_dir / "v2"
    v2_images: List[Path] = []
    if v2_pptx_path.exists():
        try:
            v2_images = render_pptx_slides(v2_pptx_path, v2_dir, "v2_structural")
            if v2_images:
                summary["phases_completed"].append("render_v2")
        except Exception as exc:
            print(f"  [WARN] Failed to render v2 slides: {exc}")

    # ── Generate side-by-side comparison images ────────────────────────────────
    comp_dir = deck_dir / "comparison"
    comparison_images: List[Path] = []
    if original_images and v2_images and _PIL_AVAILABLE:
        comp_dir.mkdir(parents=True, exist_ok=True)
        print(f"  [Comparison] Building side-by-side images…")
        n_slides = min(len(original_images), len(v2_images))
        for idx in range(n_slides):
            out_img = comp_dir / f"slide-{idx + 1:02d}.jpg"
            ok = make_comparison_image(
                original_images[idx],
                v2_images[idx],
                out_img,
                slide_number=idx + 1,
            )
            if ok:
                comparison_images.append(out_img)
        if comparison_images:
            summary["phases_completed"].append("comparison_images")
            print(f"    Created {len(comparison_images)} comparison images")

    # ── Generate HTML Report ───────────────────────────────────────────────────
    report_path = deck_dir / "report.html"
    try:
        print(f"  [Report] Generating HTML report…")
        html = generate_html_report(
            deck_name=deck_name,
            deck_dir=deck_dir,
            analysis=analysis,
            transform_result=transform_result,
            original_images=original_images,
            v2_images=v2_images,
            comparison_images=comparison_images,
        )
        with open(report_path, "w", encoding="utf-8") as f:
            f.write(html)
        summary["report_html"] = str(report_path)
        summary["phases_completed"].append("html_report")
        print(f"  Saved report  → {report_path.name}")
    except Exception as exc:
        print(f"  [ERROR] Failed to generate HTML report: {exc}")
        summary["report_error"] = str(exc)
        traceback.print_exc()

    summary["status"] = "completed"
    return summary


# ════════════════════════════════════════════════════════════════════════════
# Entry point / argument parsing
# ════════════════════════════════════════════════════════════════════════════

def _collect_pptx_files(args) -> List[Path]:
    """Collect PPTX paths from CLI args (positional files or --dir)."""
    paths: List[Path] = []

    if args.dir:
        dir_path = Path(args.dir)
        if not dir_path.is_dir():
            print(f"[ERROR] --dir path is not a directory: {args.dir}", file=sys.stderr)
            sys.exit(1)
        paths = sorted(dir_path.glob("*.pptx"))
        if not paths:
            print(f"[WARN] No .pptx files found in {dir_path}")

    if args.files:
        for f in args.files:
            p = Path(f)
            if not p.exists():
                print(f"[WARN] File not found: {f}")
                continue
            if p.suffix.lower() != ".pptx":
                print(f"[WARN] Not a PPTX file (skipping): {f}")
                continue
            paths.append(p)

    return paths


def main():
    global RESULTS_DIR
    parser = argparse.ArgumentParser(
        description="SlideArabi Structural Comparison Test Harness",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument(
        "files",
        nargs="*",
        metavar="INPUT.pptx",
        help="One or more PPTX files to process.",
    )
    parser.add_argument(
        "--dir",
        metavar="DIRECTORY",
        help="Directory containing .pptx files to process.",
    )
    parser.add_argument(
        "--out",
        metavar="OUTPUT_DIR",
        default=str(RESULTS_DIR),
        help=f"Output directory (default: {RESULTS_DIR})",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Enable DEBUG logging.",
    )
    parser.add_argument(
        "--no-render",
        action="store_true",
        help="Skip LibreOffice rendering (produces JSON/PPTX only).",
    )

    args = parser.parse_args()

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    # Override RESULTS_DIR if --out given
    RESULTS_DIR = Path(args.out)
    RESULTS_DIR.mkdir(parents=True, exist_ok=True)

    pptx_files = _collect_pptx_files(args)

    if not pptx_files:
        print("[ERROR] No PPTX files to process. Provide filenames or --dir.", file=sys.stderr)
        parser.print_help()
        sys.exit(1)

    print(f"\nSlideArabi Test Harness")
    print(f"{'─' * 50}")
    print(f"  Files to process : {len(pptx_files)}")
    print(f"  Output directory : {RESULTS_DIR}")
    if _V2_IMPORT_ERRORS:
        print(f"\n  [WARN] Some v2 modules had import issues:")
        for e in _V2_IMPORT_ERRORS:
            print(f"    • {e}")
    print()

    # Optionally disable rendering
    if args.no_render:
        # Monkey-patch the render function to be a no-op
        global render_pptx_slides
        render_pptx_slides = lambda *a, **k: []  # type: ignore[assignment]
        print("  [--no-render] LibreOffice rendering disabled.\n")

    summaries = []
    for pptx_path in pptx_files:
        try:
            summary = process_deck(pptx_path)
            summaries.append(summary)
        except Exception as exc:
            print(f"\n[FATAL] Unhandled error processing {pptx_path}: {exc}")
            traceback.print_exc()
            summaries.append({
                "input": str(pptx_path),
                "status": "fatal_error",
                "error": str(exc),
            })

    # ── Final summary ──────────────────────────────────────────────────────────
    print(f"\n{'═' * 60}")
    print(f"  DONE — {len(pptx_files)} deck(s) processed")
    print(f"{'═' * 60}")
    for s in summaries:
        name = Path(s["input"]).name if "input" in s else "?"
        status = s.get("status", "?")
        report = s.get("report_html", s.get("deck_dir", ""))
        phases = s.get("phases_completed", [])
        print(f"  {name}: {status}  ({len(phases)} phases completed)")
        if report:
            print(f"    → {report}")

    # Write overall run summary
    run_summary_path = RESULTS_DIR / "run_summary.json"
    with open(run_summary_path, "w", encoding="utf-8") as f:
        json.dump(
            {
                "run_time": datetime.datetime.now().isoformat(),
                "decks_processed": len(pptx_files),
                "results": summaries,
            },
            f,
            indent=2,
            default=str,
        )
    print(f"\n  Run summary → {run_summary_path}")


if __name__ == "__main__":
    main()

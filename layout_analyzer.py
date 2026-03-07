"""
layout_analyzer.py — Slide layout type classifier.

SlideArabi: Template-First Deterministic RTL Transformation Engine.

Classifies each slide's layout into a canonical ST_SlideLayoutType for
deterministic transformation by the TemplateRegistry.

Classification strategy (priority order):
1. Read explicit `type` attribute from slideLayout XML element.
2. Infer from placeholder configuration using heuristic rules.
3. Fall back to spatial analysis (e.g., detect two-column by geometry).
4. Flag as 'cust' (custom) if no confident match — may require AI.

All 36 standard OOXML ST_SlideLayoutType values are supported.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from enum import Enum
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ST_SlideLayoutType — All 36 standard OOXML values
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

class SlideLayoutType(str, Enum):
    """All 36 standard OOXML ST_SlideLayoutType values."""
    TITLE = 'title'
    TX = 'tx'
    TWO_COL_TX = 'twoColTx'
    TBL = 'tbl'
    TX_AND_CHART = 'txAndChart'
    CHART_AND_TX = 'chartAndTx'
    DGM = 'dgm'
    CHART = 'chart'
    TX_AND_CLIP_ART = 'txAndClipArt'
    CLIP_ART_AND_TX = 'clipArtAndTx'
    TITLE_ONLY = 'titleOnly'
    BLANK = 'blank'
    TX_AND_OBJ = 'txAndObj'
    OBJ_AND_TX = 'objAndTx'
    OBJ_ONLY = 'objOnly'
    OBJ = 'obj'
    TX_AND_MEDIA = 'txAndMedia'
    MEDIA_AND_TX = 'mediaAndTx'
    OBJ_TX = 'objTx'
    TX_OBJ = 'txObj'
    OBJ_OVER_TX = 'objOverTx'
    TX_OVER_OBJ = 'txOverObj'
    TX_AND_TWO_OBJ = 'txAndTwoObj'
    TWO_OBJ_AND_TX = 'twoObjAndTx'
    TWO_OBJ_OVER_TX = 'twoObjOverTx'
    FOUR_OBJ = 'fourObj'
    TWO_TX_TWO_OBJ = 'twoTxTwoObj'
    TWO_OBJ_AND_OBJ = 'twoObjAndObj'
    SEC_HEAD = 'secHead'
    TWO_OBJ = 'twoObj'
    OBJ_AND_TWO_OBJ = 'objAndTwoObj'
    PIC_TX = 'picTx'
    VERT_TX = 'vertTx'
    VERT_TITLE_AND_TX = 'vertTitleAndTx'
    VERT_TITLE_AND_TX_OVER_CHART = 'vertTitleAndTxOverChart'
    CUST = 'cust'


# Set of all valid layout type strings for quick membership checks
_VALID_LAYOUT_TYPES = frozenset(lt.value for lt in SlideLayoutType)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Placeholder type normalisation
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

# python-pptx PP_PLACEHOLDER enum string suffixes → canonical type names.
# We normalise via .split('.')[-1].lower() to handle both raw strings
# and PP_PLACEHOLDER enums.
_PH_TYPE_MAP = {
    'title': 'title',
    'center_title': 'ctrTitle',
    'subtitle': 'subTitle',
    'body': 'body',
    'object': 'obj',
    'chart': 'chart',
    'table': 'tbl',
    'org_chart': 'dgm',
    'slide_number': 'sldNum',
    'date': 'dt',
    'footer': 'ftr',
    'slide_image': 'pic',
    'media_clip': 'media',
    'clip_art': 'clipArt',
    'bitmap': 'pic',
    'picture': 'pic',
    # Fallback — raw XML type values map to themselves
    'ctrtitle': 'ctrTitle',
    'subtitle': 'subTitle',
    'sldnum': 'sldNum',
    'dt': 'dt',
    'ftr': 'ftr',
    'tbl': 'tbl',
    'dgm': 'dgm',
    'pic': 'pic',
    'obj': 'obj',
    'media': 'media',
    'clipart': 'clipArt',
}

# Placeholder types that are "structural" (used in layout classification)
# vs "decorative" (date, footer, slide number — ignored in heuristics)
_DECORATIVE_PH_TYPES = frozenset({'dt', 'ftr', 'sldNum'})


def _normalise_ph_type(raw_type) -> str:
    """Normalise a python-pptx placeholder type to a canonical string.

    Handles PP_PLACEHOLDER enum members (e.g., PP_PLACEHOLDER.TITLE),
    raw XML type attribute strings (e.g., 'ctrTitle'), and None.

    Returns one of the canonical strings: 'title', 'ctrTitle', 'subTitle',
    'body', 'obj', 'chart', 'tbl', 'dgm', 'pic', 'media', 'clipArt',
    'dt', 'ftr', 'sldNum', or the raw string lowered if unknown.
    """
    if raw_type is None:
        return 'body'  # OOXML default when type attr is absent
    s = str(raw_type).split('.')[-1].lower()
    return _PH_TYPE_MAP.get(s, s)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Layout classification result
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

@dataclass
class LayoutClassification:
    """Classification result for a single slide's layout.

    Attributes:
        slide_number: 1-based slide index.
        layout_name: Display name from the layout element.
        explicit_type: The raw `type` attribute from the slideLayout XML,
            or None if absent / set to 'cust'.
        resolved_type: Final canonical layout type string (always populated).
        confidence: 1.0 for explicit XML types, 0.0–1.0 for inferred.
        placeholder_summary: Counts of structural placeholder types found
            (e.g., {'title': 1, 'body': 2}).
        requires_ai_classification: True when confidence is below threshold
            and the layout should be sent to an AI classifier.
    """
    slide_number: int
    layout_name: str
    explicit_type: Optional[str]
    resolved_type: str
    confidence: float
    placeholder_summary: Dict[str, int]
    requires_ai_classification: bool


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# LayoutAnalyzer
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

# Confidence threshold below which we flag for AI classification
_AI_CONFIDENCE_THRESHOLD = 0.7


class LayoutAnalyzer:
    """Classify each slide's layout into a canonical ST_SlideLayoutType.

    Accepts a python-pptx ``Presentation`` object and produces a mapping of
    slide number → LayoutClassification for every slide.

    Classification precedence:
    1. Explicit XML ``type`` attribute on the ``<p:sldLayout>`` element.
    2. Heuristic inference from placeholder configuration.
    3. Spatial analysis for ambiguous cases (e.g., two-column detection).
    4. Fallback to 'cust' with low confidence.

    Usage::

        from pptx import Presentation
        prs = Presentation('input.pptx')
        analyzer = LayoutAnalyzer(prs)
        classifications = analyzer.analyze_all()
        for slide_num, cls in classifications.items():
            print(f"Slide {slide_num}: {cls.resolved_type} ({cls.confidence:.0%})")
    """

    def __init__(self, presentation):
        """Initialise with a python-pptx Presentation object.

        Args:
            presentation: A ``pptx.Presentation`` instance.
        """
        self._prs = presentation
        self._slide_width = presentation.slide_width
        # Cache layout classifications to avoid re-computing for shared layouts
        self._layout_cache: Dict[int, Tuple[str, float]] = {}

    # ── public API ──────────────────────────────────────────────────────────

    def analyze_all(self) -> Dict[int, LayoutClassification]:
        """Classify every slide in the presentation.

        Returns:
            Mapping of 1-based slide number → ``LayoutClassification``.
        """
        results: Dict[int, LayoutClassification] = {}
        for idx, slide in enumerate(self._prs.slides):
            slide_number = idx + 1
            classification = self.classify_slide(slide, slide_number)
            results[slide_number] = classification
            logger.debug(
                'Slide %d: layout=%r type=%s (%.0f%% confidence)',
                slide_number,
                classification.layout_name,
                classification.resolved_type,
                classification.confidence * 100,
            )
        logger.info(
            'LayoutAnalyzer: classified %d slides', len(results),
        )
        return results

    def classify_slide(self, slide, slide_number: int) -> LayoutClassification:
        """Classify a single slide.

        Args:
            slide: A python-pptx ``Slide`` object.
            slide_number: 1-based slide index.

        Returns:
            ``LayoutClassification`` for the slide.
        """
        layout = slide.slide_layout
        layout_name = layout.name or ''

        resolved_type, confidence = self.classify_layout(layout)

        # Build placeholder summary from the slide itself
        # (may differ from layout if user added/removed placeholders)
        ph_summary = self._get_placeholder_summary(slide)

        explicit_type = self._get_explicit_type(layout)

        requires_ai = confidence < _AI_CONFIDENCE_THRESHOLD

        return LayoutClassification(
            slide_number=slide_number,
            layout_name=layout_name,
            explicit_type=explicit_type,
            resolved_type=resolved_type,
            confidence=confidence,
            placeholder_summary=ph_summary,
            requires_ai_classification=requires_ai,
        )

    def classify_layout(self, layout) -> Tuple[str, float]:
        """Classify a slide layout element.

        Uses caching so that shared layouts (multiple slides using the same
        layout) are only classified once.

        Args:
            layout: A python-pptx ``SlideLayout`` object.

        Returns:
            Tuple of (layout_type_string, confidence).
        """
        cache_key = id(layout)
        if cache_key in self._layout_cache:
            return self._layout_cache[cache_key]

        # Strategy 1: Read explicit type from XML
        explicit = self._get_explicit_type(layout)
        if explicit is not None and explicit != 'cust':
            result = (explicit, 1.0)
            self._layout_cache[cache_key] = result
            return result

        # Strategy 2: Infer from placeholder configuration
        inferred_type, confidence = self._infer_type_from_placeholders(layout)
        result = (inferred_type, confidence)
        self._layout_cache[cache_key] = result
        return result

    # ── private helpers ─────────────────────────────────────────────────────

    def _get_explicit_type(self, layout) -> Optional[str]:
        """Read the ``type`` attribute from the ``<p:sldLayout>`` XML element.

        Args:
            layout: A python-pptx ``SlideLayout`` object.

        Returns:
            The type string if present and valid, else None.
        """
        try:
            raw = layout._element.get('type')
            if raw is None:
                return None
            # Validate that it's a known type
            if raw in _VALID_LAYOUT_TYPES:
                return raw
            # Unknown value — log and return as-is (may be vendor extension)
            logger.debug('Unknown layout type attribute: %r', raw)
            return raw
        except Exception:
            return None

    def _infer_type_from_placeholders(self, layout) -> Tuple[str, float]:
        """Infer layout type from its placeholder configuration.

        Applies ordered heuristic rules matching placeholder type counts
        and spatial arrangement to canonical layout types.

        Args:
            layout: A python-pptx ``SlideLayout`` object.

        Returns:
            Tuple of (inferred_type_string, confidence).
        """
        ph_summary = self._get_placeholder_summary(layout)

        # Count structural placeholders (exclude decorative)
        ctr_title = ph_summary.get('ctrTitle', 0)
        sub_title = ph_summary.get('subTitle', 0)
        title_count = ph_summary.get('title', 0)
        body_count = ph_summary.get('body', 0)
        obj_count = ph_summary.get('obj', 0)
        chart_count = ph_summary.get('chart', 0)
        tbl_count = ph_summary.get('tbl', 0)
        pic_count = ph_summary.get('pic', 0)
        dgm_count = ph_summary.get('dgm', 0)
        media_count = ph_summary.get('media', 0)
        clip_art_count = ph_summary.get('clipArt', 0)

        total_structural = sum(
            v for k, v in ph_summary.items() if k not in _DECORATIVE_PH_TYPES
        )

        # ── Rule 0: No structural placeholders → blank
        if total_structural == 0:
            return ('blank', 0.95)

        # ── Rule 1: Center title + subtitle → title slide
        if ctr_title >= 1 and sub_title >= 1:
            return ('title', 0.95)

        # ── Rule 2: Title only (no body, no objects, no content)
        if (title_count >= 1 and body_count == 0 and obj_count == 0
                and chart_count == 0 and tbl_count == 0 and pic_count == 0
                and dgm_count == 0 and media_count == 0 and clip_art_count == 0
                and ctr_title == 0):
            return ('titleOnly', 0.9)

        # ── Rule 3: Title + table
        if title_count >= 1 and tbl_count >= 1:
            return ('tbl', 0.9)

        # ── Rule 4: Title + chart only
        if title_count >= 1 and chart_count >= 1 and body_count == 0:
            return ('chart', 0.9)

        # ── Rule 5: Title + diagram
        if title_count >= 1 and dgm_count >= 1:
            return ('dgm', 0.85)

        # ── Rule 6: Title + body + chart → text and chart
        if title_count >= 1 and body_count >= 1 and chart_count >= 1:
            return ('txAndChart', 0.85)

        # ── Rule 7: Title + body + media → text and media
        if title_count >= 1 and body_count >= 1 and media_count >= 1:
            return ('txAndMedia', 0.8)

        # ── Rule 8: Title + body + clip art → text and clip art
        if title_count >= 1 and body_count >= 1 and clip_art_count >= 1:
            return ('txAndClipArt', 0.8)

        # ── Rule 9: Title + picture → picture + text
        if title_count >= 1 and pic_count >= 1:
            return ('picTx', 0.85)

        # ── Rule 10: Title + 2 body — check spatial for two-column
        if title_count >= 1 and body_count == 2:
            placeholders = self._collect_body_placeholders(layout)
            if self._detect_two_column_spatial(placeholders, self._slide_width):
                return ('twoColTx', 0.85)
            return ('twoColTx', 0.75)

        # ── Rule 11: Title + 2 objects
        if title_count >= 1 and obj_count == 2 and body_count == 0:
            return ('twoObj', 0.85)

        # ── Rule 12: Title + 1 body + 2 objects
        if title_count >= 1 and body_count >= 1 and obj_count == 2:
            return ('txAndTwoObj', 0.8)

        # ── Rule 13: Title + 1 body + 1 object → text and object
        if title_count >= 1 and body_count >= 1 and obj_count >= 1:
            return ('txAndObj', 0.85)

        # ── Rule 14: 4 objects
        if obj_count == 4:
            return ('fourObj', 0.85)

        # ── Rule 15: Title + single body → standard text layout
        if title_count >= 1 and body_count == 1:
            return ('tx', 0.9)

        # ── Rule 16: Title + single object → object layout
        if title_count >= 1 and obj_count == 1:
            return ('obj', 0.85)

        # ── Rule 17: Object(s) only, no title
        if obj_count >= 1 and title_count == 0 and ctr_title == 0:
            return ('objOnly', 0.8)

        # ── Rule 18: Section header pattern (title + body, no other content)
        # This is lower priority than tx because sec_head is layout-name driven
        # We only reach here if previous rules didn't match

        # ── Rule 19: Body only (no title, no center title)
        if body_count >= 1 and title_count == 0 and ctr_title == 0:
            return ('tx', 0.6)

        # ── Fallback: custom / unclassified
        logger.debug(
            'Could not confidently classify layout. '
            'Placeholder summary: %s', ph_summary,
        )
        return ('cust', 0.4)

    def _get_placeholder_summary(self, layout_or_slide) -> Dict[str, int]:
        """Count structural placeholders by normalised type.

        Decorative placeholders (date, footer, slide number) are still
        counted but are excluded from heuristic rules by the caller.

        Args:
            layout_or_slide: A python-pptx SlideLayout or Slide object
                with a ``placeholders`` attribute.

        Returns:
            Dict mapping placeholder type string → count.
        """
        counts: Dict[str, int] = {}
        try:
            placeholders = layout_or_slide.placeholders
        except Exception:
            return counts

        for ph in placeholders:
            try:
                raw_type = ph.placeholder_format.type
                ph_type = _normalise_ph_type(raw_type)
                counts[ph_type] = counts.get(ph_type, 0) + 1
            except Exception:
                # Defensive: if placeholder_format access fails, skip
                logger.debug(
                    'Could not read placeholder type for shape %r',
                    getattr(ph, 'name', '?'),
                )
        return counts

    def _collect_body_placeholders(self, layout) -> List[Any]:
        """Collect body-type placeholders with position info.

        Returns a list of placeholder shape objects that are classified as
        'body' (for spatial analysis).
        """
        bodies = []
        try:
            for ph in layout.placeholders:
                try:
                    raw_type = ph.placeholder_format.type
                    ph_type = _normalise_ph_type(raw_type)
                    if ph_type == 'body':
                        bodies.append(ph)
                except Exception:
                    continue
        except Exception:
            pass
        return bodies

    def _detect_two_column_spatial(
        self, placeholders: List[Any], slide_width: int
    ) -> bool:
        """Detect two-column layout by spatial analysis of body placeholders.

        Two body placeholders are considered "two columns" if:
        - One is primarily in the left half of the slide.
        - The other is primarily in the right half.
        - They don't significantly overlap horizontally.

        Args:
            placeholders: List of body placeholder shape objects.
            slide_width: Slide width in EMU.

        Returns:
            True if spatial analysis indicates a two-column arrangement.
        """
        if len(placeholders) < 2:
            return False

        try:
            midpoint = slide_width // 2
            left_found = False
            right_found = False

            for ph in placeholders:
                ph_left = ph.left
                ph_width = ph.width
                if ph_left is None or ph_width is None:
                    continue
                ph_center = ph_left + ph_width // 2

                if ph_center < midpoint:
                    left_found = True
                else:
                    right_found = True

            return left_found and right_found
        except Exception:
            return False

    # ── batch helpers ───────────────────────────────────────────────────────

    def get_layout_type_for_slide(self, slide_number: int) -> Optional[str]:
        """Convenience: get resolved type for a slide by number.

        Must call ``analyze_all()`` first to populate the cache.

        Args:
            slide_number: 1-based slide index.

        Returns:
            Layout type string, or None if not yet analyzed.
        """
        # This is a convenience accessor; the caller should use analyze_all()
        # and keep the returned dict.
        try:
            slides_list = list(self._prs.slides)
            if slide_number < 1 or slide_number > len(slides_list):
                return None
            slide = slides_list[slide_number - 1]
            layout = slide.slide_layout
            resolved_type, _ = self.classify_layout(layout)
            return resolved_type
        except Exception:
            return None

    def get_all_layout_types(self) -> Dict[str, str]:
        """Get a mapping of layout name → layout type for all layouts.

        Useful for debugging and reporting.

        Returns:
            Dict mapping layout display name → resolved type string.
        """
        result: Dict[str, str] = {}
        try:
            for master in self._prs.slide_masters:
                for layout in master.slide_layouts:
                    name = layout.name or '(unnamed)'
                    lt, _ = self.classify_layout(layout)
                    result[name] = lt
        except Exception:
            pass
        return result
